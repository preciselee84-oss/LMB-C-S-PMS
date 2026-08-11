from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


BASE = Path(__file__).resolve().parent
ASSET_DIR = BASE / "erp_casebook_assets"
OUTPUT = BASE / "통합CMS_ERP연계_사례집.pptx"

FONT = "Malgun Gothic"
NAVY = RGBColor(20, 42, 72)
BLUE = RGBColor(38, 114, 196)
TEAL = RGBColor(0, 143, 130)
LIGHT_BLUE = RGBColor(232, 242, 252)
LIGHT_GRAY = RGBColor(244, 247, 250)
MID_GRAY = RGBColor(113, 128, 150)
TEXT = RGBColor(31, 41, 55)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, x, y, w, h, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.9, 0.45, 15, 0.55, "통합CMS ERP연계 사례집", 10, True, BLUE)
    add_textbox(slide, 0.9, 1.05, 20, 0.9, title, 25, True, NAVY)
    if subtitle:
        add_textbox(slide, 0.95, 1.9, 20, 0.55, subtitle, 10.5, False, MID_GRAY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0.9), Cm(2.62), Cm(23.6), Cm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(210, 218, 226)
    line.line.fill.background()


def add_footer(slide, page):
    add_textbox(slide, 0.9, 13.55, 5, 0.3, "통합CMS 프로그램", 7.5, False, MID_GRAY)
    add_textbox(slide, 23.4, 13.55, 1, 0.3, str(page), 7.5, False, MID_GRAY, PP_ALIGN.RIGHT)


def add_chip(slide, x, y, text, color=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(3.0), Cm(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_textbox(slide, x, y + 0.08, 3.0, 0.36, text, 9, True, WHITE, PP_ALIGN.CENTER)
    return shape


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(213, 221, 230)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(y), Cm(0.12), Cm(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, x + 0.35, y + 0.35, w - 0.7, 0.45, title, 13, True, NAVY)
    box = slide.shapes.add_textbox(Cm(x + 0.35), Cm(y + 1.0), Cm(w - 0.7), Cm(h - 1.25))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, line in enumerate(body):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(9.3)
        p.font.color.rgb = TEXT
        p.space_after = Pt(4)
    return card


def add_table(slide, x, y, w, h, rows, cols, data, header_fill=NAVY):
    table = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(w), Cm(h)).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.margin_left = Cm(0.12)
            cell.margin_right = Cm(0.12)
            cell.margin_top = Cm(0.06)
            cell.margin_bottom = Cm(0.06)
            fill = header_fill if r == 0 else (LIGHT_GRAY if r % 2 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(8.5 if rows > 7 else 9)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else TEXT
            cell.vertical_anchor = 3
    return table


def fit_picture(slide, path, x, y, w, h):
    pic = slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w))
    if pic.height > Cm(h):
        pic.height = Cm(h)
    pic.left = Cm(x + (w - pic.width.cm) / 2)
    pic.top = Cm(y + (h - pic.height.cm) / 2)
    return pic


def add_flow(slide, y, items):
    x = 1.2
    for idx, (title, desc, color) in enumerate(items):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(4.2), Cm(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        add_textbox(slide, x, y + 0.18, 4.2, 0.35, title, 10.5, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.15, y + 0.7, 3.9, 0.36, desc, 7.8, False, WHITE, PP_ALIGN.CENTER)
        if idx < len(items) - 1:
            add_textbox(slide, x + 4.35, y + 0.42, 0.6, 0.35, "→", 18, True, MID_GRAY, PP_ALIGN.CENTER)
        x += 5.1


def add_screen_slide(prs, title, image_paths, caption, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, caption)
    x = 0.9
    for idx, path in enumerate(image_paths):
        box_w = 7.6 if len(image_paths) == 3 else 11.4
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(3.05), Cm(box_w), Cm(8.6))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RGBColor(198, 210, 224)
        fit_picture(slide, path, x + 0.25, 3.35, box_w - 0.5, 7.8)
        x += box_w + 0.45
    add_footer(slide, page)


def add_case_slide(prs, page, erp, protocol, jobs, flow, notes, image_path=None, accent=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, f"{erp} 연계 사례", f"{protocol} 기반으로 통합CMS 조회/이체 데이터를 ERP 업무 화면에 반영")
    add_chip(slide, 1.0, 3.0, protocol, accent)
    add_textbox(slide, 1.0, 3.85, 8.4, 0.5, "대표 연계업무", 13, True, NAVY)
    for i, job in enumerate(jobs):
        add_card(slide, 1.0 + (i % 2) * 4.65, 4.5 + (i // 2) * 1.8, 4.25, 1.25, job[0], [job[1]], accent)
    add_textbox(slide, 10.6, 3.0, 13.4, 0.55, "처리 흐름", 13, True, NAVY)
    add_flow(slide, 3.85, flow)
    add_card(slide, 10.6, 5.75, 6.0, 3.3, "적용 포인트", notes[:3], accent)
    add_card(slide, 17.1, 5.75, 6.9, 3.3, "고객 기대효과", notes[3:], TEAL)
    if image_path:
        fit_picture(slide, image_path, 10.6, 9.45, 13.4, 3.5)
    add_footer(slide, page)


def build():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.288)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(14.288))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(246, 249, 252)
    bg.line.fill.background()
    add_textbox(slide, 1.4, 1.2, 12, 0.5, "통합CMS 프로그램", 13, True, BLUE)
    add_textbox(slide, 1.4, 2.2, 16, 1.1, "ERP연계 사례집", 34, True, NAVY)
    add_textbox(slide, 1.45, 3.45, 16.5, 0.8, "더존, 영림원, SAP, 자체개발 ERP별 연계 방식과 활용 화면 구성", 15, False, TEXT)
    add_flow(
        slide,
        5.0,
        [
            ("통합CMS", "은행/카드/세금계산서 수집", BLUE),
            ("연계 Agent", "스케줄·검증·전송", TEAL),
            ("ERP", "회계·자금 업무 반영", NAVY),
            ("업무 활용", "조회·이체·정산 자동화", RGBColor(93, 95, 239)),
        ],
    )
    fit_picture(slide, ASSET_DIR / "corp_card_1.png", 1.4, 7.5, 11.2, 4.7)
    fit_picture(slide, ASSET_DIR / "bulk_transfer.png", 13.0, 7.5, 11.2, 4.7)
    add_footer(slide, 1)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "사례집 구성", "ERP사별 프로토콜과 대표 업무를 한눈에 비교")
    matrix = [
        ["ERP사", "연계 프로토콜", "대표 업무", "구성 방향"],
        ["더존 아마란스 10", "API", "계좌거래내역, 법인카드, 세금계산서", "표준 API 호출 및 응답값 검증"],
        ["더존 옴니이솔 10", "DB TO DB", "계좌원장, 집금내역, 대량이체", "ERP DB 중계 테이블 적재"],
        ["영림원", "DB TO DB", "계좌거래내역, 법인카드 승인/이용", "ERP 기준 테이블 매핑"],
        ["SAP", "RFC", "지급이체, 법인카드, 전표처리", "RFC Function 파라미터 송수신"],
        ["자체개발", "DB TO DB / API", "고객사 맞춤 조회·이체·정산", "고객 개발규격 기반 유연 연계"],
    ]
    add_table(slide, 1.0, 3.05, 23.4, 5.0, len(matrix), 4, matrix)
    add_card(
        slide,
        1.0,
        8.6,
        7.2,
        2.9,
        "조회 업무",
        ["계좌거래내역, 법인카드 승인/이용, 세금계산서 데이터를 ERP 화면에서 확인", "수기 다운로드와 업로드 업무를 줄이는 구조"],
        BLUE,
    )
    add_card(
        slide,
        9.1,
        8.6,
        7.2,
        2.9,
        "이체 업무",
        ["대량이체·급여이체·단건이체 데이터를 통합CMS로 전송", "승인·전송·결과 회신 흐름을 ERP와 연결"],
        TEAL,
    )
    add_card(
        slide,
        17.2,
        8.6,
        7.2,
        2.9,
        "운영 관리",
        ["스케줄러, 서버 IP, 처리방식, 전송 특이사항을 기준으로 운영 체크", "장애 발생 시 조회/전송 구간을 분리해 점검"],
        RGBColor(93, 95, 239),
    )
    add_footer(slide, 2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "공통 연계 아키텍처", "통합CMS가 금융 데이터를 수집하고 ERP별 프로토콜에 맞춰 전달")
    add_flow(
        slide,
        3.0,
        [
            ("금융기관", "계좌·카드·이체 결과", BLUE),
            ("통합CMS", "조회/수집/검증", TEAL),
            ("연계 Agent", "스케줄·변환·로그", RGBColor(93, 95, 239)),
            ("ERP", "API·DB·RFC 반영", NAVY),
        ],
    )
    add_card(slide, 1.1, 5.3, 7.0, 4.2, "API 방식", ["ERP가 제공하는 API 엔드포인트로 데이터 전송", "인증키, 응답코드, 재전송 정책을 기준으로 운영", "대표 ERP: 더존 아마란스 10, 자체개발"], BLUE)
    add_card(slide, 9.2, 5.3, 7.0, 4.2, "DB TO DB 방식", ["ERP DB의 연계 테이블 또는 View에 데이터 적재", "컬럼 매핑, PK, 중복 스킵 규칙을 사전에 정의", "대표 ERP: 옴니이솔 10, 영림원, 자체개발"], TEAL)
    add_card(slide, 17.3, 5.3, 7.0, 4.2, "RFC 방식", ["SAP RFC Function을 호출해 파라미터 송수신", "Function별 입력값과 Return 메시지 기준으로 검증", "대표 ERP: SAP"], RGBColor(93, 95, 239))
    add_footer(slide, 3)

    add_screen_slide(
        prs,
        "ERP 예시 화면: 법인카드 내역 조회",
        [ASSET_DIR / "corp_card_1.png", ASSET_DIR / "corp_card_2.png", ASSET_DIR / "corp_card_3.png"],
        "카드번호, 승인일자, 금액, 가맹점 정보 등 법인카드 연계 데이터가 ERP 조회 그리드에 반영되는 형태",
        4,
    )
    add_screen_slide(
        prs,
        "ERP 예시 화면: 지급 이체 관리",
        [ASSET_DIR / "bulk_transfer.png"],
        "ERP 지급 대상 데이터를 통합CMS로 전송하고 이체 결과를 다시 확인하는 대량이체 업무 화면 예시",
        5,
    )

    flow_api = [("ERP 요청", "조회 조건/전송 데이터", BLUE), ("통합CMS API", "검증·변환", TEAL), ("금융 처리", "조회/이체 실행", RGBColor(93, 95, 239)), ("ERP 반영", "응답/결과 저장", NAVY)]
    flow_db = [("통합CMS", "데이터 수집", BLUE), ("중계 테이블", "Insert/Update", TEAL), ("ERP 배치", "정합성 체크", RGBColor(93, 95, 239)), ("ERP 화면", "업무 활용", NAVY)]
    flow_rfc = [("통합CMS", "전송 데이터 구성", BLUE), ("RFC 호출", "Function 파라미터", TEAL), ("SAP 처리", "전표/조회 반영", RGBColor(93, 95, 239)), ("결과 회신", "Return 메시지", NAVY)]

    add_case_slide(
        prs,
        6,
        "더존 아마란스 10",
        "API",
        [("계좌거래내역", "기간별 거래내역 조회 후 ERP 수납/자금 화면 반영"), ("법인카드", "승인/이용 데이터 자동 수집 및 비용 처리 지원"), ("세금계산서", "매입·매출 세금계산서 조회 데이터 연계"), ("대량이체", "ERP 지급대상 데이터를 CMS 이체로 전송")],
        flow_api,
        ["API 인증정보와 호출 URL을 고객사 환경에 맞게 설정", "응답코드 기준 성공/실패 로그를 남겨 재처리 가능", "ERP 화면에서 조회 조건 입력 후 통합CMS 데이터를 확인", "수기 엑셀 업로드 감소", "전표 처리 전 데이터 정합성 확보", "운영자가 장애 구간을 API 응답 기준으로 빠르게 확인"],
        ASSET_DIR / "corp_card_1.png",
        BLUE,
    )
    add_case_slide(
        prs,
        7,
        "더존 옴니이솔 10",
        "DB TO DB",
        [("계좌원장", "계좌별 잔액/거래 데이터를 ERP DB 기준으로 적재"), ("집금내역", "가상계좌·집금 결과를 ERP 수납 업무에 반영"), ("대량이체", "지급대상 데이터를 중계 테이블로 전달"), ("환율", "일자별 환율 데이터 ERP 기준정보 반영")],
        flow_db,
        ["고객 ERP DB 접속정보와 중계 테이블 권한을 사전 확인", "PK, 처리상태, 전송일시 컬럼으로 중복 반영 방지", "스케줄러가 지정 시간에 자동 적재", "조회/이체 담당자의 반복 입력 감소", "마감 시점 데이터 누락 점검 용이", "ERP 화면 수정 없이 기존 업무 흐름 유지"],
        ASSET_DIR / "bulk_transfer.png",
        TEAL,
    )
    add_case_slide(
        prs,
        8,
        "영림원",
        "DB TO DB",
        [("계좌거래내역", "은행 거래내역을 ERP 자금관리 화면에 연결"), ("법인카드 승인내역", "승인번호·공급가액·부가세 기준으로 비용 처리 지원"), ("법인카드 이용내역", "가맹점 정보와 카드 소유자 기준 조회"), ("예금주조회", "이체 전 계좌 실명 검증 결과 활용")],
        flow_db,
        ["영림원 기준 테이블의 필수 컬럼과 코드값 매핑", "승인/취소 데이터의 중복 처리 규칙 정의", "조회 데이터와 ERP 비용처리 항목 연결", "카드 데이터 누락 문의 감소", "회계 담당자의 증빙 확인 시간 단축", "ERP 표준 화면에서 통합 조회 가능"],
        ASSET_DIR / "corp_card_1.png",
        RGBColor(82, 122, 183),
    )
    add_case_slide(
        prs,
        9,
        "SAP",
        "RFC",
        [("지급이체", "SAP 지급 데이터를 RFC로 받아 통합CMS 이체 실행"), ("법인카드", "승인/이용 내역을 SAP Function으로 전달"), ("계좌거래내역", "기간별 은행 거래내역을 SAP에 반영"), ("전자어음", "어음 조회/처리 결과를 SAP 업무에 연결")],
        flow_rfc,
        ["SAP Function Module별 입력/출력 파라미터 정의", "처리일자, 회사코드, 계좌번호 등 필수값 검증", "Return 메시지 기준으로 성공/오류 내역 관리", "SAP 업무 화면에서 외부 금융 데이터 활용", "이체 결과 회신으로 후속 전표 처리 지원", "RFC 로그를 기준으로 장애 원인 추적 가능"],
        ASSET_DIR / "bulk_transfer.png",
        RGBColor(93, 95, 239),
    )
    add_case_slide(
        prs,
        10,
        "자체개발 ERP",
        "DB TO DB / API",
        [("법인카드", "고객사 카드 모니터링 화면에 승인/이용 데이터 제공"), ("대량이체", "고객 개발 시스템의 지급 데이터를 CMS 이체로 연계"), ("세금계산서", "매입·매출 조회 데이터를 자체 정산 화면에 반영"), ("국내/외화송금", "지급 요청과 처리 결과를 개발 규격에 맞춰 송수신")],
        [("고객 규격", "API/DB 정의서", BLUE), ("통합CMS", "데이터 변환", TEAL), ("검증/테스트", "샘플 데이터 확인", RGBColor(93, 95, 239)), ("운영 반영", "스케줄/로그 관리", NAVY)],
        ["고객사 개발 규격서 기준으로 인터페이스를 설계", "API와 DB 방식을 업무별로 혼합 적용 가능", "테스트 데이터로 필드 길이·코드·금액 정합성 확인", "기존 ERP 기능을 유지하면서 금융 데이터 확장", "고객 맞춤 화면에 필요한 컬럼만 선별 제공", "추가 업무 연계 시 재사용 가능한 구조"],
        ASSET_DIR / "corp_card_3.png",
        NAVY,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "업무별 활용 시나리오", "통합CMS 연계 데이터가 ERP 업무에서 사용되는 대표 흐름")
    scenarios = [
        ("법인카드 승인/이용", "카드사 데이터 수집 → ERP 비용처리/증빙 검토 → 누락·중복 점검"),
        ("대량이체", "ERP 지급대상 생성 → 통합CMS 전송 → 은행 이체 실행 → 결과 회신"),
        ("계좌거래내역", "은행 거래 조회 → ERP 자금일보/수납 화면 반영 → 마감 검증"),
        ("세금계산서", "홈택스 매입·매출 조회 → ERP 증빙/정산 업무 반영 → 담당자 검토"),
        ("환율/예금주조회", "외부 조회 결과 수집 → ERP 기준정보 또는 지급 검증에 활용"),
        ("전자어음/B2B", "은행/금융기관 자료 조회 → ERP 채권·채무 관리 업무 반영"),
    ]
    for i, (title, body) in enumerate(scenarios):
        add_card(slide, 1.0 + (i % 3) * 7.8, 3.2 + (i // 3) * 3.0, 7.0, 2.3, title, [body], [BLUE, TEAL, RGBColor(93, 95, 239)][i % 3])
    add_footer(slide, 11)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "도입 전 확인 항목", "연계 방식이 달라도 사전 확인 구조를 통일하면 구축과 운영이 안정적입니다.")
    checks = [
        ["구분", "확인 내용", "담당"],
        ["접속정보", "ERP 서버 IP, DB 계정, API URL, SAP RFC Destination", "고객사/개발사"],
        ["데이터", "업무별 필수 컬럼, 코드값, 금액/일자 포맷, PK 기준", "통합CMS/개발사"],
        ["스케줄", "조회·전송 주기, 업무 마감 시간, 재처리 기준", "통합CMS/고객사"],
        ["검증", "샘플 데이터 반영, ERP 화면 확인, 전송 성공/실패 로그", "전체"],
        ["운영", "장애 연락체계, 변경 관리, 추가 연계 요청 절차", "전체"],
    ]
    add_table(slide, 1.0, 3.15, 23.4, 4.9, len(checks), 3, checks)
    add_card(slide, 1.0, 8.7, 11.2, 2.6, "권장 진행 순서", ["1. 고객 ERP와 연계 방식 확정", "2. 업무별 컬럼 매핑표 작성", "3. 샘플 데이터 전송 및 ERP 화면 검증", "4. 스케줄 적용 후 운영 로그 확인"], BLUE)
    add_card(slide, 13.2, 8.7, 11.2, 2.6, "사례집 활용 포인트", ["ERP사별 표준 설명자료로 사용", "영업/구축 미팅에서 프로토콜 차이를 간단히 안내", "화면 예시와 업무 시나리오를 함께 제시"], TEAL)
    add_footer(slide, 12)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build()
