from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


BASE = Path(__file__).resolve().parent
ASSET_DIR = BASE / "erp_casebook_assets"
OUTPUT = BASE / "통합CMS_ERP연계_고객사_안내자료.pptx"

FONT = "Malgun Gothic"
DEEP = RGBColor(18, 54, 63)
GREEN = RGBColor(0, 96, 82)
HANA = RGBColor(0, 168, 136)
MINT = RGBColor(229, 247, 243)
PALE = RGBColor(245, 250, 249)
LINE = RGBColor(198, 220, 216)
TEXT = RGBColor(35, 45, 52)
MUTED = RGBColor(95, 113, 121)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 200, 77)


def set_run(run, size, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def text_box(slide, x, y, w, h, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return shape


def fill(shape, color, line_color=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def rect(slide, x, y, w, h, color, line_color=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Cm(x), Cm(y), Cm(w), Cm(h))
    return fill(shape, color, line_color)


def title(slide, heading, sub=None):
    rect(slide, 0, 0, 25.4, 14.288, PALE)
    rect(slide, 0, 0, 25.4, 0.22, HANA)
    text_box(slide, 1.0, 0.8, 20.5, 0.78, heading, 24, True, DEEP)
    if sub:
        text_box(slide, 1.02, 1.66, 20.5, 0.45, sub, 10.5, False, MUTED)
    rect(slide, 1.0, 2.55, 23.4, 0.03, LINE)


def footer(slide, page):
    text_box(slide, 1.05, 13.45, 8, 0.28, "하나은행 통합CMS ERP연계 안내", 7.5, False, MUTED)
    text_box(slide, 23.5, 13.45, 0.8, 0.28, str(page), 7.5, True, MUTED, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, label, body, accent=HANA, number=None, body_size=8.7):
    rect(slide, x, y, w, h, WHITE, LINE, True)
    rect(slide, x, y, w, 0.13, accent)
    if number:
        badge = rect(slide, x + 0.35, y + 0.35, 0.72, 0.72, accent, None, True)
        text_box(slide, x + 0.35, y + 0.43, 0.72, 0.35, str(number), 11, True, WHITE, PP_ALIGN.CENTER)
        tx = x + 1.25
        tw = w - 1.6
    else:
        tx = x + 0.45
        tw = w - 0.9
    text_box(slide, tx, y + 0.34, tw, 0.48, label, 13, True, DEEP)
    box = slide.shapes.add_textbox(Cm(x + 0.45), Cm(y + 1.03), Cm(w - 0.9), Cm(h - 1.23))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    lines = body if isinstance(body, list) else [body]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(body_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(3)
    return box


def pill(slide, x, y, w, text, color=GREEN):
    rect(slide, x, y, w, 0.56, color, None, True)
    text_box(slide, x, y + 0.04, w, 0.38, text, 9.3, True, WHITE, PP_ALIGN.CENTER)


def table(slide, x, y, w, h, rows):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Cm(x), Cm(y), Cm(w), Cm(h)).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.margin_left = Cm(0.1)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.06)
            cell.margin_bottom = Cm(0.06)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DEEP if r == 0 else (WHITE if r % 2 else MINT)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_run(run, 8.8, r == 0, WHITE if r == 0 else TEXT)
    return tbl


def flow(slide, y, items):
    x = 1.1
    for idx, (head, desc, color) in enumerate(items):
        rect(slide, x, y, 4.25, 1.25, color, None, True)
        text_box(slide, x, y + 0.16, 4.25, 0.35, head, 10.5, True, WHITE, PP_ALIGN.CENTER)
        text_box(slide, x + 0.2, y + 0.68, 3.85, 0.32, desc, 7.8, False, WHITE, PP_ALIGN.CENTER)
        if idx < len(items) - 1:
            text_box(slide, x + 4.35, y + 0.45, 0.65, 0.28, "→", 17, True, GREEN, PP_ALIGN.CENTER)
        x += 5.05


def fit_pic(slide, path, x, y, w, h):
    pic = slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w))
    if pic.height > Cm(h):
        pic.height = Cm(h)
    pic.left = Cm(x + (w - pic.width.cm) / 2)
    pic.top = Cm(y + (h - pic.height.cm) / 2)
    return pic


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 25.4, 14.288, DEEP)
    rect(slide, 0, 0, 25.4, 0.26, HANA)
    text_box(slide, 1.2, 1.65, 14.2, 1.45, "통합CMS ERP연계\n고객사 안내자료", 30, True, WHITE)
    text_box(slide, 1.25, 3.95, 14.0, 0.55, "은행 업무 결과를 ERP 회계·자금 업무로 연결합니다.", 13.2, False, MINT)
    card(slide, 1.25, 5.3, 5.1, 1.95, "조회", ["계좌·카드·세금계산서\n데이터 자동 수집"], HANA, "1", 8.0)
    card(slide, 6.7, 5.3, 5.1, 1.95, "전송", ["API·DB·RFC 방식으로\nERP에 전달"], HANA, "2", 8.0)
    card(slide, 12.15, 5.3, 5.1, 1.95, "활용", ["전표·정산·마감 업무에서\n바로 확인"], HANA, "3", 8.0)
    rect(slide, 18.15, 1.75, 5.85, 8.75, MINT, None, True)
    fit_pic(slide, ASSET_DIR / "corp_card_1.png", 18.55, 2.2, 5.05, 3.0)
    fit_pic(slide, ASSET_DIR / "bulk_transfer.png", 18.55, 5.95, 5.05, 3.0)
    text_box(slide, 1.25, 12.7, 10, 0.36, "하나은행 통합CMS", 10, True, MINT)


def erp_slide(prs, page, name, protocol, summary, bullets, image, accent=HANA):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, f"{name} 연계 안내", f"{protocol} 방식으로 통합CMS 데이터를 ERP 업무에 연결")
    pill(slide, 1.05, 3.05, 3.2, protocol, accent)
    text_box(slide, 1.1, 3.85, 13.4, 0.9, summary, 13.5, True, DEEP)
    card(slide, 1.1, 5.1, 6.7, 2.55, "연계 대상 업무", bullets[:3], accent, None, 8.4)
    card(slide, 8.2, 5.1, 6.7, 2.55, "운영 포인트", bullets[3:6], GREEN, None, 8.4)
    card(slide, 1.1, 8.05, 13.8, 2.55, "고객사 기대효과", bullets[6:], HANA, None, 8.4)
    rect(slide, 15.7, 3.25, 8.4, 7.55, WHITE, LINE, True)
    fit_pic(slide, image, 16.05, 3.7, 7.7, 6.5)
    footer(slide, page)


def build():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.288)

    cover(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "ERP연계를 도입하면 업무 흐름이 짧아집니다", "은행 조회·이체 결과가 ERP 화면과 마감 데이터로 이어집니다.")
    flow(slide, 3.2, [("은행/카드", "거래·승인·이체 결과", DEEP), ("통합CMS", "수집·검증·스케줄", HANA), ("연계 방식", "API·DB·RFC", GREEN), ("ERP 업무", "전표·정산·마감", DEEP)])
    card(slide, 1.1, 5.7, 7.0, 3.1, "수기 입력 감소", ["엑셀 다운로드, 복사/붙여넣기, ERP 재입력 과정을 줄입니다.", "담당자는 입력보다 검토와 승인 업무에 집중할 수 있습니다."], HANA, "1")
    card(slide, 9.2, 5.7, 7.0, 3.1, "데이터 정합성 확보", ["사업자번호, 계좌번호, 승인번호, 금액 등 핵심값을 기준으로 중복과 누락을 점검합니다.", "전송 로그로 오류 구간을 빠르게 확인합니다."], GREEN, "2")
    card(slide, 17.3, 5.7, 7.0, 3.1, "마감시간 단축", ["계좌거래내역, 법인카드, 이체 결과가 ERP에 연결되어 월말 확인 시간이 줄어듭니다.", "반복 조회와 수동 대조 부담을 낮춥니다."], DEEP, "3")
    footer(slide, 2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "ERP사별 권장 연계 방식", "고객사 ERP 환경에 맞춰 안정적인 프로토콜을 선택합니다.")
    rows = [
        ["ERP사", "권장 프로토콜", "주요 연계업무", "고객사 확인사항"],
        ["더존 아마란스 10", "API", "계좌거래내역, 법인카드, 세금계산서", "API URL, 인증정보, 응답코드"],
        ["더존 옴니이솔 10", "DB TO DB", "계좌원장, 집금내역, 대량이체", "DB 접속정보, 중계 테이블"],
        ["영림원", "DB TO DB", "법인카드 승인/이용, 거래내역", "컬럼 매핑, 코드값, 처리상태"],
        ["SAP", "RFC", "지급이체, 법인카드, 전표처리", "RFC Destination, Function 파라미터"],
        ["자체개발", "DB TO DB / API", "고객사 맞춤 조회·이체·정산", "개발 규격서, 테스트 데이터"],
    ]
    table(slide, 1.0, 3.2, 23.4, 5.1, rows)
    card(slide, 1.1, 9.0, 11.1, 2.2, "안내 포인트", ["ERP 제품명보다 중요한 것은 고객사 운영 방식입니다.", "조회 업무와 이체 업무를 나누어 프로토콜을 정하면 장애 대응이 쉬워집니다."], HANA)
    card(slide, 13.1, 9.0, 11.1, 2.2, "사전 준비자료", ["ERP 접속정보, 연계 테이블/API 명세, 업무별 샘플 데이터", "운영 담당자와 개발 담당자 연락처를 함께 준비하면 구축 기간을 줄일 수 있습니다."], GREEN)
    footer(slide, 3)

    erp_slide(
        prs,
        4,
        "더존 아마란스 10",
        "API",
        "표준 API 호출로 통합CMS 조회/전송 결과를 아마란스 업무 화면에 반영합니다.",
        ["계좌거래내역", "법인카드 승인·이용내역", "세금계산서 매입·매출", "API 인증정보와 호출 URL 확인", "응답코드별 성공/실패 로그 관리", "재전송 기준 사전 정의", "ERP 화면에서 최신 금융 데이터를 빠르게 확인", "수기 업로드와 입력 오류 감소", "마감 전 데이터 검토 시간이 단축"],
        ASSET_DIR / "corp_card_1.png",
        HANA,
    )
    erp_slide(
        prs,
        5,
        "더존 옴니이솔 10",
        "DB TO DB",
        "ERP DB 중계 테이블에 데이터를 적재해 기존 업무 화면에서 활용합니다.",
        ["계좌원장", "집금내역", "대량이체", "DB 접속 권한과 IP 허용 확인", "PK와 처리상태 컬럼 기준 정의", "스케줄러 실행 시간 협의", "ERP 표준 업무 흐름을 유지", "중복 반영과 누락 점검 용이", "운영자가 DB 처리상태로 장애 확인"],
        ASSET_DIR / "bulk_transfer.png",
        GREEN,
    )
    erp_slide(
        prs,
        6,
        "영림원",
        "DB TO DB",
        "영림원 기준 테이블 구조에 맞춰 거래내역과 법인카드 데이터를 연결합니다.",
        ["계좌거래내역", "법인카드 승인내역", "법인카드 이용내역", "ERP 기준 코드와 필수 컬럼 확인", "승인/취소 데이터 처리 규칙 정의", "전송 후 ERP 조회 화면 검증", "회계 담당자의 증빙 확인 시간 감소", "카드 데이터 누락 문의 대응이 쉬움", "조회부터 비용처리까지 흐름 단순화"],
        ASSET_DIR / "corp_card_2.png",
        HANA,
    )
    erp_slide(
        prs,
        7,
        "SAP",
        "RFC",
        "SAP RFC Function을 호출해 지급·카드·거래 데이터를 송수신합니다.",
        ["지급이체", "법인카드 승인·이용내역", "계좌거래내역", "RFC Destination과 권한 확인", "Function별 입력/출력 파라미터 정의", "Return 메시지 기준 오류 관리", "SAP 기준 업무 승인 흐름과 연동", "이체 결과 회신으로 후속 전표 처리 지원", "대형 고객사 표준 통제 구조에 적합"],
        ASSET_DIR / "bulk_transfer.png",
        DEEP,
    )
    erp_slide(
        prs,
        8,
        "자체개발 ERP",
        "DB TO DB / API",
        "고객사 개발 규격에 맞춰 업무별로 DB와 API 방식을 유연하게 적용합니다.",
        ["법인카드 모니터링", "대량이체·급여이체", "세금계산서·정산 데이터", "고객사 인터페이스 정의서 확인", "샘플 데이터로 필드 길이와 코드 검증", "업무별 API/DB 혼합 적용 가능", "기존 ERP 화면과 프로세스를 유지", "필요한 컬럼만 맞춤 제공", "추가 업무 확장이 쉬운 구조"],
        ASSET_DIR / "corp_card_3.png",
        GREEN,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "ERP 예시 화면", "고객사 ERP 화면에서 통합CMS 데이터를 조회하고 이체 결과를 확인합니다.")
    rect(slide, 1.0, 3.0, 11.4, 7.6, WHITE, LINE, True)
    fit_pic(slide, ASSET_DIR / "corp_card_1.png", 1.35, 3.35, 10.7, 5.7)
    text_box(slide, 1.4, 9.45, 10.6, 0.4, "법인카드 내역 조회", 12, True, DEEP, PP_ALIGN.CENTER)
    rect(slide, 13.0, 3.0, 11.4, 7.6, WHITE, LINE, True)
    fit_pic(slide, ASSET_DIR / "bulk_transfer.png", 13.35, 3.35, 10.7, 5.7)
    text_box(slide, 13.4, 9.45, 10.6, 0.4, "지급 이체 관리", 12, True, DEEP, PP_ALIGN.CENTER)
    card(slide, 1.0, 11.0, 23.4, 1.4, "화면 활용 방식", ["조회 화면에서는 카드번호·승인일자·금액·가맹점 정보를 확인하고, 이체 화면에서는 ERP 지급대상을 통합CMS로 전송해 처리 결과를 회신받습니다."], HANA)
    footer(slide, 9)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "도입 진행 절차", "자료 확인부터 운영 반영까지 단계별로 진행합니다.")
    flow(slide, 3.2, [("1. 연계 범위", "업무·ERP·프로토콜 확정", DEEP), ("2. 환경 확인", "서버·DB·API·RFC 정보", HANA), ("3. 개발/설정", "컬럼 매핑·스케줄 등록", GREEN), ("4. 검증/오픈", "샘플 반영·운영 전환", DEEP)])
    card(slide, 1.1, 5.75, 7.0, 3.1, "고객사 준비", ["ERP 담당자/개발사 연락처", "서버 IP, DB/API/RFC 접속정보", "업무별 테스트 데이터"], HANA)
    card(slide, 9.2, 5.75, 7.0, 3.1, "통합CMS 준비", ["업무별 연계 설정", "스케줄러 및 로그 확인", "조회/전송 테스트 지원"], GREEN)
    card(slide, 17.3, 5.75, 7.0, 3.1, "오픈 후 운영", ["전송 성공/실패 모니터링", "월마감 전 누락 점검", "추가 연계업무 확장 협의"], DEEP)
    rect(slide, 1.1, 10.1, 23.2, 1.25, MINT, LINE, True)
    text_box(slide, 1.4, 10.45, 22.4, 0.35, "통합CMS ERP연계는 고객사 업무 방식에 맞춰 조회·이체·정산 데이터를 연결하는 고객 맞춤형 구축 서비스입니다.", 12.5, True, GREEN, PP_ALIGN.CENTER)
    footer(slide, 10)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build()
