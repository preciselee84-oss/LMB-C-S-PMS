# VERSION: 20260518_staff_filter_deadline
import streamlit as st
import pandas as pd
import numpy as np
import time
import json
import os
import html
import re
import base64
import hashlib
import requests as _requests
from io import BytesIO
from datetime import datetime, timedelta
from streamlit_cookies_controller import CookieController

st.set_page_config(page_title="내부 관리", layout="wide", initial_sidebar_state="expanded")


def resolve_template_file(file_name):
    base_dir = os.path.dirname(__file__)
    candidates = [
        os.path.join(base_dir, "templates", file_name),
        os.path.join(base_dir, "github-work", "templates", file_name),
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return candidates[0]


DB_FILE = "users.json"
PERF_FILE = "manual_perf.json"
SAVED_STATE_FILE = "saved_state.json"
SERVER_CONNECTION_FILE = "server_connection.json"
DELEGATED_WORKPLACE_FILE = "delegated_workplaces.json"
APPROVAL_FILE = "approvals.json"
BANK_ACCOUNT_FILE = "bank_accounts.json"
USAGE_REPORT_FILE = "usage_reports.json"
COMPANY_PROFILE_FILE = "company_profile.json"
EXCEL_SAMPLE_FILE = resolve_template_file("LMB월간 활동실적_000000(샘플).xlsx")
PPT_TEMPLATE_FILE = resolve_template_file("LMB활동실적보고서_202605_하나지사.pptx")
WEEKLY_PPT_TEMPLATE_FILE = resolve_template_file("주간보고_통합CMS고객_개설운영_주간보고_템플릿.pptx")

DEFAULT_URL_ANALYSIS = "https://docs.google.com/spreadsheets/d/e/2PACX-1vT9XPHqrqcaFf9bCOVya7yHORr-c1R4KCF0eEpdE3ESn8qJELP0BkqTOslur9bsGcVabRUIcyOa877R/pub?output=csv"
DEFAULT_URL_SYNC = "https://docs.google.com/spreadsheets/d/1yS4gaES-iuzt1NSRTSdj9Ivg1fjbN5mIyX4pGnvEYN0/export?format=csv&gid=1533424484"
DEFAULT_URL_HANA = "https://docs.google.com/spreadsheets/d/e/2PACX-1vQgRHnTZD4eDW2UeODQuGxmxFrflKpbQda3sBsVjj1s3qAFWMKcpke2U58UuT6VEDlkbXveZlaroTCr/pub?gid=0&single=true&output=csv"
DEFAULT_URL_HANA_BILLING = "https://docs.google.com/spreadsheets/d/e/2PACX-1vQgRHnTZD4eDW2UeODQuGxmxFrflKpbQda3sBsVjj1s3qAFWMKcpke2U58UuT6VEDlkbXveZlaroTCr/pub?gid=1172734914&single=true&output=csv"
DEFAULT_URL_HANA_PERFORMANCE = "https://docs.google.com/spreadsheets/d/e/2PACX-1vS2CUE3No1cptBOTehN8r1xoTQyUni07sjbut-f1Teo9mpB-rcJgpE5xfI6dTy0M4IUxSg8Mv5_uT4l/pub?gid=1749034066&single=true&output=csv"
URL_EDUCATION_WAITING = "https://docs.google.com/spreadsheets/d/1btxxNGPw-SLEnhvWSyyDbSHf-SniXZDS4IJit3wiqXE/export?format=csv&gid=0"


@st.cache_data(ttl=300, show_spinner=False)
def read_csv_cached(url, header=0, cache_buster=0):
    url = normalize_google_sheet_csv_url(url)
    return pd.read_csv(url, header=header)


def read_google_csv(url, header=0, force_refresh=False):
    cache_buster = int(datetime.utcnow().timestamp()) if force_refresh else 0
    return read_csv_cached(url, header=header, cache_buster=cache_buster).copy()


def normalize_google_sheet_csv_url(url):
    """Google Sheets pubhtml/pub links are normalized to direct CSV export."""
    if not isinstance(url, str):
        return url
    normalized = url.strip()
    if "docs.google.com/spreadsheets" not in normalized:
        return normalized
    normalized = normalized.replace("/pubhtml?", "/pub?")
    normalized = normalized.replace("/pubhtml", "/pub")
    if "/pub?" in normalized and "output=csv" not in normalized:
        separator = "&" if "?" in normalized else "?"
        normalized = f"{normalized}{separator}output=csv"
    return normalized


GITHUB_REPO = "preciselee84-oss/LMB-C-S-PMS"
GITHUB_BRANCH = "main"
GITHUB_DATA_DIR = "data"
SESSION_UID_COOKIE = "auto_login_uid"
LAST_MENU_COOKIE = "last_menu"
BILLING_MENU = "청구자료 생성"
ACTIVITY_TEMPLATE_CONVERT_MENU = "활동이력 템플릿 변환"
OPERATION_TARGET_MENU = "운영관리 활동고객 선정"

CRM_MENU_LABELS = {
    "대시보드": "홈 대시보드",
    "업로드 및 실적 확인": "활동 이력",
    "이번달 활동 대상고객 추천": "대상 고객 추천",
    OPERATION_TARGET_MENU: "운영관리 타깃",
    "주간보고 이력 작성": "주간 리포트",
    "관리자용 실적 확인": "성과 관리",
    "실적 분석/계산": "성과 분석",
    "실적 보고서": "성과 리포트",
    "주간보고 취합": "팀 리포트",
    "운영계획": "운영 플랜",
    BILLING_MENU: "청구자료 생성",
    "청구자료 작성": "청구자료 작성",
    "직원 및 권한설정": "사용자/권한",
    "구글 스트레드시트 연동": "데이터 연동",
    ACTIVITY_TEMPLATE_CONVERT_MENU: "이력 템플릿 변환",
}


def _get_github_token():
    try:
        return st.secrets.get("GITHUB_TOKEN", "")
    except Exception:
        return os.environ.get("GITHUB_TOKEN", "")


# ── DART 전자공시 API ─────────────────────────────────────────────

def _get_dart_api_key():
    # 1) 설정 메뉴에서 입력한 값 우선
    session_key = st.session_state.get("dart_api_key", "")
    if session_key:
        return session_key
    # 2) Streamlit secrets → 환경변수
    try:
        return st.secrets.get("DART_API_KEY", "")
    except Exception:
        return os.environ.get("DART_API_KEY", "")


@st.cache_data(ttl=86400)
def _load_dart_corp_map(api_key):
    """DART 기업코드 목록 로드 ({정규화기업명: corp_code}). 24시간 캐시."""
    if not api_key:
        return {}
    try:
        import zipfile, io
        import xml.etree.ElementTree as ET
        resp = _requests.get(
            f"https://opendart.fss.or.kr/api/corpCode.xml?crtfc_key={api_key}",
            timeout=30,
        )
        if resp.status_code != 200:
            return {}
        with zipfile.ZipFile(io.BytesIO(resp.content)) as zf:
            xml_bytes = zf.read("CORPCODE.xml")
        root = ET.fromstring(xml_bytes.decode("utf-8"))
        return {
            item.findtext("corp_name", "").strip().replace(" ", "").lower(): item.findtext("corp_code", "").strip()
            for item in root.findall("list")
            if item.findtext("corp_name", "").strip() and item.findtext("corp_code", "").strip()
        }
    except Exception:
        return {}


@st.cache_data(ttl=3600)
def _dart_get_filings(api_key, corp_code, months=6):
    """DART 최근 N개월 공시 목록. 1시간 캐시."""
    if not api_key or not corp_code:
        return []
    try:
        today = datetime.utcnow() + timedelta(hours=9)
        bgn = (today - timedelta(days=months * 31)).strftime("%Y%m%d")
        end = today.strftime("%Y%m%d")
        resp = _requests.get(
            "https://opendart.fss.or.kr/api/list.json",
            params={
                "crtfc_key": api_key, "corp_code": corp_code,
                "bgn_de": bgn, "end_de": end, "page_count": 10,
            },
            timeout=10,
        )
        d = resp.json()
        return d.get("list", []) if d.get("status") == "000" else []
    except Exception:
        return []


# 유통활동 관련 공시 키워드 → (가산점수, 사유 설명)
_DART_KW_SCORE = [
    ("유상증자",  30, "유상증자 공시 → 자금 유입, 이체 활성화 기대"),
    ("무상증자",  15, "무상증자 공시 → 재무구조 개선"),
    ("합병",      25, "합병 공시 → 계좌 통합 니즈 발생"),
    ("분할",      20, "기업분할 공시 → 신규 계좌 가능성"),
    ("수주",      20, "수주 공시 → 매출 증가·자금 흐름 활성화"),
    ("신규사업",  15, "신규사업 공시 → 향후 거래 확대 가능"),
    ("투자",      10, "투자 공시 → 자금 유출입 예상"),
]


def _dart_enrich(api_key, corp_map, company_name):
    """기업명으로 DART 공시 조회 후 (가산점수, 사유 목록) 반환."""
    key = str(company_name).strip().replace(" ", "").lower()
    corp_code = corp_map.get(key)
    if not corp_code:
        return 0, []
    filings = _dart_get_filings(api_key, corp_code)
    if not filings:
        return 0, []
    bonus, reasons, seen = 0, [], set()
    for f in filings:
        nm = f.get("report_nm", "")
        dt = str(f.get("rcept_dt", ""))[:8]
        for kw, pts, desc in _DART_KW_SCORE:
            if kw in nm and kw not in seen:
                bonus = max(bonus, pts)
                reasons.append(f"[DART공시] {desc} ({dt[:4]}.{dt[4:6]}.{dt[6:]})")
                seen.add(kw)
    return bonus, reasons


# ── 국세청 사업자등록 상태조회 API ─────────────────────────────────

def _get_nts_api_key():
    session_key = st.session_state.get("nts_api_key", "")
    if session_key:
        return session_key
    try:
        return st.secrets.get("NTS_API_KEY", "")
    except Exception:
        return os.environ.get("NTS_API_KEY", "")


def _check_business_status(api_key, biz_numbers):
    """사업자등록번호 목록의 휴폐업 상태를 조회. {사업자번호(숫자만): 상태} 반환."""
    clean_numbers = sorted({re.sub(r"\D", "", str(b)) for b in biz_numbers if str(b).strip()})
    if not api_key or not clean_numbers:
        return {}
    try:
        resp = _requests.post(
            f"https://api.odcloud.kr/api/nts-businessman/v1/status?serviceKey={api_key}",
            json={"b_no": clean_numbers},
            headers={"Content-Type": "application/json"},
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json().get("data", [])
        return {item.get("b_no", ""): item.get("b_stt", "") or "조회결과 없음" for item in data}
    except Exception:
        return {}


# ── 예금주조회 (계좌 실명조회 API 연동 전 임시 구현) ────────────────

def _lookup_account_holder_name(bank_name, account_number, expected_holder=""):
    """계좌 실명조회 API 연동 전까지, 등록된 입금은행/계좌번호가 있으면 예상예금주를 조회 결과로 반환."""
    if not bank_name or not account_number:
        return ""
    return expected_holder


def _github_save(file_path, data):
    token = _get_github_token()
    if not token:
        return False
    try:
        headers = {"Authorization": f"token {token}", "Accept": "application/vnd.github.v3+json"}
        gh_path = f"{GITHUB_DATA_DIR}/{file_path}"
        url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{gh_path}"
        resp = _requests.get(url, headers=headers, params={"ref": GITHUB_BRANCH}, timeout=10)
        sha = resp.json().get("sha") if resp.status_code == 200 else None
        content = base64.b64encode(json.dumps(data, ensure_ascii=False, indent=4).encode()).decode()
        payload = {"message": f"auto-save {file_path}", "content": content, "branch": GITHUB_BRANCH}
        if sha:
            payload["sha"] = sha
        resp = _requests.put(url, json=payload, headers=headers, timeout=15)
        return resp.status_code in (200, 201)
    except Exception:
        return False


def _github_load(file_path, default_data):
    token = _get_github_token()
    if not token:
        return default_data
    try:
        headers = {"Authorization": f"token {token}", "Accept": "application/vnd.github.v3+json"}
        gh_path = f"{GITHUB_DATA_DIR}/{file_path}"
        url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{gh_path}"
        resp = _requests.get(url, headers=headers, params={"ref": GITHUB_BRANCH}, timeout=10)
        if resp.status_code == 200:
            content = base64.b64decode(resp.json()["content"]).decode()
            return json.loads(content)
    except Exception:
        pass
    return default_data


def load_db(file_path, default_data):
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    data = _github_load(file_path, default_data)
    if data != default_data:
        try:
            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=4)
        except Exception:
            pass
    return data


def save_db(file_path, data, allow_shrink=False):
    if file_path == DB_FILE and not allow_shrink:
        existing = load_user_db()
        if len(_real_users(existing)) > len(_real_users(data)):
            data = merge_user_db(existing, data)
    with open(file_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=4)
    # 직원DB는 .bak 백업 파일도 항상 동기화 (로컬 데이터 유실 방지)
    if file_path == DB_FILE:
        try:
            with open(file_path + ".bak", "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=4)
        except Exception:
            pass
    _github_save(file_path, data)


def safe_cookie_controller():
    try:
        return CookieController()
    except Exception:
        return None


def cookie_get(cookie_manager, key, default=""):
    if cookie_manager is None:
        return default
    try:
        return cookie_manager.get(key) or default
    except Exception:
        return default


def cookie_set(cookie_manager, key, value, **kwargs):
    if cookie_manager is None:
        return False
    try:
        cookie_manager.set(key, value, **kwargs)
        return True
    except TypeError:
        try:
            cookie_manager.set(key, value)
            return True
        except Exception:
            return False
    except Exception:
        return False


def cookie_remove(cookie_manager, key):
    if cookie_manager is None:
        return False
    try:
        cookie_manager.remove(key)
        return True
    except Exception:
        return False


def restore_saved_work_state(user_name):
    saved_db = load_db(SAVED_STATE_FILE, {})
    user_saved = saved_db.get(user_name)
    if user_saved:
        if user_saved.get("user_excel_data"):
            st.session_state.user_excel_data = pd.DataFrame.from_dict(user_saved["user_excel_data"])
        if user_saved.get("user_prev_month_sel"):
            st.session_state.user_prev_month_sel = user_saved["user_prev_month_sel"]

    admin_analysis = saved_db.get("admin_analysis")
    if admin_analysis and admin_analysis.get("sent_df"):
        st.session_state.analysis_result = pd.DataFrame.from_dict(admin_analysis["sent_df"])
    if admin_analysis and admin_analysis.get("adm_prev_month"):
        st.session_state.adm_prev_month = admin_analysis["adm_prev_month"]
    deadline_info = saved_db.get("deadline")
    if deadline_info:
        st.session_state.deadline_time = deadline_info.get("time", "")
    report_closed_info = saved_db.get("report_closed")
    if report_closed_info:
        st.session_state.report_closed = report_closed_info.get("time", "")


def restore_login_from_cookie():
    if st.session_state.get("logged_in"):
        return
    cookie_manager = safe_cookie_controller()
    uid = str(cookie_get(cookie_manager, SESSION_UID_COOKIE, "")).strip()
    if not uid:
        return

    db = st.session_state.get("user_db", {})
    is_super = uid == "1"
    is_user = uid in db and db[uid].get("access") == "허용"
    if not (is_super or is_user):
        cookie_remove(cookie_manager, SESSION_UID_COOKIE)
        return

    user = db.get(uid, {"role": "관리자", "name": "최고관리자"})
    st.session_state.logged_in = True
    st.session_state.user_role = user.get("role", "관리자")
    st.session_state.user_name = user.get("name", "최고관리자")
    if not st.session_state.get("login_time"):
        st.session_state.login_time = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M")

    last_menu = str(cookie_get(cookie_manager, LAST_MENU_COOKIE, "")).strip()
    if last_menu:
        st.session_state.current_menu = last_menu

    restore_saved_work_state(st.session_state.user_name)


def persist_login_session(user_id):
    cookie_manager = safe_cookie_controller()
    cookie_set(cookie_manager, SESSION_UID_COOKIE, str(user_id), max_age=60 * 60 * 24 * 30)


def persist_current_menu():
    cookie_manager = safe_cookie_controller()
    menu = st.session_state.get("current_menu", "")
    if menu:
        cookie_set(cookie_manager, LAST_MENU_COOKIE, str(menu), max_age=60 * 60 * 24 * 30)


def send_kakao_notify(message):
    try:
        token = ""
        try:
            token = st.secrets.get("KAKAO_TOKEN", "")
        except Exception:
            token = os.environ.get("KAKAO_TOKEN", "")
        if not token:
            return False
        headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/x-www-form-urlencoded"}
        data = {"template_object": json.dumps({"object_type": "text", "text": message, "link": {"web_url": "", "mobile_web_url": ""}, "button_title": "확인"})}
        resp = _requests.post("https://kapi.kakao.com/v2/api/talk/memo/default/send", headers=headers, data=data, timeout=10)
        return resp.status_code == 200
    except Exception:
        return False


def get_staff_names():
    return {
        info.get("name")
        for uid, info in st.session_state.get("user_db", {}).items()
        if uid != "1"
        and info.get("name")
        and info.get("dept_type", "사업부") == "C&S"
    }


def get_current_user_rank():
    name = st.session_state.get("user_name", "")
    for uid, info in st.session_state.get("user_db", {}).items():
        if info.get("name") == name:
            return info.get("rank", "직원")
    return "직원"


def filter_by_staff(df, name_col="담당자"):
    if df.empty or name_col not in df.columns:
        return df
    staff_names = get_staff_names()
    if not staff_names:
        return df
    return df[df[name_col].isin(staff_names)].reset_index(drop=True)


def hide_department_heads(df):
    if df is None or df.empty or "직급" not in df.columns:
        return df
    return df[df["직급"].astype(str).str.strip() != "부서장"].reset_index(drop=True)


def dataframe_to_upload_payload(df):
    safe_df = strip_activity_time_columns(df).copy()
    safe_df = safe_df.replace({np.nan: ""})
    return {
        "columns": [str(c) for c in safe_df.columns],
        "rows": safe_df.astype(str).values.tolist(),
        "saved_at": (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S"),
    }


def upload_payload_to_dataframe(payload):
    if not payload:
        return pd.DataFrame()
    return pd.DataFrame(payload.get("rows", []), columns=payload.get("columns", []))


_USER_DB_DEFAULT = {
    "1": {
        "pw": "1",
        "name": "최고관리자",
        "email": "",
        "access": "허용",
        "role": "관리자",
        "dept_type": "사업부",
        "staff_type": "정규직",
        "outsource": "아니오",
        "outsource_period": "해당없음",
    },
    "T": {
        "pw": "1111",
        "name": "이성환",
        "email": "",
        "access": "허용",
        "role": "관리자",
        "dept_type": "사업부",
        "staff_type": "정규직",
        "outsource": "아니오",
        "outsource_period": "해당없음",
    },
    "iminjee": {
        "pw": "1111",
        "name": "임인지",
        "email": "",
        "access": "허용",
        "role": "사용자",
        "dept_type": "사업부",
        "staff_type": "정규직",
        "outsource": "아니오",
        "outsource_period": "해당없음",
    },
    "kangkt": {
        "pw": "1111",
        "name": "강경태",
        "email": "",
        "access": "허용",
        "role": "사용자",
        "dept_type": "사업부",
        "staff_type": "정규직",
        "outsource": "아니오",
        "outsource_period": "해당없음",
    },
    "leesh": {
        "pw": "1111",
        "name": "이수현",
        "email": "",
        "access": "허용",
        "role": "사용자",
        "dept_type": "사업부",
        "staff_type": "정규직",
        "outsource": "아니오",
        "outsource_period": "해당없음",
    },
    "gilmj": {
        "pw": "1111",
        "name": "길민종",
        "email": "",
        "access": "허용",
        "role": "사용자",
        "dept_type": "사업부",
        "staff_type": "정규직",
        "outsource": "아니오",
        "outsource_period": "해당없음",
    },
    "maengks": {
        "pw": "1111",
        "name": "맹국성",
        "email": "",
        "access": "허용",
        "role": "사용자",
        "dept_type": "사업부",
        "staff_type": "정규직",
        "outsource": "아니오",
        "outsource_period": "해당없음",
    },
}


def _real_users(d):
    """_deleted 같은 메타키를 제외한 실제 사용자 수 반환."""
    return {k: v for k, v in d.items() if not k.startswith("_")}


def load_user_db():
    """users.json → users.json.bak → GitHub → {} 순서로 로드 후 기본 계정과 병합하여 반환.
    _deleted 목록에 있는 ID는 기본 계정이라도 복원하지 않는다."""
    candidates = []

    # 1) 메인 파일
    if os.path.exists(DB_FILE):
        try:
            with open(DB_FILE, "r", encoding="utf-8") as f:
                d = json.load(f)
            if isinstance(d, dict):
                candidates.append(d)
        except Exception:
            pass

    # 2) 백업 파일
    bak = DB_FILE + ".bak"
    if os.path.exists(bak):
        try:
            with open(bak, "r", encoding="utf-8") as f:
                d = json.load(f)
            if isinstance(d, dict):
                candidates.append(d)
        except Exception:
            pass

    # 3) GitHub
    gh = _github_load(DB_FILE, None)
    if isinstance(gh, dict):
        candidates.append(gh)

    loaded = {}
    for candidate in candidates:
        if len(_real_users(candidate)) > len(_real_users(loaded)):
            loaded = candidate

    # 삭제된 ID 목록 추출 (파일에만 존재하는 메타 키)
    deleted_ids = set(loaded.get("_deleted", []))

    if len(_real_users(loaded)) > 1:
        try:
            with open(DB_FILE, "w", encoding="utf-8") as f:
                json.dump(loaded, f, ensure_ascii=False, indent=4)
            with open(DB_FILE + ".bak", "w", encoding="utf-8") as f:
                json.dump(loaded, f, ensure_ascii=False, indent=4)
        except Exception:
            pass

    # 기본 계정 베이스 + 파일 데이터 병합, _deleted ID는 제외
    merged = {**_USER_DB_DEFAULT, **_real_users(loaded)}
    for uid in deleted_ids:
        merged.pop(uid, None)
    return merged


def merge_user_db(existing, incoming):
    deleted_ids = set(incoming.get("_deleted", [])) | set(existing.get("_deleted", []))
    merged = {**_USER_DB_DEFAULT}
    if isinstance(existing, dict):
        merged.update(_real_users(existing))
    if isinstance(incoming, dict):
        for uid, info in _real_users(incoming).items():
            if uid not in merged or len(_real_users(incoming)) >= len(_real_users(existing or {})):
                merged[uid] = info
    for uid in deleted_ids:
        merged.pop(uid, None)
    return merged


def init_state():
    if "user_db" not in st.session_state:
        st.session_state.user_db = load_user_db()

    defaults = {
        "logged_in": False,
        "user_role": "사용자",
        "user_name": "",
        "auth_mode": "login",
        "current_menu": "업로드 및 실적 확인",
        "url_analysis": DEFAULT_URL_ANALYSIS,
        "url_sync": DEFAULT_URL_SYNC,
        "url_hana": DEFAULT_URL_HANA,
        "url_hana_billing": DEFAULT_URL_HANA_BILLING,
        "url_hana_performance": DEFAULT_URL_HANA_PERFORMANCE,
        "dart_api_key": "",
        "hana_sheet_df": None,
        "hana_billing_df": None,
        "analysis_lookup_df": None,
        "cloud_sheet_df": None,
        "analysis_result": None,
        "user_excel_data": None,
        "temp_cloud_df": None,
        "auto_prev_df": None,
        "deadline_time": "",
        "report_closed": "",
        "dark_mode": False,
        "login_time": "",
        "_prev_menu": None,
    }

    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

    if "1yS4gaES-iuzt1NSRTSdj9Ivg1fjbN5mIyX4pGnvEYN0" not in str(st.session_state.get("url_sync", "")):
        st.session_state.url_sync = DEFAULT_URL_SYNC
        st.session_state.cloud_sheet_df = None

    if st.session_state.current_menu == "구글 스트레드시트 연동":
        st.session_state.current_menu = "서버 접속 정보"

    transfer_result_upload_menus = {
        "이체 결과 엑셀 업로드",
        "이체결과 엑셀 업로드",
        "이체 결과 업로드",
        "지급 결과 엑셀 업로드",
    }
    if st.session_state.current_menu in transfer_result_upload_menus:
        st.session_state.current_menu = "지급 결과 확인"

    removed_menus = {
        "이력확인 및 작성",
        "은행 이력 업로드",
        "최종 실적 확인",
    }
    if st.session_state.current_menu in removed_menus:
        st.session_state.current_menu = "전도금 요청"

    RENAMED_MENUS = {
        "사업장 정보 등록": "위탁 사업장 관리",
        "사업장 예측/보고": "대시보드",
        "보고서": "대시보드",
        "계좌 관리": "위탁 사업장 관리",
        "담당자 관리": "위탁 사업장 관리",
        "사업장 관리": "위탁 사업장 관리",
        "이체 자료 생성": "이체 자료 확정",
        "사용품의서 보고": "전도금 사용 결의 보고",
    }
    if st.session_state.current_menu in RENAMED_MENUS:
        st.session_state.current_menu = RENAMED_MENUS[st.session_state.current_menu]


init_state()


def clean_header_logic(df):
    try:
        df = df.copy()
        df.columns = [str(c).strip() for c in df.columns]

        if len(df.columns) > 0 and str(df.columns[0]).startswith("Unnamed"):
            for i in range(min(len(df), 10)):
                row_text = " ".join(df.iloc[i].astype(str).tolist())
                if any(k in row_text for k in ["등록자", "담당자", "성명", "업체명", "사업자번호", "고객번호"]):
                    df.columns = [str(c).strip() for c in df.iloc[i]]
                    df = df.iloc[i + 1:].reset_index(drop=True)
                    break

        keep = ~pd.Series(df.columns).astype(str).str.contains("^Unnamed|^nan", case=False, na=False).values
        df = df.loc[:, keep]

        # 중요한 컬럼은 값이 비어있어도 유지
        important_keywords = ["업무번호", "플로우", "식권", "비즈플레이"]
        important_cols = [col for col in df.columns if any(keyword in str(col) for keyword in important_keywords)]

        # 빈 컬럼 제거 (단, 중요 컬럼은 제외)
        empty_cols = df.columns[df.isna().all()]
        cols_to_drop = [col for col in empty_cols if col not in important_cols]
        df = df.drop(columns=cols_to_drop, errors="ignore")

        # 빈 행 제거
        df = df.dropna(how="all", axis=0)
        return strip_activity_time_columns(df)
    except Exception:
        return df


def strip_activity_time_columns(df):
    try:
        df = df.copy()
        for col in df.columns:
            col_name = str(col).replace(" ", "").replace("　", "")
            if "활동일" not in col_name:
                continue

            converted = pd.to_datetime(df[col], errors="coerce")
            valid = converted.notna()
            if valid.any():
                df.loc[valid, col] = converted.loc[valid].dt.strftime("%Y-%m-%d")
        return df
    except Exception:
        return df


def find_col(df, keys, fallback=None):
    for c in df.columns:
        c_normalized = str(c).replace(" ", "").replace("　", "").lower()
        if any(k.replace(" ", "").replace("　", "").lower() in c_normalized for k in keys):
            return c
    return fallback


def is_blank_value(value):
    text = "" if pd.isna(value) else str(value).strip()
    return text == "" or text.lower() in ["nan", "nat", "none", "-", "null"]


def parse_sheet_date(value):
    if is_blank_value(value):
        return pd.NaT
    text = str(value).strip()
    digits = re.sub(r"\D", "", text)
    if len(digits) >= 8:
        parsed = pd.to_datetime(digits[:8], format="%Y%m%d", errors="coerce")
        if pd.notna(parsed):
            return parsed
    return pd.to_datetime(text, errors="coerce")


def filter_visit_rows(df):
    if df is None:
        return pd.DataFrame()
    if df.empty:
        return df.copy()
    visit_col = find_col(df, ["활동구분", "접수유형"])
    if not visit_col or visit_col not in df.columns:
        return df.copy()
    mask = df[visit_col].astype(str).str.strip().str.contains("방문", na=False)
    return df[mask].copy()


def render_history_search_filters(source_df, key_prefix):
    st.markdown("##### 검색 조건")
    c1, c2, c3, c4, c5 = st.columns(5)
    with c1:
        company = st.text_input("업체명", key=f"{key_prefix}_company")

    staff_options = ["전체"]
    date_options = ["전체"]
    category_options = ["전체"]
    detail_options = ["전체"]
    if source_df is not None and not source_df.empty:
        staff_col = find_col(source_df, ["등록자", "담당자", "성명"])
        date_col = find_col(source_df, ["활동일자", "활동일", "초과일자", "일자"])
        category_col = find_col(source_df, ["활동구분", "접수유형"])
        detail_col = find_col(source_df, ["활동상세", "활동내용"])
        if staff_col and staff_col in source_df.columns:
            staff_options += sorted(v for v in source_df[staff_col].astype(str).str.strip().unique() if v)
        if date_col and date_col in source_df.columns:
            _dates = pd.to_datetime(source_df[date_col], errors="coerce").dropna().dt.strftime("%Y-%m-%d").unique()
            date_options += sorted(_dates, reverse=True)
        if category_col and category_col in source_df.columns:
            category_options += sorted(v for v in source_df[category_col].astype(str).str.strip().unique() if v)
        if detail_col and detail_col in source_df.columns:
            detail_options += sorted(v for v in source_df[detail_col].astype(str).str.strip().unique() if v)

    with c2:
        staff = st.selectbox("담당자", staff_options, key=f"{key_prefix}_staff")
    with c3:
        activity_date = st.selectbox("활동일자", date_options, key=f"{key_prefix}_date")
    with c4:
        activity_category = st.selectbox("활동구분", category_options, key=f"{key_prefix}_category")
    with c5:
        activity_detail = st.selectbox("활동상세", detail_options, key=f"{key_prefix}_detail")

    return {
        "company": company.strip(),
        "staff": staff,
        "date": activity_date,
        "category": activity_category,
        "detail": activity_detail,
    }


def apply_history_search_filters(df, filters):
    if df is None or df.empty:
        return df
    if not filters:
        return df
    result = df.copy()
    company_col = find_col(result, ["업체명", "상호", "고객명"])
    staff_col = find_col(result, ["등록자", "담당자", "성명"])
    date_col = find_col(result, ["활동일자", "활동일", "초과일자", "일자"])
    category_col = find_col(result, ["활동구분", "접수유형"])
    detail_col = find_col(result, ["활동상세", "활동내용"])

    if filters.get("company") and company_col in result.columns:
        result = result[result[company_col].astype(str).str.contains(filters["company"], case=False, na=False)]
    if filters.get("staff") and filters["staff"] != "전체" and staff_col in result.columns:
        target_staff = re.sub(r"\s+", "", str(filters["staff"]).strip())
        result = result[result[staff_col].apply(lambda value: re.sub(r"\s+", "", str(value).strip())) == target_staff]
    if filters.get("date") and filters["date"] != "전체" and date_col in result.columns:
        _date_values = pd.to_datetime(result[date_col], errors="coerce").dt.strftime("%Y-%m-%d")
        result = result[_date_values == filters["date"]]
    if filters.get("category") and filters["category"] != "전체" and category_col in result.columns:
        result = result[result[category_col].astype(str).str.strip() == filters["category"]]
    if filters.get("detail") and filters["detail"] != "전체" and detail_col in result.columns:
        result = result[result[detail_col].astype(str).str.strip() == filters["detail"]]
    return result


def has_active_history_filters(filters):
    if not filters:
        return False
    return (
        bool(filters.get("company"))
        or filters.get("staff") != "전체"
        or filters.get("date") != "전체"
        or filters.get("category") != "전체"
        or filters.get("detail") != "전체"
    )


def history_filter_signature(filters):
    if not filters:
        return "all"
    raw = "|".join(str(filters.get(key, "")) for key in ["company", "staff", "date", "category", "detail"])
    return hashlib.md5(raw.encode("utf-8")).hexdigest()[:10]


def normalize_biz(series):
    def normalize_one(value):
        if pd.isna(value):
            return ""

        # 숫자 타입이면 정수로 변환 (과학적 표기법 방지)
        if isinstance(value, (int, float, np.integer, np.floating)):
            try:
                value = int(float(value))
            except (ValueError, OverflowError):
                pass

        text = str(value).strip().replace(",", "")
        if re.fullmatch(r"\d+\.0+", text):
            text = text.split(".", 1)[0]

        digits = re.sub(r"[^0-9]", "", text)
        if "." in str(value) and digits.endswith("0") and len(digits) == 11:
            digits = digits[:-1]
        return digits

    if isinstance(series, pd.Series):
        return series.apply(normalize_one)
    return normalize_one(series)


def get_uploaded_month(df):
    try:
        d_col = find_col(df, ["활동일", "일자"])
        if d_col and d_col in df.columns:
            dates = pd.to_datetime(df[d_col], errors="coerce").dropna()
            if not dates.empty:
                return dates.dt.strftime("%Y-%m").value_counts().idxmax()
    except Exception:
        pass
    return ""


def attach_cloud_dates(user_df):
    df = user_df.copy()
    cloud = st.session_state.get("cloud_sheet_df")

    if cloud is None or df.empty:
        if cloud is None and not df.empty:
            st.warning("⚠️ 본사 구글시트 데이터를 불러올 수 없습니다. 관리자 > 본사 URL 설정에서 URL을 확인해주세요.")
        return df

    cloud = clean_header_logic(cloud.copy())

    biz_col_user = find_col(df, ["사업자번호"])
    biz_col_cloud = find_col(cloud, ["사업자번호"])
    open_col = find_col(cloud, ["개설완료일자", "개설일"])
    erp_col = find_col(cloud, ["ERP연계일자", "연계일자"])

    if not biz_col_user or not biz_col_cloud:
        return df

    div_col = find_col(cloud, ["신규/이행구분", "이행구분", "신규이행"])
    add_col = find_col(cloud, ["이행추가연계"])

    cloud_cols = [biz_col_cloud]
    rename_map = {}

    if open_col:
        cloud_cols.append(open_col)
        rename_map[open_col] = "본사 개설완료일자"
    if erp_col:
        cloud_cols.append(erp_col)
        rename_map[erp_col] = "본사 ERP연계일자"
    if div_col:
        cloud_cols.append(div_col)
        rename_map[div_col] = "본사 신규이행구분"
    if add_col:
        cloud_cols.append(add_col)
        rename_map[add_col] = "본사 이행추가연계"

    if len(cloud_cols) == 1:
        return df

    df["_biz_key"] = normalize_biz(df[biz_col_user])
    cloud["_biz_key"] = normalize_biz(cloud[biz_col_cloud])

    cloud_map = cloud[["_biz_key"] + [c for c in cloud_cols if c != biz_col_cloud]].rename(columns=rename_map)
    cloud_map = cloud_map.drop_duplicates("_biz_key")

    df = pd.merge(df, cloud_map, on="_biz_key", how="left", suffixes=("", "_cloud"))
    for col in ["본사 개설완료일자", "본사 ERP연계일자", "본사 신규이행구분", "본사 이행추가연계"]:
        cloud_col = f"{col}_cloud"
        if cloud_col in df.columns:
            if col in df.columns:
                df[col] = df[col].where(df[col].notna() & (df[col].astype(str).str.strip() != ""), df[cloud_col])
            else:
                df[col] = df[cloud_col]
            df = df.drop(columns=[cloud_col], errors="ignore")
    return df.drop(columns=["_biz_key"], errors="ignore")


def build_other_validation_errors(df):
    other_errors = []
    if df is None or df.empty:
        return pd.DataFrame()

    biz_col = find_col(df, ["사업자번호"], "사업자번호")
    comp_col = find_col(df, ["업체명", "상호"], "업체명")
    user_col = find_col(df, ["등록자", "담당자", "성명"], "등록자")
    open_col = "본사 개설완료일자"
    erp_col = "본사 ERP연계일자"

    if open_col not in df.columns or erp_col not in df.columns:
        return pd.DataFrame()

    def parse_compare_date(value):
        if pd.isna(value):
            return pd.NaT
        text = str(value).strip()
        digits = re.sub(r"[^0-9]", "", text)
        if len(digits) >= 8:
            parsed = pd.to_datetime(digits[:8], format="%Y%m%d", errors="coerce")
            if pd.notna(parsed):
                return parsed
        return pd.to_datetime(text, errors="coerce")

    df_check = df.copy()
    df_check["_open_date"] = df_check[open_col].apply(parse_compare_date)
    df_check["_erp_date"] = df_check[erp_col].apply(parse_compare_date)

    for _, row in df_check.iterrows():
        open_date = row["_open_date"]
        erp_date = row["_erp_date"]
        if pd.notna(open_date) and pd.notna(erp_date) and erp_date.normalize() < open_date.normalize():
            other_errors.append({
                "업체명": row.get(comp_col, "") if comp_col else "",
                "사업자번호": row.get(biz_col, "") if biz_col else "",
                "등록자": row.get(user_col, "") if user_col else "",
                open_col: open_date.strftime("%Y-%m-%d"),
                erp_col: erp_date.strftime("%Y-%m-%d"),
                "오류 사유": "본사 ERP연계일자가 개설완료일자보다 이전",
            })

    return pd.DataFrame(other_errors)


def normalize_customer_no(value):
    if pd.isna(value):
        return ""
    text = str(value).strip().replace(",", "")
    if re.fullmatch(r"\d+\.0+", text):
        text = text.split(".", 1)[0]
    digits = re.sub(r"[^0-9]", "", text)
    if digits and len(digits) < 8:
        digits = digits.zfill(8)
    return digits


def read_excel_history_file(uploaded_file):
    df = pd.read_excel(uploaded_file, sheet_name=0)
    df = clean_header_logic(df)
    return df.replace({np.nan: ""})


def parse_history_date(value):
    if pd.isna(value) or str(value).strip() == "":
        return ""
    if isinstance(value, (datetime, pd.Timestamp)):
        return pd.to_datetime(value).strftime("%Y-%m-%d")
    if isinstance(value, (int, float, np.integer, np.floating)):
        num = float(value)
        if 20000 <= num <= 60000:
            return (datetime(1899, 12, 30) + timedelta(days=int(num))).strftime("%Y-%m-%d")
    text = str(value).strip()
    digits = re.sub(r"[^0-9]", "", text)
    if len(digits) >= 8:
        parsed = pd.to_datetime(digits[:8], format="%Y%m%d", errors="coerce")
        if pd.notna(parsed):
            return parsed.strftime("%Y-%m-%d")
    parsed = pd.to_datetime(text, errors="coerce")
    return parsed.strftime("%Y-%m-%d") if pd.notna(parsed) else text


def infer_activity_detail(row, source_cols):
    text = " ".join(str(row.get(col, "")) for col in source_cols if col)
    if "개설" in text:
        return "개설"
    if "연계" in text or "ERP" in text.upper():
        return "연계"
    return "운영"


def title_from_activity_detail(value):
    detail = str(value).strip()
    if "연계" in detail:
        return "연계 방문"
    if "개설" in detail:
        return "개설 방문"
    return "운영방문"


def normalize_converted_history_df(df):
    if df is None or df.empty:
        return pd.DataFrame()
    result = df.copy()
    for col in list(result.columns):
        if "업무번호" in str(col):
            result = result.rename(columns={col: "업무번호"})
            break
    result["지사"] = "HANA지사"
    if "_is_manual" not in result.columns:
        result["_is_manual"] = False
    if "활동구분" in result.columns:
        result["활동구분"] = result["활동구분"].astype(str).map(
            lambda v: "방문" if "방문" in v else ("원격" if "원격" in v else ("상담" if v == "상담" else "상담"))
        )
    if "활동상세" in result.columns:
        result["활동상세"] = result["활동상세"].astype(str).map(
            lambda v: "연계" if "연계" in v else ("개설" if "개설" in v else "운영")
        )
        default_titles = result["활동상세"].map(title_from_activity_detail)
        if "제목" in result.columns:
            result["제목"] = result["제목"].where(
                result["제목"].notna() & result["제목"].astype(str).str.strip().ne(""),
                default_titles,
            )
        else:
            result["제목"] = default_titles
    if "활동일" in result.columns and "활동일자" not in result.columns:
        result = result.rename(columns={"활동일": "활동일자"})
    if "활동내용" in result.columns and "활동내역" not in result.columns:
        result = result.rename(columns={"활동내용": "활동내역"})
    return result


def hana_customer_biz_map():
    hana = st.session_state.get("hana_sheet_df")
    if hana is None or hana.empty:
        hana = read_google_csv(st.session_state.get("url_hana", DEFAULT_URL_HANA), header=2)
        st.session_state.hana_sheet_df = hana

    hana = clean_header_logic(hana.copy())
    customer_col = find_col(hana, ["고객번호", "고개번호", "고객NO", "고객 No", "고객"])
    biz_col = find_col(hana, ["사업자번호"])
    if not customer_col or customer_col not in hana.columns or not biz_col or biz_col not in hana.columns:
        return {}, "하나은행 구글 시트에서 고객번호 또는 사업자번호 컬럼을 찾을 수 없습니다."

    mapping_df = hana[[customer_col, biz_col]].copy()
    mapping_df["_customer_key"] = mapping_df[customer_col].apply(normalize_customer_no)
    mapping_df["_biz_value"] = normalize_biz(mapping_df[biz_col])
    mapping_df = mapping_df[(mapping_df["_customer_key"] != "") & (mapping_df["_biz_value"] != "")]
    mapping_df = mapping_df.drop_duplicates("_customer_key")
    return dict(zip(mapping_df["_customer_key"], mapping_df["_biz_value"])), ""


def convert_history_to_sample_df(history_df, user_name):
    history_df = clean_header_logic(history_df.copy()).replace({np.nan: ""})
    customer_col = find_col(history_df, ["고객번호", "고개번호", "고객NO", "고객 No", "고객"])
    staff_col = find_col(history_df, ["담당자"]) or find_col(history_df, ["처리자", "접수자", "등록자", "성명"])
    date_col = find_col(history_df, ["활동일", "처리일자", "접수일자", "일자"])
    company_col = find_col(history_df, ["업체명", "고객명", "상호", "회사명"])
    product_col = find_col(history_df, ["상품", "서비스", "제품"])
    category_col = find_col(history_df, ["활동구분", "상담구분", "접수구분", "접수유형"])
    detail_col = find_col(history_df, ["활동상세", "처리유형", "진행상태", "접수유형"])
    location_col = find_col(history_df, ["방문장소", "주소", "지역"])
    work_no_col = find_col(history_df, ["업무번호", "플로우", "작성번호"])
    title_col = find_col(history_df, ["제목", "접수내용", "문의내용", "진행상태"])
    content_col = find_col(history_df, ["활동내용", "처리내용", "상담내용", "내용"])

    missing = []
    if not customer_col or customer_col not in history_df.columns:
        missing.append("고객번호")
    if not staff_col or staff_col not in history_df.columns:
        missing.append("담당자")
    if missing:
        return pd.DataFrame(), {"error": f"이력 파일에서 {', '.join(missing)} 컬럼을 찾을 수 없습니다."}

    filtered = history_df[history_df[staff_col].astype(str).str.strip() == str(user_name).strip()].copy()
    if filtered.empty:
        return pd.DataFrame(), {"error": f"로그인 사용자({user_name})와 담당자가 일치하는 이력이 없습니다."}

    biz_map, map_error = hana_customer_biz_map()
    if map_error:
        return pd.DataFrame(), {"error": map_error}

    rows = []
    unmatched = 0
    for _, row in filtered.iterrows():
        customer_key = normalize_customer_no(row.get(customer_col, ""))
        biz_no = biz_map.get(customer_key, "")
        if not biz_no:
            unmatched += 1

        activity_detail = infer_activity_detail(row, [detail_col, title_col, content_col, category_col])
        activity_category = str(row.get(category_col, "")).strip() if category_col else ""
        activity_text = " ".join(str(row.get(col, "")) for col in [category_col, title_col, content_col, detail_col] if col)
        if "방문" in activity_text:
            activity_category = "방문"
        elif "원격" in activity_category:
            activity_category = "원격"
        elif activity_category not in ["방문", "상담", "원격"]:
            activity_category = "상담"

        rows.append({
            "지사": "HANA지사",
            "상품": row.get(product_col, "통합CMS") if product_col else "통합CMS",
            "업체명": row.get(company_col, "") if company_col else "",
            "사업자번호": biz_no,
            "등록자": user_name,
            "활동일자": parse_history_date(row.get(date_col, "")) if date_col else "",
            "방문장소 (시, 군, 구까지)": row.get(location_col, "") if location_col else "",
            "활동구분": activity_category,
            "활동상세": activity_detail,
            "업무번호": row.get(work_no_col, "") if work_no_col else "",
            "제목": title_from_activity_detail(activity_detail),
            "활동내역": row.get(content_col, "") if content_col else "",
        })

    return normalize_converted_history_df(pd.DataFrame(rows)), {"total": len(rows), "unmatched": unmatched}


def sample_format_excel_bytes(df):
    from openpyxl import load_workbook

    output = BytesIO()
    if os.path.exists(EXCEL_SAMPLE_FILE):
        wb = load_workbook(EXCEL_SAMPLE_FILE)
        ws = wb.worksheets[0]
        if ws.max_row > 2:
            ws.delete_rows(3, ws.max_row - 2)
        def header_key(value):
            return re.sub(r"\s+", "", str(value or ""))

        header_map = {header_key(ws.cell(row=2, column=col).value): col for col in range(1, ws.max_column + 1)}
        aliases = {
            "활동내역": "활동내용",
            "활동일자": "활동일",
            "업무번호": "업무번호\n(플로우에 식권, 비즈플레이, 플로우 작성번호)",
        }
        for row_idx, (_, row) in enumerate(df.iterrows(), start=3):
            for header, value in row.items():
                col_idx = header_map.get(header_key(header)) or header_map.get(header_key(aliases.get(header, "")))
                if not col_idx:
                    continue
                ws.cell(row=row_idx, column=col_idx, value=value)
        wb.save(output)
    else:
        output.write(dataframe_to_excel_bytes({"Sheet0": df}))
    output.seek(0)
    return output.getvalue()


def _activity_template_text(value):
    if pd.isna(value):
        return ""
    return str(value).strip()


def _activity_template_customer_key(value):
    text = _activity_template_text(value)
    if re.fullmatch(r"\d+\.0+", text):
        text = text.split(".", 1)[0]
    digits = re.sub(r"\D", "", text)
    return digits.lstrip("0") if digits else ""


def _activity_template_time(value):
    if pd.isna(value):
        return ""
    if isinstance(value, (datetime, pd.Timestamp)):
        return pd.to_datetime(value).strftime("%H%M")

    if isinstance(value, (int, float, np.integer, np.floating)) and not pd.isna(value):
        number = float(value)
        if 0 <= number < 1:
            total_minutes = int(round(number * 24 * 60)) % (24 * 60)
            return f"{total_minutes // 60:02d}{total_minutes % 60:02d}"
        if number.is_integer():
            text = str(int(number))
        else:
            text = str(value)
    else:
        text = _activity_template_text(value)

    if not text:
        return ""

    time_match = re.search(r"(\d{1,2})\s*:\s*(\d{1,2})", text)
    if time_match:
        return f"{int(time_match.group(1)):02d}{int(time_match.group(2)):02d}"[:4]

    if re.fullmatch(r"\d+\.0+", text):
        text = text.split(".", 1)[0]

    try:
        number = float(text)
        if 0 <= number < 1:
            total_minutes = int(round(number * 24 * 60)) % (24 * 60)
            return f"{total_minutes // 60:02d}{total_minutes % 60:02d}"
    except ValueError:
        pass

    digits = re.sub(r"\D", "", text)
    if not digits:
        return ""
    if len(digits) in (3, 4):
        return digits.zfill(4)[:4]
    if len(digits) >= 12:
        return digits[8:12]
    if len(digits) == 5:
        return digits.zfill(6)[:4]
    if len(digits) > 4:
        return digits[:4]
    if len(digits) < 3:
        return digits.zfill(4)
    return ""


def _activity_template_datetime(date_value, time_value):
    date_digits = re.sub(r"\D", "", _activity_template_text(date_value))
    if len(date_digits) >= 8:
        date_part = date_digits[:8]
    else:
        parsed = pd.to_datetime(date_value, errors="coerce")
        date_part = parsed.strftime("%Y%m%d") if pd.notna(parsed) else date_digits
    time_part = _activity_template_time(time_value)
    if not time_part:
        parsed = pd.to_datetime(date_value, errors="coerce")
        if pd.notna(parsed) and (parsed.hour or parsed.minute):
            time_part = parsed.strftime("%H%M")
    if time_part:
        return f"{date_part}{time_part}"
    return date_part


def _activity_template_type(value):
    text = _activity_template_text(value)
    if "원격" in text:
        return "원격활동"
    if "유선" in text or "전화" in text or "통화" in text:
        return "통화"
    if "메일" in text or "이메일" in text:
        return "이메일"
    if "방문" in text or "미팅" in text:
        return "미팅"
    return "메모"


def _activity_template_title(row, request_col, result_col, work_type_col):
    request_text = _activity_template_text(row.get(request_col, "")) if request_col else ""
    result_text = _activity_template_text(row.get(result_col, "")) if result_col else ""
    work_type = _activity_template_text(row.get(work_type_col, "")) if work_type_col else ""
    title = request_text or result_text or work_type or "활동이력"
    first_line = next((line.strip() for line in title.splitlines() if line.strip()), title)
    return first_line[:250]


def _activity_template_purpose(value):
    text = _activity_template_text(value)
    if "교육" in text:
        return "사용자교육"
    if "ERP" in text.upper() or "연계" in text:
        return "ERP연계진행"
    return "운영활동"


def _activity_template_issue_type(value):
    text = _activity_template_text(value)
    if any(keyword in text for keyword in ["오류", "장애", "에러", "실패"]):
        return "장애발생"
    return "단순요청"


def _activity_customer_map_from_df(sheet):
    sheet = clean_header_logic(sheet.copy()).replace({np.nan: ""})
    customer_col = find_col(sheet, ["고객번호", "고객NO", "고객 No", "고객"])
    biz_col = find_col(sheet, ["사업자번호", "사업자등록번호"])
    company_col = find_col(sheet, ["업체명", "회사명", "고객명", "고객사명", "상호"])
    if not customer_col or customer_col not in sheet.columns or not biz_col or biz_col not in sheet.columns:
        return {}, "고객번호 또는 사업자번호 컬럼을 찾을 수 없습니다."

    mapping = {}
    for _, row in sheet.iterrows():
        customer_key = _activity_template_customer_key(row.get(customer_col, ""))
        biz_no = re.sub(r"\D", "", _activity_template_text(row.get(biz_col, "")))
        company = _activity_template_text(row.get(company_col, "")) if company_col else ""
        if customer_key and biz_no:
            mapping[customer_key] = {"biz_no": biz_no, "company": company}
    return mapping, ""


@st.cache_data(ttl=600, show_spinner=False)
def read_activity_google_customer_map():
    try:
        hana_sheet = read_google_csv(DEFAULT_URL_HANA, header=2).fillna("")
        return _activity_customer_map_from_df(hana_sheet)
    except Exception as exc:
        return {}, f"고객원장 CSV 조회 실패: {exc}"


def convert_history_to_activity_template_df(history_df, template_bytes):
    raw_history_df = history_df.copy()
    for col in list(raw_history_df.columns):
        if not str(col).startswith("Unnamed"):
            continue
        sample = raw_history_df[col].dropna().astype(str).str.strip()
        numeric_ratio = sample.str.match(r"^\d+(\.0+)?$").mean() if not sample.empty else 0
        if numeric_ratio >= 0.7:
            raw_history_df = raw_history_df.rename(columns={col: "고객번호"})
            break
    history_df = clean_header_logic(raw_history_df).replace({np.nan: ""})
    template_df = pd.read_excel(BytesIO(template_bytes), sheet_name="Activities", nrows=0, dtype=str)
    template_columns = list(template_df.columns)
    google_customer_map, google_map_error = read_activity_google_customer_map()

    date_col = find_col(history_df, ["접수일자", "일자", "날짜"])
    time_col = find_col(history_df, ["접수시간", "시간"])
    product_col = find_col(history_df, ["CMS구분", "제품", "상품"])
    customer_col = find_col(history_df, ["고객번호", "고객NO", "고객 No"])
    if not customer_col and len(history_df.columns) > 4:
        customer_col = history_df.columns[4]
    company_col = find_col(history_df, ["업체명", "회사명", "고객명", "상호"])
    type_col = find_col(history_df, ["접수유형", "활동유형", "유형"])
    work_type_col = find_col(history_df, ["업무유형", "처리업무"])
    receiver_col = find_col(history_df, ["접수자"])
    owner_col = find_col(history_df, ["담당자", "활동자"])
    request_col = find_col(history_df, ["요청사항", "문의내용", "접수내용"])
    result_col = find_col(history_df, ["처리내용", "활동내역", "조치내용"])

    missing = []
    for label, col in [("접수일자", date_col), ("업체명", company_col), ("담당자", owner_col), ("요청사항 또는 처리내용", request_col or result_col)]:
        if not col or col not in history_df.columns:
            missing.append(label)
    if missing:
        return pd.DataFrame(columns=template_columns), {"error": f"원본 파일에서 {', '.join(missing)} 컬럼을 찾을 수 없습니다."}

    converted_rows = []
    unmatched = 0
    for _, row in history_df.iterrows():
        if str(row.get("NO", "")).strip().lower() in ["", "nan"]:
            continue
        customer_key = _activity_template_customer_key(row.get(customer_col, "")) if customer_col else ""
        mapped = google_customer_map.get(customer_key, {}).copy()
        if customer_key and not mapped.get("biz_no"):
            unmatched += 1

        company = _activity_template_text(row.get(company_col, ""))
        if not company and mapped.get("company"):
            company = mapped["company"]
        work_type = _activity_template_text(row.get(work_type_col, "")) if work_type_col else ""
        request_text = _activity_template_text(row.get(request_col, "")) if request_col else ""
        result_text = _activity_template_text(row.get(result_col, "")) if result_col else ""
        activity_text = result_text or request_text
        owner = _activity_template_text(row.get(owner_col, "")) if owner_col else ""
        if not owner and receiver_col:
            owner = _activity_template_text(row.get(receiver_col, ""))

        converted_rows.append({
            "유형(필수: 통화/미팅/원격활동/이메일/메모)": _activity_template_type(row.get(type_col, "")) if type_col else "메모",
            "제목(필수)": _activity_template_title(row, request_col, result_col, work_type_col),
            "일시(필수: 202606110930 또는 2026-06-11 09:30)": _activity_template_datetime(row.get(date_col, ""), row.get(time_col, "")) if date_col else "",
            "활동자(필수: 이름)": owner,
            "소요시간(분)": "",
            "회사명": company,
            "고객사명": company,
            "사업자번호": mapped.get("biz_no", ""),
            "고객": "",
            "관련제품(코드 또는 이름)": _activity_template_text(row.get(product_col, "")) if product_col else "통합CMS",
            "활동내역": activity_text,
            "활동 상세(마케팅/개설/운영/연계/기타)": "운영",
            "활동 구분(IN/OUT)": "IN",
            "활동 목적(운영활동/신규구축/사용자교육/ERP연계진행/ERP연계완료/ERP추가연계/상품전환/재구축/사전협의/수수료연체안내)": _activity_template_purpose(work_type),
            "팀구분(MANAGE/HOTLINE/TECH)": "HOTLINE",
            "장애 구분(단순요청/장애발생)": _activity_template_issue_type(f"{work_type} {request_text} {result_text}"),
            "HOTLINE 처리업무(통합자금관리/이체/B2B/부가서비스/기초정보/수수료관리/외화무역/ERP/클라이언트설치/인증서/결재함/기타)": "",
            "HOTLINE 장애유형(사용자 단순 문의/오류/이체 관련 문의/오류/전자확인증 관련 문의/오류/자금집금 오류/문의/시재 관련 문의/오류/보고서 관련 문의/오류/클라이언트 설치/기초정보관련 문의/오류/인증서 관련 문의/오류/보안모듈 관련 문의/오류/ERP 관련 문의/오류/법인카드 관련 문의/오류/Ibase 관련 문의/오류/수수료 관련 문의/오류/접속 관련 문의/오류/서버 관련 문의/오류/서버 재설치/요건 관련 문의/기타 문의/오류)": "",
        })

    converted_df = pd.DataFrame(converted_rows)
    for col in template_columns:
        if col not in converted_df.columns:
            converted_df[col] = ""
    return converted_df[template_columns], {
        "total": len(converted_df),
        "unmatched": unmatched,
        "google_map_error": google_map_error,
        "google_map_count": len(google_customer_map),
    }


def activity_template_excel_bytes(converted_df, template_bytes):
    from openpyxl import load_workbook

    output = BytesIO()
    wb = load_workbook(BytesIO(template_bytes))
    ws = wb["Activities"] if "Activities" in wb.sheetnames else wb.worksheets[0]
    if ws.max_row > 1:
        ws.delete_rows(2, ws.max_row - 1)
    header_map = {str(ws.cell(row=1, column=col).value or "").strip(): col for col in range(1, ws.max_column + 1)}
    for row_idx, (_, row) in enumerate(converted_df.iterrows(), start=2):
        for header, col_idx in header_map.items():
            if header in converted_df.columns:
                ws.cell(row=row_idx, column=col_idx, value=row.get(header, ""))
    wb.save(output)
    output.seek(0)
    return output.getvalue()


def show_activity_template_converter():
    st.markdown("#### 활동이력 템플릿 변환")
    st.caption("은행/핫라인 활동이력 xls 파일을 activities_template (9).xlsx의 Activities 시트 형식으로 변환합니다.")

    col_history, col_template = st.columns(2)
    with col_history:
        history_file = st.file_uploader("202607활동이력(07-08).xls 파일 업로드", type=["xls", "xlsx"], key="activity_template_history_upload")
    with col_template:
        template_file = st.file_uploader("activities_template (9).xlsx 파일 업로드", type=["xlsx"], key="activity_template_template_upload")

    if history_file is None or template_file is None:
        st.info("원본 활동이력 파일과 변환 기준 템플릿 파일을 모두 업로드해주세요.")
        return

    try:
        template_bytes = template_file.getvalue()
        with st.spinner("활동이력을 템플릿 형식으로 변환하는 중입니다."):
            history_df = pd.read_excel(history_file, sheet_name=0)
            converted_df, info = convert_history_to_activity_template_df(history_df, template_bytes)
        if converted_df.empty:
            st.warning(info.get("error", "변환할 데이터가 없습니다."))
            return

        st.success(f"변환 완료: {int(info.get('total', len(converted_df))):,}건")
        unmatched = int(info.get("unmatched", 0))
        google_map_error = info.get("google_map_error", "")
        google_map_count = int(info.get("google_map_count", 0) or 0)
        if google_map_error:
            st.warning(google_map_error)
        elif google_map_count:
            st.caption(f"고객번호 기준 사업자번호 매핑 {google_map_count:,}건을 적용했습니다.")
        if unmatched:
            st.caption(f"사업자번호 매핑 실패 {unmatched:,}건은 사업자번호가 빈칸으로 저장됩니다.")
        st.dataframe(converted_df.head(30), use_container_width=True, hide_index=True)

        output_bytes = activity_template_excel_bytes(converted_df, template_bytes)
        ym = ""
        date_col = find_col(converted_df, ["일시"])
        if date_col:
            date_digits = converted_df[date_col].astype(str).str.extract(r"(\d{6})", expand=False).dropna()
            ym = date_digits.iloc[0] if not date_digits.empty else ""
        file_month = ym or (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m")
        st.download_button(
            "변환 파일 다운로드",
            data=output_bytes,
            file_name=f"activities_template_converted_{file_month}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
        )
    except ImportError:
        st.error("xls 파일 변환을 위해 xlrd 패키지가 필요합니다.")
    except Exception as exc:
        st.error(f"변환 중 오류가 발생했습니다: {exc}")


def criteria_df():
    data = [
        ["교차판매", "타겟고객선별", 5, 100],
        ["교차판매", "메일발송", 2, 100],
        ["교차판매", "제안서전달", 5, 100],
        ["교차판매", "방문설명회&견적발송", 30, "-"],
        ["교차판매", "계약진행 시", 50, "-"],
        ["교차판매", "유선 (해피콜)", 5, 100],
        ["교차판매", "활성화 (조회업무)", 10, 100],
        ["교차판매", "이체, 집금 활성화", 30, 150],
        ["교차판매", "계열사 추가도입", 30, "-"],
        ["교차판매", "신규연계도입", 60, "-"],
        ["추가활동", "문서 작성 (본사)", 100, 100],
        ["추가활동", "문서 작성 (가이드)", 50, "-"],
        ["추가활동", "문서 작성 (기타)", 20, "-"],
        ["추가활동", "VOC (아이디어)", 10, 50],
        ["추가활동", "운영활동(원격)", 10, 200],
    ]
    return pd.DataFrame(data, columns=["활동구분", "구분", "단위 점수", "월 최대점수"])


def render_manual_perf_input_table(base):
    """표는 보고서 표와 동일하게 렌더링하고, 실적 입력은 표 아래에서 처리한다."""
    result_df = base.copy()
    result_df["입력(건)"] = pd.to_numeric(result_df["입력(건)"], errors="coerce").fillna(0).astype(int)

    style_report_logic(result_df.drop(columns=["입력(건)"], errors="ignore"))

    st.markdown("#### 실적 입력")
    st.info("월 최대점수가 있는 항목은 `월 최대점수 ÷ 단위 점수`까지만 입력할 수 있습니다. 예: 타겟고객선별은 100 ÷ 5 = 최대 20건입니다.")
    st.markdown(
        """<style>
        .manual-input-row {
            color:#4A5568;
            font-size:13px;
            font-weight:700;
            padding-top:9px;
            white-space:nowrap;
        }
        .manual-input-help {
            color:#718096;
            font-size:12px;
            margin-top:-4px;
            margin-bottom:4px;
        }
        body:has(#pms-d:checked) .manual-input-row,
        body:has(#pms-d:checked) .manual-input-help {
            color:#ffffff !important;
        }
        </style>""",
        unsafe_allow_html=True,
    )

    results = {}
    for i in range(0, len(result_df), 3):
        cols = st.columns(3)
        for col, (_, row) in zip(cols, result_df.iloc[i:i + 3].iterrows()):
            item = str(row["구분"])
            unit_score = int(row["단위 점수"])
            monthly_limit = row["월 최대점수"]
            max_count = None
            if str(monthly_limit).strip() != "-":
                try:
                    max_count = int(float(monthly_limit) // unit_score)
                except Exception:
                    max_count = None

            value = int(row["입력(건)"])
            if max_count is not None:
                value = min(value, max_count)

            with col:
                st.markdown(f"<div class='manual-input-row'>{html.escape(item)}</div>", unsafe_allow_html=True)
                if max_count is not None:
                    st.markdown(f"<div class='manual-input-help'>최대 {max_count}건 입력 가능</div>", unsafe_allow_html=True)
                else:
                    st.markdown("<div class='manual-input-help'>월 최대점수 제한 없음</div>", unsafe_allow_html=True)
                widget_key = f"perf_{item}"
                kwargs = {
                    "label": "입력(건)",
                    "min_value": 0,
                    "value": value,
                    "key": widget_key,
                    "label_visibility": "collapsed",
                    "step": 1,
                }
                if max_count is not None:
                    kwargs["max_value"] = max_count
                st.number_input(**kwargs)
                results[item] = int(st.session_state.get(widget_key, value) or 0)

    result_df["입력(건)"] = result_df["구분"].map(results).fillna(0).astype(int)
    return result_df


def manual_points_for_user(name):
    override_db = st.session_state.get("manual_perf_preview_override", {})
    saved = override_db.get(name)
    if saved is None:
        perf_db = load_db(PERF_FILE, {})
        saved = perf_db.get(name, {})

    total = 0
    criteria = criteria_df()
    criteria["_key"] = criteria["구분"].astype(str).str.strip()
    point_map = criteria.set_index("_key")["단위 점수"].to_dict()
    limit_map = criteria.set_index("_key")["월 최대점수"].to_dict()
    for item, count in saved.items():
        try:
            key = str(item).strip()
            score = int(point_map.get(key, 0)) * int(count)
            limit_value = limit_map.get(key, "-")
            if str(limit_value).strip() != "-":
                score = min(score, int(float(limit_value)))
            total += score
        except Exception:
            pass
    return total


def calculate_manual_perf_total(edited_df):
    if edited_df is None or edited_df.empty:
        return 0
    total = 0
    for _, row in edited_df.iterrows():
        try:
            score = int(float(row.get("단위 점수", 0))) * int(float(row.get("입력(건)", 0)))
            limit_value = row.get("월 최대점수", "-")
            if str(limit_value).strip() != "-":
                score = min(score, int(float(limit_value)))
            total += score
        except Exception:
            pass
    return int(total)


def apply_rs_allowance_formula(perf_df, user_db, return_debug=False):
    if perf_df is None or perf_df.empty or "합계포인트" not in perf_df.columns:
        return (perf_df, {}) if return_debug else perf_df

    def num(value):
        try:
            if pd.isna(value):
                return 0
            return int(float(str(value).replace(",", "").strip() or 0))
        except Exception:
            return 0

    def payable_points(row):
        open_link_points = num(row.get("개설포인트", 0)) + num(row.get("연계포인트", 0))
        operation_points = num(row.get("운영포인트 (실제 활동)", 0))
        operation_points += num(row.get("운영포인트 (추가 활동)", 0))
        operation_points += num(row.get("운영포인트(추가 활동)", 0))
        return min(open_link_points, 1000) + min(operation_points, 1800)

    result = perf_df.copy()
    result["지급포인트"] = 0
    result["지급예상금액"] = 0

    work_df = result[result["담당자"].astype(str) != "합계"].copy() if "담당자" in result.columns else result.copy()
    if work_df.empty:
        return (result, {}) if return_debug else result

    # 직원 정보 매핑 (이름 -> 정보)
    name_to_info = {}
    for uid, info in user_db.items():
        if isinstance(info, dict) and info.get("name"):
            name_to_info[info["name"]] = {
                "rank": info.get("rank", "직원"),
                "outsource": info.get("outsource", "아니오"),
                "staff_type": info.get("staff_type", "정규직"),
                "period": info.get("outsource_period", "해당없음"),
                "dept_type": info.get("dept_type", "사업부")
            }

    # C&S 부서만 대상
    cs_staff = []
    outsource_staff = []
    regular_staff = []

    for idx, row in work_df.iterrows():
        staff_name = str(row.get("담당자", ""))
        if staff_name in name_to_info:
            staff_info = name_to_info[staff_name]
            if staff_info["dept_type"] == "C&S":
                cs_staff.append((idx, row, staff_info))
                if staff_info["outsource"] == "예" or staff_info["staff_type"] == "외주":
                    outsource_staff.append((idx, row, staff_info))
                else:
                    regular_staff.append((idx, row, staff_info))

    if not cs_staff:
        return (result, {}) if return_debug else result

    # BU 평균 계산 (C&S 전체 지급 산정 포인트 / C&S 인원)
    total_bu_points = sum(payable_points(row) for _, row, _ in cs_staff)
    bu_count = len(cs_staff)
    bu_average = total_bu_points / bu_count if bu_count > 0 else 0

    # 외주직원 가감포인트 계산
    outsource_total_points = 0
    for idx, row, staff_info in outsource_staff:
        staff_points = payable_points(row)

        # 근무기간 계수
        period = staff_info["period"]
        if period == "1년 미만":
            period_factor = 0.8
        elif period == "1년 이상":
            period_factor = 0.9
        else:  # 2년 이상
            period_factor = 1.0

        # 외주직원 가감포인트 = 외주직원포인트 - (BU평균 × 근무기간)
        adjustment = staff_points - (bu_average * period_factor)
        outsource_total_points += adjustment

    # 일반직원 수
    regular_count = len(regular_staff)
    outsource_point_per_regular = outsource_total_points / regular_count if regular_count > 0 else 0

    # 팀장수당
    team_leader_bonus = 267

    # 디버깅 정보
    debug_info = {
        "BU합산": total_bu_points,
        "BU인원": bu_count,
        "BU평균": bu_average,
        "외주직원수": len(outsource_staff),
        "일반직원수": regular_count,
        "외주가감총합": outsource_total_points,
        "일반직원1인당분배": outsource_point_per_regular,
        "팀장수당": team_leader_bonus
    }

    # 최종 지급액 계산
    for idx, row, staff_info in cs_staff:
        total_points = payable_points(row)

        # 팀장 여부 확인
        is_team_leader = staff_info["rank"] == "팀장"

        # 외주직원 여부 확인
        is_outsource = staff_info["outsource"] == "예" or staff_info["staff_type"] == "외주"

        if is_outsource:
            # 외주직원: 외주실적 가감포인트 분배 대상에서는 제외
            final_pay_point = total_points - 1000
        else:
            # 팀장/직원 공식: (개설+연계+운영(+팀장수당)-1000)+(외주실적가감포인트/일반직원수)
            final_pay_point = total_points - 1000 + outsource_point_per_regular
            if is_team_leader:
                final_pay_point += team_leader_bonus

        final_pay_point = int(max(0, final_pay_point))
        result.at[idx, "지급포인트"] = final_pay_point
        result.at[idx, "지급예상금액"] = int(final_pay_point * 500)

    if return_debug:
        return result, debug_info
    return result


def may_2026_business_dates():
    holidays = {"2026-05-01", "2026-05-05", "2026-05-24", "2026-05-25"}
    days = pd.date_range("2026-05-01", "2026-05-31", freq="D")
    return [d.strftime("%Y-%m-%d") for d in days if d.weekday() < 5 and d.strftime("%Y-%m-%d") not in holidays]


def korean_public_holidays(year):
    holidays_by_year = {
        2026: {
            "2026-01-01",
            "2026-02-16", "2026-02-17", "2026-02-18",
            "2026-03-01", "2026-03-02",
            "2026-05-01", "2026-05-05", "2026-05-24", "2026-05-25",
            "2026-06-03", "2026-06-06",
            "2026-08-15", "2026-08-17",
            "2026-09-24", "2026-09-25", "2026-09-26", "2026-09-27", "2026-09-28",
            "2026-10-03", "2026-10-05", "2026-10-09",
            "2026-12-25",
        }
    }
    fixed = {
        f"{year}-01-01",
        f"{year}-03-01",
        f"{year}-05-01",
        f"{year}-05-05",
        f"{year}-06-06",
        f"{year}-08-15",
        f"{year}-10-03",
        f"{year}-10-09",
        f"{year}-12-25",
    }
    return fixed | holidays_by_year.get(int(year), set())


def is_korean_business_day(value):
    date_value = pd.to_datetime(value, errors="coerce")
    if pd.isna(date_value):
        return False
    date_text = date_value.strftime("%Y-%m-%d")
    return date_value.weekday() < 5 and date_text not in korean_public_holidays(date_value.year)


def month_business_dates(target_ym):
    month_start = pd.to_datetime(f"{target_ym}-01", errors="coerce")
    if pd.isna(month_start):
        return []
    days = pd.date_range(month_start, month_start + pd.offsets.MonthEnd(0), freq="D")
    return [d.strftime("%Y-%m-%d") for d in days if is_korean_business_day(d)]


def prepare_random_history_source_df(df):
    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        return pd.DataFrame()
    result = df.copy()
    result.columns = [str(col).strip() for col in result.columns]
    result = result.loc[:, ~pd.Index(result.columns).duplicated()]
    for col in result.columns:
        if not isinstance(result[col], pd.Series):
            result[col] = pd.Series(result[col], index=result.index)
    return clean_header_logic(result)


def prepare_display_dataframe(df):
    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        return pd.DataFrame()
    result = df.copy()
    result.columns = [str(col).strip() for col in result.columns]
    result = result.loc[:, ~pd.Index(result.columns).duplicated()]
    result = result[
        [
            col for col in result.columns
            if not str(col).strip().lower().startswith("unnamed")
            and not str(col).strip().startswith("_")
        ]
    ]
    if result.empty:
        return result
    return result.loc[
        :,
        result.apply(lambda col: col.astype(str).str.strip().replace("nan", "").ne("").any(), axis=0)
    ]


def limit_history_to_total_point_cap(history_df, perf_df):
    if history_df is None or not isinstance(history_df, pd.DataFrame) or history_df.empty:
        return history_df
    if perf_df is None or not isinstance(perf_df, pd.DataFrame) or perf_df.empty:
        return history_df

    df = clean_header_logic(history_df.copy())
    user_col = find_col(df, ["등록자", "담당자", "성명"])
    detail_col = find_col(df, ["활동상세", "활동내용"])
    category_col = find_col(df, ["활동구분", "접수유형"])
    if not user_col or user_col not in df.columns:
        return history_df

    base_points_by_user = {}
    for _, row in perf_df.iterrows():
        name = str(row.get("담당자", "")).strip()
        if not name or name == "합계":
            continue
        open_points = int(float(row.get("개설포인트", 0) or 0))
        link_points = int(float(row.get("연계포인트", 0) or 0))
        base_points_by_user[name] = max(0, 2800 - open_points - link_points)

    used_points_by_user = {name: 0 for name in base_points_by_user}
    keep_indices = []
    for idx, row in df.iterrows():
        name = str(row.get(user_col, "")).strip()
        if name not in base_points_by_user:
            keep_indices.append(idx)
            continue

        detail = str(row.get(detail_col, "")).strip() if detail_col else ""
        category = str(row.get(category_col, "")).strip() if category_col else ""
        is_operation = "운영" in detail or "방문" in category or "원격" in category
        if not is_operation:
            keep_indices.append(idx)
            continue

        point = 10 if "원격" in category else 30
        if used_points_by_user[name] + point > base_points_by_user[name]:
            continue
        used_points_by_user[name] += point
        keep_indices.append(idx)

    return df.loc[keep_indices].reset_index(drop=True)


def build_random_admin_extra_history(source_df, existing_df, target_count):
    if source_df is None or not isinstance(source_df, pd.DataFrame) or source_df.empty:
        return pd.DataFrame(), "하나지사 활동이력 데이터가 없습니다."

    source = prepare_random_history_source_df(source_df)
    existing = prepare_random_history_source_df(existing_df) if isinstance(existing_df, pd.DataFrame) else pd.DataFrame()

    detail_col = find_col(source, ["활동상세", "활동내용"])
    if not detail_col or detail_col not in source.columns:
        return pd.DataFrame(), "하나지사 활동이력에서 활동상세 컬럼을 찾을 수 없습니다."

    source = source[source[detail_col].astype(str).str.contains("운영", na=False)].copy()
    if source.empty:
        return pd.DataFrame(), "활동상세가 운영인 데이터가 없습니다."

    company_col = find_col(source, ["업체명", "상호", "고객명"])
    biz_col = find_col(source, ["사업자번호"])
    user_col = find_col(source, ["등록자", "담당자", "성명"])
    place_col = find_col(source, ["방문장소", "주소", "지역"])
    category_col = find_col(source, ["활동구분", "접수유형"])
    work_col = find_col(source, ["업무번호"])
    note_col = find_col(source, ["활동내역", "활동내용", "비고", "제목"])

    if not user_col or user_col not in source.columns:
        return pd.DataFrame(), "하나지사 활동이력에서 담당자 컬럼을 찾을 수 없습니다."

    business_dates = may_2026_business_dates()
    if not business_dates:
        return pd.DataFrame(), "2026년 5월 영업일을 만들 수 없습니다."

    existing_visit_counts = {}
    existing_keys = set()
    if not existing.empty:
        e_user_col = find_col(existing, ["등록자", "담당자", "성명"])
        e_biz_col = find_col(existing, ["사업자번호"])
        e_date_col = find_col(existing, ["활동일자", "활동일", "일자"])
        e_detail_col = find_col(existing, ["활동상세", "활동내용"])
        e_category_col = find_col(existing, ["활동구분", "접수유형"])
        if e_user_col and e_date_col and e_user_col in existing.columns and e_date_col in existing.columns:
            for _, row in existing.iterrows():
                user = str(row.get(e_user_col, "")).strip()
                date = str(row.get(e_date_col, "")).strip()
                category = str(row.get(e_category_col, "")) if e_category_col else ""
                if user and date and "방문" in category:
                    existing_visit_counts[(user, date)] = existing_visit_counts.get((user, date), 0) + 1
                biz = normalize_biz(row.get(e_biz_col, "")) if e_biz_col else ""
                detail = str(row.get(e_detail_col, "")).strip() if e_detail_col else ""
                if biz and user and date and detail:
                    existing_keys.add((biz, user, date, detail))

    rng_state = int((datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m%d%H%M%S")) % (2**32 - 1)
    source = source.sample(frac=1, random_state=rng_state).reset_index(drop=True)
    target_count = max(1, int(target_count or 1))
    rows = []

    for _, row in source.iterrows():
        if len(rows) >= target_count:
            break

        user = str(row.get(user_col, "")).strip()
        if not user:
            continue

        raw_category = str(row.get(category_col, "")).strip() if category_col else ""
        activity_category = "원격" if "원격" in raw_category else "방문"
        biz = normalize_biz(row.get(biz_col, "")) if biz_col else ""

        assigned_date = ""
        for date in np.random.default_rng().permutation(business_dates):
            date = str(date)
            if activity_category == "방문" and existing_visit_counts.get((user, date), 0) >= 5:
                continue
            if biz and (biz, user, date, "운영") in existing_keys:
                continue
            assigned_date = date
            break

        if not assigned_date:
            continue

        if activity_category == "방문":
            existing_visit_counts[(user, assigned_date)] = existing_visit_counts.get((user, assigned_date), 0) + 1
        if biz:
            existing_keys.add((biz, user, assigned_date, "운영"))

        rows.append({
            "업체명": row.get(company_col, "") if company_col else "",
            "사업자번호": row.get(biz_col, "") if biz_col else "",
            "등록자": user,
            "활동일": assigned_date,
            "방문장소 (시, 군, 구까지)": row.get(place_col, "") if place_col else "",
            "활동구분": activity_category,
            "활동상세": "운영",
            "업무번호": row.get(work_col, "") if work_col else "",
            "제목": "운영방문" if activity_category == "방문" else "운영원격",
            "활동내역": row.get(note_col, "") if note_col else "",
            "본사 개설완료일자": assigned_date,
            "본사 ERP연계일자": assigned_date,
            "_is_manual": False,
        })

    if not rows:
        return pd.DataFrame(), "중복/초과방문 조건을 피해서 추가할 수 있는 데이터가 없습니다."

    return pd.DataFrame(rows), ""


def save_manual_perf_override_for_current_user():
    name = st.session_state.get("user_name")
    override = st.session_state.get("manual_perf_preview_override", {}).get(name)
    if name and override is not None:
        db = load_db(PERF_FILE, {})
        db[name] = override
        save_db(PERF_FILE, db)


def has_performance_required_columns(df):
    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        return False
    return all(
        find_col(df, keys) in df.columns
        for keys in [["등록자", "담당자", "성명"], ["활동상세", "활동내용"], ["활동일자", "활동일", "일자"]]
    )


def prepare_history_analysis_df(raw_df):
    if raw_df is None or not isinstance(raw_df, pd.DataFrame) or raw_df.empty:
        return pd.DataFrame()

    df = normalize_converted_history_df(clean_header_logic(raw_df.copy()))
    if has_performance_required_columns(df):
        return df

    try:
        converted_df, _ = convert_history_to_sample_df(df, st.session_state.get("user_name", ""))
        converted_df = normalize_converted_history_df(converted_df)
        if has_performance_required_columns(converted_df):
            return converted_df
    except Exception:
        pass

    return df


def current_history_analysis_df():
    excel_df = st.session_state.get("user_excel_data")
    if st.session_state.get("user_excel_source") == "hq" and isinstance(excel_df, pd.DataFrame):
        excel_df = prepare_history_analysis_df(excel_df)
        if has_performance_required_columns(excel_df):
            return excel_df
    preview_df = st.session_state.get("history_convert_preview_data")
    if isinstance(preview_df, pd.DataFrame):
        preview_df = prepare_history_analysis_df(preview_df)
        if has_performance_required_columns(preview_df):
            return preview_df
    if isinstance(excel_df, pd.DataFrame):
        excel_df = prepare_history_analysis_df(excel_df)
        if has_performance_required_columns(excel_df):
            return excel_df
    return pd.DataFrame()


def apply_rank_from_user_db(df):
    if "담당자" not in df.columns or "직급" not in df.columns:
        return df
    rank_map = {
        info.get("name", ""): info.get("rank", "직원")
        for uid, info in st.session_state.get("user_db", {}).items()
        if uid != "1" and info.get("name")
    }
    if not rank_map:
        return df
    df = df.copy()
    df["직급"] = df["담당자"].astype(str).str.strip().map(lambda n: rank_map.get(n, "직원"))
    return df


_RANK_ORDER = {"부서장": 0, "팀장": 1, "과장": 2, "대리": 3, "주임": 4, "직원": 5}


def sort_by_rank_name(df):
    if df.empty:
        return df
    name_col = next((c for c in ["담당자", "성명"] if c in df.columns), None)
    if "직급" not in df.columns or name_col is None:
        return df
    df = df.copy()
    df["_rank_order"] = df["직급"].map(_RANK_ORDER).fillna(9)
    df = df.sort_values(["_rank_order", name_col]).drop(columns=["_rank_order"]).reset_index(drop=True)
    return df


def process_performance_analysis(curr_df_raw, prev_df_raw=None):
    try:
        if curr_df_raw is None:
            return pd.DataFrame(), pd.DataFrame(), pd.DataFrame()
        df = normalize_converted_history_df(clean_header_logic(curr_df_raw))
        if df.empty and len(df.columns) == 0:
            return pd.DataFrame(), pd.DataFrame(), pd.DataFrame()

        u_col = find_col(df, ["등록자", "담당자", "성명"])
        d_col = find_col(df, ["활동상세", "활동내용"])
        date_col = find_col(df, ["활동일자", "활동일", "일자"])
        biz_col = find_col(df, ["사업자번호"], "사업자번호")
        comp_col = find_col(df, ["업체명", "상호"], "업체명")

        if not all([u_col, d_col, date_col]):
            prepared_df = prepare_history_analysis_df(curr_df_raw)
            if not prepared_df.empty:
                df = prepared_df
                u_col = find_col(df, ["등록자", "담당자", "성명"])
                d_col = find_col(df, ["활동상세", "활동내용"])
                date_col = find_col(df, ["활동일자", "활동일", "일자"])
                biz_col = find_col(df, ["사업자번호"], "사업자번호")
                comp_col = find_col(df, ["업체명", "상호"], "업체명")

        required_cols = [("등록자", u_col), ("활동상세", d_col), ("활동일", date_col)]
        missing = [label for label, col in required_cols if not col or col not in df.columns]
        if missing:
            return f"필수 컬럼이 없습니다: {', '.join(missing)}", None, None

        _udb = st.session_state.get("user_db", {})
        member_db = {
            info.get("name", ""): info.get("rank", "직원")
            for uid, info in _udb.items()
            if uid != "1" and info.get("name")
        }

        df_clean = df.dropna(subset=[u_col, date_col]).copy()
        df_clean[u_col] = df_clean[u_col].astype(str).str.strip()
        _staff = get_staff_names()
        if _staff and u_col in df_clean.columns:
            staff_keys = {re.sub(r"\s+", "", str(name).strip()) for name in _staff}
            df_clean = df_clean[df_clean[u_col].apply(lambda value: re.sub(r"\s+", "", str(value).strip())).isin(staff_keys)]
        df_clean[date_col] = pd.to_datetime(df_clean[date_col], errors="coerce")
        df_clean = df_clean.dropna(subset=[date_col])
        df_clean[date_col] = df_clean[date_col].dt.strftime("%Y-%m-%d")

        daily_counts = df_clean.groupby([u_col, date_col]).size().reset_index(name="일방문횟수")
        monthly_counts = df_clean.groupby(u_col).size().reset_index(name="월총방문")

        if comp_col in df_clean.columns and biz_col in df_clean.columns:
            error_raw = pd.merge(
                daily_counts,
                df_clean[[u_col, date_col, comp_col, biz_col]].drop_duplicates(),
                on=[u_col, date_col],
                how="left",
            )
            error_df = pd.merge(error_raw, monthly_counts, on=u_col, how="left")
            error_df = error_df[[comp_col, biz_col, u_col, date_col, "일방문횟수", "월총방문"]].drop_duplicates()
            error_df.columns = ["업체명", "사업자번호", "담당자", "초과일자", "일방문", "월총방문"]
        else:
            error_df = pd.DataFrame(columns=["업체명", "사업자번호", "담당자", "초과일자", "일방문", "월총방문"])

        # 은행 이력(실제 활동)과 이력 추가(추가 활동) 분리
        _flag = "_is_manual"
        if _flag in df.columns:
            df_real = df[~df[_flag].fillna(False).astype(bool)].copy()
            df_added = df[df[_flag].fillna(False).astype(bool)].copy()
        else:
            df_real = df
            df_added = pd.DataFrame()

        # 중복 이력: 업로드 파일 기준 동일일자에 동일사업자번호 방문 이력이 2건 이상인 행
        if biz_col in df_clean.columns and date_col in df_clean.columns:
            _dup_df = df_clean.copy()
            _dup_df["_dup_biz"] = normalize_biz(_dup_df[biz_col])
            _dup_df["_dup_date"] = pd.to_datetime(_dup_df[date_col], errors="coerce").dt.strftime("%Y-%m-%d")
            _dup_keys = ["_dup_biz", "_dup_date"]
            _dup_df = _dup_df[(_dup_df[_dup_keys] != "").all(axis=1)]
            dup_biz_df = _dup_df[_dup_df.duplicated(subset=_dup_keys, keep=False)].drop(columns=_dup_keys, errors="ignore")
            sort_cols = [col for col in [date_col, biz_col, u_col, d_col] if col in dup_biz_df.columns]
            dup_biz_df = dup_biz_df.sort_values(by=sort_cols) if sort_cols else dup_biz_df
        else:
            dup_biz_df = pd.DataFrame()

        # 이력 추가 행의 포인트를 담당자별로 집계
        added_pts_by_name = {}
        if not df_added.empty and u_col in df_added.columns and d_col in df_added.columns:
            for _n, _g in df_added.dropna(subset=[u_col, d_col]).groupby(u_col):
                _det = _g[d_col].astype(str)
                _pts = _det.str.contains("운영|방문|점검").sum() * 30
                _pts += _det.str.contains("개설").sum() * 90
                _pts += _det.str.contains("연계").sum() * 120
                added_pts_by_name[_n] = int(_pts)

        summary = (
            df_real.dropna(subset=[u_col, d_col])
            .groupby(u_col)
            .agg(
                o=(d_col, lambda x: x.astype(str).str.contains("개설").sum()),
                l=(d_col, lambda x: x.astype(str).str.contains("연계").sum()),
                v=(d_col, lambda x: x.astype(str).str.contains("운영|방문|점검").sum()),
            )
            .reset_index()
        )

        if summary.empty:
            return pd.DataFrame(), error_df, dup_biz_df

        name_to_info = {
            info.get("name"): {
                "staff_type": info.get("staff_type", "정규직"),
                "period": info.get("outsource_period", "해당없음"),
            }
            for uid, info in st.session_state.user_db.items()
            if uid != "1"
        }

        member_stats = {}
        total_points = 0

        for _, row in summary.iterrows():
            name = row[u_col]
            o_p = int(row["o"]) * 90
            l_p = int(row["l"]) * 120
            v_actual_p = int(row["v"]) * 30
            manual_p = manual_points_for_user(name) + added_pts_by_name.get(name, 0)
            p_sum = min(2800, o_p + l_p + v_actual_p + manual_p)

            member_stats[name] = {
                "o_p": o_p,
                "l_p": l_p,
                "v_actual_p": v_actual_p,
                "manual_p": manual_p,
                "p_sum": p_sum,
            }
            total_points += p_sum

        bu_avg = total_points / len(summary) if len(summary) > 0 else 0
        outsource_adj_pool = 0
        regular_count = 0

        for name, stats in member_stats.items():
            info = name_to_info.get(name, {})
            if info.get("staff_type", "정규직") == "외주":
                tenure = {"1년 미만": 0.8, "1년 이상": 0.9, "2년 이상": 1.0}.get(info.get("period"), 1.0)
                outsource_adj_pool += stats["p_sum"] - (bu_avg * tenure)
            else:
                regular_count += 1

        adj_per_regular = outsource_adj_pool / regular_count if regular_count else 0

        rows = []
        for _, row in summary.iterrows():
            name = row[u_col]
            stats = member_stats[name]
            rank = member_db.get(name, "직원")
            is_outsource = name_to_info.get(name, {}).get("staff_type", "정규직") == "외주"

            rows.append(
                {
                    "담당자": name,
                    "직급": rank,
                    "개설건수": int(row["o"]),
                    "개설포인트": stats["o_p"],
                    "연계건수": int(row["l"]),
                    "연계포인트": stats["l_p"],
                    "운영건수 (실제 활동)": int(row["v"]),
                    "운영포인트 (실제 활동)": stats["v_actual_p"],
                    "운영건수 (추가 활동)": round(manual_points_for_user(name) / 30),
                    "운영포인트(추가 활동)": manual_points_for_user(name),
                    "합계포인트": stats["p_sum"],
                    "지급포인트": 0,
                    "지급예상금액": 0,
                    "전월대비": 0,
                }
            )

        res_df = pd.DataFrame(rows)
        res_df = apply_rank_from_user_db(res_df)
        res_df = apply_rs_allowance_formula(res_df, st.session_state.user_db)
        res_df = hide_department_heads(res_df)
        res_df = sort_by_rank_name(res_df)

        if prev_df_raw is not None:
            try:
                prev_res_df, _, _ = process_performance_analysis(prev_df_raw)
                if isinstance(prev_res_df, pd.DataFrame) and not prev_res_df.empty:
                    p_map = prev_res_df.set_index("담당자")["지급예상금액"].to_dict()
                    res_df["전월대비"] = res_df.apply(
                        lambda r: int(r["지급예상금액"] - p_map.get(r["담당자"], 0)),
                        axis=1,
                    )
            except Exception:
                pass

        return res_df, error_df, dup_biz_df

    except Exception as e:
        return f"ERR: {str(e)}", None, None


def build_upload_over_visit_df(curr_df_raw):
    if curr_df_raw is None:
        return pd.DataFrame(columns=["업체명", "사업자번호", "담당자", "초과일자", "일방문", "월총방문"])
    df = normalize_converted_history_df(clean_header_logic(curr_df_raw))
    if df.empty:
        return pd.DataFrame(columns=["업체명", "사업자번호", "담당자", "초과일자", "일방문", "월총방문"])

    u_col = find_col(df, ["등록자", "담당자", "성명"])
    date_col = find_col(df, ["활동일자", "활동일", "일자"])
    biz_col = find_col(df, ["사업자번호"], "사업자번호")
    comp_col = find_col(df, ["업체명", "상호", "고객명"], "업체명")
    if not u_col or not date_col or u_col not in df.columns or date_col not in df.columns:
        return pd.DataFrame(columns=["업체명", "사업자번호", "담당자", "초과일자", "일방문", "월총방문"])

    work = df.dropna(subset=[u_col, date_col]).copy()
    work[u_col] = work[u_col].astype(str).str.strip()
    work[date_col] = pd.to_datetime(work[date_col], errors="coerce")
    work = work.dropna(subset=[date_col])
    work[date_col] = work[date_col].dt.strftime("%Y-%m-%d")

    daily_counts = work.groupby([u_col, date_col]).size().reset_index(name="일방문")
    monthly_counts = work.groupby(u_col).size().reset_index(name="월총방문")
    over_days = daily_counts[daily_counts["일방문"] >= 6]
    if over_days.empty:
        return pd.DataFrame(columns=["업체명", "사업자번호", "담당자", "초과일자", "일방문", "월총방문"])

    detail_cols = [col for col in [u_col, date_col, comp_col, biz_col] if col and col in work.columns]
    result = over_days.merge(work[detail_cols].drop_duplicates(), on=[u_col, date_col], how="left")
    result = result.merge(monthly_counts, on=u_col, how="left")
    rename_map = {u_col: "담당자", date_col: "초과일자"}
    if comp_col and comp_col in result.columns:
        rename_map[comp_col] = "업체명"
    if biz_col and biz_col in result.columns:
        rename_map[biz_col] = "사업자번호"
    result = result.rename(columns=rename_map)
    for col in ["업체명", "사업자번호"]:
        if col not in result.columns:
            result[col] = ""
    return result[["업체명", "사업자번호", "담당자", "초과일자", "일방문", "월총방문"]].drop_duplicates()


def build_daily_visit_matrix_df(curr_df_raw):
    if curr_df_raw is None:
        return pd.DataFrame()
    df = normalize_converted_history_df(clean_header_logic(curr_df_raw))
    if df.empty:
        return pd.DataFrame()

    u_col = find_col(df, ["등록자", "담당자", "성명"])
    date_col = find_col(df, ["활동일자", "활동일", "일자"])
    if not u_col or not date_col or u_col not in df.columns or date_col not in df.columns:
        return pd.DataFrame()

    work = df.dropna(subset=[u_col, date_col]).copy()
    work[u_col] = work[u_col].astype(str).str.strip()
    work[date_col] = pd.to_datetime(work[date_col], errors="coerce")
    work = work.dropna(subset=[date_col])
    if work.empty:
        return pd.DataFrame()

    target_ym = get_uploaded_month(work)
    if not target_ym:
        target_ym = work[date_col].dt.strftime("%Y-%m").value_counts().idxmax()

    month_start = pd.to_datetime(f"{target_ym}-01", errors="coerce")
    if pd.isna(month_start):
        return pd.DataFrame()
    month_end = month_start + pd.offsets.MonthEnd(0)
    all_dates = pd.date_range(month_start, month_end, freq="D")
    work = work[work[date_col].dt.strftime("%Y-%m") == target_ym].copy()
    if work.empty:
        return pd.DataFrame()

    work["_visit_date"] = work[date_col].dt.strftime("%Y-%m-%d")
    grouped = work.groupby([u_col, "_visit_date"]).size().unstack(fill_value=0)
    date_keys = [d.strftime("%Y-%m-%d") for d in all_dates]
    grouped = grouped.reindex(columns=date_keys, fill_value=0)
    grouped = grouped.sort_index()
    display_cols = [f"{d.month}/{d.day}" for d in all_dates]
    grouped.columns = display_cols
    result = grouped.reset_index().rename(columns={u_col: "담당자"})
    result.attrs["date_map"] = {label: key for label, key in zip(display_cols, date_keys)}
    return result


def render_daily_visit_matrix(matrix_df):
    if matrix_df is None or matrix_df.empty:
        st.info("담당자별 일 방문횟수 데이터가 없습니다.")
        return

    date_cols = [c for c in matrix_df.columns if c != "담당자"]

    def color_over_limit(value):
        try:
            return "background-color:#FEE2E2;color:#B91C1C;font-weight:800;" if int(value) >= 6 else ""
        except Exception:
            return ""

    styler = matrix_df.style
    if hasattr(styler, "map"):
        styler = styler.map(color_over_limit, subset=date_cols)
    else:
        styler = styler.applymap(color_over_limit, subset=date_cols)
    styler = styler.format({c: "{:.0f}" for c in date_cols})
    st.dataframe(styler, use_container_width=True, hide_index=True)


def build_visit_change_guide_df(matrix_df):
    if matrix_df is None or matrix_df.empty or "담당자" not in matrix_df.columns:
        return pd.DataFrame()

    date_cols = [c for c in matrix_df.columns if c != "담당자"]
    date_map = matrix_df.attrs.get("date_map", {})
    business_date_cols = [
        col for col in date_cols
        if is_korean_business_day(date_map.get(col, col))
    ]
    guide_rows = []
    shortage_rows = []

    for _, row in matrix_df.iterrows():
        staff = row.get("담당자", "")
        counts = {}
        for col in date_cols:
            try:
                counts[col] = int(float(row.get(col, 0) or 0))
            except Exception:
                counts[col] = 0

        spare_slots = []
        for col in business_date_cols:
            spare = max(0, 5 - counts.get(col, 0))
            spare_slots.extend([col] * spare)

        for from_date in date_cols:
            over_count = max(0, counts.get(from_date, 0) - 5)
            if over_count <= 0:
                continue

            used_targets = []
            while over_count > 0 and spare_slots:
                to_date = spare_slots.pop(0)
                if to_date == from_date:
                    continue
                used_targets.append(to_date)
                over_count -= 1

            if used_targets:
                guide_rows.append({
                    "담당자": staff,
                    "변경 필요일": from_date,
                    "현재 방문횟수": counts.get(from_date, 0),
                    "권장 변경일": ", ".join(used_targets),
                    "변경 권장건수": len(used_targets),
                })
            if over_count > 0:
                shortage_rows.append({
                    "담당자": staff,
                    "변경 필요일": from_date,
                    "현재 방문횟수": counts.get(from_date, 0),
                    "권장 변경일": "여유 일자 부족",
                    "변경 권장건수": over_count,
                })

    result = pd.DataFrame(guide_rows + shortage_rows)
    if result.empty:
        return result
    return result[["담당자", "변경 필요일", "현재 방문횟수", "권장 변경일", "변경 권장건수"]]


def render_upload_ppt_download_button(report_df, upload_df, key_suffix="upload"):
    btn_label = "실적보고서 PPT 다운로드"
    try:
        if not isinstance(report_df, pd.DataFrame) or report_df.empty:
            raise ValueError("PPT로 만들 실적 데이터가 없습니다.")
        ym = get_uploaded_month(upload_df) if isinstance(upload_df, pd.DataFrame) else ""
        if ym:
            year_month = ym.replace("-", "")
            curr_month_label = f"{int(ym.split('-')[1])}월"
        else:
            now_kst = datetime.utcnow() + timedelta(hours=9)
            year_month = now_kst.strftime("%Y%m")
            curr_month_label = f"{int(now_kst.strftime('%m'))}월"

        user_sel = st.session_state.get("user_prev_month_sel", "선택안함")
        if user_sel and user_sel != "선택안함":
            prev_month_label = f"{int(str(user_sel).split('-')[1])}월"
        else:
            prev_month_label = "전월"

        ppt_bytes = build_report_ppt_bytes(
            report_df.copy(),
            pd.DataFrame(),
            curr_month_label,
            prev_month_label,
            upload_df,
        )
        st.download_button(
            btn_label,
            data=ppt_bytes,
            file_name=f"LMB활동실적보고서_{year_month}_하나지사.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            key=f"ppt_download_{key_suffix}",
        )
    except Exception as e:
        st.button(btn_label, use_container_width=True, disabled=True, key=f"ppt_download_disabled_{key_suffix}")
        st.caption(f"PPT 생성 준비 중 오류: {e}")


def build_adjusted_history_download_df(upload_df):
    if not isinstance(upload_df, pd.DataFrame) or upload_df.empty:
        return pd.DataFrame()

    result = upload_df.copy()
    df = normalize_converted_history_df(clean_header_logic(upload_df.copy()))
    if df.empty or len(df) != len(result):
        return result

    u_col = find_col(df, ["등록자", "담당자", "성명"])
    date_col = find_col(df, ["활동일자", "활동일", "일자"])
    original_date_col = find_col(result, ["활동일자", "활동일", "일자"]) or date_col
    if not u_col or not date_col or not original_date_col or u_col not in df.columns or date_col not in df.columns:
        return result

    work = df.dropna(subset=[u_col, date_col]).copy()
    work[u_col] = work[u_col].astype(str).str.strip()
    work[date_col] = pd.to_datetime(work[date_col], errors="coerce")
    work = work.dropna(subset=[date_col])
    if work.empty:
        return result

    target_ym = get_uploaded_month(work)
    if not target_ym:
        target_ym = work[date_col].dt.strftime("%Y-%m").value_counts().idxmax()
    month_start = pd.to_datetime(f"{target_ym}-01", errors="coerce")
    if pd.isna(month_start):
        return result

    all_dates = month_business_dates(target_ym)
    if not all_dates:
        return result
    work = work[work[date_col].dt.strftime("%Y-%m") == target_ym].copy()
    work["_visit_date"] = work[date_col].dt.strftime("%Y-%m-%d")

    for staff, staff_df in work.groupby(u_col):
        counts = staff_df["_visit_date"].value_counts().to_dict()
        spare_slots = []
        for day in all_dates:
            spare_slots.extend([day] * max(0, 5 - int(counts.get(day, 0))))

        for over_day, count in sorted(counts.items()):
            move_count = max(0, int(count) - 5)
            if move_count <= 0:
                continue
            move_indices = staff_df[staff_df["_visit_date"] == over_day].index.tolist()[5:]
            for idx in move_indices[:move_count]:
                while spare_slots and spare_slots[0] == over_day:
                    spare_slots.pop(0)
                if not spare_slots:
                    break
                result.at[idx, original_date_col] = spare_slots.pop(0)

    if original_date_col in result.columns:
        result[original_date_col] = pd.to_datetime(result[original_date_col], errors="coerce").dt.strftime("%Y-%m-%d").fillna("")
    result = result.drop(columns=["_is_manual"], errors="ignore")
    return result


def render_adjusted_history_download_button(upload_df, key_suffix="upload"):
    btn_label = "최종 실적 다운로드"
    try:
        adjusted_df = build_adjusted_history_download_df(upload_df)
        if adjusted_df.empty:
            raise ValueError("다운로드할 이력 데이터가 없습니다.")
        ym = get_uploaded_month(adjusted_df) or (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m")
        year_month = ym.replace("-", "")
        excel_bytes = dataframe_to_excel_bytes({"실적파일": adjusted_df})
        st.download_button(
            btn_label,
            data=excel_bytes,
            file_name=f"LMB월간 활동실적__{year_month}_최종.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
            key=f"adjusted_history_download_{key_suffix}",
        )
    except Exception as e:
        st.button(btn_label, use_container_width=True, disabled=True, key=f"adjusted_history_download_disabled_{key_suffix}")
        st.caption(f"최종 실적 파일 생성 준비 중 오류: {e}")


def render_plain_html_table(
    df, max_rows=500, center_align=True, merge_cols=None, stretch=True, max_width=None, border=True
):
    """AG Grid 없이 순수 HTML 테이블로 렌더링 — 다크모드 완전 호환."""
    if df is None or df.empty:
        st.info("표시할 데이터가 없습니다.")
        return
    df = df.head(max_rows).reset_index(drop=True)
    merge_cols = [c for c in (merge_cols or []) if c in df.columns]

    rowspan_map = {}
    for col in merge_cols:
        rowspan_map[col] = {}
        last_value = None
        start_idx = None
        span = 0
        for i, value in enumerate(df[col].tolist()):
            text = "" if pd.isna(value) else str(value).strip()
            if text and start_idx is not None and text == last_value:
                span += 1
                rowspan_map[col][i] = 0
            elif text:
                if start_idx is not None:
                    rowspan_map[col][start_idx] = span
                last_value = text
                start_idx = i
                span = 1
            elif start_idx is not None and last_value:
                span += 1
                rowspan_map[col][i] = 0
            else:
                rowspan_map[col][i] = 1
        if start_idx is not None:
            rowspan_map[col][start_idx] = span

    th = "background:#EDF2F7;color:#4A5568;font-weight:700;font-size:12px;padding:6px 10px;white-space:nowrap;border-bottom:2px solid #E2E8F0;text-align:center;"
    headers = "".join(f"<th style='{th}'>{html.escape(str(c))}</th>" for c in df.columns)
    body = ""
    for i, row in df.iterrows():
        bg = "#FFFFFF" if i % 2 == 0 else "#F7FAFC"
        tds = ""
        for col in df.columns:
            if col in rowspan_map and rowspan_map[col].get(i, 1) == 0:
                continue
            val = "" if pd.isna(row[col]) else html.escape(str(row[col]))
            align = "center" if center_align else "left"
            td_align = f"text-align:{align};"
            rowspan = ""
            if col in rowspan_map and rowspan_map[col].get(i, 1) > 1:
                rowspan = f" rowspan='{rowspan_map[col][i]}'"
            tds += f"<td{rowspan} style='background:{bg};padding:5px 10px;border-bottom:1px solid #EDF2F7;font-size:12px;color:#2D3748;white-space:nowrap;vertical-align:middle;{td_align}'>{val}</td>"
        body += f"<tr>{tds}</tr>"
    wrapper_display = "display:block;" if stretch else "display:inline-block;max-width:100%;"
    if max_width:
        wrapper_display += f"max-width:{max_width};"
    table_width = "100%" if stretch else "auto"
    wrapper_border = "border:1px solid #E2E8F0;box-shadow:0 2px 6px rgba(0,0,0,0.05);" if border else "border:none;box-shadow:none;"
    st.markdown(
        f"""<div class="pms-report-table" style="overflow-x:auto;{wrapper_border}border-radius:8px;margin-bottom:1rem;{wrapper_display}">
        <table style="width:{table_width};border-collapse:collapse;">
            <thead><tr>{headers}</tr></thead>
            <tbody>{body}</tbody>
        </table></div>""",
        unsafe_allow_html=True,
    )


def style_report_logic(df, compact=False, align_overrides=None, default_align=None):
    if df is None or df.empty:
        st.info("표시할 데이터가 없습니다.")
        return

    df = strip_activity_time_columns(df)
    df = hide_department_heads(df)
    if df.empty:
        st.info("표시할 데이터가 없습니다.")
        return

    # 지사, 상품 컬럼 제거
    df = df.drop(columns=["지사", "상품"], errors="ignore")

    diff_cols = [c for c in df.columns if "전월대비" in str(c) or "증감" in str(c) or "대비" in str(c)]
    exclude_from_num = ["담당자", "직급", "전송시각", "등록월", "항목", "활동구분", "구분", "일치여부"] + [c for c in df.columns if "사업자번호" in str(c)]
    num_cols = [c for c in df.columns if c not in exclude_from_num + diff_cols]

    def fmt_diff(x):
        try:
            v = int(float(x))
            if v > 0:
                return f"<span style='color:#E53E3E;font-weight:800;'>▲ {v:,}</span>"
            if v < 0:
                return f"<span style='color:#3182CE;font-weight:800;'>▼ {abs(v):,}</span>"
            return f"{v:,}"
        except Exception:
            return ""

    def fmt_num(x):
        try:
            if pd.isna(x):
                return ""
            if isinstance(x, (int, float, np.integer, np.floating)):
                return f"{int(x):,}"
            return html.escape(str(x))
        except Exception:
            return html.escape(str(x))

    def fmt_match(x):
        text = "" if pd.isna(x) else str(x)
        if text == "불일치":
            return "<span style='color:#E53E3E;font-weight:900;'>불일치</span>"
        if text == "일치":
            return "<span style='color:#2F855A;font-weight:900;'>일치</span>"
        return html.escape(text)

    # 전월대비 컬럼명 치환
    user_sel = st.session_state.get("user_prev_month_sel") or st.session_state.get("adm_prev_month")
    if user_sel and user_sel != "선택안함":
        prev_m = str(int(user_sel.split("-")[1])) + "월"
        비교월_label = f"{prev_m} 대비"
    else:
        비교월_label = "전월 대비"

    def format_col_name(col):
        if "전월대비" in str(col):
            return str(col).replace("전월대비", 비교월_label)
        return str(col)

    th_pad = "5px 6px" if compact else "9px 12px"
    th_font = "12px" if compact else "13px"
    td_pad = "5px 6px" if compact else "8px 12px"
    td_font = "12px" if compact else "13px"
    th = f"background:#EDF2F7;color:#4A5568;font-weight:800;font-size:{th_font};padding:{th_pad};text-align:center;border-bottom:2px solid #E2E8F0;white-space:nowrap;"
    headers = "".join(f"<th style='{th}'>{html.escape(format_col_name(c))}</th>" for c in df.columns)
    align_overrides = align_overrides or {}

    body = ""
    for i, row in df.reset_index(drop=True).iterrows():
        tds = ""
        for col in df.columns:
            align = default_align or ("right" if col in num_cols or col in diff_cols else "center")
            if str(col).strip() == "활동내역":
                align = "left"
            align = align_overrides.get(col, align)
            bg = "#FFFFFF" if i % 2 == 0 else "#F7FAFC"

            if col == "일치여부":
                value = fmt_match(row[col])
            elif col in diff_cols:
                value = fmt_diff(row[col])
            elif col in num_cols:
                value = fmt_num(row[col])
            else:
                value = "" if pd.isna(row[col]) else html.escape(str(row[col]))

            _ws = "white-space:nowrap;"
            tds += (
                f"<td style='background:{bg};padding:{td_pad};border-bottom:1px solid #EDF2F7;"
                f"font-size:{td_font};color:#2D3748;text-align:{align};{_ws}'>{value}</td>"
            )
        body += f"<tr>{tds}</tr>"

    _ov = "visible" if compact else "auto"
    _tbl_style = "width:100%;border-collapse:collapse;table-layout:fixed;" if compact else "width:100%;border-collapse:collapse;"
    st.markdown(
        f"""
        <div class="pms-report-table" style="overflow-x:{_ov};border:1px solid #E2E8F0;border-radius:10px;box-shadow:0 2px 8px rgba(0,0,0,0.06);margin-bottom:1rem;">
            <table style="{_tbl_style}">
                <thead><tr>{headers}</tr></thead>
                <tbody>{body}</tbody>
            </table>
        </div>
        """,
        unsafe_allow_html=True,
    )


def show_auth_page():
    inject_theme_toggle()
    st.markdown(
        """
        <style>
        .stAppHeader, header, [data-testid="stHeader"], .stDecoration {
            display: none !important;
            height: 0 !important;
            visibility: hidden !important;
        }
        #MainMenu, footer,
        [data-testid="stToolbar"],
        [data-testid="stStatusWidget"],
        [data-testid="stDecoration"],
        .stDeployButton,
        ._profileContainer_gzau3_53,
        ._profilePreview_gzau3_63,
        [class*="viewerBadge"],
        [class*="StatusWidget"] {
            display: none !important;
            visibility: hidden !important;
        }
        [data-testid="stSidebar"] { display: none !important; }
        body:not(:has(#pms-d:checked)) .stApp,
        body:not(:has(#pms-d:checked)) [data-testid="stAppViewContainer"],
        body:not(:has(#pms-d:checked)) .main {
            background: #FFFFFF !important;
        }
        .main .block-container {
            max-width: 100% !important;
            padding: 64px 0 48px !important;
        }
        .auth-logo-card {
            background: transparent;
            text-align: center;
            padding: 22px 32px 18px;
            margin-bottom: 16px;
        }
        [data-testid="stColumn"] {
            background: transparent !important;
            box-shadow: none !important;
        }
        .auth-logo-title {
            color: #1E1A3A;
            font-size: 56px;
            font-weight: 900;
            letter-spacing: 0;
            line-height: 1.05;
            margin-bottom: 5px;
        }
        .auth-logo-sub {
            color: #7B79AA;
            font-size: 16px;
            font-weight: 600;
            letter-spacing: 0.3px;
        }
        .auth-small {
            text-align: center;
            color: #A09AC5;
            font-size: 14px;
            margin: 16px 0 8px;
        }
        div[data-testid="stTextInput"] { margin-bottom: 14px !important; }
        div[data-testid="stTextInput"] label {
            color: #3A3660 !important;
            font-size: 14px !important;
            font-weight: 700 !important;
            margin-bottom: 4px !important;
        }
        div[data-testid="stTextInput"] > div[data-baseweb="input"] {
            min-height: 52px !important;
            border-radius: 12px !important;
            border: none !important;
            background: #F2F2F5 !important;
            box-shadow: none !important;
        }
        div[data-testid="stTextInput"] input {
            color: #1E1A3A !important;
            caret-color: #1E1A3A !important;
            font-size: 15px !important;
            background: #F2F2F5 !important;
        }
        div[data-testid="stTextInput"] > div[data-baseweb="input"]:focus-within {
            border: 1.5px solid #7B6FD4 !important;
            box-shadow: 0 0 0 2px rgba(123,111,212,0.18) !important;
        }
        div[data-testid="stTextInput"] input::placeholder {
            color: #AEACC8 !important;
        }
        [data-testid="InputInstructions"] { display: none !important; }
        div[data-testid="stCheckbox"] {
            margin-bottom: 16px !important;
        }
        div[data-testid="stCheckbox"] label {
            color: #4B466D !important;
            font-size: 14px !important;
            font-weight: 600 !important;
        }
        div.stButton > button {
            height: 54px !important;
            border-radius: 12px !important;
            font-size: 15px !important;
            font-weight: 700 !important;
            width: 100% !important;
            transition: all 0.18s !important;
        }
        [data-testid="stBaseButton-primary"] {
            background: #655CF0 !important;
            color: #FFFFFF !important;
            border: 1.5px solid #655CF0 !important;
            box-shadow: 0 8px 18px rgba(101,92,240,0.22) !important;
        }
        [data-testid="stBaseButton-primary"]:hover {
            background: #5A52DF !important;
            border-color: #5A52DF !important;
            color: #FFFFFF !important;
        }
        [data-testid="stBaseButton-secondary"] {
            background: #FAFAFC !important;
            color: #3D3580 !important;
            border: 1.5px solid #EEEFF5 !important;
            box-shadow: 0 2px 8px rgba(30,26,58,0.05) !important;
        }
        [data-testid="stBaseButton-secondary"]:hover {
            background: #FFFFFF !important;
            border-color: #D8D3F5 !important;
            color: #493E9A !important;
        }
        [data-testid="stBaseButton-primary"] p,
        [data-testid="stBaseButton-primary"] span {
            color: #FFFFFF !important;
        }
        [data-testid="stBaseButton-secondary"] p,
        [data-testid="stBaseButton-secondary"] span {
            color: #3D3580 !important;
        }
        hr {
            border: 0 !important;
            border-top: 1px solid #EDEBF8 !important;
            margin: 18px 0 0 !important;
        }
        /* 로그인 페이지 다크모드 */
        body:has(#pms-d:checked) .auth-logo-title { color: #ffffff !important; }
        body:has(#pms-d:checked) .auth-logo-sub   { color: #a6adc8 !important; }
        body:has(#pms-d:checked) .auth-small       { color: #a6adc8 !important; }
        body:has(#pms-d:checked) div[data-testid="stTextInput"] label { color: #cdd6f4 !important; }
        body:has(#pms-d:checked) div[data-testid="stTextInput"] > div[data-baseweb="input"] {
            background: #252535 !important; border: 1.5px solid #45475a !important;
        }
        body:has(#pms-d:checked) div[data-testid="stTextInput"] input {
            background: #252535 !important; color: #ffffff !important;
            caret-color: #ffffff !important;
        }
        body:has(#pms-d:checked) div[data-testid="stTextInput"] > div[data-baseweb="input"]:focus-within {
            border-color: #818cf8 !important;
            box-shadow: 0 0 0 2px rgba(129,140,248,0.22) !important;
        }
        body:has(#pms-d:checked) div[data-testid="stCheckbox"] label { color: #cdd6f4 !important; }
        body:has(#pms-d:checked) hr { border-top-color: #45475a !important; }
        /* 모바일: 로그인 화면 */
        @media (max-width: 768px) {
            .main .block-container {
                padding: 24px 16px 32px !important;
                max-width: 100% !important;
            }
            /* 컬럼 가로 배치 해제 → 전체 폭 사용 */
            [data-testid="stHorizontalBlock"] {
                flex-direction: column !important;
            }
            [data-testid="stHorizontalBlock"] > [data-testid="stColumn"] {
                width: 100% !important;
                min-width: 100% !important;
                flex: 1 1 100% !important;
            }
            .auth-logo-title {
                font-size: 36px !important;
            }
            .auth-logo-sub {
                font-size: 14px !important;
                word-break: keep-all !important;
            }
            .auth-logo-card {
                padding: 12px 16px 10px !important;
                margin-bottom: 8px !important;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    _, center, _ = st.columns([1.5, 1.1, 1.5])

    with center:
        st.markdown(
            """
            <div class="auth-logo-card">
                <div class="auth-logo-title">내부 관리</div>
                <div class="auth-logo-sub">Webcash We · 360° Control</div>
            </div>
            """,
            unsafe_allow_html=True,
        )

        st.session_state.auth_mode = "login"

        if st.session_state.auth_mode == "login":
            cookie_manager = safe_cookie_controller()
            sid = cookie_get(cookie_manager, "saved_id", "")

            with st.form("login_form", border=False):
                u_id = st.text_input("아이디", value=sid, placeholder="아이디를 입력하세요", key="l_id")
                u_pw = st.text_input("비밀번호", type="password", placeholder="비밀번호를 입력하세요", key="l_pw")
                save_id_cb = st.checkbox("아이디 저장", value=bool(sid))
                _login_submitted = st.form_submit_button("로그인", use_container_width=True, type="primary")

            if _login_submitted:
                db = st.session_state.user_db
                u_id_str = str(u_id).strip() if u_id else ""
                u_pw_str = str(u_pw).strip() if u_pw else ""

                is_super = u_id_str == "1" and u_pw_str == "1"
                is_user = (
                    u_id_str in db
                    and db[u_id_str].get("pw", "") == u_pw_str
                    and db[u_id_str].get("access") == "허용"
                )

                if is_super or is_user:
                    if save_id_cb:
                        cookie_set(cookie_manager, "saved_id", u_id_str)
                    else:
                        try:
                            cookie_remove(cookie_manager, "saved_id")
                        except Exception:
                            pass
                    persist_login_session(u_id_str)

                    user = db.get(u_id_str, {"role": "관리자", "name": "최고관리자"})
                    st.session_state.logged_in = True
                    st.session_state.user_role = user.get("role", "관리자")
                    st.session_state.user_name = user.get("name", "최고관리자")
                    st.session_state.login_time = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M")

                    # 저장된 데이터 로드
                    saved_db = load_db(SAVED_STATE_FILE, {})
                    user_saved = saved_db.get(st.session_state.user_name)
                    if user_saved:
                        if user_saved.get("user_excel_data"):
                            st.session_state.user_excel_data = pd.DataFrame.from_dict(user_saved["user_excel_data"])
                        if user_saved.get("user_prev_month_sel"):
                            st.session_state.user_prev_month_sel = user_saved["user_prev_month_sel"]

                    admin_analysis = saved_db.get("admin_analysis")
                    if admin_analysis and admin_analysis.get("sent_df"):
                        st.session_state.analysis_result = pd.DataFrame.from_dict(admin_analysis["sent_df"])
                    if admin_analysis and admin_analysis.get("adm_prev_month"):
                        st.session_state.adm_prev_month = admin_analysis["adm_prev_month"]
                    deadline_info = saved_db.get("deadline")
                    if deadline_info:
                        st.session_state.deadline_time = deadline_info.get("time", "")
                    report_closed_info = saved_db.get("report_closed")
                    if report_closed_info:
                        st.session_state.report_closed = report_closed_info.get("time", "")

                    st.session_state.current_menu = "업로드 및 실적 확인"
                    persist_current_menu()
                    st.rerun()
                elif not u_id_str:
                    st.error("아이디를 입력해주세요.")
                elif not u_pw_str:
                    st.error("비밀번호를 입력해주세요.")
                else:
                    st.error("아이디 또는 비밀번호가 올바르지 않습니다.")

            st.divider()
            st.caption("계정은 관리자가 발급한 ID와 비밀번호로만 로그인할 수 있습니다.")


def show_sidebar():
    # 모바일: 메뉴 클릭 후 사이드바를 CSS로 숨김
    if st.session_state.pop("_close_sidebar_mobile", False):
        st.markdown("""
        <style>
        @media (max-width: 768px) {
            [data-testid="stSidebar"] {
                transform: translateX(-100vw) !important;
                visibility: hidden !important;
                pointer-events: none !important;
                transition: none !important;
            }
            [data-testid="stSidebarCollapsedControl"] {
                display: flex !important;
                visibility: visible !important;
            }
        }
        </style>
        """, unsafe_allow_html=True)
    with st.sidebar:
        st.markdown(
            """
            <style>
            #MainMenu, footer,
            [data-testid="stStatusWidget"],
            [data-testid="stDecoration"],
            [data-testid="stAppViewBlockContainer"] > div:last-child,
            .stDeployButton,
            ._profileContainer_gzau3_53,
            ._profilePreview_gzau3_63,
            .viewerBadge_container__r5tak,
            .viewerBadge_link__qRIco,
            [class*="viewerBadge"],
            [class*="StatusWidget"],
            [data-testid="stToolbar"] > * {
                display: none !important;
                visibility: hidden !important;
            }
            header[data-testid="stHeader"] {
                background: transparent !important;
                height: 0px !important;
                min-height: 0px !important;
                overflow: hidden !important;
            }
            [data-testid="stSidebar"] {
                display: flex !important;
                transform: none !important;
                visibility: visible !important;
                min-width: 272px !important;
                width: 272px !important;
            }
            [data-testid="stSidebarCollapseButton"],
            [data-testid="collapsedControl"],
            [data-testid="stSidebarCollapsedControl"] {
                display: none !important;
            }
            /* CRM navigation palette */
            [data-testid="stSidebar"] {
                background-color: #0F172A !important;
                border-right: 1px solid #1E293B !important;
            }
            [data-testid="stSidebar"] > div:first-child {
                padding-top: 0 !important;
            }
            [data-testid="stSidebarContent"] {
                padding: 14px 12px 16px !important;
                background-color: #0F172A !important;
            }
            [data-testid="stSidebar"] * {
                color: #E8EEF8 !important;
            }
            [data-testid="stSidebar"] .gpt-side-shell {
                display: flex;
                flex-direction: column;
                gap: 10px;
                padding: 4px 2px 8px;
            }
            [data-testid="stSidebar"] .gpt-brand {
                display: flex;
                align-items: center;
                gap: 10px;
                min-height: 44px;
                padding: 6px 8px;
                border-radius: 8px;
                border: 1px solid #1E293B;
                background: #111C33;
            }
            [data-testid="stSidebar"] .gpt-brand-mark {
                display: inline-flex;
                align-items: center;
                justify-content: center;
                width: 40px;
                height: 40px;
                border: none;
                border-radius: 8px;
                color: #ffffff !important;
                font-size: 16px;
                font-weight: 900;
                letter-spacing: 0;
                background: #2563EB;
                box-shadow: none;
            }
            [data-testid="stSidebar"] .gpt-brand-text {
                display: flex;
                flex-direction: column;
                line-height: 1.2;
            }
            [data-testid="stSidebar"] .gpt-brand-title {
                font-size: 20px;
                font-weight: 800;
                color: #ffffff !important;
                line-height: 1.1;
                word-break: keep-all;
            }
            [data-testid="stSidebar"] .gpt-brand-subtitle {
                margin-top: 2px;
                font-size: 12px;
                color: #94A3B8 !important;
                font-weight: 500;
            }
            [data-testid="stSidebar"] .gpt-user-card {
                margin: 2px 2px 8px;
                padding: 10px 12px;
                border: 1px solid #1E293B;
                border-radius: 8px;
                background: #111827;
                box-shadow: none;
            }
            [data-testid="stSidebar"] .gpt-user-name {
                font-size: 14px;
                font-weight: 800;
                color: #ffffff !important;
                line-height: 1.25;
            }
            [data-testid="stSidebar"] .gpt-user-meta {
                margin-top: 4px;
                font-size: 11px;
                color: #9FB3D6 !important;
                font-weight: 500;
                line-height: 1.35;
            }
            [data-testid="stSidebar"] .gpt-section {
                margin: 10px 8px 4px;
                font-size: 11px;
                color: #64748B !important;
                font-weight: 800;
                letter-spacing: 0.08em;
                text-transform: uppercase;
            }
            [data-testid="stSidebar"] .gpt-section:first-of-type {
                margin-top: 2px;
            }
            [data-testid="stSidebar"] [data-testid="stMarkdownContainer"] {
                margin: 0 !important;
                padding: 0 !important;
            }
            [data-testid="stSidebar"] .gpt-sidebar-spacer {
                height: 10px;
            }
            [data-testid="stSidebar"] .gpt-sidebar-divider {
                height: 1px;
                margin: 12px 8px;
                background: #1E293B;
            }
            [data-testid="stSidebar"] div.stButton {
                margin: 0 !important;
                padding: 0 !important;
            }
            [data-testid="stSidebar"] div.stButton > button {
                margin-bottom: 0 !important;
                min-height: 32px !important;
                height: auto !important;
            }
            [data-testid="stSidebar"] [data-testid="stElementContainer"] {
                padding: 0 !important;
                margin: 0 !important;
            }
            [data-testid="stSidebar"] [data-testid="stVerticalBlock"] {
                gap: 2px !important;
            }
            [data-testid="stSidebar"] hr {
                border-color: #24395E !important;
                margin: 10px 0 !important;
            }
            [data-testid="stSidebar"] div.stButton > button {
                background: transparent !important;
                border: none !important;
                box-shadow: none !important;
                color: #CBD5E1 !important;
                font-size: 13px !important;
                font-weight: 600 !important;
                text-align: left !important;
                padding: 7px 10px !important;
                border-radius: 6px !important;
                justify-content: flex-start !important;
                transition: background 0.12s ease, color 0.12s ease !important;
            }
            [data-testid="stSidebar"] div.stButton > button p,
            [data-testid="stSidebar"] div.stButton > button span,
            [data-testid="stSidebar"] [data-testid="stBaseButton-secondary"] p,
            [data-testid="stSidebar"] [data-testid="stBaseButton-secondary"] span {
                font-size: 13px !important;
                font-weight: 600 !important;
                color: #CBD5E1 !important;
                white-space: normal !important;
                line-height: 1.25 !important;
            }
            [data-testid="stSidebar"] div.stButton > button:hover {
                background: #1E293B !important;
                color: #ffffff !important;
            }
            [data-testid="stSidebar"] [data-testid="stElementContainer"]:has(.gpt-nav-active) + div [data-testid="stButton"] button,
            [data-testid="stSidebar"] [data-testid="stElementContainer"]:has(.gpt-nav-active) + [data-testid="stElementContainer"] [data-testid="stButton"] button {
                background: #2563EB !important;
                color: #ffffff !important;
                box-shadow: none;
            }
            [data-testid="stSidebar"] [data-testid="stElementContainer"]:has(.gpt-nav-active) + div [data-testid="stButton"] button p,
            [data-testid="stSidebar"] [data-testid="stElementContainer"]:has(.gpt-nav-active) + [data-testid="stElementContainer"] [data-testid="stButton"] button p {
                color: #ffffff !important;
                font-weight: 700 !important;
            }
            [data-testid="stSidebar"] .gpt-logout-marker + div [data-testid="stButton"] button,
            [data-testid="stSidebar"] .gpt-logout-marker + [data-testid="stElementContainer"] [data-testid="stButton"] button {
                color: #B9C8E4 !important;
            }
            [data-testid="stSidebar"] .gpt-logout-marker + div [data-testid="stButton"] button:hover,
            [data-testid="stSidebar"] .gpt-logout-marker + [data-testid="stElementContainer"] [data-testid="stButton"] button:hover {
                background: #7F1D1D !important;
                color: #ffffff !important;
            }
            </style>
            """,
            unsafe_allow_html=True,
        )

        def render_nav_button(menu_name):
            active_class = " gpt-nav-active" if st.session_state.current_menu == menu_name else ""
            label = CRM_MENU_LABELS.get(menu_name, menu_name)
            st.markdown(f"<div class='gpt-nav-marker{active_class}'></div>", unsafe_allow_html=True)
            if st.button(label, use_container_width=True, key=f"nav_{menu_name}"):
                st.session_state.current_menu = menu_name
                persist_current_menu()
                if menu_name == "업로드 및 실적 확인":
                    st.session_state["show_add_history_form"] = False
                st.session_state["_close_sidebar_mobile"] = True
                st.rerun()

        st.markdown(
            f"<div class='gpt-side-shell'>"
            f"<div class='gpt-brand'>"
            f"<div class='gpt-brand-mark'>CRM</div>"
            f"<div class='gpt-brand-text'>"
            f"<div class='gpt-brand-title'>내부관리 CRM</div>"
            f"<div class='gpt-brand-subtitle'>Customer Operations</div>"
            f"</div></div>"
            f"</div>",
            unsafe_allow_html=True,
        )

        if st.session_state.user_role == "관리자":
            st.markdown("<div class='gpt-section'>Workspace</div>", unsafe_allow_html=True)
            for menu_name in ["대시보드", "운영계획"]:
                render_nav_button(menu_name)

            st.markdown("<div class='gpt-section'>Performance</div>", unsafe_allow_html=True)
            for menu_name in ["관리자용 실적 확인", "실적 분석/계산", "실적 보고서", "주간보고 취합"]:
                render_nav_button(menu_name)

            st.markdown("<div class='gpt-section'>Billing</div>", unsafe_allow_html=True)
            for menu_name in [BILLING_MENU, "청구자료 작성"]:
                render_nav_button(menu_name)

            st.markdown("<div class='gpt-section'>Admin</div>", unsafe_allow_html=True)
            for menu_name in ["직원 및 권한설정", "구글 스트레드시트 연동", ACTIVITY_TEMPLATE_CONVERT_MENU]:
                render_nav_button(menu_name)

        st.markdown("<div class='gpt-section'>Customer CRM</div>", unsafe_allow_html=True)
        for menu_name in ["업로드 및 실적 확인", OPERATION_TARGET_MENU, "이번달 활동 대상고객 추천", "주간보고 이력 작성", ACTIVITY_TEMPLATE_CONVERT_MENU]:
            render_nav_button(menu_name)

        if st.session_state.user_role != "관리자":
            st.markdown("<div class='gpt-section'>Billing</div>", unsafe_allow_html=True)
            render_nav_button(BILLING_MENU)

        st.markdown("<div class='gpt-sidebar-divider'></div><div class='gpt-logout-marker'></div>", unsafe_allow_html=True)
        if st.button("로그아웃", use_container_width=True, key="nav_logout"):
            st.session_state.logged_in = False
            st.session_state.auth_mode = "login"
            try:
                _cm_logout = safe_cookie_controller()
                cookie_remove(_cm_logout, SESSION_UID_COOKIE)
                cookie_remove(_cm_logout, LAST_MENU_COOKIE)
            except Exception:
                pass
            st.rerun()


def _load_delegated_workplaces():
    default_data = {"workplaces": [], "requests": []}
    data = load_db(DELEGATED_WORKPLACE_FILE, default_data)
    if not isinstance(data, dict):
        return default_data
    data.setdefault("workplaces", [])
    data.setdefault("requests", [])
    return data


def _save_delegated_workplaces(data):
    data.setdefault("workplaces", [])
    data.setdefault("requests", [])
    save_db(DELEGATED_WORKPLACE_FILE, data, allow_shrink=True)


def _load_approvals():
    default_data = {"documents": []}
    data = load_db(APPROVAL_FILE, default_data)
    if not isinstance(data, dict):
        return default_data
    data.setdefault("documents", [])
    return data


def _save_approvals(data):
    data.setdefault("documents", [])
    save_db(APPROVAL_FILE, data, allow_shrink=True)


def _load_usage_reports():
    default_data = {"reports": []}
    data = load_db(USAGE_REPORT_FILE, default_data)
    if not isinstance(data, dict):
        return default_data
    data.setdefault("reports", [])
    return data


def _save_usage_reports(data):
    data.setdefault("reports", [])
    save_db(USAGE_REPORT_FILE, data, allow_shrink=True)


def _load_bank_accounts():
    default_data = {"accounts": []}
    data = load_db(BANK_ACCOUNT_FILE, default_data)
    if not isinstance(data, dict):
        return default_data
    data.setdefault("accounts", [])
    return data


def _save_bank_accounts(data):
    data.setdefault("accounts", [])
    save_db(BANK_ACCOUNT_FILE, data, allow_shrink=True)


def _account_company_label(account, workplaces_by_id):
    site = workplaces_by_id.get(account.get("linked_workplace_id"))
    if site:
        return site.get("workplace_name", "")
    return account.get("account_type", "")


def _my_workplaces(workplaces):
    """로그인 사용자가 관리자가 아닐 경우, 본인 담당/소속 사업장만 반환한다."""
    if st.session_state.user_role == "관리자":
        return workplaces

    user_name = st.session_state.get("user_name", "")
    user_dept = ""
    for info in _real_users(st.session_state.user_db).values():
        if info.get("name") == user_name:
            user_dept = info.get("dept_type", "")
            break
    return [
        row for row in workplaces
        if row.get("manager_name") == user_name or row.get("workplace_name") == user_dept
    ]


def _load_company_profile():
    default_data = {
        "name": "",
        "business_number": "",
        "ceo_name": "",
        "address": "",
        "contact": "",
        "memo": "",
        "updated_at": "",
    }
    data = load_db(COMPANY_PROFILE_FILE, default_data)
    if not isinstance(data, dict):
        return default_data
    for key, value in default_data.items():
        data.setdefault(key, value)
    return data


def _save_company_profile(data):
    save_db(COMPANY_PROFILE_FILE, data, allow_shrink=True)


def _normalize_name(value):
    return re.sub(r"\s+", "", str(value or "").strip().lower())


def _format_won(value):
    try:
        return f"{int(float(value or 0)):,}원"
    except Exception:
        return "0원"


def _current_kst():
    return datetime.utcnow() + timedelta(hours=9)


def _parse_kst_datetime(value):
    if not value:
        return None
    parsed = pd.to_datetime(str(value), errors="coerce")
    if pd.isna(parsed):
        return None
    return parsed.to_pydatetime().replace(tzinfo=None)


def _normalize_account_number(value):
    return re.sub(r"[^0-9A-Za-z]", "", str(value or ""))


def _coerce_int_amount(value):
    try:
        cleaned = re.sub(r"[^0-9.-]", "", str(value or "0"))
        return int(float(cleaned or 0))
    except Exception:
        return 0


def _history_amount_value(history_row, previous_balance=None):
    for key in ["amount", "deposit_amount", "transaction_amount", "입금액", "거래금액"]:
        if key in history_row:
            amount = _coerce_int_amount(history_row.get(key))
            if amount:
                return amount
    if previous_balance is not None and "balance" in history_row:
        balance = _coerce_int_amount(history_row.get("balance"))
        diff = balance - previous_balance
        if diff != 0:
            return diff
    if previous_balance is None and "balance" in history_row:
        return _coerce_int_amount(history_row.get("balance"))
    return 0


def _deposit_match_label(account, target_amount):
    history = account.get("balance_history", []) or []
    if not history:
        return ""

    sorted_history = sorted(history, key=lambda item: str(item.get("at", "")))
    previous_balance = None
    for item in sorted_history:
        amount = _history_amount_value(item, previous_balance)
        if "balance" in item:
            previous_balance = _coerce_int_amount(item.get("balance"))
        if amount != int(target_amount or 0):
            continue
        deposited_at = item.get("at") or item.get("date") or item.get("transaction_at") or ""
        deposited_date = str(deposited_at).split(" ")[0] if deposited_at else "입금일자 확인"
        return f"{deposited_date} 입금"
    return ""


def _reconcile_transfer_deposits(pending_rows, bank_accounts):
    accounts_by_number = {
        _normalize_account_number(row.get("account_number")): row
        for row in bank_accounts
        if row.get("account_number")
    }
    matched_count = 0
    for request_row in pending_rows:
        account_number = request_row.get("transfer_account_number", "")
        if not account_number:
            linked_account = next(
                (
                    row for row in bank_accounts
                    if row.get("linked_workplace_id") == request_row.get("workplace_id")
                ),
                {},
            )
            account_number = linked_account.get("account_number", "")
        account = accounts_by_number.get(_normalize_account_number(account_number))
        label = _deposit_match_label(account or {}, request_row.get("request_amount", 0))
        request_row["deposit_reconciliation"] = label or "미확인"
        request_row["deposit_reconciled_at"] = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
        if label:
            matched_count += 1
    return matched_count


def _is_recent_request(value, days=7):
    requested_at = _parse_kst_datetime(value)
    if not requested_at:
        return False
    cutoff = _current_kst() - timedelta(days=days)
    return requested_at >= cutoff


def _month_key(value):
    if not value:
        return ""
    return str(value)[:7]


def _workplace_request_metrics(requests):
    current_month = _current_kst().strftime("%Y-%m")
    month_requests = [row for row in requests if _month_key(row.get("requested_at")) == current_month]
    paid = [row for row in month_requests if row.get("status") == "이체 완료"]
    pending = [row for row in requests if row.get("status") == "요청"]
    approved = [row for row in requests if row.get("status") == "품의 확정"]
    return {
        "pending": len(pending),
        "approved": len(approved),
        "paid_amount": sum(int(row.get("request_amount", 0) or 0) for row in paid),
        "month_requests": len(month_requests),
    }


def _workplace_forecast_rows(workplaces, requests):
    current_month = _current_kst().strftime("%Y-%m")
    today_day = int(_current_kst().strftime("%d"))
    rows = []
    for site in workplaces:
        site_requests = [
            row for row in requests
            if row.get("workplace_id") == site.get("id") and row.get("status") == "이체 완료"
        ]
        month_totals = {}
        for row in site_requests:
            key = _month_key(row.get("paid_at") or row.get("requested_at"))
            if not key:
                continue
            month_totals[key] = month_totals.get(key, 0) + int(row.get("request_amount", 0) or 0)
        historical = [amount for key, amount in month_totals.items() if key != current_month]
        current_amount = month_totals.get(current_month, 0)
        average_amount = int(sum(historical) / len(historical)) if historical else current_amount
        recommended = max(average_amount, current_amount)
        regular_day = int(site.get("regular_payment_day", 0) or 0)
        days_to_payment = regular_day - today_day if regular_day else None
        risk_notes = []
        if average_amount and current_amount > average_amount * 1.3:
            risk_notes.append("평상시 대비 사용 증가")
        if days_to_payment is not None and 0 <= days_to_payment <= 3:
            risk_notes.append("정기 지급일 도래")
        if not site_requests:
            risk_notes.append("지급 이력 없음")
        rows.append(
            {
                "사업장명": site.get("workplace_name", ""),
                "담당자": site.get("manager_name", ""),
                "정기 지급일": f"매월 {regular_day}일" if regular_day else "미등록",
                "당월 지급액": _format_won(current_amount),
                "평균 지급액": _format_won(average_amount),
                "추천 지급액": _format_won(recommended),
                "리스크": ", ".join(risk_notes) if risk_notes else "정상",
            }
        )
    return rows


def _render_page_chrome(breadcrumb_items):
    crumbs = " &gt; ".join(html.escape(item) for item in breadcrumb_items[:-1])
    current = html.escape(breadcrumb_items[-1])
    st.markdown(
        f"""
        <style>
        .block-container {{
            max-width: 1180px !important;
            padding-left: 2rem !important;
            padding-right: 2rem !important;
        }}
        .pms-breadcrumb {{
            text-align: right;
            color: #64748b;
            font-size: 12px;
            font-weight: 600;
            margin: -4px 0 14px;
        }}
        .pms-breadcrumb span {{
            color: #008c78;
        }}
        [data-testid="stCheckbox"] label {{
            white-space: nowrap !important;
        }}
        </style>
        <div class="pms-breadcrumb">{crumbs} &gt; <span>{current}</span></div>
        """,
        unsafe_allow_html=True,
    )


def _workplace_dashboard_css():
    st.markdown(
        """
        <style>
        .sales-home-shell {
            position: relative;
            overflow: hidden;
            border-radius: 22px;
            padding: 14px 18px 22px;
            background:
                radial-gradient(circle at 48% 42%, rgba(16, 185, 173, 0.30) 0 32%, transparent 33%),
                linear-gradient(135deg, #f7f8fb 0%, #ffffff 52%, #f4f6fa 100%);
            box-shadow: 0 18px 50px rgba(15, 23, 42, 0.12);
            border: 1px solid rgba(226, 232, 240, 0.9);
        }
        .sales-topbar {
            height: 46px;
            display: flex;
            align-items: center;
            gap: 18px;
            border-radius: 24px;
            padding: 0 18px;
            background: rgba(255, 255, 255, 0.86);
            box-shadow: inset 0 0 0 1px rgba(226, 232, 240, 0.85);
            margin-bottom: 18px;
        }
        .sales-cloud {
            width: 34px;
            height: 22px;
            border-radius: 999px;
            background: #1d9bd7;
            position: relative;
            flex: 0 0 auto;
        }
        .sales-cloud::before, .sales-cloud::after {
            content: "";
            position: absolute;
            background: #1d9bd7;
            border-radius: 999px;
        }
        .sales-cloud::before { width: 18px; height: 18px; left: 5px; top: -7px; }
        .sales-cloud::after { width: 18px; height: 18px; right: 4px; top: -5px; }
        .sales-nav-pill {
            display: inline-flex;
            align-items: center;
            gap: 8px;
            font-weight: 800;
            color: #111827;
            white-space: nowrap;
        }
        .sales-app-icon {
            display: inline-grid;
            place-items: center;
            width: 18px;
            height: 18px;
            border-radius: 5px;
            background: #14b8a6;
            color: #ffffff;
            font-size: 12px;
            font-weight: 900;
        }
        .sales-home-title {
            font-size: 22px;
            font-weight: 850;
            color: #111827;
            margin: 8px 0 14px;
            letter-spacing: 0;
        }
        .work-card {
            min-height: 170px;
            border-radius: 16px;
            background: rgba(255, 255, 255, 0.90);
            border: 1px solid rgba(226, 232, 240, 0.8);
            box-shadow: 0 16px 35px rgba(15, 23, 42, 0.10);
            padding: 18px 18px 14px;
        }
        .work-card-title {
            color: #111827;
            font-size: 15px;
            font-weight: 850;
            margin-bottom: 12px;
        }
        .donut {
            width: 118px;
            height: 118px;
            border-radius: 50%;
            display: grid;
            place-items: center;
            margin: 0 auto 10px;
            background: conic-gradient(#8b5cf6 var(--p), #2f8eea var(--p) 72%, #7bb9f4 72% 100%);
        }
        .donut-inner {
            width: 76px;
            height: 76px;
            border-radius: 50%;
            background: #ffffff;
            display: grid;
            place-items: center;
            text-align: center;
            color: #111827;
            line-height: 1.12;
            font-size: 13px;
            box-shadow: inset 0 0 0 1px rgba(226, 232, 240, 0.9);
        }
        .donut-value {
            display: block;
            font-size: 20px;
            font-weight: 900;
            margin-bottom: 2px;
        }
        .card-dots {
            display: flex;
            gap: 6px;
            justify-content: flex-end;
            margin-top: 4px;
        }
        .card-dots span {
            width: 5px;
            height: 5px;
            border-radius: 50%;
            background: #2f8eea;
            display: block;
        }
        .card-dots span:nth-child(3), .card-dots span:nth-child(4) { background: #8b5cf6; }
        .suggestion-panel {
            min-height: 386px;
            border-radius: 18px;
            background: rgba(255, 255, 255, 0.88);
            border: 1px solid rgba(226, 232, 240, 0.85);
            box-shadow: 0 16px 35px rgba(15, 23, 42, 0.10);
            padding: 18px;
        }
        .suggestion-title {
            color: #6b7280;
            font-size: 15px;
            font-weight: 850;
            margin-bottom: 10px;
        }
        .suggestion-row {
            display: grid;
            grid-template-columns: 28px 1fr;
            gap: 12px;
            padding: 13px 0;
        }
        .suggestion-icon {
            width: 24px;
            height: 24px;
            border-radius: 5px;
            background: #8b7cf6;
            color: #fff;
            display: grid;
            place-items: center;
            font-size: 13px;
            font-weight: 900;
        }
        .skeleton {
            height: 6px;
            border-radius: 99px;
            background: #d7d7d7;
            margin: 8px 0;
        }
        .skeleton.short { width: 38%; }
        .skeleton.mid { width: 62%; }
        .skeleton.long { width: 96%; }
        .workflow-card {
            border-radius: 16px;
            background: rgba(255, 255, 255, 0.92);
            border: 1px solid rgba(226, 232, 240, 0.85);
            box-shadow: 0 12px 30px rgba(15, 23, 42, 0.10);
            padding: 16px;
            margin-top: 14px;
        }
        .workflow-title {
            font-size: 15px;
            font-weight: 850;
            color: #111827;
            margin-bottom: 12px;
        }
        .request-card {
            border-radius: 14px;
            background: rgba(255, 255, 255, 0.92);
            border: 1px solid rgba(226, 232, 240, 0.9);
            box-shadow: 0 10px 24px rgba(15, 23, 42, 0.08);
            padding: 13px 14px;
            min-height: 122px;
            margin-bottom: 12px;
        }
        .request-title {
            color: #111827;
            font-size: 15px;
            font-weight: 850;
            margin-bottom: 4px;
        }
        .request-meta {
            color: #6b7280;
            font-size: 13px;
            line-height: 1.4;
        }
        .status-chip {
            display: inline-flex;
            align-items: center;
            border-radius: 999px;
            padding: 4px 9px;
            font-size: 12px;
            font-weight: 800;
            background: #eef2ff;
            color: #4f46e5;
            margin-top: 8px;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def _donut_card(title, value, label, percent):
    pct = max(0, min(100, int(percent)))
    return f"""
    <div class="work-card">
        <div class="work-card-title">{html.escape(str(title))}</div>
        <div class="donut" style="--p:{pct}%;">
            <div class="donut-inner">
                <div>
                    <span class="donut-value">{html.escape(str(value))}</span>
                    {html.escape(str(label))}
                </div>
            </div>
        </div>
        <div class="card-dots"><span></span><span></span><span></span><span></span></div>
    </div>
    """


def show_advance_payment_request():
    _render_page_chrome(["홈", "업무", "전도금 요청", "전도금 요청"])
    st.markdown("### 전도금 요청")
    st.caption("위탁 사업장의 전도금을 요청하고 처리 현황을 확인합니다.")

    data = _load_delegated_workplaces()
    workplaces = data.get("workplaces", [])
    requests = data.get("requests", [])

    workplaces = _my_workplaces(workplaces)

    my_requests = [
        row
        for row in requests
        if row.get("requested_by") == st.session_state.get("user_name", "")
        and _is_recent_request(row.get("requested_at"), days=7)
    ]
    if my_requests:
        st.markdown("#### 최근 요청현황")
        _workplace_dashboard_css()
        request_cols = st.columns(3)
        for index, row in enumerate(sorted(my_requests, key=lambda item: item.get("requested_at", ""), reverse=True)):
            with request_cols[index % 3]:
                reject_line = ""
                if row.get("status") == "반려" and row.get("reject_reason"):
                    reject_line = f"<div class='request-meta'>반려 사유: {html.escape(str(row.get('reject_reason')))}</div>"
                st.markdown(
                    (
                        "<div class='request-card'>"
                        f"<div class='request-title'>{html.escape(str(row.get('workplace_name', '')))}</div>"
                        f"<div class='request-meta'>{html.escape(str(row.get('request_reason') or '요청 사유 없음'))}</div>"
                        f"<div class='request-meta'>{html.escape(_format_won(row.get('request_amount')))} · {html.escape(str(row.get('requested_at', '')))}</div>"
                        f"{reject_line}"
                        f"<span class='status-chip'>{html.escape(str(row.get('status', '')))}</span>"
                        "</div>"
                    ),
                    unsafe_allow_html=True,
                )
    else:
        st.info("최근 7일간 요청한 전도금 내역이 없습니다.")

    if not workplaces:
        if st.session_state.user_role != "관리자":
            st.info("담당으로 등록된 사업장이 없습니다. 관리자에게 문의해주세요.")
        else:
            st.info("관리자 메뉴의 [위탁 사업장 관리]에서 사업장을 먼저 등록해주세요.")
    else:
        workplace_options = {row.get("workplace_name", ""): row for row in workplaces}

        if st.session_state.get("_show_advance_payment_form"):
            with st.form("delegated_fund_request_form", clear_on_submit=True):
                selected_name = st.selectbox("사업장", list(workplace_options.keys()))
                requested_by = st.text_input(
                    "요청자명",
                    value=st.session_state.get("user_name", ""),
                    placeholder="요청자명",
                )
                requester_phone = st.text_input("핸드폰번호", placeholder="010-0000-0000")
                request_amount = st.number_input("요청 금액", min_value=0, step=100000, format="%d")
                request_reason = st.text_area("요청 사유", placeholder="전도금 사용 목적 및 필요 사유")
                col_submit, col_cancel = st.columns(2)
                requested = col_submit.form_submit_button("전도금 요청 등록", type="primary", use_container_width=True)
                cancelled = col_cancel.form_submit_button("취소", use_container_width=True)

            if requested:
                site = workplace_options[selected_name]
                if request_amount <= 0:
                    st.warning("요청 금액을 입력해주세요.")
                else:
                    now_str = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
                    request_id = int(time.time() * 1000)
                    requests.append(
                        {
                            "id": request_id,
                            "workplace_id": site.get("id"),
                            "workplace_name": site.get("workplace_name"),
                            "request_amount": int(request_amount),
                            "requested_by": requested_by.strip(),
                            "requester_phone": requester_phone.strip(),
                            "request_reason": request_reason.strip(),
                            "status": "요청",
                            "requested_at": now_str,
                            "approved_at": "",
                            "paid_at": "",
                            "approval_doc_id": request_id,
                            "reject_reason": "",
                            "transfer_file_generated_at": "",
                        }
                    )
                    _save_delegated_workplaces(data)

                    approvals = _load_approvals()
                    approvals["documents"].append(
                        {
                            "id": request_id,
                            "doc_type": "전도금요청",
                            "title": f"{site.get('workplace_name', '')} 전도금 요청",
                            "amount": int(request_amount),
                            "requester": requested_by.strip(),
                            "requester_phone": requester_phone.strip(),
                            "reason": request_reason.strip(),
                            "ref_request_id": request_id,
                            "ref_workplace_id": site.get("id"),
                            "status": "결재대기",
                            "requested_at": now_str,
                            "processed_at": "",
                            "processed_by": "",
                            "reject_reason": "",
                        }
                    )
                    _save_approvals(approvals)

                    st.session_state["_show_advance_payment_form"] = False
                    st.success("전도금 요청이 등록되었습니다. [전자결재]에서 처리 현황을 확인할 수 있습니다.")
                    st.rerun()
            if cancelled:
                st.session_state["_show_advance_payment_form"] = False
                st.rerun()
        else:
            col_add, _ = st.columns([35, 165])
            if col_add.button("+ 전도금 요청 등록", key="show_advance_payment_form_btn", use_container_width=True):
                st.session_state["_show_advance_payment_form"] = True
                st.rerun()


def show_e_approval():
    st.markdown("### 전자결재")
    st.caption("전도금 요청 등 결재 문서를 승인/반려 처리합니다.")

    approvals = _load_approvals()
    documents = approvals.get("documents", [])
    wp_data = _load_delegated_workplaces()
    wp_requests = wp_data.get("requests", [])
    requests_by_id = {row.get("id"): row for row in wp_requests}

    # 결재함 정보 표시
    pending_count = len([doc for doc in documents if doc.get("status") == "결재대기"])
    done_count = len([doc for doc in documents if doc.get("status") in ["승인", "반려"]])

    st.markdown(
        f"""
        <div style='background: linear-gradient(135deg, #2F6FED 0%, #1B4FC4 100%);
                    padding: 20px; border-radius: 12px; margin-bottom: 20px; color: white;'>
            <div style='font-size: 18px; font-weight: 700; margin-bottom: 10px;'>
                {html.escape(st.session_state.user_name)}님 반갑습니다
            </div>
            <div style='display: flex; gap: 30px; font-size: 14px;'>
                <div><strong>결재대기:</strong> {pending_count}건</div>
                <div><strong>처리완료:</strong> {done_count}건</div>
            </div>
        </div>
        """,
        unsafe_allow_html=True
    )

    pending_tab, done_tab, stats_tab = st.tabs(["결재대기", "처리 완료", "처리 통계"])

    with pending_tab:
        pending_docs = [doc for doc in documents if doc.get("status") == "결재대기"]
        if not pending_docs:
            st.info("결재 대기 중인 문서가 없습니다.")
        else:
            _workplace_dashboard_css()
            cols = st.columns(3)
            for index, doc in enumerate(sorted(pending_docs, key=lambda item: item.get("requested_at", ""), reverse=True)):
                with cols[index % 3]:
                    st.markdown(
                        (
                            "<div class='request-card'>"
                            f"<div class='request-title'>{html.escape(str(doc.get('title', '')))}</div>"
                            f"<div class='request-meta'>{html.escape(str(doc.get('reason') or '사유 없음'))}</div>"
                            f"<div class='request-meta'>{html.escape(_format_won(doc.get('amount')))} · {html.escape(str(doc.get('requester', '')))}</div>"
                            f"<div class='request-meta'>{html.escape(str(doc.get('requested_at', '')))}</div>"
                            f"<span class='status-chip'>{html.escape(str(doc.get('status', '')))}</span>"
                            "</div>"
                        ),
                        unsafe_allow_html=True,
                    )
                    if st.session_state.user_role == "관리자":
                        now_str = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
                        approve_col, reject_col = st.columns(2)
                        if approve_col.button("승인", key=f"approve_doc_{doc.get('id')}", use_container_width=True, type="primary"):
                            doc["status"] = "승인"
                            doc["processed_at"] = now_str
                            doc["processed_by"] = st.session_state.user_name
                            linked = requests_by_id.get(doc.get("ref_request_id"))
                            if doc.get("doc_type") == "사용품의서":
                                usage_data = _load_usage_reports()
                                usage_by_id = {row.get("id"): row for row in usage_data.get("reports", [])}
                                usage_report = usage_by_id.get(doc.get("ref_usage_report_id"))
                                if usage_report:
                                    usage_report["status"] = "승인"
                                    usage_report["processed_at"] = now_str
                                    usage_report["processed_by"] = st.session_state.user_name
                                _save_usage_reports(usage_data)
                                if linked:
                                    linked["usage_report_status"] = "승인"
                            elif linked:
                                linked["status"] = "품의 확정"
                                linked["approved_at"] = now_str
                            _save_approvals(approvals)
                            _save_delegated_workplaces(wp_data)
                            st.rerun()
                        if reject_col.button("반려", key=f"reject_doc_{doc.get('id')}", use_container_width=True):
                            st.session_state[f"_reject_target_{doc.get('id')}"] = True
                        if st.session_state.get(f"_reject_target_{doc.get('id')}"):
                            reason = st.text_input("반려 사유", key=f"reject_reason_{doc.get('id')}")
                            if st.button("반려 확정", key=f"reject_confirm_{doc.get('id')}", use_container_width=True):
                                doc["status"] = "반려"
                                doc["processed_at"] = now_str
                                doc["processed_by"] = st.session_state.user_name
                                doc["reject_reason"] = reason.strip()
                                linked = requests_by_id.get(doc.get("ref_request_id"))
                                if doc.get("doc_type") == "사용품의서":
                                    usage_data = _load_usage_reports()
                                    usage_by_id = {row.get("id"): row for row in usage_data.get("reports", [])}
                                    usage_report = usage_by_id.get(doc.get("ref_usage_report_id"))
                                    if usage_report:
                                        usage_report["status"] = "반려"
                                        usage_report["processed_at"] = now_str
                                        usage_report["processed_by"] = st.session_state.user_name
                                        usage_report["reject_reason"] = reason.strip()
                                    _save_usage_reports(usage_data)
                                    if linked:
                                        linked["usage_report_status"] = "반려"
                                elif linked:
                                    linked["status"] = "반려"
                                    linked["reject_reason"] = reason.strip()
                                    linked["processed_at"] = now_str
                                _save_approvals(approvals)
                                _save_delegated_workplaces(wp_data)
                                st.session_state.pop(f"_reject_target_{doc.get('id')}", None)
                                st.rerun()
                    else:
                        st.caption("관리자 승인 대기 중입니다.")

    with done_tab:
        done_docs = [doc for doc in documents if doc.get("status") != "결재대기"]
        if not done_docs:
            st.info("처리된 결재 문서가 없습니다.")
        else:
            done_df = pd.DataFrame(done_docs)
            done_df["amount_won"] = done_df["amount"].apply(_format_won)
            for col in ["processed_at", "processed_by", "reject_reason"]:
                if col not in done_df.columns:
                    done_df[col] = ""
            done_view = done_df[
                ["requested_at", "title", "amount_won", "requester", "status", "processed_at", "processed_by", "reject_reason"]
            ].rename(
                columns={
                    "requested_at": "요청일시",
                    "title": "문서명",
                    "amount_won": "금액",
                    "requester": "요청자",
                    "status": "처리결과",
                    "processed_at": "처리일시",
                    "processed_by": "처리자",
                    "reject_reason": "반려사유",
                }
            )
            render_plain_html_table(done_view.sort_values("요청일시", ascending=False), center_align=True)

    with stats_tab:
        if documents:
            pending_n = len([doc for doc in documents if doc.get("status") == "결재대기"])
            approved_n = len([doc for doc in documents if doc.get("status") == "승인"])
            rejected_n = len([doc for doc in documents if doc.get("status") == "반려"])
            s1, s2, s3 = st.columns(3)
            s1.metric("결재대기", f"{pending_n:,}건")
            s2.metric("승인", f"{approved_n:,}건")
            s3.metric("반려", f"{rejected_n:,}건")

            month_counts = {}
            for doc in documents:
                key = _month_key(doc.get("requested_at"))
                if key:
                    month_counts[key] = month_counts.get(key, 0) + 1
            if month_counts:
                import plotly.graph_objects as go
                months = sorted(month_counts.keys())
                fig = go.Figure(go.Bar(x=months, y=[month_counts[m] for m in months]))
                fig.update_layout(**_chart_layout(height=280))
                st.plotly_chart(fig, use_container_width=True, theme=None)
        else:
            st.info("전자결재 이력이 없습니다.")


def show_approval_result():
    _render_page_chrome(["홈", "조회", "품의 결과", "품의 결과"])
    st.markdown("### 품의 결과")
    st.caption("전도금 요청의 품의(승인/반려) 처리 결과를 조회합니다.")

    data = _load_delegated_workplaces()
    requests = data.get("requests", [])
    workplaces = data.get("workplaces", [])

    if st.session_state.user_role != "관리자":
        my_ids = {row.get("id") for row in _my_workplaces(workplaces)}
        requests = [row for row in requests if row.get("workplace_id") in my_ids]

    processed = [row for row in requests if row.get("status") != "요청"]
    if not processed:
        st.info("처리된 품의 결과가 없습니다.")
        return

    def _processed_at(row):
        return row.get("approved_at") or row.get("paid_at") or row.get("processed_at") or ""

    requester_names = ["전체"] + sorted({row.get("requested_by", "") for row in processed if row.get("requested_by")})
    status_options = ["전체"] + sorted({row.get("status", "") for row in processed if row.get("status")})
    today = _current_kst().date()
    if "approval_result_start_date" not in st.session_state:
        st.session_state.approval_result_start_date = today - timedelta(days=6)
    if "approval_result_end_date" not in st.session_state:
        st.session_state.approval_result_end_date = today
    period_presets = {
        "오늘": (today, today),
        "어제": (today - timedelta(days=1), today - timedelta(days=1)),
        "1주일": (today - timedelta(days=6), today),
        "1개월": (today - timedelta(days=30), today),
    }

    with st.container():
        request_date_label_col, request_date_input_col, _ = st.columns([0.1, 0.54, 0.36])
        with request_date_label_col:
            st.markdown("**조회기간**")
        with request_date_input_col:
            preset_cols = st.columns([0.13, 0.13, 0.17, 0.17, 0.40])
            for index, (label, date_range) in enumerate(period_presets.items()):
                button_type = "primary" if label == st.session_state.get("approval_result_preset", "1주일") else "secondary"
                if preset_cols[index].button(label, key=f"approval_result_preset_{label}", type=button_type, use_container_width=True):
                    st.session_state.approval_result_preset = label
                    st.session_state.approval_result_start_date = date_range[0]
                    st.session_state.approval_result_end_date = date_range[1]
                    st.rerun()

            date_start_col, tilde_col, date_end_col, month_col, request_check_col, process_check_col = st.columns(
                [0.18, 0.03, 0.18, 0.17, 0.17, 0.17]
            )
            with date_start_col:
                request_start_date = st.date_input(
                    "조회 시작일",
                    key="approval_result_start_date",
                    label_visibility="collapsed",
                )
            with tilde_col:
                st.markdown("<div style='padding-top:8px;text-align:center;'>~</div>", unsafe_allow_html=True)
            with date_end_col:
                request_end_date = st.date_input(
                    "조회 종료일",
                    key="approval_result_end_date",
                    label_visibility="collapsed",
                )
            with month_col:
                month_options = ["월별 선택"] + [
                    (today.replace(day=1) - pd.DateOffset(months=idx)).strftime("%Y-%m")
                    for idx in range(12)
                ]
                selected_month = st.selectbox(
                    "월별 선택",
                    month_options,
                    key="approval_result_month",
                    label_visibility="collapsed",
                )
            with request_check_col:
                date_by_request = st.checkbox("요청일자", value=True, key="approval_result_date_by_request")
            with process_check_col:
                date_by_processed = st.checkbox("처리일자", value=False, key="approval_result_date_by_processed")

        requester_label_col, requester_input_col, _ = st.columns([0.1, 0.24, 0.66])
        with requester_label_col:
            st.markdown("**요청자**")
        with requester_input_col:
            selected_requester = st.selectbox(
                "요청자",
                requester_names,
                key="approval_result_requester",
                label_visibility="collapsed",
            )

        status_label_col, status_input_col, _ = st.columns([0.1, 0.24, 0.66])
        with status_label_col:
            st.markdown("**처리결과**")
        with status_input_col:
            selected_status = st.selectbox(
                "처리결과",
                status_options,
                key="approval_result_status",
                label_visibility="collapsed",
            )

    button_left, button_center, button_right = st.columns([0.25, 0.1, 0.65])
    with button_center:
        search_clicked = st.button("조회", key="approval_result_search", type="primary", use_container_width=True)

    if selected_month != "월별 선택":
        month_start = pd.to_datetime(f"{selected_month}-01").date()
        month_end = (pd.to_datetime(f"{selected_month}-01") + pd.offsets.MonthEnd(0)).date()
        request_start_date, request_end_date = month_start, month_end

    if search_clicked or "_approval_result_query" not in st.session_state:
        st.session_state["_approval_result_query"] = {
            "request_date_range": (request_start_date, request_end_date),
            "date_by_request": date_by_request,
            "date_by_processed": date_by_processed,
            "selected_requester": selected_requester,
            "selected_status": selected_status,
        }

    query = st.session_state.get("_approval_result_query", {})
    request_date_range = query.get("request_date_range", ())
    date_by_request = query.get("date_by_request", True)
    date_by_processed = query.get("date_by_processed", False)
    selected_requester = query.get("selected_requester", "전체")
    selected_status = query.get("selected_status", "전체")

    def _date_in_range(value, date_range):
        if not isinstance(date_range, tuple) or len(date_range) != 2:
            return True
        parsed = _parse_kst_datetime(value)
        if not parsed:
            return False
        start, end = date_range
        return start <= parsed.date() <= end

    approvals_data = _load_approvals()
    docs_by_request_id = {
        doc.get("ref_request_id"): doc
        for doc in approvals_data.get("documents", [])
        if doc.get("ref_request_id") is not None
    }

    filtered = processed
    if request_date_range and (date_by_request or date_by_processed):
        if date_by_request:
            filtered = [row for row in filtered if _date_in_range(row.get("requested_at"), request_date_range)]
        if date_by_processed:
            filtered = [row for row in filtered if _date_in_range(_processed_at(row), request_date_range)]
    if selected_requester != "전체":
        filtered = [row for row in filtered if row.get("requested_by") == selected_requester]
    if selected_status != "전체":
        filtered = [row for row in filtered if row.get("status") == selected_status]

    if not filtered:
        st.info("조건에 맞는 품의 결과가 없습니다.")
        return

    result_df = pd.DataFrame(filtered)
    result_df["request_amount_won"] = result_df["request_amount"].apply(_format_won)
    for col in ["approved_at", "paid_at", "processed_at", "reject_reason"]:
        if col not in result_df.columns:
            result_df[col] = ""
    result_df["처리일시"] = result_df.apply(_processed_at, axis=1)
    result_df["문서구분"] = result_df.apply(
        lambda row: "전도금 사용 결의 보고"
        if (docs_by_request_id.get(row.get("id"), {}).get("doc_type") == "사용품의서" or row.get("usage_report_id"))
        else "전도금 요청",
        axis=1,
    )
    type_col1, type_col2, _ = st.columns([0.13, 0.2, 0.67])
    with type_col1:
        show_advance_requests = st.checkbox("전도금 요청", value=True, key="approval_result_show_advance_requests")
    with type_col2:
        show_usage_reports = st.checkbox("전도금 사용 결의 보고", value=True, key="approval_result_show_usage_reports")

    visible_doc_types = []
    if show_advance_requests:
        visible_doc_types.append("전도금 요청")
    if show_usage_reports:
        visible_doc_types.append("전도금 사용 결의 보고")
    result_df = result_df[result_df["문서구분"].isin(visible_doc_types)]
    if result_df.empty:
        st.info("선택한 문서구분에 해당하는 품의 결과가 없습니다.")
        return

    result_view = result_df[
        ["문서구분", "requested_at", "workplace_name", "request_amount_won", "requested_by", "status", "처리일시", "reject_reason"]
    ].rename(
        columns={
            "requested_at": "요청일시",
            "workplace_name": "사업장명",
            "request_amount_won": "요청 금액",
            "requested_by": "요청자",
            "status": "처리결과",
            "reject_reason": "반려사유",
        }
    )
    render_plain_html_table(
        result_view.sort_values("요청일시", ascending=False),
        center_align=True,
        stretch=True,
        max_width="760px",
        border=False,
    )


def show_usage_report():
    st.markdown("### 전도금 사용 결의 보고")
    st.caption("계좌 거래내역을 기반으로 출금 내역의 사용 사유를 작성하여 본사에 보고합니다.")

    tab1, tab2 = st.tabs(["보고서 작성", "제출 내역"])

    with tab1:
        _render_usage_report_form()

    with tab2:
        _render_usage_report_history()


def _render_usage_report_form():
    """전도금 사용 결의 보고서 작성 폼"""
    st.markdown("#### 보고서 작성")

    # 폼 너비 제한
    st.markdown(
        """
        <style>
        [data-testid="stVerticalBlock"] > [style*="flex-direction: column;"] > [data-testid="stVerticalBlock"] {
            max-width: 1000px;
            margin: 0 auto;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

    # 계좌 선택
    bank_data = _load_bank_accounts()
    accounts = bank_data.get("accounts", [])

    if not accounts:
        st.warning("등록된 계좌가 없습니다. [계좌 관리]에서 계좌를 먼저 등록해주세요.")
        return

    wp_data = _load_delegated_workplaces()
    workplaces_by_id = {site.get("id"): site for site in wp_data.get("workplaces", [])}

    account_options = {
        f"{row.get('account_name')} ({row.get('bank_name')} {row.get('account_number')})": row
        for row in accounts
    }

    with st.container(border=False):
        # 계좌번호
        account_label_col, account_input_col = st.columns([0.12, 0.88])
        with account_label_col:
            st.markdown("**계좌번호 <span style='color:#008c78'>*</span>**", unsafe_allow_html=True)
        with account_input_col:
            account_select_col, _ = st.columns([0.48, 0.52])
            with account_select_col:
                selected_account_label = st.selectbox(
                    "계좌번호",
                    list(account_options.keys()),
                    key="usage_report_account",
                    label_visibility="collapsed"
                )
                selected_account = account_options[selected_account_label]

        # 조회기간
        today = _current_kst().date()
        if "usage_report_start_date" not in st.session_state:
            st.session_state.usage_report_start_date = today.replace(day=1)
        if "usage_report_end_date" not in st.session_state:
            st.session_state.usage_report_end_date = today

        period_presets = {
            "오늘": (today, today),
            "어제": (today - timedelta(days=1), today - timedelta(days=1)),
            "1주일": (today - timedelta(days=6), today),
            "1개월": (today - timedelta(days=30), today),
        }

        period_label_col, period_input_col = st.columns([0.12, 0.88])
        with period_label_col:
            st.markdown("**조회기간**")
        with period_input_col:
            preset_cols = st.columns([0.13, 0.13, 0.17, 0.17, 0.40])
            for index, (label, date_range) in enumerate(period_presets.items()):
                button_type = "primary" if label == st.session_state.get("usage_report_preset", "1주일") else "secondary"
                if preset_cols[index].button(label, key=f"usage_report_preset_{label}", type=button_type, use_container_width=True):
                    st.session_state.usage_report_preset = label
                    st.session_state.usage_report_start_date = date_range[0]
                    st.session_state.usage_report_end_date = date_range[1]
                    st.rerun()

            date_start_col, tilde_col, date_end_col, month_col, _ = st.columns([0.22, 0.03, 0.22, 0.18, 0.35])
            with date_start_col:
                start_date = st.date_input("시작일", key="usage_report_start_date", label_visibility="collapsed")
            with tilde_col:
                st.markdown("<div style='padding-top:8px;text-align:center;'>~</div>", unsafe_allow_html=True)
            with date_end_col:
                end_date = st.date_input("종료일", key="usage_report_end_date", label_visibility="collapsed")
            with month_col:
                month_options = ["월별 선택"] + [
                    (today.replace(day=1) - pd.DateOffset(months=idx)).strftime("%Y-%m")
                    for idx in range(12)
                ]
                selected_month = st.selectbox(
                    "월별 선택",
                    month_options,
                    key="usage_report_month",
                    label_visibility="collapsed",
                )

    # 조회 버튼
    button_left, button_center, button_right = st.columns([0.42, 0.16, 0.42])
    with button_center:
        if st.button("조회", key="usage_report_search", type="primary", use_container_width=True):
            if selected_month != "월별 선택":
                month_start = pd.to_datetime(f"{selected_month}-01").date()
                month_end = (pd.to_datetime(f"{selected_month}-01") + pd.offsets.MonthEnd(0)).date()
                st.session_state.usage_report_start_date = month_start
                st.session_state.usage_report_end_date = month_end

            st.session_state.withdrawal_query_done = True
            st.rerun()

    if not st.session_state.get("withdrawal_query_done"):
        st.info("계좌와 조회기간을 선택한 후 '조회' 버튼을 클릭해주세요.")
        return

    st.markdown("---")
    st.markdown("#### 출금 내역 및 사용 사유 입력")

    # 계좌 거래내역에서 출금 데이터 가져오기
    # 실제 거래내역이 있다면 여기서 필터링
    # 임시로 샘플 데이터 또는 balance_history 활용
    transactions = selected_account.get("transactions", [])

    # 조회기간에 해당하는 출금 내역만 필터링
    query_start = st.session_state.usage_report_start_date
    query_end = st.session_state.usage_report_end_date

    withdrawals = [
        t for t in transactions
        if t.get("type") == "출금" and query_start <= pd.to_datetime(t.get("date")).date() <= query_end
    ]

    if not withdrawals:
        st.warning(
            "선택한 기간에 출금 내역이 없습니다.\n\n"
            "💡 참고: 거래내역은 [계좌 잔고 확인] 메뉴의 거래내역조회에서 확인할 수 있습니다."
        )
        st.info(
            "**임시 입력 모드**\n\n"
            "아래 버튼을 클릭하여 출금 내역을 직접 입력할 수 있습니다."
        )

    # 출금 항목별 사용 사유 입력
    if "usage_report_items" not in st.session_state:
        st.session_state.usage_report_items = {}

    # 거래내역이 있으면 표시하고 사용 사유 입력 받기
    if withdrawals:
        st.markdown(f"**조회된 출금 내역: {len(withdrawals)}건**")

        for idx, withdrawal in enumerate(withdrawals):
            with st.expander(
                f"📤 {withdrawal.get('date')} | {_format_won(withdrawal.get('amount'))} | {withdrawal.get('recipient', '-')}"
            ):
                trans_id = f"{withdrawal.get('date')}_{withdrawal.get('amount')}"

                usage_reason = st.text_area(
                    "사용 사유",
                    key=f"reason_{idx}",
                    placeholder="출금 사용 목적을 상세히 작성해주세요.",
                    value=st.session_state.usage_report_items.get(trans_id, {}).get("reason", "")
                )

                usage_file = st.file_uploader(
                    "증빙서류 첨부 (선택)",
                    key=f"file_{idx}",
                    type=["pdf", "jpg", "jpeg", "png"]
                )

                if st.button("저장", key=f"save_{idx}", type="primary"):
                    file_name = usage_file.name if usage_file else ""
                    st.session_state.usage_report_items[trans_id] = {
                        "date": withdrawal.get("date"),
                        "amount": withdrawal.get("amount"),
                        "recipient": withdrawal.get("recipient", ""),
                        "reason": usage_reason.strip(),
                        "file_name": file_name,
                    }
                    st.success("저장되었습니다.")

    # 수동 입력 옵션
    with st.expander("➕ 출금 내역 수동 추가"):
        with st.form("add_withdrawal_manual"):
            col_a, col_b, col_c = st.columns(3)
            with col_a:
                manual_date = st.date_input("출금일자")
            with col_b:
                manual_amount = st.number_input("출금금액", min_value=0, step=1000, format="%d")
            with col_c:
                manual_recipient = st.text_input("거래처")

            manual_reason = st.text_area("사용 사유", placeholder="출금 사용 목적을 상세히 작성해주세요.")
            manual_file = st.file_uploader("증빙서류 첨부 (선택)", type=["pdf", "jpg", "jpeg", "png"])

            if st.form_submit_button("추가", type="primary", use_container_width=True):
                if not manual_reason.strip():
                    st.warning("사용 사유를 입력해주세요.")
                else:
                    trans_id = f"{manual_date.strftime('%Y-%m-%d')}_{manual_amount}_manual"
                    file_name = manual_file.name if manual_file else ""

                    st.session_state.usage_report_items[trans_id] = {
                        "date": manual_date.strftime("%Y-%m-%d"),
                        "amount": int(manual_amount),
                        "recipient": manual_recipient.strip(),
                        "reason": manual_reason.strip(),
                        "file_name": file_name,
                    }
                    st.success("출금 항목이 추가되었습니다.")
                    st.rerun()

    st.markdown("---")

    # 입력된 항목 표시
    if st.session_state.usage_report_items:
        st.markdown("#### 작성된 사용 사유")

        for trans_id, item in st.session_state.usage_report_items.items():
            with st.container():
                col_info, col_del = st.columns([0.9, 0.1])
                with col_info:
                    file_badge = f" 📎 {item['file_name']}" if item['file_name'] else ""
                    st.markdown(
                        f"""
                        **{item['date']}** | {_format_won(item['amount'])} | {item['recipient']}
                        사유: {item['reason']}{file_badge}
                        """
                    )
                with col_del:
                    if st.button("🗑️", key=f"del_{trans_id}"):
                        del st.session_state.usage_report_items[trans_id]
                        st.rerun()
                st.markdown("---")

        # 보고서 제출
        total_amount = sum(item['amount'] for item in st.session_state.usage_report_items.values())
        st.markdown(f"**총 출금액:** {_format_won(total_amount)}")
        st.markdown(f"**작성 완료 건수:** {len(st.session_state.usage_report_items)}건")

        if st.button("보고서 제출", type="primary", use_container_width=True):
            # 사용 사유가 입력되지 않은 항목 확인
            incomplete = [item for item in st.session_state.usage_report_items.values() if not item.get('reason', '').strip()]

            if incomplete:
                st.warning(f"사용 사유가 입력되지 않은 항목이 {len(incomplete)}건 있습니다. 모든 항목의 사용 사유를 입력해주세요.")
            else:
                # 보고서 저장
                usage_data = _load_usage_reports()
                if "reports" not in usage_data:
                    usage_data["reports"] = []

                workplace_name = workplaces_by_id.get(selected_account.get("linked_workplace_id"), {}).get("workplace_name", "")

                query_start = st.session_state.usage_report_start_date
                query_end = st.session_state.usage_report_end_date

                usage_data["reports"].append({
                    "id": int(time.time() * 1000),
                    "workplace_id": selected_account.get("linked_workplace_id"),
                    "workplace_name": workplace_name,
                    "account_id": selected_account.get("id"),
                    "report_period": f"{query_start.strftime('%Y-%m-%d')} ~ {query_end.strftime('%Y-%m-%d')}",
                    "total_amount": total_amount,
                    "items": list(st.session_state.usage_report_items.values()),
                    "status": "제출완료",
                    "requested_at": _current_kst().strftime("%Y-%m-%d %H:%M:%S"),
                })

                _save_usage_reports(usage_data)
                st.session_state.usage_report_items = {}
                st.session_state.withdrawal_query_done = False
                st.success("전도금 사용 결의 보고서가 제출되었습니다.")
                st.rerun()
    else:
        st.info("출금 내역의 사용 사유를 입력해주세요.")


def _render_usage_report_history():
    """제출된 전도금 사용 결의 보고 내역"""
    wp_data = _load_delegated_workplaces()
    workplaces = _my_workplaces(wp_data.get("workplaces", []))
    my_ids = {row.get("id") for row in workplaces}

    usage_data = _load_usage_reports()
    usage_reports = usage_data.get("reports", [])
    my_reports = [row for row in usage_reports if row.get("workplace_id") in my_ids]

    if not my_reports:
        st.info("제출된 전도금 사용 결의 보고가 없습니다.")
        return

    _workplace_dashboard_css()
    cols = st.columns(3)
    for index, row in enumerate(sorted(my_reports, key=lambda item: item.get("requested_at", ""), reverse=True)):
        with cols[index % 3]:
            reject_line = ""
            if row.get("status") == "반려" and row.get("reject_reason"):
                reject_line = f"<div class='request-meta'>반려 사유: {html.escape(str(row.get('reject_reason')))}</div>"

            items_count = len(row.get("items", []))
            report_period = row.get('report_period', row.get('report_month', ''))  # 하위 호환성
            st.markdown(
                (
                    "<div class='request-card'>"
                    f"<div class='request-title'>{html.escape(str(row.get('workplace_name', '')))}</div>"
                    f"<div class='request-meta'>{html.escape(report_period)} · {items_count}건</div>"
                    f"<div class='request-meta'>{html.escape(_format_won(row.get('total_amount')))} · {html.escape(str(row.get('requested_at', '')))}</div>"
                    f"{reject_line}"
                    f"<span class='status-chip'>{html.escape(str(row.get('status', '')))}</span>"
                    "</div>"
                ),
                unsafe_allow_html=True,
            )


def show_approval_line_settings():
    st.markdown("### 결재선 설정")
    st.caption("결재 문서의 승인자와 결재 단계를 관리합니다.")

    # 결재선 데이터 로드 (임시로 세션 스테이트 사용)
    if "approval_lines" not in st.session_state:
        st.session_state.approval_lines = []

    approval_lines = st.session_state.approval_lines

    # 등록된 결재선 표시
    if approval_lines:
        st.markdown("#### 등록된 결재선")
        line_df = pd.DataFrame(approval_lines)
        line_view = line_df[["line_name", "approver1", "approver2", "approver3"]].rename(
            columns={
                "line_name": "결재선명",
                "approver1": "1차 승인자",
                "approver2": "2차 승인자",
                "approver3": "3차 승인자",
            }
        )
        st.dataframe(line_view, use_container_width=True, hide_index=True)
    else:
        st.info("등록된 결재선이 없습니다.")

    # 결재선 추가 폼
    if st.session_state.get("_show_approval_line_add_form"):
        with st.form("approval_line_add_form"):
            line_name = st.text_input("결재선명", placeholder="예: 전도금 결재선")
            col1, col2, col3 = st.columns(3)
            approver1 = col1.text_input("1차 승인자", placeholder="필수")
            approver2 = col2.text_input("2차 승인자", placeholder="선택")
            approver3 = col3.text_input("3차 승인자", placeholder="선택")

            col_submit, col_cancel = st.columns(2)
            submitted = col_submit.form_submit_button("등록", type="primary", use_container_width=True)
            cancelled = col_cancel.form_submit_button("취소", use_container_width=True)

        if submitted:
            if not line_name.strip() or not approver1.strip():
                st.warning("결재선명과 1차 승인자는 필수입니다.")
            else:
                approval_lines.append({
                    "id": int(time.time() * 1000),
                    "line_name": line_name.strip(),
                    "approver1": approver1.strip(),
                    "approver2": approver2.strip(),
                    "approver3": approver3.strip(),
                })
                st.session_state["_show_approval_line_add_form"] = False
                st.success("결재선이 등록되었습니다.")
                st.rerun()
        if cancelled:
            st.session_state["_show_approval_line_add_form"] = False
            st.rerun()

    # 결재선 삭제
    elif st.session_state.get("_show_approval_line_delete_select"):
        if approval_lines:
            delete_options = {line.get("line_name"): line.get("id") for line in approval_lines}
            target_label = st.selectbox("삭제 대상 결재선", list(delete_options.keys()), key="approval_line_delete_target")
            target_id = delete_options[target_label]

            col_confirm, col_cancel = st.columns(2)
            if col_confirm.button("삭제 확인", type="primary", use_container_width=True):
                st.session_state.approval_lines = [line for line in approval_lines if line.get("id") != target_id]
                st.session_state["_show_approval_line_delete_select"] = False
                st.success("결재선이 삭제되었습니다.")
                st.rerun()
            if col_cancel.button("취소", use_container_width=True):
                st.session_state["_show_approval_line_delete_select"] = False
                st.rerun()

    # 버튼들
    else:
        col_add, col_delete, _ = st.columns([20, 20, 160])
        if col_add.button("+ 결재선 추가", key="show_approval_line_add_form_btn", use_container_width=True):
            st.session_state["_show_approval_line_add_form"] = True
            st.rerun()

        if col_delete.button("삭제", key="show_approval_line_delete_btn", use_container_width=True, disabled=not approval_lines):
            st.session_state["_show_approval_line_delete_select"] = True
            st.rerun()


def show_company_profile():
    st.markdown("### 회사 관리")
    st.caption("당사(우리 회사)의 기본 정보를 관리합니다.")

    profile = _load_company_profile()
    profile.setdefault("withdrawal_accounts", [])

    # 등록된 회사 정보 표시
    if profile.get("name"):
        st.markdown("#### 등록된 회사 정보")
        info_data = {
            "회사명": profile.get("name", "-"),
            "사업자번호": profile.get("business_number", "-"),
            "대표자": profile.get("ceo_name", "-"),
            "연락처": profile.get("contact", "-"),
            "주소": profile.get("address", "-"),
            "비고": profile.get("memo", "-"),
        }
        info_df = pd.DataFrame([info_data]).T.reset_index()
        info_df.columns = ["항목", "내용"]
        render_plain_html_table(info_df, center_align=True)
        if profile.get("updated_at"):
            st.caption(f"최종 수정: {profile.get('updated_at')}")
    else:
        st.info("등록된 회사 정보가 없습니다.")

    st.markdown("#### 출금계좌")
    withdrawal_accounts = profile.get("withdrawal_accounts", [])
    withdrawal_accounts.sort(key=lambda row: (not row.get("is_main", False), row.get("created_at", "")))

    if withdrawal_accounts:
        account_df = pd.DataFrame(withdrawal_accounts)
        for col in ["account_name", "bank_name", "account_number", "holder_name", "is_main", "memo"]:
            if col not in account_df.columns:
                account_df[col] = "" if col != "is_main" else False
        account_df["main_label"] = account_df["is_main"].apply(lambda value: "메인" if value else "")
        account_view = account_df[
            ["main_label", "account_name", "bank_name", "account_number", "holder_name", "memo"]
        ].rename(
            columns={
                "main_label": "메인 출금계좌",
                "account_name": "계좌명",
                "bank_name": "은행명",
                "account_number": "계좌번호",
                "holder_name": "예금주",
                "memo": "비고",
            }
        ).reset_index(drop=True)
        account_view.insert(0, "순번", range(1, len(account_view) + 1))
        render_plain_html_table(account_view, center_align=True)
    else:
        st.info("등록된 출금계좌가 없습니다.")

    if st.session_state.get("_show_withdrawal_add_form"):
        with st.form("withdrawal_account_add_form", clear_on_submit=True):
            col_a, col_b = st.columns(2)
            account_name = col_a.text_input("계좌명", placeholder="예: 본사 운영계좌")
            bank_name = col_b.text_input("은행명", placeholder="예: 국민은행")
            col_c, col_d = st.columns(2)
            account_number = col_c.text_input("계좌번호")
            holder_name = col_d.text_input("예금주", value=profile.get("name", ""))
            is_main = st.checkbox("메인 출금계좌로 설정", value=not bool(withdrawal_accounts))
            memo = st.text_input("비고")
            col_submit, col_cancel = st.columns(2)
            submitted = col_submit.form_submit_button("등록", type="primary", use_container_width=True)
            cancelled = col_cancel.form_submit_button("취소", use_container_width=True)

        if submitted:
            if not account_name.strip() or not bank_name.strip() or not account_number.strip():
                st.warning("계좌명, 은행명, 계좌번호를 입력해주세요.")
            else:
                if is_main:
                    for row in withdrawal_accounts:
                        row["is_main"] = False
                withdrawal_accounts.append(
                    {
                        "id": int(time.time() * 1000),
                        "account_name": account_name.strip(),
                        "bank_name": bank_name.strip(),
                        "account_number": account_number.strip(),
                        "holder_name": holder_name.strip(),
                        "is_main": bool(is_main),
                        "memo": memo.strip(),
                        "created_at": _current_kst().strftime("%Y-%m-%d %H:%M:%S"),
                    }
                )
                profile["withdrawal_accounts"] = withdrawal_accounts
                profile["updated_at"] = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
                _save_company_profile(profile)
                st.session_state["_show_withdrawal_add_form"] = False
                st.success("출금계좌가 등록되었습니다.")
                st.rerun()
        if cancelled:
            st.session_state["_show_withdrawal_add_form"] = False
            st.rerun()

    elif st.session_state.get("_show_withdrawal_manage_form"):
        if withdrawal_accounts:
            account_options = {
                f"{'[메인] ' if row.get('is_main') else ''}{row.get('account_name', '')} ({row.get('bank_name', '')} {row.get('account_number', '')})": row.get("id")
                for row in withdrawal_accounts
            }
            selected_label = st.selectbox("대상 출금계좌", list(account_options.keys()), key="withdrawal_manage_target")
            selected_id = account_options[selected_label]
            target = next((row for row in withdrawal_accounts if row.get("id") == selected_id), None)

            col_main, col_delete, col_cancel = st.columns(3)
            if col_main.button("메인 출금계좌 설정", type="primary", use_container_width=True):
                for row in withdrawal_accounts:
                    row["is_main"] = row.get("id") == selected_id
                profile["withdrawal_accounts"] = withdrawal_accounts
                profile["updated_at"] = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
                _save_company_profile(profile)
                st.session_state["_show_withdrawal_manage_form"] = False
                st.success("메인 출금계좌가 설정되었습니다.")
                st.rerun()
            if col_delete.button("삭제", use_container_width=True, disabled=not target):
                profile["withdrawal_accounts"] = [row for row in withdrawal_accounts if row.get("id") != selected_id]
                if profile["withdrawal_accounts"] and not any(row.get("is_main") for row in profile["withdrawal_accounts"]):
                    profile["withdrawal_accounts"][0]["is_main"] = True
                profile["updated_at"] = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
                _save_company_profile(profile)
                st.session_state["_show_withdrawal_manage_form"] = False
                st.success("출금계좌가 삭제되었습니다.")
                st.rerun()
            if col_cancel.button("취소", use_container_width=True):
                st.session_state["_show_withdrawal_manage_form"] = False
                st.rerun()

    else:
        col_add, col_manage, _ = st.columns([0.18, 0.22, 0.6])
        if col_add.button("+ 출금계좌 등록", key="show_withdrawal_add_btn", use_container_width=True):
            st.session_state["_show_withdrawal_add_form"] = True
            st.rerun()
        if col_manage.button("메인/삭제 관리", key="show_withdrawal_manage_btn", use_container_width=True, disabled=not withdrawal_accounts):
            st.session_state["_show_withdrawal_manage_form"] = True
            st.rerun()

    # 수정 폼
    if st.session_state.get("_show_company_edit_form"):
        with st.form("company_profile_form"):
            col_a, col_b = st.columns(2)
            name = col_a.text_input("회사명", value=profile.get("name", ""), placeholder="예: (주)회사명")
            business_number = col_b.text_input("사업자번호", value=profile.get("business_number", ""), placeholder="예: 123-45-67890")
            col_c, col_d = st.columns(2)
            ceo_name = col_c.text_input("대표자", value=profile.get("ceo_name", ""), placeholder="예: 홍길동")
            contact = col_d.text_input("연락처", value=profile.get("contact", ""), placeholder="예: 02-1234-5678")
            address = st.text_input("주소", value=profile.get("address", ""), placeholder="예: 서울특별시 강남구 ...")
            memo = st.text_area("비고", value=profile.get("memo", ""), placeholder="회사 관련 특이사항")
            col_submit, col_cancel = st.columns(2)
            submitted = col_submit.form_submit_button("저장", type="primary", use_container_width=True)
            cancelled = col_cancel.form_submit_button("취소", use_container_width=True)

        if submitted:
            if not name.strip():
                st.warning("회사명을 입력해주세요.")
            else:
                profile.update(
                    {
                        "name": name.strip(),
                        "business_number": business_number.strip(),
                        "ceo_name": ceo_name.strip(),
                        "address": address.strip(),
                        "contact": contact.strip(),
                        "memo": memo.strip(),
                        "updated_at": _current_kst().strftime("%Y-%m-%d %H:%M:%S"),
                    }
                )
                _save_company_profile(profile)
                st.session_state["_show_company_edit_form"] = False
                st.success("회사 정보가 저장되었습니다.")
                st.rerun()
        if cancelled:
            st.session_state["_show_company_edit_form"] = False
            st.rerun()

    # 수정 버튼
    else:
        if st.button("수정", key="show_company_edit_btn", use_container_width=False):
            st.session_state["_show_company_edit_form"] = True
            st.rerun()


def show_workplace_admin():
    tab1, tab2, tab3 = st.tabs(["사업장 정보 관리", "계좌 관리", "사용자 관리"])
    with tab1:
        _render_workplace_info_admin()
    with tab2:
        _render_bank_account_management()
    with tab3:
        _render_staff_admin()


def _render_bank_account_management():
    st.caption("위탁 사업장의 계좌 정보를 관리합니다.")

    data = _load_bank_accounts()
    accounts = data.get("accounts", [])

    wp_data = _load_delegated_workplaces()
    workplaces = wp_data.get("workplaces", [])
    workplaces_by_id = {site.get("id"): site for site in workplaces}

    # 등록된 계좌 표시
    if accounts:
        st.markdown("#### 등록된 계좌")
        account_df = pd.DataFrame(accounts)
        account_df["balance_won"] = account_df["balance"].apply(_format_won)
        for col in ["holder_name", "balance_updated_at", "memo"]:
            if col not in account_df.columns:
                account_df[col] = ""
        account_df["회사"] = [_account_company_label(row, workplaces_by_id) for row in accounts]
        account_view = account_df[
            ["회사", "account_name", "bank_name", "account_number", "holder_name", "balance_won", "balance_updated_at", "memo"]
        ].rename(
            columns={
                "account_name": "계좌명",
                "bank_name": "은행명",
                "account_number": "계좌번호",
                "holder_name": "예금주",
                "balance_won": "현재잔고",
                "balance_updated_at": "최종업데이트",
                "memo": "메모",
            }
        )
        render_plain_html_table(account_view, center_align=True)
    else:
        st.info("등록된 계좌가 없습니다.")

    # 계좌 수정 폼
    if st.session_state.get("_show_account_edit_form"):
        edit_target_id = st.session_state.get("_account_edit_target_id")
        target = next((row for row in accounts if row.get("id") == edit_target_id), None)
        if target is None:
            st.session_state["_show_account_edit_form"] = False
            st.rerun()

        company_options = {site.get("workplace_name", ""): site for site in workplaces}
        current_workplace_id = target.get("linked_workplace_id")
        current_workplace = workplaces_by_id.get(current_workplace_id, {})
        current_workplace_name = current_workplace.get("workplace_name", "")

        with st.form("bank_account_edit_form", clear_on_submit=False):
            company_label = st.selectbox("회사 선택", list(company_options.keys()),
                                        index=list(company_options.keys()).index(current_workplace_name) if current_workplace_name in company_options.keys() else 0)
            account_name = st.text_input("계좌명", value=target.get("account_name", ""))
            bank_name = st.text_input("은행명", value=target.get("bank_name", ""))
            account_number = st.text_input("계좌번호", value=target.get("account_number", ""))
            holder_name = st.text_input("예금주", value=target.get("holder_name", ""))
            memo = st.text_area("메모", value=target.get("memo", ""))
            col_submit, col_cancel = st.columns(2)
            submitted = col_submit.form_submit_button("수정 완료", type="primary", use_container_width=True)
            cancelled = col_cancel.form_submit_button("취소", use_container_width=True)

        if submitted:
            if not account_name.strip() or not account_number.strip():
                st.warning("계좌명과 계좌번호를 입력해주세요.")
            else:
                site = company_options[company_label]
                target["account_name"] = account_name.strip()
                target["bank_name"] = bank_name.strip()
                target["account_number"] = account_number.strip()
                target["holder_name"] = holder_name.strip()
                target["linked_workplace_id"] = site.get("id")
                target["memo"] = memo.strip()
                _save_bank_accounts(data)
                st.session_state["_show_account_edit_form"] = False
                st.success("계좌 정보가 수정되었습니다.")
                st.rerun()
        if cancelled:
            st.session_state["_show_account_edit_form"] = False
            st.rerun()

    # 계좌 추가 폼
    elif st.session_state.get("_show_account_add_form"):
        if not workplaces:
            st.warning("[사업장 정보 관리]에서 사업장을 먼저 등록해주세요.")
            st.session_state["_show_account_add_form"] = False
            st.rerun()
        else:
            company_options = {site.get("workplace_name", ""): site for site in workplaces}
            with st.form("bank_account_form", clear_on_submit=True):
                col_a, col_b = st.columns(2)
                company_label = col_a.selectbox("회사 선택", list(company_options.keys()))
                account_name = col_b.text_input("계좌명", placeholder="예: 본사 지급계좌")
                col_c, col_d, col_e = st.columns([1, 1, 0.4])
                bank_name = col_c.text_input("은행명", placeholder="예: 하나은행")
                account_number = col_d.text_input("계좌번호", placeholder="예: 123-456789-01234")
                col_e.markdown("<div style='height:1.8em'></div>", unsafe_allow_html=True)
                holder_lookup = col_e.form_submit_button("예금주조회", use_container_width=True)
                holder_name = st.text_input("예금주", placeholder="예: (주)회사명")
                memo = st.text_area("메모", placeholder="계좌 관련 특이사항")
                col_submit, col_cancel = st.columns(2)
                submitted = col_submit.form_submit_button("등록", type="primary", use_container_width=True)
                cancelled = col_cancel.form_submit_button("취소", use_container_width=True)

            if holder_lookup:
                st.info("예금주조회 기능은 준비 중입니다. (API 연동 예정)")
            if submitted:
                if not account_name.strip() or not account_number.strip():
                    st.warning("계좌명과 계좌번호를 입력해주세요.")
                else:
                    site = company_options[company_label]
                    accounts.append(
                        {
                            "id": int(time.time() * 1000),
                            "account_type": "위탁사업장",
                            "account_name": account_name.strip(),
                            "bank_name": bank_name.strip(),
                            "account_number": account_number.strip(),
                            "holder_name": holder_name.strip(),
                            "linked_workplace_id": site.get("id"),
                            "balance": 0,
                            "balance_updated_at": "",
                            "balance_history": [],
                            "memo": memo.strip(),
                            "created_at": _current_kst().strftime("%Y-%m-%d %H:%M:%S"),
                        }
                    )
                    _save_bank_accounts(data)
                    st.session_state["_show_account_add_form"] = False
                    st.success("계좌가 등록되었습니다.")
                    st.rerun()
            if cancelled:
                st.session_state["_show_account_add_form"] = False
                st.rerun()

    # 삭제 선택
    elif st.session_state.get("_show_account_delete_select"):
        if accounts:
            delete_options = {f"{row.get('account_name')} ({row.get('bank_name')} {row.get('account_number')})": row.get("id") for row in accounts}
            target_label = st.selectbox("삭제 대상 계좌", list(delete_options.keys()), key="bank_account_delete_target")
            target_id = delete_options[target_label]

            col_confirm, col_cancel = st.columns(2)
            if col_confirm.button("삭제 확인", type="primary", use_container_width=True):
                data["accounts"] = [row for row in accounts if row.get("id") != target_id]
                _save_bank_accounts(data)
                st.session_state["_show_account_delete_select"] = False
                st.success("계좌가 삭제되었습니다.")
                st.rerun()
            if col_cancel.button("취소", use_container_width=True):
                st.session_state["_show_account_delete_select"] = False
                st.rerun()

    # 수정 선택
    elif st.session_state.get("_show_account_edit_select"):
        if accounts:
            edit_options = {f"{row.get('account_name')} ({row.get('bank_name')} {row.get('account_number')})": row.get("id") for row in accounts}
            target_label = st.selectbox("수정 대상 계좌", list(edit_options.keys()), key="bank_account_edit_target")
            target_id = edit_options[target_label]

            col_confirm, col_cancel = st.columns(2)
            if col_confirm.button("수정 진행", type="primary", use_container_width=True):
                st.session_state["_show_account_edit_select"] = False
                st.session_state["_show_account_edit_form"] = True
                st.session_state["_account_edit_target_id"] = target_id
                st.rerun()
            if col_cancel.button("취소", use_container_width=True):
                st.session_state["_show_account_edit_select"] = False
                st.rerun()

    # 버튼들
    else:
        col_add, col_edit, col_delete, _ = st.columns([20, 20, 20, 140])
        if col_add.button("+ 계좌 추가", key="show_account_add_form_btn", use_container_width=True):
            st.session_state["_show_account_add_form"] = True
            st.rerun()

        if col_edit.button("수정", key="show_account_edit_form_btn", use_container_width=True, disabled=not accounts):
            st.session_state["_show_account_edit_select"] = True
            st.rerun()

        if col_delete.button("삭제", key="show_account_delete_btn", use_container_width=True, disabled=not accounts):
            st.session_state["_show_account_delete_select"] = True
            st.rerun()


def _render_transfer_confirmed_section(confirmed_rows, accounts_by_workplace_id):
    st.markdown("#### 이체 자료 확정 건")
    if not confirmed_rows:
        st.info("이체 자료 확정 건이 없습니다.")
        return
    st.caption("이체 자료가 확정되어 [지급 결과 확인]에서 처리 대기 중인 건입니다.")
    confirmed_view = pd.DataFrame(
        [
            {
                "사업장명": row.get("workplace_name", ""),
                "입금은행": accounts_by_workplace_id.get(row.get("workplace_id"), {}).get("bank_name", ""),
                "입금계좌번호": accounts_by_workplace_id.get(row.get("workplace_id"), {}).get("account_number", ""),
                "입금액": _format_won(row.get("request_amount")),
                "이체자료확정일시": row.get("transfer_file_generated_at", ""),
            }
            for row in sorted(confirmed_rows, key=lambda item: item.get("transfer_file_generated_at", ""), reverse=True)
        ]
    )
    render_plain_html_table(confirmed_view, center_align=True, stretch=True, max_width="760px", border=False)


def show_transfer_file_generation():
    _render_page_chrome(["홈", "지급관리", "이체 자료 확정", "이체 자료 확정"])
    st.markdown("### 이체 자료 확정")
    st.caption("품의가 확정된 전도금 요청을 선택하여 이체 자료를 확정하고 엑셀로 다운로드합니다.")

    # 다크모드 스타일
    st.markdown(
        """
        <style>
        /* 이체 대상 선택 data_editor 다크모드 */
        body:has(#pms-d:checked) [data-testid="stDataEditor"],
        body:has(#pms-d:checked) [data-testid="stDataEditor"] > div,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="grid"] {
            background-color: #1e1e2e !important;
            color: #ffffff !important;
            border-color: #45475a !important;
        }
        body:has(#pms-d:checked) [data-testid="stDataEditor"] th,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="columnheader"] {
            background-color: #252535 !important;
            color: #ffffff !important;
            border-color: #45475a !important;
        }
        body:has(#pms-d:checked) [data-testid="stDataEditor"] td,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="gridcell"] {
            background-color: #1e1e2e !important;
            color: #ffffff !important;
            border-color: #313244 !important;
        }
        body:has(#pms-d:checked) [data-testid="stDataEditor"] input {
            background-color: #252535 !important;
            color: #ffffff !important;
            border-color: #45475a !important;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

    wp_data = _load_delegated_workplaces()
    requests = wp_data.get("requests", [])

    bank_data = _load_bank_accounts()
    accounts_by_workplace_id = {
        row.get("linked_workplace_id"): row
        for row in bank_data.get("accounts", [])
        if row.get("linked_workplace_id")
    }

    company_profile = _load_company_profile()
    withdrawal_accounts = company_profile.get("withdrawal_accounts", [])
    withdrawal_accounts = sorted(
        withdrawal_accounts,
        key=lambda row: (not row.get("is_main", False), row.get("created_at", "")),
    )

    targets = [row for row in requests if row.get("status") == "품의 확정"]
    confirmed_rows = [row for row in requests if row.get("status") == "이체 대상"]

    if not targets:
        st.info("이체 자료를 생성할 품의 확정 건이 없습니다.")
        _render_transfer_confirmed_section(confirmed_rows, accounts_by_workplace_id)
        return

    if not withdrawal_accounts:
        st.warning("[회사 관리]에서 출금계좌를 먼저 등록해주세요.")
        _render_transfer_confirmed_section(confirmed_rows, accounts_by_workplace_id)
        return

    source_options = {
        f"{'[메인] ' if row.get('is_main') else ''}{row.get('account_name')} ({row.get('bank_name')} {row.get('account_number')})": row
        for row in withdrawal_accounts
    }
    source_label = st.selectbox("출금 계좌", list(source_options.keys()))

    st.markdown("#### 이체 대상 선택")
    lookup_map = st.session_state.get("_transfer_holder_lookup", {})
    target_rows = []
    lookup_inputs = {}
    for row in sorted(targets, key=lambda item: item.get("requested_at", "")):
        account = accounts_by_workplace_id.get(row.get("workplace_id"), {})
        workplace_name = row.get("workplace_name", "")
        passbook_label = f"{workplace_name} 전도금"
        expected_holder = account.get("holder_name", "") or workplace_name
        lookup_inputs[row.get("id")] = (account.get("bank_name", ""), account.get("account_number", ""), expected_holder)
        target_rows.append(
            {
                "id": row.get("id"),
                "사업장명": workplace_name,
                "선택": True,
                "입금은행": account.get("bank_name", ""),
                "입금계좌번호": account.get("account_number", ""),
                "입금액": _format_won(row.get("request_amount")),
                "예상예금주": expected_holder,
                "조회한예금주": lookup_map.get(row.get("id"), ""),
                "입금통장표시": passbook_label,
                "출금통장표시": passbook_label,
            }
        )

    target_df = pd.DataFrame(target_rows)
    column_order = [
        "id", "선택", "사업장명", "입금은행", "입금계좌번호", "입금액",
        "예상예금주", "조회한예금주", "입금통장표시", "출금통장표시",
    ]
    edited_df = st.data_editor(
        target_df[column_order],
        column_config={
            "id": None,
            "선택": st.column_config.CheckboxColumn("선택", width="small"),
            "사업장명": st.column_config.TextColumn("사업장명", disabled=True, width="small"),
            "입금은행": st.column_config.TextColumn("입금은행", disabled=True, width="small"),
            "입금계좌번호": st.column_config.TextColumn("입금계좌번호", disabled=True, width="small"),
            "입금액": st.column_config.TextColumn("입금액", disabled=True, width="small"),
            "예상예금주": st.column_config.TextColumn("예상예금주", disabled=True, width="small"),
            "조회한예금주": st.column_config.TextColumn("조회한예금주", disabled=True, width="small"),
            "입금통장표시": st.column_config.TextColumn("입금통장표시", width="small"),
            "출금통장표시": st.column_config.TextColumn("출금통장표시", width="small"),
        },
        hide_index=True,
        use_container_width=True,
        key="transfer_target_editor",
    )
    selected_ids = edited_df[edited_df["선택"]]["id"].tolist()

    col_spacer, col_holder_lookup = st.columns([165, 35])
    if col_holder_lookup.button("예금주조회", key="transfer_holder_lookup_btn", use_container_width=True):
        new_lookup = {}
        for row_id, (bank_name, account_number, expected_holder) in lookup_inputs.items():
            holder = _lookup_account_holder_name(bank_name, account_number, expected_holder)
            if holder:
                new_lookup[row_id] = holder
        st.session_state["_transfer_holder_lookup"] = new_lookup
        if new_lookup:
            st.success(f"{len(new_lookup)}건의 예금주 정보를 조회했습니다.")
        else:
            st.warning("조회 가능한 입금계좌 정보가 없습니다. [계좌 관리]에서 계좌를 먼저 등록해주세요.")
        st.rerun()

    lookup_map = st.session_state.get("_transfer_holder_lookup", {})
    lookup_complete = bool(selected_ids) and all(row_id in lookup_map for row_id in selected_ids)

    source = source_options[source_label]
    edited_by_id = {item["id"]: item for item in edited_df.to_dict("records")}
    transfer_rows = []
    for row in targets:
        if row.get("id") not in selected_ids:
            continue
        account = accounts_by_workplace_id.get(row.get("workplace_id"), {})
        edited_row = edited_by_id.get(row.get("id"), {})
        transfer_rows.append(
            {
                "출금은행": source.get("bank_name", ""),
                "출금계좌번호": source.get("account_number", ""),
                "출금통장표시": edited_row.get("출금통장표시", ""),
                "입금은행": account.get("bank_name", ""),
                "입금계좌번호": account.get("account_number", ""),
                "입금액": row.get("request_amount", 0),
                "예상예금주": lookup_map.get(row.get("id")) or account.get("holder_name", "") or row.get("workplace_name", ""),
                "입금통장표시": edited_row.get("입금통장표시", ""),
                "사업장명": row.get("workplace_name", ""),
            }
        )

    col_confirm, col_download = st.columns(2)
    confirm_clicked = col_confirm.button(
        "이체 자료 확정", type="primary", disabled=not lookup_complete, use_container_width=True
    )
    excel_bytes = dataframe_to_excel_bytes({"이체자료": pd.DataFrame(transfer_rows)}) if transfer_rows else b""
    col_download.download_button(
        "📥 엑셀 다운로드",
        data=excel_bytes,
        file_name=f"이체자료_{_current_kst().strftime('%Y%m%d_%H%M%S')}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True,
        disabled=not transfer_rows,
    )

    if confirm_clicked:
        now_str = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
        for row in targets:
            if row.get("id") in selected_ids:
                account = accounts_by_workplace_id.get(row.get("workplace_id"), {})
                edited_row = edited_by_id.get(row.get("id"), {})
                row["status"] = "이체 대상"
                row["transfer_file_generated_at"] = now_str
                row["transfer_expected_holder"] = (
                    lookup_map.get(row.get("id")) or account.get("holder_name", "") or row.get("workplace_name", "")
                )
                row["transfer_verified_holder"] = lookup_map.get(row.get("id"), "")
                row["transfer_deposit_passbook"] = edited_row.get("입금통장표시", "")
                row["transfer_withdrawal_passbook"] = edited_row.get("출금통장표시", "")
                row["transfer_bank_name"] = account.get("bank_name", "")
                row["transfer_account_number"] = account.get("account_number", "")
                row["transfer_amount"] = row.get("request_amount", 0)
                row["deposit_reconciliation"] = ""
                row["deposit_reconciled_at"] = ""

        st.session_state["_transfer_holder_lookup"] = {
            row_id: name for row_id, name in lookup_map.items() if row_id not in selected_ids
        }
        _save_delegated_workplaces(wp_data)
        st.success(f"{len(transfer_rows)}건이 이체 대상으로 확정되었습니다. [지급 결과 확인]에서 처리 결과를 확정해주세요.")
        st.rerun()

    _render_transfer_confirmed_section(confirmed_rows, accounts_by_workplace_id)


def show_transfer_result_confirmation():
    _render_page_chrome(["홈", "지급관리", "지급 결과 확인", "지급 결과 확인"])
    st.markdown("### 지급 결과 확인")
    st.caption("이체 자료가 확정된 건의 이체 완료 여부를 확인하고, 완료 이력을 관리합니다.")

    wp_data = _load_delegated_workplaces()
    requests = wp_data.get("requests", [])

    bank_data = _load_bank_accounts()
    accounts_by_workplace_id = {
        row.get("linked_workplace_id"): row
        for row in bank_data.get("accounts", [])
        if row.get("linked_workplace_id")
    }

    pending = [row for row in requests if row.get("status") == "이체 대상"]

    title_col, action_col = st.columns([0.78, 0.22])
    with title_col:
        st.markdown("#### 이체 대상")
    with action_col:
        check_clicked = st.button(
            "이체 결과 확인",
            type="primary",
            disabled=not pending,
            use_container_width=True,
            key="transfer_deposit_reconcile",
        )

    if pending:
        if check_clicked:
            matched_count = _reconcile_transfer_deposits(pending, bank_data.get("accounts", []))
            _save_delegated_workplaces(wp_data)
            st.success(f"입금 대사 완료: {matched_count}건 일치, {len(pending) - matched_count}건 미확인")
            st.rerun()

        pending_sorted = sorted(pending, key=lambda item: item.get("transfer_file_generated_at", ""), reverse=True)
        pending_rows = []
        for row in pending_sorted:
            account = accounts_by_workplace_id.get(row.get("workplace_id"), {})
            bank_name = row.get("transfer_bank_name") or account.get("bank_name", "")
            account_number = row.get("transfer_account_number") or account.get("account_number", "")
            pending_rows.append(
                {
                    "id": row.get("id"),
                    "선택": False,
                    "사업장명": row.get("workplace_name", ""),
                    "입금은행": bank_name,
                    "입금계좌번호": account_number,
                    "입금액": _format_won(row.get("request_amount")),
                    "예상예금주": row.get("transfer_expected_holder", "") or account.get("holder_name", "") or row.get("workplace_name", ""),
                    "조회한예금주": row.get("transfer_verified_holder", ""),
                    "입금통장표시": row.get("transfer_deposit_passbook", ""),
                    "출금통장표시": row.get("transfer_withdrawal_passbook", ""),
                    "입금 대사": row.get("deposit_reconciliation", ""),
                }
            )

        pending_df = pd.DataFrame(pending_rows)
        column_order = [
            "id", "선택", "사업장명", "입금은행", "입금계좌번호", "입금액",
            "예상예금주", "조회한예금주", "입금통장표시", "출금통장표시", "입금 대사",
        ]
        st.markdown(
            """
            <style>
            [data-testid="stDataFrame"] [role="columnheader"],
            [data-testid="stDataFrame"] [role="gridcell"],
            [data-testid="stDataEditor"] [role="columnheader"],
            [data-testid="stDataEditor"] [role="gridcell"] {
                text-align: center !important;
                justify-content: center !important;
            }
            </style>
            """,
            unsafe_allow_html=True,
        )
        edited_confirm = st.data_editor(
            pending_df[column_order],
            column_config={
                "id": None,
                "선택": st.column_config.CheckboxColumn("선택", width="small"),
                "사업장명": st.column_config.TextColumn("사업장명", disabled=True, width="small"),
                "입금은행": st.column_config.TextColumn("입금은행", disabled=True, width="small"),
                "입금계좌번호": st.column_config.TextColumn("입금계좌번호", disabled=True, width="small"),
                "입금액": st.column_config.TextColumn("입금액", disabled=True, width="small"),
                "예상예금주": st.column_config.TextColumn("예상예금주", disabled=True, width="small"),
                "조회한예금주": st.column_config.TextColumn("조회한예금주", disabled=True, width="small"),
                "입금통장표시": st.column_config.TextColumn("입금통장표시", disabled=True, width="small"),
                "출금통장표시": st.column_config.TextColumn("출금통장표시", disabled=True, width="small"),
                "입금 대사": st.column_config.TextColumn(
                    "입금 대사",
                    help="계좌번호와 입금액이 일치하는 입금 내역이 확인된 날짜입니다.",
                    disabled=True,
                    width="small",
                ),
            },
            hide_index=True,
            use_container_width=True,
            key="transfer_confirm_editor",
        )
        confirm_ids = edited_confirm[edited_confirm["선택"]]["id"].tolist()
        st.caption("입금 대사는 이체 대상의 입금계좌번호와 입금액이 실제 입금 내역과 일치하는지 확인한 결과입니다. 일치하면 입금일이 표시되고, 없으면 미확인으로 표시됩니다.")
        if st.button("선택 건 이체 완료 확인", type="primary", disabled=not confirm_ids):
            now_str = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
            for row in pending:
                if row.get("id") in confirm_ids:
                    row["status"] = "이체 완료"
                    row["paid_at"] = now_str
            _save_delegated_workplaces(wp_data)
            st.success(f"{len(confirm_ids)}건을 이체 완료로 처리했습니다.")
            st.rerun()
    else:
        st.info("이체 확정 대기 중인 건이 없습니다.")

    st.markdown("#### 이체 완료 이력")
    done = [row for row in requests if row.get("status") == "이체 완료"]
    if done:
        done_df = pd.DataFrame(done)
        done_df["request_amount_won"] = done_df["request_amount"].apply(_format_won)
        if "paid_at" not in done_df.columns:
            done_df["paid_at"] = ""
        done_view = done_df[
            ["paid_at", "workplace_name", "request_amount_won", "requested_by", "requested_at"]
        ].rename(
            columns={
                "paid_at": "이체완료일시",
                "workplace_name": "사업장명",
                "request_amount_won": "이체금액",
                "requested_by": "요청자",
                "requested_at": "요청일시",
            }
        )
        render_plain_html_table(
            done_view.sort_values("이체완료일시", ascending=False),
            center_align=True,
            stretch=True,
            max_width="760px",
            border=False,
        )
    else:
        st.info("이체 완료 이력이 없습니다.")


def show_account_balance_check():
    _render_page_chrome(["홈", "조회", "계좌 잔고 확인", "계좌 잔고 확인"])
    st.markdown("### 계좌 잔고 확인")
    st.caption("위탁 사업장 계좌의 잔고를 수동으로 관리하고 변동 추이를 확인합니다.")

    if not st.session_state.get("_balance_update_form_initialized"):
        st.session_state.show_balance_update_form = False
        st.session_state["_balance_update_form_initialized"] = True

    data = _load_bank_accounts()
    accounts = data.get("accounts", [])

    wp_data = _load_delegated_workplaces()
    workplaces = wp_data.get("workplaces", [])
    workplaces_by_id = {site.get("id"): site for site in workplaces}

    if st.session_state.user_role != "관리자":
        my_ids = {row.get("id") for row in _my_workplaces(workplaces)}
        accounts = [row for row in accounts if row.get("linked_workplace_id") in my_ids]

    if not accounts:
        st.info("등록된 계좌가 없습니다. [계좌 관리]에서 계좌를 먼저 등록해주세요.")
        return

    account_df = pd.DataFrame(accounts)
    account_df["balance_won"] = account_df["balance"].apply(_format_won)
    if "balance_updated_at" not in account_df.columns:
        account_df["balance_updated_at"] = ""
    account_df["회사"] = [_account_company_label(row, workplaces_by_id) for row in accounts]
    account_view = account_df[
        ["회사", "account_name", "bank_name", "account_number", "balance_won", "balance_updated_at"]
    ].rename(
        columns={
            "account_name": "계좌명",
            "bank_name": "은행명",
            "account_number": "계좌번호",
            "balance_won": "현재잔고",
            "balance_updated_at": "최종업데이트",
        }
    )
    title_col, button_col = st.columns([0.78, 0.22])
    with title_col:
        st.markdown("#### 계좌 리스트")
    with button_col:
        if st.button("잔액 정보 업데이트", key="open_balance_update_form", use_container_width=True, type="primary"):
            st.session_state.show_balance_update_form = True
            st.rerun()

    render_plain_html_table(account_view, center_align=True, stretch=True, max_width="100%", border=False)

    st.markdown("#### 거래내역조회")

    history_options = {f"{row.get('account_name')} ({row.get('bank_name')} {row.get('account_number')})": row for row in accounts}
    today = _current_kst().date()
    if "balance_history_start_date" not in st.session_state:
        st.session_state.balance_history_start_date = today - timedelta(days=6)
    if "balance_history_end_date" not in st.session_state:
        st.session_state.balance_history_end_date = today

    period_presets = {
        "오늘": (today, today),
        "어제": (today - timedelta(days=1), today - timedelta(days=1)),
        "1주일": (today - timedelta(days=6), today),
        "1개월": (today - timedelta(days=30), today),
    }

    st.markdown(
        """
        <style>
        .balance-history-help {
            color:#7a8599;
            font-size:12px;
            margin-top:-6px;
        }
        div[data-testid="stRadio"] > label,
        div[data-testid="stSelectbox"] > label,
        div[data-testid="stDateInput"] > label,
        div[data-testid="stTextInput"] > label {
            font-weight:700 !important;
            color:#0f172a !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    with st.container(border=False):
        account_label_col, account_input_col = st.columns([0.12, 0.88])
        with account_label_col:
            st.markdown("**계좌번호 <span style='color:#008c78'>*</span>**", unsafe_allow_html=True)
        with account_input_col:
            account_select_col, _ = st.columns([0.48, 0.52])
            with account_select_col:
                history_label = st.selectbox(
                    "계좌번호",
                    list(history_options.keys()),
                    key="balance_history_account",
                    label_visibility="collapsed",
                    placeholder="계좌 선택",
                )

        period_label_col, period_input_col = st.columns([0.12, 0.88])
        with period_label_col:
            st.markdown("**조회기간**")
        with period_input_col:
            preset_cols = st.columns([0.13, 0.13, 0.17, 0.17, 0.40])
            for index, (label, date_range) in enumerate(period_presets.items()):
                button_type = "primary" if label == st.session_state.get("balance_history_preset", "1주일") else "secondary"
                if preset_cols[index].button(label, key=f"balance_history_preset_{label}", type=button_type, use_container_width=True):
                    st.session_state.balance_history_preset = label
                    st.session_state.balance_history_start_date = date_range[0]
                    st.session_state.balance_history_end_date = date_range[1]
                    st.rerun()

            date_start_col, tilde_col, date_end_col, month_col, _ = st.columns([0.22, 0.03, 0.22, 0.18, 0.35])
            with date_start_col:
                start_date = st.date_input("시작일", key="balance_history_start_date", label_visibility="collapsed")
            with tilde_col:
                st.markdown("<div style='padding-top:8px;text-align:center;'>~</div>", unsafe_allow_html=True)
            with date_end_col:
                end_date = st.date_input("종료일", key="balance_history_end_date", label_visibility="collapsed")
            with month_col:
                month_options = ["월별 선택"] + [
                    (today.replace(day=1) - pd.DateOffset(months=idx)).strftime("%Y-%m")
                    for idx in range(12)
                ]
                selected_month = st.selectbox(
                    "월별 선택",
                    month_options,
                    key="balance_history_month",
                    label_visibility="collapsed",
                )
            st.markdown("<div class='balance-history-help'>ㆍ 직접입력 예시 : YYYYMMDD</div>", unsafe_allow_html=True)

        content_label_col, content_input_col = st.columns([0.12, 0.88])
        with content_label_col:
            st.markdown("**조회내용**")
        with content_input_col:
            history_type = st.radio(
                "조회내용",
                ["전체(입금+출금)", "입금내역", "출금내역"],
                horizontal=True,
                key="balance_history_type",
                label_visibility="collapsed",
            )

        sort_label_col, sort_input_col = st.columns([0.12, 0.88])
        with sort_label_col:
            st.markdown("**정렬방식**")
        with sort_input_col:
            sort_col, count_col = st.columns([0.22, 0.5])
            with sort_col:
                sort_order = st.selectbox(
                    "정렬방식",
                    ["최근거래먼저", "과거거래먼저"],
                    key="balance_history_sort",
                    label_visibility="collapsed",
                )
            with count_col:
                result_count = st.radio(
                    "조회건수",
                    [15, 30, 50, 100],
                    index=1,
                    horizontal=True,
                    key="balance_history_limit",
                    label_visibility="collapsed",
                    format_func=lambda value: f"{value}건",
                )

        search_label_col, search_input_col = st.columns([0.12, 0.88])
        with search_label_col:
            st.markdown("**검색조건**")
        with search_input_col:
            search_type_col, keyword_col, _ = st.columns([0.22, 0.40, 0.16])
            with search_type_col:
                search_type = st.selectbox(
                    "검색조건",
                    ["적요"],
                    key="balance_history_search_type",
                    label_visibility="collapsed",
                )
            with keyword_col:
                search_keyword = st.text_input(
                    "검색어",
                    key="balance_history_keyword",
                    label_visibility="collapsed",
                    placeholder="적요(통장 메모) 최대 25자까지 입력가능",
                    max_chars=25,
                )

    button_left, button_center, button_right = st.columns([0.42, 0.16, 0.42])
    with button_center:
        history_clicked = st.button("조회", key="balance_history_search", type="primary", use_container_width=True)

    if selected_month != "월별 선택":
        month_start = pd.to_datetime(f"{selected_month}-01").date()
        month_end = (pd.to_datetime(f"{selected_month}-01") + pd.offsets.MonthEnd(0)).date()
        start_date, end_date = month_start, month_end

    if history_clicked:
        st.session_state["_balance_history_target"] = history_label
        st.session_state["_balance_history_query"] = {
            "start_date": start_date,
            "end_date": end_date,
            "history_type": history_type,
            "sort_order": sort_order,
            "result_count": result_count,
            "search_type": search_type,
            "search_keyword": search_keyword,
        }

    history_target_label = st.session_state.get("_balance_history_target")
    history_query = st.session_state.get("_balance_history_query", {})
    if history_target_label and history_target_label in history_options:
        target_account = history_options[history_target_label]
        history_rows = target_account.get("balance_history", [])
        start_date = history_query.get("start_date", start_date)
        end_date = history_query.get("end_date", end_date)
        history_type = history_query.get("history_type", history_type)
        sort_order = history_query.get("sort_order", sort_order)
        result_count = history_query.get("result_count", result_count)
        search_keyword = str(history_query.get("search_keyword", search_keyword) or "").strip()

        enriched_history = []
        previous_balance = None
        for item in sorted(history_rows, key=lambda row: str(row.get("at", ""))):
            amount = _history_amount_value(item, previous_balance)
            if "balance" in item:
                previous_balance = _coerce_int_amount(item.get("balance"))
            enriched = dict(item)
            enriched["_amount"] = amount
            enriched_history.append(enriched)

        filtered_history = []
        for item in enriched_history:
            parsed = _parse_kst_datetime(item.get("at"))
            if parsed and not (start_date <= parsed.date() <= end_date):
                continue
            if history_type == "입금내역" and item.get("_amount", 0) <= 0:
                continue
            if history_type == "출금내역" and item.get("_amount", 0) >= 0:
                continue
            if search_keyword and search_keyword not in str(item.get("memo", "")):
                continue
            filtered_history.append(item)

        reverse_sort = sort_order == "최근거래먼저"
        filtered_history = sorted(filtered_history, key=lambda row: str(row.get("at", "")), reverse=reverse_sort)
        filtered_history = filtered_history[: int(result_count)]

        if filtered_history:
            history_df = pd.DataFrame(filtered_history)
            history_df["거래금액"] = history_df["_amount"].apply(_format_won)
            history_df["잔고"] = history_df["balance"].apply(_format_won)
            if "memo" not in history_df.columns:
                history_df["memo"] = ""
            history_view = history_df[["at", "거래금액", "잔고", "memo"]].rename(
                columns={"at": "거래일시", "memo": "적요"}
            )
            render_plain_html_table(history_view, center_align=True, stretch=True, max_width="760px", border=False)
        else:
            st.info("조건에 맞는 거래내역이 없습니다.")

    if not st.session_state.get("show_balance_update_form"):
        return

    st.markdown("#### 잔액 정보 업데이트")
    account_options = {f"{row.get('account_name')} ({row.get('bank_name')} {row.get('account_number')})": row for row in accounts}
    selected_label = st.selectbox("계좌 선택", list(account_options.keys()), key="balance_target_account")
    selected_account = account_options[selected_label]

    with st.form("balance_update_form", clear_on_submit=True):
        new_balance = st.number_input(
            "현재 잔고", min_value=0, step=10000, format="%d", value=int(selected_account.get("balance", 0) or 0)
        )
        memo = st.text_input("메모", placeholder="예: 정기 확인")
        form_col1, form_col2 = st.columns(2)
        with form_col1:
            submitted = st.form_submit_button("업데이트", type="primary", use_container_width=True)
        with form_col2:
            cancelled = st.form_submit_button("취소", use_container_width=True)

    if submitted:
        now_str = _current_kst().strftime("%Y-%m-%d %H:%M:%S")
        selected_account["balance"] = int(new_balance)
        selected_account["balance_updated_at"] = now_str
        selected_account.setdefault("balance_history", []).append(
            {"at": now_str, "balance": int(new_balance), "memo": memo.strip()}
        )
        _save_bank_accounts(data)
        st.session_state.show_balance_update_form = False
        st.success("잔고가 업데이트되었습니다.")
        st.rerun()

    if cancelled:
        st.session_state.show_balance_update_form = False
        st.rerun()

    history = selected_account.get("balance_history", [])
    if history:
        st.markdown("#### 잔고 변동 추이")
        import plotly.graph_objects as go
        fig = go.Figure(go.Scatter(
            x=[item.get("at") for item in history],
            y=[item.get("balance") for item in history],
            mode="lines+markers",
        ))
        fig.update_layout(**_chart_layout(height=300))
        st.plotly_chart(fig, use_container_width=True, theme=None)


def _render_workplace_info_admin():
    st.caption("위탁 사업장의 기본 정보를 관리합니다.")

    data = _load_delegated_workplaces()
    workplaces = data.get("workplaces", [])

    if workplaces:
        st.markdown("#### 등록된 사업장")
        site_df = pd.DataFrame(workplaces)
        if "business_number" not in site_df.columns:
            site_df["business_number"] = ""
        if "business_alias" not in site_df.columns:
            site_df["business_alias"] = ""
        if "regular_payment_day" not in site_df.columns:
            site_df["regular_payment_day"] = 0
        if "manager_name" not in site_df.columns:
            site_df["manager_name"] = ""
        if "memo" not in site_df.columns:
            site_df["memo"] = ""
        site_df["regular_payment_day_label"] = site_df["regular_payment_day"].apply(
            lambda value: f"매월 {int(float(value or 0))}일" if int(float(value or 0)) else "미등록"
        )
        site_view = site_df[
            ["workplace_name", "business_number", "business_alias", "regular_payment_day_label", "manager_name", "memo"]
        ].rename(
            columns={
                "workplace_name": "사업장명",
                "business_number": "사업자번호",
                "business_alias": "사업자별칭",
                "regular_payment_day_label": "정기지급일자",
                "manager_name": "담당자",
                "memo": "비고",
            }
        ).reset_index(drop=True)
        site_view.insert(0, "순번", range(1, len(site_view) + 1))

        status_map = st.session_state.get("_biz_status_map", {})
        if status_map:
            site_view["휴폐업상태"] = site_df["business_number"].apply(
                lambda v: status_map.get(re.sub(r"\D", "", str(v)), "")
            )

        render_plain_html_table(site_view, center_align=True)

        col_spacer, col_status_btn = st.columns([165, 35])
        if col_status_btn.button("휴폐업조회", key="biz_status_check_btn", use_container_width=True):
            api_key = _get_nts_api_key()
            biz_numbers = site_df["business_number"].tolist()
            if not api_key:
                st.warning("국세청 API 키(NTS_API_KEY)가 설정되지 않았습니다.")
            elif not any(str(b).strip() for b in biz_numbers):
                st.warning("조회할 사업자번호가 없습니다.")
            else:
                result = _check_business_status(api_key, biz_numbers)
                if result:
                    st.session_state["_biz_status_map"] = result
                    st.rerun()
                else:
                    st.warning("휴폐업 조회 결과를 가져오지 못했습니다.")
    else:
        st.info("등록된 사업장이 없습니다.")

    staff_names = sorted({
        info.get("name") for info in _real_users(load_user_db()).values()
        if isinstance(info, dict) and info.get("name")
    })
    manager_options = ["미지정"] + staff_names

    if st.session_state.get("_show_workplace_edit_form"):
        edit_target_id = st.session_state.get("_workplace_edit_target_id")
        target = next((row for row in workplaces if row.get("id") == edit_target_id), None)
        if target is None:
            st.session_state["_show_workplace_edit_form"] = False
            st.rerun()

        current_manager = target.get("manager_name", "")
        manager_index = manager_options.index(current_manager) if current_manager in manager_options else 0

        with st.form("delegated_workplace_edit_form", clear_on_submit=False):
            workplace_name = st.text_input("사업장명", value=target.get("workplace_name", ""))
            business_number = st.text_input("사업자번호", value=target.get("business_number", ""))
            business_alias = st.text_input("사업자별칭", value=target.get("business_alias", ""))
            regular_payment_day = st.number_input(
                "정기지급일자", min_value=0, max_value=31, step=1, value=int(target.get("regular_payment_day", 0) or 0)
            )
            manager_name = st.selectbox("담당자", manager_options, index=manager_index)
            memo = st.text_area("비고", value=target.get("memo", ""))
            col_submit, col_cancel = st.columns(2)
            submitted = col_submit.form_submit_button("수정 완료", type="primary", use_container_width=True)
            cancelled = col_cancel.form_submit_button("취소", use_container_width=True)

        if submitted:
            normalized = _normalize_name(workplace_name)
            duplicate = any(
                row.get("id") != edit_target_id and _normalize_name(row.get("workplace_name")) == normalized
                for row in workplaces
            )
            if not workplace_name.strip():
                st.warning("사업장명을 입력해주세요.")
            elif duplicate:
                st.error("이미 등록된 사업장입니다.")
            else:
                target["workplace_name"] = workplace_name.strip()
                target["business_number"] = business_number.strip()
                target["business_alias"] = business_alias.strip()
                target["regular_payment_day"] = int(regular_payment_day)
                target["manager_name"] = "" if manager_name == "미지정" else manager_name
                target["memo"] = memo.strip()
                _save_delegated_workplaces(data)
                st.session_state["_show_workplace_edit_form"] = False
                st.success("사업장 정보가 수정되었습니다.")
                st.rerun()
        if cancelled:
            st.session_state["_show_workplace_edit_form"] = False
            st.rerun()
    elif st.session_state.get("_show_workplace_add_form"):
        with st.form("delegated_workplace_form", clear_on_submit=True):
            workplace_name = st.text_input("사업장명", placeholder="예: 강남 위탁사업장")
            business_number = st.text_input("사업자번호", placeholder="예: 123-45-67890")
            business_alias = st.text_input("사업자별칭", placeholder="예: 강남점")
            regular_payment_day = st.number_input("정기지급일자", min_value=0, max_value=31, step=1, value=0)
            manager_name = st.selectbox("담당자", manager_options)
            memo = st.text_area("비고", placeholder="사업장 관련 특이사항")
            col_submit, col_cancel = st.columns(2)
            submitted = col_submit.form_submit_button("등록", type="primary", use_container_width=True)
            cancelled = col_cancel.form_submit_button("취소", use_container_width=True)

        if submitted:
            normalized = _normalize_name(workplace_name)
            duplicate = any(_normalize_name(row.get("workplace_name")) == normalized for row in workplaces)
            if not workplace_name.strip():
                st.warning("사업장명을 입력해주세요.")
            elif duplicate:
                st.error("이미 등록된 사업장입니다.")
            else:
                workplaces.append(
                    {
                        "id": int(time.time() * 1000),
                        "workplace_name": workplace_name.strip(),
                        "business_number": business_number.strip(),
                        "business_alias": business_alias.strip(),
                        "bank_name": "",
                        "account_number": "",
                        "regular_payment_day": int(regular_payment_day),
                        "manager_name": "" if manager_name == "미지정" else manager_name,
                        "manager_contact": "",
                        "memo": memo.strip(),
                        "created_at": _current_kst().strftime("%Y-%m-%d %H:%M:%S"),
                    }
                )
                _save_delegated_workplaces(data)
                st.session_state["_show_workplace_add_form"] = False
                st.success("사업장 정보가 등록되었습니다.")
                st.rerun()
        if cancelled:
            st.session_state["_show_workplace_add_form"] = False
            st.rerun()
    elif st.session_state.get("_show_workplace_delete_select"):
        if workplaces:
            site_options = {
                f"{row.get('workplace_name', '')} ({row.get('business_number', '') or '-'})": row.get("id")
                for row in workplaces
            }
            target_label = st.selectbox("삭제 대상 사업장", list(site_options.keys()), key="workplace_delete_target")
            target_id = site_options[target_label]

            col_confirm, col_cancel = st.columns(2)
            if col_confirm.button("삭제 확인", type="primary", use_container_width=True):
                data["workplaces"] = [row for row in workplaces if row.get("id") != target_id]
                _save_delegated_workplaces(data)
                st.session_state["_show_workplace_delete_select"] = False
                st.success("사업장 정보가 삭제되었습니다.")
                st.rerun()
            if col_cancel.button("취소", use_container_width=True):
                st.session_state["_show_workplace_delete_select"] = False
                st.rerun()
    elif st.session_state.get("_show_workplace_edit_select"):
        if workplaces:
            site_options = {
                f"{row.get('workplace_name', '')} ({row.get('business_number', '') or '-'})": row.get("id")
                for row in workplaces
            }
            target_label = st.selectbox("수정 대상 사업장", list(site_options.keys()), key="workplace_modify_target")
            target_id = site_options[target_label]

            col_confirm, col_cancel = st.columns(2)
            if col_confirm.button("수정 진행", type="primary", use_container_width=True):
                st.session_state["_show_workplace_edit_select"] = False
                st.session_state["_show_workplace_edit_form"] = True
                st.session_state["_workplace_edit_target_id"] = target_id
                st.rerun()
            if col_cancel.button("취소", use_container_width=True):
                st.session_state["_show_workplace_edit_select"] = False
                st.rerun()
    else:
        col_add, col_edit, col_delete, _ = st.columns([20, 20, 20, 140])
        if col_add.button("+ 사업장 추가", key="show_workplace_add_form_btn", use_container_width=True):
            st.session_state["_show_workplace_add_form"] = True
            st.rerun()

        if col_edit.button("수정", key="show_workplace_edit_form_btn", use_container_width=True, disabled=not workplaces):
            st.session_state["_show_workplace_edit_select"] = True
            st.rerun()

        if col_delete.button("삭제", key="delete_workplace_btn", use_container_width=True, disabled=not workplaces):
            st.session_state["_show_workplace_delete_select"] = True
            st.rerun()


def _dashboard_ai_messages(metrics, forecast_rows):
    """대시보드 상단에 표시할 AI 인사이트 메시지 목록을 생성."""
    messages = []
    if metrics["pending"]:
        messages.append(("📥", f"결재 대기 중인 전도금 요청이 {metrics['pending']}건 있습니다. [전자결재]에서 확인해주세요."))
    if metrics["approved"]:
        messages.append(("💸", f"품의 확정되어 이체 대상인 요청이 {metrics['approved']}건 있습니다. [이체 자료 확정]에서 처리해주세요."))
    for row in forecast_rows:
        risk = row.get("리스크", "")
        if risk and risk != "정상":
            messages.append(("⚠️", f"{row.get('사업장명', '')}: {risk} (추천 지급액 {row.get('추천 지급액', '')})"))
    if not messages:
        messages.append(("✅", "현재 특이사항 없이 정상적으로 운영되고 있습니다."))
    return messages


def show_dashboard():
    st.markdown("### 대시보드")
    st.caption("위탁 사업장 운영 현황과 전도금 지급 이력, 계좌 잔고 현황을 확인합니다.")

    wp_data = _load_delegated_workplaces()
    workplaces = wp_data.get("workplaces", [])
    requests = wp_data.get("requests", [])
    metrics = _workplace_request_metrics(requests)
    paid_count = len([row for row in requests if row.get("status") == "이체 완료"])
    total_count = max(len(requests), 1)
    forecast_rows = _workplace_forecast_rows(workplaces, requests)

    _workplace_dashboard_css()
    c1, c2, c3, c4 = st.columns(4)
    c1.markdown(_donut_card("Workplaces", f"{len(workplaces):,}", "Sites", min(len(workplaces) * 12, 100)), unsafe_allow_html=True)
    c2.markdown(_donut_card("결재대기", f"{metrics['pending']:,}", "Pending", min(metrics["pending"] * 18, 100)), unsafe_allow_html=True)
    c3.markdown(_donut_card("품의확정", f"{metrics['approved']:,}", "Approved", min(metrics["approved"] * 18, 100)), unsafe_allow_html=True)
    c4.markdown(_donut_card("이체완료", f"{paid_count:,}", "Transfers", int((paid_count / total_count) * 100)), unsafe_allow_html=True)

    ai_messages = _dashboard_ai_messages(metrics, forecast_rows)
    rows_html = "".join(
        f"<div class='suggestion-row'><div class='suggestion-icon'>{html.escape(icon)}</div><div>{html.escape(text)}</div></div>"
        for icon, text in ai_messages
    )
    st.markdown(
        f"<div class='suggestion-panel'><div class='suggestion-title'>🤖 AI 인사이트</div>{rows_html}</div>",
        unsafe_allow_html=True,
    )

    st.markdown("#### 전도금 지급 이력")
    if requests:
        request_df = pd.DataFrame(requests)
        request_df["request_amount_won"] = request_df["request_amount"].apply(_format_won)
        for col in ["approved_at", "paid_at", "reject_reason"]:
            if col not in request_df.columns:
                request_df[col] = ""
        request_view = request_df[
            ["requested_at", "workplace_name", "request_amount_won", "requested_by", "status", "approved_at", "paid_at", "reject_reason"]
        ].rename(
            columns={
                "requested_at": "요청일시",
                "workplace_name": "사업장명",
                "request_amount_won": "요청 금액",
                "requested_by": "요청자",
                "status": "상태",
                "approved_at": "품의 확정일시",
                "paid_at": "이체 완료일시",
                "reject_reason": "반려사유",
            }
        )
        render_plain_html_table(request_view.sort_values("요청일시", ascending=False), center_align=True)
    else:
        st.info("전도금 지급 이력이 없습니다.")

    if forecast_rows:
        st.markdown("#### AI 예측 안내")
        st.caption("현재 MVP는 지급 이력 기반의 규칙형 예측입니다. 충분한 이력이 쌓이면 모델 기반 예측으로 확장할 수 있습니다.")
        render_plain_html_table(pd.DataFrame(forecast_rows), center_align=True)
    elif not workplaces:
        st.info("예측을 위해 먼저 사업장 정보를 등록해주세요.")

    st.markdown("#### 계좌별 잔고 현황")
    bank_data = _load_bank_accounts()
    accounts = bank_data.get("accounts", [])
    if accounts:
        workplaces_by_id = {site.get("id"): site for site in workplaces}
        account_df = pd.DataFrame(accounts)
        account_df["balance_won"] = account_df["balance"].apply(_format_won)
        if "balance_updated_at" not in account_df.columns:
            account_df["balance_updated_at"] = ""
        account_df["회사"] = [_account_company_label(row, workplaces_by_id) for row in accounts]
        account_view = account_df[
            ["회사", "account_name", "bank_name", "balance_won", "balance_updated_at"]
        ].rename(
            columns={
                "account_name": "계좌명",
                "bank_name": "은행명",
                "balance_won": "현재잔고",
                "balance_updated_at": "최종업데이트",
            }
        )
        render_plain_html_table(account_view, center_align=True)
    else:
        st.info("등록된 계좌가 없습니다.")


def _chart_layout(height=300, **overrides):
    """Plotly 차트 기본 레이아웃. 배경은 앱 테마가 보이도록 투명 처리."""
    bg = "rgba(0,0,0,0)"
    plot = "rgba(0,0,0,0)"
    txt = "#1e293b"
    grid = "#e2e8f0"
    axis = "#cbd5e1"
    legend_bg = "rgba(255,255,255,0)"
    layout = dict(
        paper_bgcolor=bg, plot_bgcolor=plot, height=height,
        font=dict(color=txt, size=12),
        xaxis=dict(gridcolor=grid, linecolor=axis,
                   tickfont=dict(color=txt), title_font=dict(color=txt),
                   showgrid=True, zeroline=False),
        yaxis=dict(gridcolor=grid, linecolor=axis,
                   tickfont=dict(color=txt), title_font=dict(color=txt),
                   showgrid=True, zeroline=False),
        legend=dict(font=dict(color=txt), bgcolor=legend_bg,
                    bordercolor=axis, borderwidth=1),
        margin=dict(t=20, b=20, l=10, r=10),
    )
    layout.update(overrides)
    return layout


def apply_global_table_css():
    st.markdown(
        """
        <style>
        .stApp {
            background: #F8FAFC;
        }
        .block-container {
            max-width: 100% !important;
            padding: 0.9rem 1.2rem 2rem !important;
        }
        h1, h2, h3, h4 {
            letter-spacing: 0 !important;
            color: #0F172A;
        }
        [data-testid="stMetric"] {
            background: #FFFFFF;
            border: 1px solid #E2E8F0;
            border-radius: 6px;
            padding: 10px 12px;
            box-shadow: none;
        }
        [data-testid="stMetricLabel"] p {
            color: #64748B !important;
            font-size: 12px !important;
            font-weight: 700 !important;
        }
        [data-testid="stMetricValue"] {
            color: #0F172A !important;
            font-size: 22px !important;
            font-weight: 800 !important;
        }
        /* 전체 버튼 색상 통일 */
        div.stButton > button,
        [data-testid="stFormSubmitButton"] button,
        [data-testid="stDownloadButton"] button {
            background-color: #2563EB !important;
            color: #FFFFFF !important;
            border: 1px solid #2563EB !important;
            border-radius: 6px !important;
            box-shadow: none !important;
        }
        div.stButton > button p, div.stButton > button span,
        [data-testid="stFormSubmitButton"] button p, [data-testid="stFormSubmitButton"] button span,
        [data-testid="stDownloadButton"] button p, [data-testid="stDownloadButton"] button span {
            color: #FFFFFF !important;
        }
        div.stButton > button:hover,
        [data-testid="stFormSubmitButton"] button:hover,
        [data-testid="stDownloadButton"] button:hover {
            background-color: #1D4ED8 !important;
            border-color: #1D4ED8 !important;
            color: #FFFFFF !important;
        }
        div.stButton > button:disabled,
        [data-testid="stFormSubmitButton"] button:disabled,
        [data-testid="stDownloadButton"] button:disabled {
            background-color: #A0AEC0 !important;
            border-color: #A0AEC0 !important;
            color: #F1F5F9 !important;
            opacity: 0.7;
        }
        div[data-testid="stDataFrame"],
        div[data-testid="stDataEditor"] {
            border: 1px solid #CBD5E1 !important;
            border-radius: 6px !important;
            overflow: hidden !important;
            background: #FFFFFF !important;
        }
        div[data-testid="stDataFrame"] th,
        div[data-testid="stDataFrame"] td,
        div[data-testid="stDataFrame"] [role="columnheader"],
        div[data-testid="stDataFrame"] [role="gridcell"],
        div[data-testid="stDataEditor"] [role="columnheader"],
        div[data-testid="stDataEditor"] [role="gridcell"] {
            white-space: nowrap !important;
            word-break: keep-all !important;
            font-size: 12px !important;
            min-height: 28px !important;
            line-height: 18px !important;
            border-color: #E2E8F0 !important;
        }
        div[data-testid="stDataFrame"] [role="columnheader"],
        div[data-testid="stDataEditor"] [role="columnheader"] {
            height: 30px !important;
            min-height: 30px !important;
            max-height: 30px !important;
            line-height: 20px !important;
            background: #E2E8F0 !important;
            color: #0F172A !important;
            font-weight: 800 !important;
        }
        /* 표 전체 가운데 정렬 */
        div[data-testid="stDataFrame"] [role="columnheader"],
        div[data-testid="stDataFrame"] [role="gridcell"],
        div[data-testid="stDataEditor"] [role="columnheader"],
        div[data-testid="stDataEditor"] [role="gridcell"] {
            text-align: center !important;
            justify-content: center !important;
        }
        /* 버튼 크기 공통 (높이 30px) */
        div.stButton > button,
        [data-testid="stFormSubmitButton"] button,
        [data-testid="stDownloadButton"] button {
            min-height: 32px !important;
            height: 32px !important;
            padding: 0 0.75rem !important;
            font-size: 13px !important;
            line-height: 1.2 !important;
        }
        [data-testid="stTextInput"] input,
        [data-testid="stNumberInput"] input,
        [data-testid="stSelectbox"] div[data-baseweb="select"] > div,
        [data-testid="stTextArea"] textarea {
            border-radius: 6px !important;
            border-color: #CBD5E1 !important;
            min-height: 32px !important;
            font-size: 13px !important;
        }
        [data-testid="stTabs"] button {
            min-height: 34px !important;
            padding: 6px 12px !important;
            font-size: 13px !important;
            font-weight: 700 !important;
        }
        .action-btn button p, .action-btn button span {
            font-size: 10px !important;
        }
        /* 직원 현황 카드 (라이트 모드) */
        .pms-staff-card {
            border: 1px solid #e2e8f0; border-radius: 10px;
            padding: 14px 18px; margin-bottom: 10px; background: #f8fafc;
        }
        .pms-card-name { font-size: 16px; font-weight: 700; margin-bottom: 8px; color: #1e293b; }
        .pms-card-stats { display: flex; gap: 12px; flex-wrap: wrap; }
        .pms-card-stats span { font-size: 13px; color: #475569; }
        .pms-card-points { margin-top: 6px; font-size: 15px; font-weight: 600; color: #4F46E5; }
        .pms-delta { font-size: 12px; margin-left: 8px; }
        .pms-delta-up   { color: #16a34a; }
        .pms-delta-down { color: #dc2626; }
        /* 직원 현황 카드 (다크 모드) */
        body:has(#pms-d:checked) .pms-staff-card {
            background: #252535 !important; border-color: #45475a !important;
        }
        body:has(#pms-d:checked) .pms-card-name  { color: #ffffff !important; }
        body:has(#pms-d:checked) .pms-card-stats span { color: #ffffff !important; }
        body:has(#pms-d:checked) .pms-card-points { color: #ffffff !important; }
        body:has(#pms-d:checked) .pms-delta-up   { color: #a6e3a1 !important; }
        body:has(#pms-d:checked) .pms-delta-down { color: #f38ba8 !important; }
        /* 기타 인라인 배경 패턴 (흰/연회색 배경 다크 처리) */
        body:has(#pms-d:checked) [style*="background:#f8fafc"],
        body:has(#pms-d:checked) [style*="background: #f8fafc"],
        body:has(#pms-d:checked) [style*="background:#EBF8FF"],
        body:has(#pms-d:checked) [style*="background: #EBF8FF"],
        body:has(#pms-d:checked) [style*="background:#fff"],
        body:has(#pms-d:checked) [style*="background: #fff"],
        body:has(#pms-d:checked) [style*="background:white"],
        body:has(#pms-d:checked) [style*="background: white"] {
            background: #252535 !important;
        }
        body:has(#pms-d:checked) [style*="color:#2B6CB0"] { color: #89dceb !important; }
        body:has(#pms-d:checked) [style*="color: #2B6CB0"] { color: #89dceb !important; }

        /* 실적 보고 테이블 (style_report_logic 래퍼) 다크모드 */
        body:has(#pms-d:checked) .pms-report-table {
            border-color: #45475a !important;
            box-shadow: none !important;
            background: #252535 !important;
        }
        body:has(#pms-d:checked) .pms-report-table th {
            background: #0f0f1f !important;
            color: #ffffff !important;
            border-color: #45475a !important;
        }
        body:has(#pms-d:checked) .pms-report-table td {
            background: #252535 !important;
            color: #ffffff !important;
            border-color: #313244 !important;
        }
        body:has(#pms-d:checked) .pms-report-table tr:nth-child(odd) td {
            background: #1e1e30 !important;
        }
        body:has(#pms-d:checked) .pms-report-table tr:hover td {
            background: #2d2d45 !important;
        }

        /* 다운로드 버튼 다크모드 */
        body:has(#pms-d:checked) [data-testid="stDownloadButton"] button,
        body:has(#pms-d:checked) [data-testid="stDownloadButton"] a {
            background-color: #2a2a3e !important;
            color: #ffffff !important;
            border-color: #45475a !important;
        }
        body:has(#pms-d:checked) [data-testid="stDownloadButton"] button:hover,
        body:has(#pms-d:checked) [data-testid="stDownloadButton"] a:hover {
            background-color: #313244 !important;
            border-color: #ffffff !important;
        }
        body:has(#pms-d:checked) [data-testid="stDownloadButton"] button p,
        body:has(#pms-d:checked) [data-testid="stDownloadButton"] button span,
        body:has(#pms-d:checked) [data-testid="stDownloadButton"] a p,
        body:has(#pms-d:checked) [data-testid="stDownloadButton"] a span {
            color: #ffffff !important;
        }
        /* 익스팬더(메뉴 이용 안내 등) 헤더 크기 공통 (높이 40px) */
        [data-testid="stExpander"] summary {
            min-height: 40px !important;
            height: 40px !important;
            padding: 0.25rem 0.75rem !important;
        }
        [data-testid="stExpander"] summary p {
            font-size: 0.85rem !important;
            margin: 0 !important;
        }
        [data-testid="stExpander"] summary svg,
        [data-testid="stExpander"] [data-testid="stExpanderToggleIcon"] {
            width: 1rem !important;
            height: 1rem !important;
        }
        [data-testid="stExpanderDetails"] {
            padding: 0.5rem 0.75rem !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


MENU_GUIDES = {
    "대시보드": [
        "🤖 AI 인사이트에서 결재 대기, 이체 대상, 리스크 사업장 등 주요 알림을 확인합니다.",
        "📊 사업장 현황, 전도금 지급 이력, AI 예측 안내를 한 화면에서 확인합니다.",
        "🏦 계좌별 잔고 현황도 함께 제공합니다.",
    ],
    "전자결재": [
        "📥 결재대기 문서를 승인/반려 처리합니다. (관리자)",
        "✅ 승인 시 해당 전도금 요청이 [품의 결과]에 '품의 확정'으로 반영됩니다.",
        "🗂️ 처리 완료 탭에서 과거 결재 이력을 조회할 수 있습니다.",
        "📈 처리 통계 탭에서 결재대기/승인/반려 건수와 월별 처리 추이를 확인합니다.",
    ],
    "회사 관리": [
        "🏬 당사(우리 회사)의 회사명, 사업자번호, 대표자, 연락처, 주소 등 기본 정보를 등록/수정합니다.",
        "🏦 여러 출금계좌를 등록하고 메인 출금계좌를 지정할 수 있습니다.",
        "💾 저장 시 최종 수정 일시가 자동으로 기록됩니다.",
        "📄 등록된 회사 정보는 문서 출력 및 엑셀 양식 등에서 활용됩니다.",
    ],
    "위탁 사업장 관리": [
        "🏢 [사업장 정보 관리] 탭에서 사업장 정보와 정기 지급일을 등록합니다.",
        "🏦 [계좌 관리] 탭에서 [사업장 정보 관리]에 등록된 사업장을 선택하여 계좌 정보를 등록/관리합니다.",
        "🔄 계좌 관리에서 '위탁 사업장 계좌 가져오기'로 사업장 계좌를 일괄 등록할 수 있습니다.",
        "👤 [담당자 관리] 탭에서 직원 계정의 접근 권한, 비밀번호, 역할을 관리합니다.",
    ],
    "서버 접속 정보": [
        "🗄️ MSSQL 서버 접속 정보를 등록합니다.",
        "🔐 비밀번호는 화면에 마스킹되어 표시됩니다.",
        "⚙️ 운영 환경에 맞게 인증 방식과 암호화 옵션을 설정합니다.",
    ],
    "이체 자료 확정": [
        "✅ 품의 확정된 전도금 요청을 선택해 이체 자료(엑셀)를 다운로드합니다.",
        "🏦 출금 계좌는 [회사 관리]에 등록된 출금계좌 중에서 선택하며, 메인 출금계좌가 먼저 표시됩니다.",
        "➡️ '이체 자료 확정' 클릭 시 해당 요청은 '이체 대상' 상태로 전환되어 [지급 결과 확인]에 표시됩니다.",
    ],
    "지급 결과 확인": [
        "💸 '이체 대상' 건을 선택해 이체 완료 여부를 확정합니다.",
        "📜 이체 완료된 지급 이력을 조회할 수 있습니다.",
    ],
    "전도금 요청": [
        "🧾 전도금 요청을 등록하면 [전자결재]로 결재 요청이 전달됩니다.",
        "📋 최근 요청현황에서 최근 7일간 요청한 내역과 처리 상태를 확인합니다.",
    ],
    "품의 결과": [
        "📊 전도금 요청의 품의(승인/반려) 처리 결과를 조회합니다.",
        "🔍 사업장/처리결과로 필터링할 수 있습니다.",
    ],
    "전도금 사용 결의 보고": [
        "🧾 제출된 전도금 사용 결의 보고 내역과 처리 상태를 확인합니다.",
    ],
    "계좌 잔고 확인": [
        "🏦 등록된 계좌의 현재 잔고를 확인하고 수동으로 업데이트합니다.",
        "📈 잔고 변동 추이를 그래프로 확인할 수 있습니다.",
    ],
}


def render_page_title(menu):
    page_label = CRM_MENU_LABELS.get(menu, menu)
    role_label = st.session_state.get("user_role", "")
    user_label = st.session_state.get("user_name", "")
    st.markdown(
        f"""
        <style>
            .crm-page-header {{
                display: flex;
                align-items: center;
                justify-content: space-between;
                gap: 16px;
                min-height: 54px;
                padding: 8px 0 14px;
                margin-bottom: 2px;
                border-bottom: 1px solid #E2E8F0;
            }}
            .crm-page-kicker {{
                color: #64748B;
                font-size: 11px;
                font-weight: 800;
                letter-spacing: 0.08em;
                text-transform: uppercase;
                margin-bottom: 3px;
            }}
            .crm-page-title {{
                color: #0F172A;
                font-size: 24px;
                font-weight: 850;
                line-height: 1.15;
            }}
            .crm-page-meta {{
                display: flex;
                align-items: center;
                gap: 8px;
                color: #475569;
                font-size: 12px;
                font-weight: 700;
                white-space: nowrap;
            }}
            .crm-user-pill {{
                display: inline-flex;
                align-items: center;
                min-height: 28px;
                padding: 4px 10px;
                border: 1px solid #CBD5E1;
                border-radius: 999px;
                background: #FFFFFF;
                color: #334155;
            }}
            @media (max-width: 768px) {{
                .crm-page-header {{
                    align-items: flex-start;
                    flex-direction: column;
                }}
                .crm-page-meta {{
                    white-space: normal;
                }}
            }}
        </style>
        <div class="crm-page-header">
            <div>
                <div class="crm-page-kicker">Internal CRM</div>
                <div class="crm-page-title">{html.escape(page_label)}</div>
            </div>
            <div class="crm-page-meta">
                <span class="crm-user-pill">{html.escape(str(user_label))}</span>
                <span class="crm-user-pill">{html.escape(str(role_label))}</span>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    _, col_guide = st.columns([0.86, 0.14])

    with col_guide:
        if menu in MENU_GUIDES:
            with st.expander("📌 메뉴 이용 안내", expanded=False):
                for line in MENU_GUIDES[menu]:
                    st.markdown(f"- {line}")

    st.markdown("<div style='height:14px'></div>", unsafe_allow_html=True)


def load_csv_to_state(url_key, state_key, force_refresh=False):
    st.session_state[state_key] = clean_header_logic(read_google_csv(st.session_state[url_key], force_refresh=force_refresh))


def refresh_google_sheets_action():
    st.session_state.cloud_sheet_df = clean_header_logic(read_google_csv(st.session_state.url_sync, force_refresh=True))
    st.session_state.analysis_lookup_df = clean_header_logic(read_google_csv(st.session_state.url_analysis, force_refresh=True))
    st.toast("구글시트 데이터를 다시 조회했습니다.")
    time.sleep(0.3)
    st.rerun()


def validation_tabs_with_refresh(key):
    st.markdown(
        """
        <style>
        /* .refresh-tab-button div과 button도 형제 — :has()+인접형제로 타겟팅 */
        [data-testid="stElementContainer"]:has(.refresh-tab-button) + [data-testid="stElementContainer"] [data-testid="stButton"] button {
            min-width: 42px !important;
            height: 38px !important;
            padding: 0 !important;
            margin-top: 2px !important;
            border-radius: 8px !important;
            font-size: 18px !important;
            font-weight: 900 !important;
        }
        body:has(#pms-d:checked) [data-testid="stElementContainer"]:has(.refresh-tab-button) + [data-testid="stElementContainer"] [data-testid="stButton"] button {
            background: #2a2a3e !important;
            background-color: #2a2a3e !important;
            color: #cdd6f4 !important;
            border-color: #45475a !important;
        }
        body:has(#pms-d:checked) [data-testid="stElementContainer"]:has(.refresh-tab-button) + [data-testid="stElementContainer"] [data-testid="stButton"] button:hover {
            background: #313244 !important;
            background-color: #313244 !important;
            color: #ffffff !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )
    tabs_col, link_col, refresh_col = st.columns([0.91, 0.045, 0.045])
    with tabs_col:
        tabs = st.tabs(["중복 이력", "초과 방문", "본사 개설완료일자 누락", "본사 ERP연계일자 누락", "기타 오류"])
    with link_col:
        _sheet_url = "https://docs.google.com/spreadsheets/d/1yS4gaES-iuzt1NSRTSdj9Ivg1fjbN5mIyX4pGnvEYN0/edit?gid=1533424484#gid=1533424484"
        st.link_button("🔗", _sheet_url, use_container_width=True, help="본사 구글시트 바로가기")
    with refresh_col:
        st.markdown("<div class='refresh-tab-button'>", unsafe_allow_html=True)
        if st.button("↻", use_container_width=True, key=key, help="구글시트 데이터 다시 조회"):
            try:
                refresh_google_sheets_action()
            except Exception as e:
                st.error(f"구글시트 갱신 실패: {e}")
        st.markdown("</div>", unsafe_allow_html=True)
    return tabs


def select_prev_month(state_key, widget_key):
    # 초기 로드 시 구글시트 자동 조회
    prev_sel_key = f"{widget_key}_prev"

    if st.session_state.analysis_lookup_df is None:
        try:
            load_csv_to_state("url_analysis", "analysis_lookup_df")
        except Exception:
            pass

    if st.session_state.analysis_lookup_df is not None:
        c_df = st.session_state.analysis_lookup_df.copy()
        d_col = find_col(c_df, ["활동일", "일자"])
        if d_col and d_col in c_df.columns:
            c_df[d_col] = pd.to_datetime(c_df[d_col], errors="coerce")
            opts = sorted(c_df[d_col].dropna().dt.strftime("%Y-%m").unique(), reverse=True)
            sel = st.selectbox("비교할 전월 선택", ["선택안함"] + list(opts), key=widget_key)

            # 선택값이 변경되었을 때 구글시트 재조회
            prev_sel = st.session_state.get(prev_sel_key)
            if prev_sel != sel and prev_sel is not None:
                try:
                    load_csv_to_state("url_analysis", "analysis_lookup_df")
                    c_df = st.session_state.analysis_lookup_df.copy()
                    c_df[d_col] = pd.to_datetime(c_df[d_col], errors="coerce")
                except Exception:
                    pass

            st.session_state[prev_sel_key] = sel
            st.session_state[state_key] = c_df[c_df[d_col].dt.strftime("%Y-%m") == sel] if sel != "선택안함" else None


def convert_bank_excel_to_activity(bank_df):
    """
    은행 엑셀 파일을 활동실적 양식으로 변환

    Args:
        bank_df: 은행 엑셀 DataFrame (접수일자, CMS구분, 고객번호, 업체명, 접수유형, 업무유형, 접수자, 담당자, 진행상태, 요청사항, 처리내용)

    Returns:
        변환된 DataFrame (등록자, 활동일, 활동상세, 업체명, 사업자번호)
    """
    try:
        # 1. 필요한 컬럼 찾기
        receipt_type_col = find_col(bank_df, ["접수유형"], "접수유형")
        receipt_date_col = find_col(bank_df, ["접수일자"], "접수일자")
        customer_no_col = find_col(bank_df, ["고객번호"], "고객번호")
        company_col = find_col(bank_df, ["업체명"], "업체명")
        manager_col = find_col(bank_df, ["담당자"], "담당자")
        request_col = find_col(bank_df, ["요청사항"], "요청사항")

        # 2. 접수유형이 "방문"인 것만 필터링
        if receipt_type_col not in bank_df.columns:
            st.error("접수유형 컬럼을 찾을 수 없습니다.")
            return pd.DataFrame()

        filtered_df = bank_df[bank_df[receipt_type_col].astype(str).str.contains("방문", na=False)].copy()

        if filtered_df.empty:
            st.warning("접수유형이 '방문'인 데이터가 없습니다.")
            return pd.DataFrame()

        # 3. 하나은행 시트 로드 (고객번호 → 사업자번호 매칭)
        try:
            load_csv_to_state("url_analysis", "analysis_lookup_df")
            hana_df = st.session_state.analysis_lookup_df.copy()
        except Exception as e:
            st.error(f"하나은행 시트를 불러올 수 없습니다: {str(e)}")
            return pd.DataFrame()

        # 하나은행 시트에서 고객번호와 사업자번호 컬럼 찾기
        hana_customer_col = find_col(hana_df, ["고객번호"], "고객번호")
        hana_biz_col = find_col(hana_df, ["사업자번호"], "사업자번호")

        if hana_customer_col not in hana_df.columns or hana_biz_col not in hana_df.columns:
            st.error("하나은행 시트에서 고객번호 또는 사업자번호 컬럼을 찾을 수 없습니다.")
            return pd.DataFrame()

        # 고객번호로 사업자번호 매칭
        hana_lookup = hana_df[[hana_customer_col, hana_biz_col]].drop_duplicates(subset=[hana_customer_col])
        hana_lookup.columns = ["고객번호_매칭", "사업자번호"]

        # 4. 본사 시트 로드 (사업자번호 → 개설완료일자, ERP연계일자 확인)
        try:
            load_csv_to_state("url_sync", "cloud_sheet_df")
            sync_df = st.session_state.cloud_sheet_df.copy()
        except Exception as e:
            st.error(f"본사 시트를 불러올 수 없습니다: {str(e)}")
            return pd.DataFrame()

        # 본사 시트에서 필요한 컬럼 찾기
        sync_biz_col = find_col(sync_df, ["사업자번호"], "사업자번호")
        sync_open_col = find_col(sync_df, ["개설완료일자"], "개설완료일자")
        sync_erp_col = find_col(sync_df, ["ERP연계일자"], "ERP연계일자")

        if sync_biz_col not in sync_df.columns:
            st.error("본사 시트에서 사업자번호 컬럼을 찾을 수 없습니다.")
            return pd.DataFrame()

        # 사업자번호 정규화
        sync_df[sync_biz_col] = normalize_biz(sync_df[sync_biz_col])
        sync_lookup = sync_df[[sync_biz_col, sync_open_col, sync_erp_col]].drop_duplicates(subset=[sync_biz_col]) if sync_open_col and sync_erp_col else sync_df[[sync_biz_col]].drop_duplicates(subset=[sync_biz_col])
        sync_lookup.columns = ["사업자번호_매칭", "개설완료일자", "ERP연계일자"] if sync_open_col and sync_erp_col else ["사업자번호_매칭"]

        # 5. 변환 로직 적용
        result_rows = []

        for idx, row in filtered_df.iterrows():
            customer_no = str(row.get(customer_no_col, "")).strip()
            receipt_date = row.get(receipt_date_col, "")
            company = row.get(company_col, "")
            manager = row.get(manager_col, "")
            request = str(row.get(request_col, "")).strip()

            # 고객번호로 사업자번호 찾기
            biz_no = ""
            matched_hana = hana_lookup[hana_lookup["고객번호_매칭"].astype(str).str.strip() == customer_no]
            if not matched_hana.empty:
                biz_no = str(matched_hana.iloc[0]["사업자번호"]).strip()

            # 사업자번호 정규화
            biz_no_normalized = normalize_biz(pd.Series([biz_no])).iloc[0]

            # 본사 시트에서 개설완료일자, ERP연계일자 확인
            open_date = ""
            erp_date = ""

            if biz_no_normalized and sync_open_col and sync_erp_col:
                matched_sync = sync_lookup[sync_lookup["사업자번호_매칭"] == biz_no_normalized]
                if not matched_sync.empty:
                    open_date = str(matched_sync.iloc[0].get("개설완료일자", "")).strip()
                    erp_date = str(matched_sync.iloc[0].get("ERP연계일자", "")).strip()

            # 활동상세 결정
            activity_detail = "운영"  # 기본값

            # 개설 조건: 요청사항에 "개설" or "구축" or "신규" 포함 AND 개설완료일자 공백
            if any(keyword in request for keyword in ["개설", "구축", "신규"]):
                if not open_date or open_date in ["", "nan", "None", "NaT"]:
                    activity_detail = "개설"

            # 연계 조건: 요청사항에 "연계" or "ERP" 포함 AND ERP연계일자 공백
            if any(keyword in request for keyword in ["연계", "ERP"]):
                if not erp_date or erp_date in ["", "nan", "None", "NaT"]:
                    activity_detail = "연계"

            # 활동일 포맷팅
            try:
                if isinstance(receipt_date, pd.Timestamp):
                    activity_date = receipt_date.strftime("%Y-%m-%d")
                else:
                    activity_date = pd.to_datetime(receipt_date).strftime("%Y-%m-%d")
            except:
                activity_date = str(receipt_date)

            result_rows.append({
                "등록자": manager,
                "활동일": activity_date,
                "활동상세": activity_detail,
                "업체명": company,
                "사업자번호": biz_no
            })

        result_df = pd.DataFrame(result_rows)

        return result_df

    except Exception as e:
        st.error(f"변환 중 오류가 발생했습니다: {str(e)}")
        return pd.DataFrame()


def render_converted_preview_editor(converted_preview_df, filters=None):
    if converted_preview_df is None or converted_preview_df.empty:
        return None

    st.markdown("#### 변환파일 미리보기")
    st.markdown(
        """
        <style>
        /* ── 변환파일 미리보기 다크모드 ───────────────────────────────────────
           canvas 에는 background-color 를 주지 않아야 filter 와 충돌하지 않음 */
        body[data-pms-theme="d"] [data-testid="stDataEditor"],
        body[data-pms-theme="d"] [data-testid="stDataEditor"] > div,
        body[data-pms-theme="d"] [data-testid="stDataEditor"] [role="grid"],
        body:has(#pms-d:checked) [data-testid="stDataEditor"],
        body:has(#pms-d:checked) [data-testid="stDataEditor"] > div,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="grid"] {
            background-color: #13131f !important;
            color: #e2e8f0 !important;
            border-color: #2d2d4a !important;
        }
        body[data-pms-theme="d"] [data-testid="stDataEditor"] canvas,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] canvas {
            filter: invert(1) hue-rotate(180deg) brightness(0.92) contrast(0.88) saturate(0.9) !important;
        }
        body[data-pms-theme="d"] [data-testid="stDataEditor"] input,
        body[data-pms-theme="d"] [data-testid="stDataEditor"] textarea,
        body[data-pms-theme="d"] [data-testid="stDataEditor"] [contenteditable="true"],
        body[data-pms-theme="d"] [data-testid="stDataEditor"] [data-baseweb="input"] input,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] input,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] textarea,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] [contenteditable="true"],
        body:has(#pms-d:checked) [data-testid="stDataEditor"] [data-baseweb="input"] input {
            background-color: #1e1e34 !important;
            color: #ffffff !important;
            caret-color: #ffffff !important;
            border-color: #6366f1 !important;
            -webkit-text-fill-color: #ffffff !important;
        }
        body[data-pms-theme="d"] [data-testid="stDataEditor"] input::selection,
        body[data-pms-theme="d"] [data-testid="stDataEditor"] textarea::selection,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] input::selection,
        body:has(#pms-d:checked) [data-testid="stDataEditor"] textarea::selection {
            background-color: #4f46e5 !important;
            color: #ffffff !important;
        }
        @media (prefers-color-scheme: dark) {
            body:has(#pms-s:checked) [data-testid="stDataEditor"],
            body:has(#pms-s:checked) [data-testid="stDataEditor"] > div,
            body:has(#pms-s:checked) [data-testid="stDataEditor"] [role="grid"] {
                background-color: #13131f !important;
                border-color: #2d2d4a !important;
            }
            body:has(#pms-s:checked) [data-testid="stDataEditor"] canvas {
                filter: invert(1) hue-rotate(180deg) brightness(0.92) contrast(0.88) saturate(0.9) !important;
            }
            body:has(#pms-s:checked) [data-testid="stDataEditor"] input,
            body:has(#pms-s:checked) [data-testid="stDataEditor"] textarea,
            body:has(#pms-s:checked) [data-testid="stDataEditor"] [contenteditable="true"],
            body:has(#pms-s:checked) [data-testid="stDataEditor"] [data-baseweb="input"] input {
                background-color: #1e1e34 !important;
                color: #ffffff !important;
                caret-color: #ffffff !important;
                border-color: #6366f1 !important;
                -webkit-text-fill-color: #ffffff !important;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )
    converted_preview_df = normalize_converted_history_df(converted_preview_df)
    data_key = "history_convert_preview_data"
    editor_source_df = st.session_state.get(data_key, converted_preview_df)
    if not isinstance(editor_source_df, pd.DataFrame) or list(editor_source_df.columns) != list(converted_preview_df.columns):
        editor_source_df = converted_preview_df
    editor_source_df = normalize_converted_history_df(editor_source_df)
    full_editor_df = editor_source_df.copy()
    active_filters = has_active_history_filters(filters)
    if active_filters:
        editor_source_df = apply_history_search_filters(editor_source_df, filters)
        st.caption(f"검색 결과 {len(editor_source_df):,}건 / 전체 {len(full_editor_df):,}건")
    editor_key = f"history_convert_preview_editor_{history_filter_signature(filters)}"
    # _sel 체크박스 컬럼 (맨 왼쪽) — 행 선택용, 저장 시 제거
    editor_source_df.insert(0, "_sel", False)
    edited_with_sel = st.data_editor(
        editor_source_df,
        key=editor_key,
        use_container_width=True,
        hide_index=True,
        num_rows="fixed",
        disabled=[col for col in editor_source_df.columns if col not in ["_sel", "활동일자", "활동구분", "활동상세"]],
        column_config={
            "_sel": st.column_config.CheckboxColumn("선택", default=False, width="small"),
            "지사": st.column_config.TextColumn("지사", disabled=True),
            "활동일자": st.column_config.TextColumn("활동일자"),
            "활동구분": st.column_config.SelectboxColumn("활동구분", options=["방문", "상담", "원격"], required=True),
            "활동상세": st.column_config.SelectboxColumn("활동상세", options=["운영", "개설", "연계"], required=True),
            "활동내역": st.column_config.TextColumn("활동내역"),
            "_is_manual": None,
        },
    )
    _selected_indices = edited_with_sel[edited_with_sel["_sel"] == True].index.tolist() if "_sel" in edited_with_sel.columns else []
    edited_preview_df = edited_with_sel.drop(columns=["_sel"], errors="ignore")
    if active_filters:
        analysis_df = full_editor_df.copy()
        common_index = [idx for idx in edited_preview_df.index if idx in analysis_df.index]
        if common_index:
            analysis_df.loc[common_index, edited_preview_df.columns] = edited_preview_df.loc[common_index, edited_preview_df.columns]
        analysis_df = normalize_converted_history_df(analysis_df)
    else:
        analysis_df = normalize_converted_history_df(edited_preview_df)
    st.session_state[data_key] = analysis_df
    if st.session_state.get("user_excel_source") in (None, "bank"):
        st.session_state.user_excel_data = analysis_df
        st.session_state.user_excel_source = "bank"

    # ── 업체명↔사업자번호 매핑 (구글시트 + 기존 미리보기 데이터) ──────────
    _comp_biz_map = {}
    _cloud_raw = st.session_state.get("cloud_sheet_df")
    if _cloud_raw is not None and not _cloud_raw.empty:
        _cc = clean_header_logic(_cloud_raw.copy())
        _ow = find_col(_cc, ["담당자", "등록자", "성명"])
        _cm = find_col(_cc, ["업체명", "고객명", "상호"])
        _bm = find_col(_cc, ["사업자번호"])
        if _cm and _bm:
            _sub = _cc if not _ow else _cc[_cc[_ow].astype(str).str.strip() == str(st.session_state.user_name).strip()]
            for _, _r in _sub.iterrows():
                _c = str(_r.get(_cm, "")).strip()
                _b = str(_r.get(_bm, "")).strip()
                if _c and _c not in _comp_biz_map:
                    _comp_biz_map[_c] = _b
    # 기존 미리보기 데이터에서도 보완
    _pv = st.session_state.get(data_key)
    if isinstance(_pv, pd.DataFrame) and not _pv.empty:
        _pc = find_col(_pv, ["업체명", "상호", "고객명"])
        _pb = find_col(_pv, ["사업자번호"])
        if _pc and _pb:
            for _, _r in _pv.iterrows():
                _c = str(_r.get(_pc, "")).strip()
                _b = str(_r.get(_pb, "")).strip()
                if _c and _c not in _comp_biz_map:
                    _comp_biz_map[_c] = _b
    _company_list = ["-- 업체명 선택 --"] + sorted(_comp_biz_map.keys())

    # ── 하단 버튼 행: [🗑️ 삭제] [spacer] [이력 추가 ＋] ──────────────────────
    st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)
    _del_col, _, _add_col = st.columns([1, 8, 2])
    with _del_col:
        _has_sel = len(_selected_indices) > 0
        if st.button("🗑️ 삭제", disabled=not _has_sel, key="del_selected_btn", use_container_width=True):
            _new_df = analysis_df.drop(index=_selected_indices).reset_index(drop=True)
            st.session_state[data_key] = normalize_converted_history_df(_new_df)
            st.rerun()
    with _add_col:
        if st.button("이력 추가 ＋", use_container_width=True, key="add_history_row"):
            st.session_state["show_add_history_form"] = not st.session_state.get("show_add_history_form", False)
            st.rerun()

    # ── 이력 추가 폼 ──────────────────────────────────────────────────────────
    if st.session_state.get("show_add_history_form", False):
        with st.container(border=True):
            st.markdown("##### 이력 추가")
            _fc1, _fc2 = st.columns([0.6, 0.4])
            with _fc1:
                _sel_comp = st.selectbox("업체명", _company_list, key="add_form_company")
                _sel_comp = "" if _sel_comp == "-- 업체명 선택 --" else _sel_comp
            with _fc2:
                _auto_biz = _comp_biz_map.get(_sel_comp, "")
                st.text_input("사업자번호 (자동입력)", value=_auto_biz, disabled=True, key="add_form_biz_display")

            _fd1, _fd2, _fd3 = st.columns(3)
            with _fd1:
                _form_date = st.text_input(
                    "활동일자",
                    value=st.session_state.get("add_form_date_val", (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d")),
                    key="add_form_date",
                )
            with _fd2:
                _form_cat = st.selectbox("활동구분", ["방문", "상담", "원격"], key="add_form_cat")
            with _fd3:
                _form_det = st.selectbox("활동상세", ["운영", "개설", "연계"], key="add_form_det")

            _form_note = st.text_input("활동내역", key="add_form_note")

            _fb1, _fb2 = st.columns(2)
            with _fb1:
                if st.button("등록", type="primary", use_container_width=True, key="add_form_submit"):
                    if not _sel_comp:
                        st.warning("업체명을 선택해주세요.")
                    else:
                        _det_title = {"운영": "운영방문", "개설": "개설 방문", "연계": "연계 방문"}.get(_form_det, "운영방문")
                        _new_row = {col: "" for col in analysis_df.columns}
                        _new_row.update({
                            "지사": "HANA지사",
                            "상품": "통합CMS",
                            "등록자": st.session_state.user_name,
                            "업체명": _sel_comp,
                            "사업자번호": _auto_biz,
                            "활동일자": _form_date,
                            "활동구분": _form_cat,
                            "활동상세": _form_det,
                            "제목": _det_title,
                            "활동내역": _form_note,
                            "_is_manual": True,
                        })
                        st.session_state[data_key] = normalize_converted_history_df(
                            pd.concat([analysis_df, pd.DataFrame([_new_row])], ignore_index=True)
                        )
                        # 등록 후 업체명·활동일자·활동내역 초기화
                        st.session_state["add_form_date_val"] = ""
                        st.session_state["add_form_company"] = "-- 업체명 선택 --"
                        for _rk in ["add_form_date", "add_form_note", "add_form_biz_display"]:
                            st.session_state.pop(_rk, None)
                        st.rerun()
            with _fb2:
                if st.button("취소", use_container_width=True, key="add_form_cancel"):
                    st.session_state["show_add_history_form"] = False
                    st.rerun()

    st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)
    return analysis_df


def dataframe_to_excel_bytes(sheets):
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        for sheet_name, df in sheets.items():
            safe_name = str(sheet_name)[:31] or "Sheet"
            df.to_excel(writer, index=False, sheet_name=safe_name)
    output.seek(0)
    return output.getvalue()


def _render_staff_admin():
    st.caption("직원 계정과 메뉴 접근 권한을 관리합니다.")

    # 파일 기준으로 재로드 (_deleted 목록 포함 반영)
    st.session_state.user_db = load_user_db()

    if st.session_state.pop("reset_staff_edit_sel", False):
        st.session_state.staff_edit_sel = "선택안함"

    if not st.session_state.get("_staff_admin_action_initialized"):
        st.session_state.staff_admin_action = ""
        st.session_state["_staff_admin_action_initialized"] = True

    workplace_data = _load_delegated_workplaces()
    workplace_names = sorted(
        {
            str(site.get("workplace_name", "")).strip()
            for site in workplace_data.get("workplaces", [])
            if str(site.get("workplace_name", "")).strip()
        }
    )

    staff_rows = []
    for uid, info in st.session_state.user_db.items():
        if uid == "1":
            continue
        staff_rows.append({
            "ID": uid,
            "성명": info.get("name", ""),
            "메일주소": info.get("email", ""),
            "핸드폰": info.get("phone", ""),
            "사업자 정보": info.get("dept_type", "사업부"),
            "마스터 구분": "관리자" if info.get("role") == "관리자" else "사용자",
        })

    if staff_rows:
        # ID 순서로 정렬
        staff_rows.sort(key=lambda x: x["ID"])

        # ── 사용자 목록 HTML 테이블 표시 (다크모드 호환, 가운데 정렬) ──
        render_plain_html_table(pd.DataFrame(staff_rows), center_align=True)
    else:
        st.info("등록된 직원이 없습니다.")

    st.markdown("---")

    action_cols = st.columns([0.15, 0.15, 0.15, 0.55])
    with action_cols[0]:
        if st.button("사용자 추가", key="staff_action_add", use_container_width=True, type="primary"):
            st.session_state.staff_admin_action = "add"
            st.rerun()
    with action_cols[1]:
        if st.button("수정", key="staff_action_edit", use_container_width=True):
            st.session_state.staff_admin_action = "edit"
            st.rerun()
    with action_cols[2]:
        if st.button("삭제", key="staff_action_delete", use_container_width=True):
            st.session_state.staff_admin_action = "delete"
            st.rerun()

    action = st.session_state.get("staff_admin_action", "")

    if action == "add":
        st.markdown("#### 사용자 추가")
        if not workplace_names:
            st.warning("[사업장 정보 관리]에서 사업장을 먼저 등록해주세요.")
            return

        with st.form("add_staff_form"):
            ac1, ac2, ac3 = st.columns(3)
            with ac1:
                new_uid = st.text_input("ID", key="add_staff_uid")
                new_pw = st.text_input("비밀번호", type="password", key="add_staff_pw")
                new_name = st.text_input("성명", key="add_staff_name")
            with ac2:
                new_email = st.text_input("메일주소", key="add_staff_email")
                new_phone = st.text_input("핸드폰", key="add_staff_phone")
            with ac3:
                new_business = st.selectbox("사업자 정보", workplace_names, key="add_staff_business")
                new_role = st.selectbox("마스터 구분", ["관리자", "사용자"], key="add_staff_role")
            add_submitted = st.form_submit_button("추가", use_container_width=True, type="primary")

        if add_submitted:
            uid_str = new_uid.strip()
            name_str = new_name.strip()
            pw_str = new_pw.strip()
            if not uid_str or not name_str or not pw_str:
                st.error("아이디, 성명, 비밀번호는 필수 입력입니다.")
            elif uid_str in st.session_state.user_db:
                st.error(f"이미 존재하는 아이디입니다: {uid_str}")
            else:
                st.session_state.user_db[uid_str] = {
                    "pw": pw_str,
                    "name": name_str,
                    "email": new_email.strip(),
                    "phone": new_phone.strip(),
                    "access": "허용",
                    "role": new_role,
                    "dept_type": new_business,
                    "staff_type": "정규직",
                    "outsource": "아니오",
                    "outsource_period": "해당없음",
                }
                # _deleted 목록 유지
                _cur_file = {}
                if os.path.exists(DB_FILE):
                    try:
                        with open(DB_FILE, "r", encoding="utf-8") as _f:
                            _cur_file = json.load(_f)
                    except Exception:
                        pass
                save_data = dict(st.session_state.user_db)
                if "_deleted" in _cur_file:
                    save_data["_deleted"] = _cur_file["_deleted"]
                save_db(DB_FILE, save_data)
                st.session_state.staff_admin_action = ""
                st.success(f"직원 '{name_str}({uid_str})'을(를) 추가했습니다.")
                st.rerun()

        cancel_col, _ = st.columns([0.15, 0.85])
        with cancel_col:
            if st.button("취소", key="cancel_add_staff", use_container_width=True):
                st.session_state.staff_admin_action = ""
                st.rerun()

        return

    if action not in {"edit", "delete"}:
        return

    st.markdown("#### 직원 정보 수정" if action == "edit" else "#### 직원 삭제")

    uid_options = [f"{r['ID']} — {r['성명']}" for r in staff_rows]
    sel_col, _ = st.columns([0.32, 0.68])
    with sel_col:
        sel = st.selectbox("대상 직원 선택", ["선택안함"] + uid_options, key="staff_edit_sel")

    if sel == "선택안함":
        return

    sel_uid = sel.split(" — ")[0].strip()
    info = st.session_state.user_db.get(sel_uid, {})

    if action == "delete":
        st.warning(f"{sel} 계정을 삭제합니다. 삭제 후 기본 계정에서도 복원되지 않도록 처리됩니다.")
        del_col, cancel_col, _ = st.columns([0.15, 0.15, 0.7])
        with del_col:
            if st.button("삭제 실행", type="primary", use_container_width=True):
                del st.session_state.user_db[sel_uid]
                # _deleted 목록을 파일에 함께 저장해 기본 계정에서도 복원되지 않도록 함
                _cur_file = {}
                if os.path.exists(DB_FILE):
                    try:
                        with open(DB_FILE, "r", encoding="utf-8") as _f:
                            _cur_file = json.load(_f)
                    except Exception:
                        pass
                _deleted_set = set(_cur_file.get("_deleted", []))
                _deleted_set.add(sel_uid)
                save_data = dict(st.session_state.user_db)
                save_data["_deleted"] = list(_deleted_set)
                save_db(DB_FILE, save_data, allow_shrink=True)
                st.session_state.reset_staff_edit_sel = True
                st.session_state.staff_admin_action = ""
                st.success(f"{sel} 삭제 완료")
                time.sleep(0.5)
                st.rerun()
        with cancel_col:
            if st.button("취소", use_container_width=True):
                st.session_state.staff_admin_action = ""
                st.rerun()
        return

    st.markdown(f"**메일주소:** {info.get('email', '—')}")

    if not workplace_names:
        st.warning("[사업장 정보 관리]에서 사업장을 먼저 등록해주세요.")
        return

    business_opts = workplace_names

    c1, c2, c3, c4, c5 = st.columns([1, 1.3, 1, 1.3, 1])
    with c1:
        new_name = st.text_input("성명", value=info.get("name", ""), key="edit_staff_name")
    with c2:
        new_email = st.text_input("메일주소", value=info.get("email", ""), key="edit_staff_email")
    with c3:
        new_phone = st.text_input("핸드폰", value=info.get("phone", ""), key="edit_staff_phone")
    with c4:
        new_dept_type = st.selectbox(
            "사업자 정보",
            business_opts,
            index=business_opts.index(info.get("dept_type", "")) if info.get("dept_type", "") in business_opts else 0,
            key="edit_dept_type",
        )
    with c5:
        new_access = st.selectbox("로그인 허용 여부", ["허용", "불가"],
                                  index=["허용", "불가"].index(info.get("access", "불가")),
                                  key="edit_access")
        new_role = st.selectbox("마스터 구분", ["관리자", "사용자"],
                                index=0 if info.get("role") == "관리자" else 1,
                                key="edit_role")

    bc1, bc2, _ = st.columns([0.15, 0.15, 0.7])
    with bc1:
        if st.button("저장", type="primary", use_container_width=True):
            st.session_state.user_db[sel_uid]["name"] = new_name.strip()
            st.session_state.user_db[sel_uid]["email"] = new_email.strip()
            st.session_state.user_db[sel_uid]["phone"] = new_phone.strip()
            st.session_state.user_db[sel_uid]["dept_type"] = new_dept_type
            st.session_state.user_db[sel_uid]["access"] = new_access
            st.session_state.user_db[sel_uid]["role"] = new_role
            # _deleted 목록 유지
            _cur_file = {}
            if os.path.exists(DB_FILE):
                try:
                    with open(DB_FILE, "r", encoding="utf-8") as _f:
                        _cur_file = json.load(_f)
                except Exception:
                    pass
            save_data = dict(st.session_state.user_db)
            if "_deleted" in _cur_file:
                save_data["_deleted"] = _cur_file["_deleted"]
            save_db(DB_FILE, save_data)
            st.session_state.reset_staff_edit_sel = True
            st.session_state.staff_admin_action = ""
            st.success("저장 완료")
            time.sleep(0.5)
            st.rerun()
    with bc2:
        if st.button("취소", use_container_width=True):
            st.session_state.staff_admin_action = ""
            st.session_state.reset_staff_edit_sel = True
            st.rerun()


def _load_server_connection():
    installed_drivers = _get_installed_odbc_drivers()
    default_data = {
        "server_host": "",
        "server_port": 1433,
        "driver_name": _preferred_mssql_driver(installed_drivers),
        "database_name": "",
        "auth_type": "SQL Server 인증",
        "username": "",
        "password": "",
        "encrypt": True,
        "trust_server_certificate": False,
        "connection_timeout": 30,
        "updated_at": "",
        "updated_by": "",
    }
    if not os.path.exists(SERVER_CONNECTION_FILE):
        return default_data
    try:
        with open(SERVER_CONNECTION_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception:
        return default_data
    if not isinstance(data, dict):
        return default_data
    merged = {**default_data, **data}
    if installed_drivers and merged.get("driver_name") not in installed_drivers:
        merged["driver_name"] = _preferred_mssql_driver(installed_drivers)
    return merged


def _save_server_connection(data):
    with open(SERVER_CONNECTION_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=4)


def _mask_secret(value):
    if not value:
        return "미등록"
    return "●" * min(max(len(str(value)), 6), 12)


def show_server_connection_info():
    st.markdown("### 서버 접속 정보")
    st.caption("내부 관리에서 사용할 MSSQL 서버 접속 정보를 관리합니다.")

    config = _load_server_connection()
    installed_drivers = _get_installed_odbc_drivers()
    driver_options = installed_drivers.copy()
    current_driver = config.get("driver_name", "")
    if current_driver and current_driver not in driver_options and not installed_drivers:
        driver_options.insert(0, current_driver)
    if not driver_options:
        driver_options = ["ODBC Driver 18 for SQL Server", "ODBC Driver 17 for SQL Server", "FreeTDS"]

    with st.form("mssql_connection_form"):
        col_a, col_b = st.columns([3, 1])
        server_host = col_a.text_input(
            "MSSQL 서버 주소",
            value=config.get("server_host", ""),
            placeholder="예: 10.10.10.20 또는 db.company.local",
        )
        server_port = col_b.number_input(
            "포트",
            min_value=1,
            max_value=65535,
            value=int(config.get("server_port", 1433) or 1433),
            step=1,
        )

        database_name = st.text_input(
            "데이터베이스명",
            value=config.get("database_name", ""),
            placeholder="예: SalesManagement",
        )
        driver_name = st.selectbox(
            "ODBC 드라이버명",
            driver_options,
            index=driver_options.index(current_driver) if current_driver in driver_options else 0,
        )
        if installed_drivers:
            st.caption(f"현재 서버에서 감지된 ODBC 드라이버: {', '.join(installed_drivers)}")
        else:
            st.warning("현재 서버에서 감지된 ODBC 드라이버가 없습니다. MSSQL 접속 전 시스템 드라이버 설치가 필요합니다.")

        col_c, col_d = st.columns(2)
        auth_type = col_c.selectbox(
            "인증 방식",
            ["SQL Server 인증", "Windows 인증"],
            index=0 if config.get("auth_type", "SQL Server 인증") == "SQL Server 인증" else 1,
        )
        connection_timeout = col_d.number_input(
            "접속 제한 시간(초)",
            min_value=1,
            max_value=300,
            value=int(config.get("connection_timeout", 30) or 30),
            step=1,
        )

        username = st.text_input("계정", value=config.get("username", ""), placeholder="예: sales_app")
        password = st.text_input(
            "비밀번호",
            value=config.get("password", ""),
            type="password",
            placeholder="MSSQL 접속 비밀번호",
        )

        col_e, col_f = st.columns(2)
        encrypt = col_e.checkbox("암호화 연결 사용", value=bool(config.get("encrypt", True)))
        trust_server_certificate = col_f.checkbox(
            "서버 인증서 신뢰",
            value=bool(config.get("trust_server_certificate", False)),
        )

        saved = st.form_submit_button("서버 접속 정보 저장", type="primary")

    if saved:
        if not server_host.strip() or not database_name.strip():
            st.warning("MSSQL 서버 주소와 데이터베이스명을 입력해주세요.")
        elif auth_type == "SQL Server 인증" and (not username.strip() or not password):
            st.warning("SQL Server 인증을 사용하려면 계정과 비밀번호를 입력해주세요.")
        else:
            payload = {
                "server_host": server_host.strip(),
                "server_port": int(server_port),
                "driver_name": str(driver_name).strip() or _preferred_mssql_driver(installed_drivers),
                "database_name": database_name.strip(),
                "auth_type": auth_type,
                "username": username.strip(),
                "password": password,
                "encrypt": bool(encrypt),
                "trust_server_certificate": bool(trust_server_certificate),
                "connection_timeout": int(connection_timeout),
                "updated_at": (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S"),
                "updated_by": st.session_state.get("user_name", ""),
            }
            _save_server_connection(payload)
            try:
                with _connect_mssql(payload) as conn:
                    _ensure_sales_registration_table(conn)
                st.session_state["_server_connection_notice"] = (
                    "success",
                    "MSSQL 서버 접속 정보가 저장되었고, 접속 및 테이블 확인까지 완료되었습니다.",
                )
            except Exception as exc:
                st.session_state["_server_connection_notice"] = (
                    "warning",
                    f"MSSQL 서버 접속 정보는 저장되었지만 연결 확인에 실패했습니다: {exc}",
                )
            st.rerun()

    st.divider()
    notice = st.session_state.pop("_server_connection_notice", None)
    if notice:
        level, text = notice
        if level == "success":
            st.success(text)
        else:
            st.warning(text)

    test_col, _ = st.columns([1, 3])
    with test_col:
        if st.button("테이블 생성/접속 테스트", use_container_width=True):
            try:
                latest_for_test = _load_server_connection()
                with _connect_mssql(latest_for_test) as conn:
                    _ensure_sales_registration_table(conn)
                st.success("MSSQL 접속 및 dbo.sales_registrations 테이블 확인이 완료되었습니다.")
            except Exception as exc:
                st.error(f"MSSQL 접속 또는 테이블 생성에 실패했습니다: {exc}")

    latest = _load_server_connection()
    st.markdown("#### 현재 등록 정보")
    summary_df = pd.DataFrame(
        [
            {"항목": "서버 주소", "값": latest.get("server_host") or "미등록"},
            {"항목": "포트", "값": latest.get("server_port") or "미등록"},
            {"항목": "ODBC 드라이버명", "값": latest.get("driver_name") or "미등록"},
            {"항목": "데이터베이스명", "값": latest.get("database_name") or "미등록"},
            {"항목": "인증 방식", "값": latest.get("auth_type") or "미등록"},
            {"항목": "계정", "값": latest.get("username") or "미등록"},
            {"항목": "비밀번호", "값": _mask_secret(latest.get("password"))},
            {"항목": "암호화 연결", "값": "사용" if latest.get("encrypt") else "미사용"},
            {"항목": "서버 인증서 신뢰", "값": "사용" if latest.get("trust_server_certificate") else "미사용"},
            {"항목": "접속 제한 시간", "값": f"{latest.get('connection_timeout') or 30}초"},
            {"항목": "최종 수정", "값": latest.get("updated_at") or "미등록"},
            {"항목": "수정자", "값": latest.get("updated_by") or "미등록"},
        ]
    )
    st.dataframe(summary_df, use_container_width=True, hide_index=True)


def inject_theme_toggle():
    # JavaScript는 st.markdown innerHTML으로 실행 불가 → CSS :has() + radio 버튼 방식 사용
    st.markdown("""
    <style>
    .pms-sw-outer {
        position: fixed; top: 14px; right: 18px; z-index: 999999;
    }
    .pms-sw-track {
        display: inline-flex; align-items: center;
        background: #e5e7eb; border-radius: 22px; padding: 3px;
        box-shadow: 0 1px 6px rgba(0,0,0,0.12); gap: 2px;
    }
    .pms-btn {
        position: relative; cursor: pointer;
        padding: 0 10px; height: 28px; border-radius: 18px;
        line-height: 28px; opacity: 0.45; color: #374151;
        display: flex; align-items: center; justify-content: center;
        transition: opacity 0.2s, background 0.15s, color 0.15s;
        user-select: none;
    }
    .pms-btn:hover { opacity: 0.75; }
    /* radio를 label 전체에 투명하게 덮어서 한 번 클릭에 즉시 체크 */
    .pms-theme-radio {
        position: absolute; inset: 0;
        opacity: 0; cursor: pointer; margin: 0;
    }
    .pms-btn:has(input:checked) {
        opacity: 1; background: #2F6FED; color: white;
    }

    /* ══ 다크 모드 컬러 팔레트 ══
       배경:    #1e1e2e  표면:   #252535  카드:    #2a2a3e
       테두리:  #45475a  텍스트: #cdd6f4  보조:    #a6adc8
       강조:    #89b4fa  성공:   #a6e3a1  경고:    #f9e2af
       오류:    #f38ba8  정보:   #89dceb  사이드:  #16162a  */

    /* 기본 배경·텍스트 — 모든 글씨 흰색 */
    body:has(#pms-d:checked) .stApp,
    body:has(#pms-d:checked) [data-testid="stAppViewContainer"],
    body:has(#pms-d:checked) [data-testid="stMain"],
    body:has(#pms-d:checked) section.main {
        background-color: #1e1e2e !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) .main .block-container { background-color: #1e1e2e !important; }

    /* 텍스트 — 전체 흰색 */
    body:has(#pms-d:checked) h1, body:has(#pms-d:checked) h2, body:has(#pms-d:checked) h3,
    body:has(#pms-d:checked) h4, body:has(#pms-d:checked) h5, body:has(#pms-d:checked) h6 {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) p, body:has(#pms-d:checked) li,
    body:has(#pms-d:checked) [data-testid="stMarkdownContainer"],
    body:has(#pms-d:checked) [data-testid="stMarkdownContainer"] p {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) label { color: #ffffff !important; }
    /* span: 인라인 style 색상이 있는 셀(대비/증감)은 제외하고 흰색 적용 */
    body:has(#pms-d:checked) span:not([style*="color"]) { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-testid="stMarkdownContainer"] span:not([style*="color"]) { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-testid="stSidebar"] span { color: #ffffff !important; }
    body:has(#pms-d:checked) .stButton span { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-baseweb="tab"] span { color: inherit !important; }
    body:has(#pms-d:checked) [data-testid="stMetricValue"] span { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-testid="stMetricLabel"] span { color: #ffffff !important; }

    /* 파일 업로더 다크모드 — 모든 자식 요소 포함 */
    body:has(#pms-d:checked) [data-testid="stFileUploader"],
    body:has(#pms-d:checked) [data-testid="stFileUploader"] > div,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] > div > div {
        background-color: #252535 !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderDropzone"],
    body:has(#pms-d:checked) [data-testid="stFileUploader"] section,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] [class*="uploadDropzone"],
    body:has(#pms-d:checked) [data-testid="stFileUploader"] [class*="drop"] {
        background-color: #1e1e30 !important;
        border: 2px dashed #45475a !important;
        border-radius: 8px !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderDropzone"] *,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] section * {
        color: #ffffff !important;
        background-color: transparent !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderDropzone"] button,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] section button {
        background-color: #313244 !important; color: #ffffff !important;
        border: 1px solid #45475a !important; border-radius: 6px !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderDropzone"] button:hover,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] section button:hover {
        background-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderFileData"],
    body:has(#pms-d:checked) [data-testid="stFileUploader"] [data-testid="stFileUploaderFileData"] {
        background-color: #252535 !important; border-color: #45475a !important;
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) small, body:has(#pms-d:checked) [data-testid="stCaptionContainer"] {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) hr { border-color: #45475a !important; }

    /* 사이드바 — Zoho 네이비 톤은 다크모드에서도 동일 유지 */
    body:has(#pms-d:checked) [data-testid="stSidebar"],
    body:has(#pms-d:checked) [data-testid="stSidebarContent"] {
        background-color: #16284A !important;
    }
    body:has(#pms-d:checked) [data-testid="stSidebar"] p,
    body:has(#pms-d:checked) [data-testid="stSidebar"] span,
    body:has(#pms-d:checked) [data-testid="stSidebar"] label {
        color: #E8EEF8 !important;
    }
    body:has(#pms-d:checked) [data-testid="stSidebar"] .stButton > button {
        background-color: transparent !important; color: #C9D6EC !important;
        border-color: transparent !important; box-shadow: none !important;
    }
    body:has(#pms-d:checked) [data-testid="stSidebar"] .stButton > button:hover {
        background-color: #233A61 !important;
    }

    /* 입력 컨트롤 */
    body:has(#pms-d:checked) input,
    body:has(#pms-d:checked) textarea,
    body:has(#pms-d:checked) [data-baseweb="input"] input,
    body:has(#pms-d:checked) [data-baseweb="textarea"] textarea {
        background-color: #252535 !important; color: #ffffff !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="input"],
    body:has(#pms-d:checked) [data-baseweb="textarea"] {
        background-color: #252535 !important; border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="select"] > div,
    body:has(#pms-d:checked) [data-baseweb="select"] div[role="combobox"] {
        background-color: #252535 !important; border-color: #45475a !important; color: #ffffff !important;
    }
    /* 셀렉트박스 드롭다운 팝업 — 전체 컨테이너 및 모든 하위 요소 */
    body:has(#pms-d:checked) [data-baseweb="popover"],
    body:has(#pms-d:checked) [data-baseweb="popover"] > div,
    body:has(#pms-d:checked) [data-baseweb="popover"] > div > div,
    body:has(#pms-d:checked) [data-baseweb="popover"] > div > div > div {
        background-color: #252535 !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="menu"],
    body:has(#pms-d:checked) [data-baseweb="menu"] ul,
    body:has(#pms-d:checked) [role="listbox"] {
        background-color: #252535 !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="menu"] li,
    body:has(#pms-d:checked) [data-baseweb="menu"] [role="option"],
    body:has(#pms-d:checked) [role="option"],
    body:has(#pms-d:checked) [role="listbox"] > div {
        background-color: #252535 !important;
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-baseweb="menu"] li span,
    body:has(#pms-d:checked) [data-baseweb="menu"] [role="option"] span,
    body:has(#pms-d:checked) [role="option"] span {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-baseweb="menu"] li:hover,
    body:has(#pms-d:checked) [data-baseweb="menu"] [aria-selected="true"],
    body:has(#pms-d:checked) [role="option"]:hover,
    body:has(#pms-d:checked) [role="option"][aria-selected="true"] {
        background-color: #3a3a5e !important;
    }
    body:has(#pms-d:checked) [data-testid="stNumberInput"] button {
        background-color: #313244 !important; color: #ffffff !important;
        border-color: #45475a !important;
    }

    /* 버튼 */
    body:has(#pms-d:checked) .stButton > button,
    body:has(#pms-d:checked) [data-testid="stButton"] button,
    body:has(#pms-d:checked) button[data-testid^="baseButton"] {
        background: #2a2a3e !important;
        background-color: #2a2a3e !important; color: #ffffff !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) .stButton > button:hover,
    body:has(#pms-d:checked) [data-testid="stButton"] button:hover,
    body:has(#pms-d:checked) button[data-testid^="baseButton"]:hover {
        background: #313244 !important;
        background-color: #313244 !important; border-color: #ffffff !important;
    }
    body:has(#pms-d:checked) .stButton > button[kind="primary"],
    body:has(#pms-d:checked) .stButton > button[data-testid="baseButton-primary"] {
        background-color: #4F46E5 !important; color: #ffffff !important;
        border-color: #6366f1 !important;
    }
    body:has(#pms-d:checked) .stButton > button[kind="primary"]:hover {
        background-color: #6366f1 !important;
    }

    /* 홈 버튼 — 아이콘 전용: 배경·테두리 없음 */
    button.pms-home-btn,
    button.pms-home-btn:hover {
        display: inline-flex !important;
        align-items: center !important;
        justify-content: center !important;
        width: 82px !important;
        height: 82px !important;
        padding: 0 !important;
        background: transparent !important;
        border: none !important;
        box-shadow: none !important;
        color: transparent !important;
        font-size: 0 !important;
    }
    button.pms-home-btn *,
    button.pms-home-btn:hover * {
        color: transparent !important;
        display: none !important;
        font-size: 0 !important;
    }
    button.pms-home-btn::before {
        content: "";
        width: 75px;
        height: 75px;
        display: block;
        flex: 0 0 auto;
        background: url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 96 96'%3E%3Cpath fill='%23000' d='M32 8h56v56c0 13.255-10.745 24-24 24H8V32C8 18.745 18.745 8 32 8Z'/%3E%3Cpath fill='%23fff' d='M20 52 48 26l28 26h-10v26H30V52H20Z'/%3E%3Cpath fill='%23fff' d='M58 34h12v18H58z'/%3E%3Cpath fill='%23000' d='M41 52h8v8h-8zM53 52h8v8h-8zM41 64h8v8h-8zM53 64h8v8h-8z'/%3E%3C/svg%3E") center / contain no-repeat;
    }
    button.pms-home-btn:hover::before {
        filter: opacity(0.86);
    }
    body:has(#pms-d:checked) button.pms-home-btn {
        background: transparent !important;
        background-color: transparent !important;
        border: none !important;
        box-shadow: none !important;
    }
    body:has(#pms-d:checked) button.pms-home-btn::before {
        filter: none;
    }
    body:has(#pms-d:checked) button.pms-home-btn:hover {
        background: transparent !important;
        background-color: transparent !important;
        border: none !important;
    }
    body:has(#pms-d:checked) button.pms-home-btn:hover::before {
        filter: opacity(0.86);
    }

    /* 새로고침 버튼 다크모드 */
    body:has(#pms-d:checked) button.pms-refresh-btn {
        background: #2a2a3e !important;
        background-color: #2a2a3e !important;
        color: #cdd6f4 !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) button.pms-refresh-btn:hover {
        background: #313244 !important;
        background-color: #313244 !important;
        color: #ffffff !important;
    }

    /* 탭 */
    body:has(#pms-d:checked) [data-baseweb="tab-list"] {
        background-color: #1e1e2e !important; border-bottom: 2px solid #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="tab"] {
        background-color: transparent !important; color: rgba(255,255,255,0.55) !important;
    }
    body:has(#pms-d:checked) [data-baseweb="tab"][aria-selected="true"] {
        color: #ffffff !important; font-weight: 700;
    }
    body:has(#pms-d:checked) [data-baseweb="tab-highlight"] {
        background-color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-baseweb="tab-border"] {
        background-color: #45475a !important;
    }

    /* 메트릭 카드 */
    body:has(#pms-d:checked) [data-testid="metric-container"] {
        background: #252535 !important; border: 1px solid #45475a !important; border-radius: 8px;
    }
    body:has(#pms-d:checked) [data-testid="stMetricValue"] {
        color: #ffffff !important; font-weight: 700;
    }
    body:has(#pms-d:checked) [data-testid="stMetricLabel"] { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-testid="stMetricDeltaIcon-Up"]   { color: #a6e3a1 !important; }
    body:has(#pms-d:checked) [data-testid="stMetricDeltaIcon-Down"] { color: #f38ba8 !important; }

    /* Plotly 차트 */
    body:has(#pms-d:checked) .js-plotly-plot,
    body:has(#pms-d:checked) .js-plotly-plot .plot-container,
    body:has(#pms-d:checked) .js-plotly-plot .svg-container {
        background: transparent !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .bglayer rect.bg {
        fill: transparent !important;
        stroke: transparent !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .legend rect.bg {
        fill: #252535 !important;
        stroke: #45475a !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot text {
        fill: #ffffff !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .gridlayer path {
        stroke: #45475a !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .zerolinelayer path,
    body:has(#pms-d:checked) .js-plotly-plot .xlines-below path,
    body:has(#pms-d:checked) .js-plotly-plot .ylines-below path {
        stroke: #585b70 !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .modebar-btn path {
        fill: #cdd6f4 !important;
    }

    /* 알림 박스 */
    body:has(#pms-d:checked) [data-testid="stNotification"],
    body:has(#pms-d:checked) .stAlert { border-radius: 6px !important; }
    body:has(#pms-d:checked) [data-testid="stNotification"][kind="info"],
    body:has(#pms-d:checked) .stAlert.stInfo {
        background-color: #1a2a3a !important; border-left: 4px solid #89dceb !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stNotification"][kind="success"],
    body:has(#pms-d:checked) .stAlert.stSuccess {
        background-color: #1a2e1a !important; border-left: 4px solid #a6e3a1 !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stNotification"][kind="warning"],
    body:has(#pms-d:checked) .stAlert.stWarning {
        background-color: #2e2a1a !important; border-left: 4px solid #f9e2af !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stNotification"][kind="error"],
    body:has(#pms-d:checked) .stAlert.stError {
        background-color: #2e1a1a !important; border-left: 4px solid #f38ba8 !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stNotification"] p,
    body:has(#pms-d:checked) .stAlert p { color: #ffffff !important; }

    /* 익스팬더 */
    body:has(#pms-d:checked) [data-testid="stExpander"] {
        background-color: #252535 !important; border: 1px solid #45475a !important; border-radius: 6px;
    }
    body:has(#pms-d:checked) [data-testid="stExpander"] summary {
        background-color: #252535 !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stExpander"] summary:hover {
        background-color: #313244 !important;
    }

    /* ─── AG Grid 완전 다크모드 ─── */

    /* 1. CSS 변수 (AG Grid가 읽는 경우 적용) */
    body:has(#pms-d:checked) [class*="ag-theme"] {
        --ag-background-color: #252535;
        --ag-foreground-color: #ffffff;
        --ag-secondary-foreground-color: #ffffff;
        --ag-border-color: #45475a;
        --ag-secondary-border-color: #313244;
        --ag-row-border-color: #313244;
        --ag-header-background-color: #0f0f1f;
        --ag-header-foreground-color: #ffffff;
        --ag-header-column-separator-color: #45475a;
        --ag-odd-row-background-color: #1e1e30;
        --ag-row-hover-color: #2d2d45;
        --ag-selected-row-background-color: #3a3a5e;
        --ag-range-selection-background-color: rgba(255,255,255,0.1);
        --ag-input-focus-border-color: #ffffff;
        --ag-cell-horizontal-border: solid #313244;
        --ag-font-size: 13px;
        --ag-data-color: #ffffff;
        --ag-alpine-active-color: #4F46E5;
    }

    /* 2. 최외곽 래퍼 + 내부 모든 div 배경 강제 */
    body:has(#pms-d:checked) [data-testid="stDataFrame"],
    body:has(#pms-d:checked) [data-testid="stDataEditor"],
    body:has(#pms-d:checked) [data-testid="stDataFrame"] > div,
    body:has(#pms-d:checked) [data-testid="stDataFrame"] > div > div,
    body:has(#pms-d:checked) [data-testid="stDataFrame"] > div > div > div,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] > div,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] > div > div,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="grid"] {
        background-color: #13131f !important;
        border-color: #2d2d4a !important;
        color: #e2e8f0 !important;
    }
    /* stDataEditor canvas — invert 필터로 라이트 렌더링을 다크로 반전 */
    body:has(#pms-d:checked) [data-testid="stDataEditor"] canvas {
        filter: invert(1) hue-rotate(180deg) brightness(0.92) contrast(0.88) saturate(0.9) !important;
    }
    /* stDataEditor 인라인 편집 입력창 */
    body:has(#pms-d:checked) [data-testid="stDataEditor"] input,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] textarea,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [contenteditable="true"],
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [data-baseweb="input"] input {
        background-color: #1e1e34 !important;
        color: #ffffff !important;
        caret-color: #ffffff !important;
        border-color: #6366f1 !important;
        -webkit-text-fill-color: #ffffff !important;
    }

    /* 3. ag-theme 컨테이너 전체 */
    body:has(#pms-d:checked) [class*="ag-theme"],
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-root-wrapper,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-root-wrapper-body,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-root,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-body,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-body-viewport,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-body-horizontal-scroll,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-body-horizontal-scroll-viewport,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-center-cols-clipper,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-center-cols-viewport,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-center-cols-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-pinned-left-cols-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-pinned-right-cols-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-full-width-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-floating-top,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-floating-bottom,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-overlay {
        background-color: #252535 !important;
        border-color: #45475a !important;
    }

    /* 4. 헤더 */
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-row,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-viewport,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-pinned-left-header,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-pinned-right-header {
        background-color: #0f0f1f !important;
        border-bottom: 2px solid #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-cell,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-cell-comp-wrapper,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-group-cell {
        background-color: #0f0f1f !important;
        color: #ffffff !important;
        border-right: 1px solid #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-cell-text,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-cell-label {
        color: #ffffff !important;
        font-weight: 700 !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-sort-indicator-icon,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-icon {
        color: #ffffff !important;
    }

    /* 5. 행 — .ag-row 자체에 배경 지정 (ag-row-even 없을 때 대비) */
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row {
        background-color: #252535 !important;
        border-bottom-color: #313244 !important;
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row-odd {
        background-color: #1e1e30 !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row:hover,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row-hover {
        background-color: #2d2d45 !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row-selected {
        background-color: #3a3a5e !important;
    }

    /* 6. 셀 */
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-cell {
        border-right-color: #313244 !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-cell:not([style*="color"]),
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-cell-value:not([style*="color"]) {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-cell span:not([style*="color"]) {
        color: #ffffff !important;
    }

    /* 7. role 속성 기반 (columnheader / gridcell) */
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="columnheader"],
    body:has(#pms-d:checked) [data-testid="stDataFrame"]  [role="columnheader"] {
        background-color: #0f0f1f !important;
        color: #ffffff !important;
        border-bottom-color: #45475a !important;
        border-right-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="columnheader"] *,
    body:has(#pms-d:checked) [data-testid="stDataFrame"]  [role="columnheader"] * {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="gridcell"]:not([style*="color"]),
    body:has(#pms-d:checked) [data-testid="stDataFrame"]  [role="gridcell"]:not([style*="color"]),
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="gridcell"]:not([style*="color"]) *,
    body:has(#pms-d:checked) [data-testid="stDataFrame"]  [role="gridcell"]:not([style*="color"]) * {
        color: #ffffff !important;
    }

    /* 8. 셀 편집 팝업 (SelectboxColumn 등) */
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-popup,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-popup-editor,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-popup-child {
        background-color: #252535 !important; border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select-list {
        background-color: #252535 !important; border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select-row,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-virtual-list-item {
        background-color: #252535 !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select-row:hover,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select-row-selected {
        background-color: #3a3a5e !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-text-field-input {
        background-color: #252535 !important; color: #ffffff !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-checkbox-input-wrapper::after {
        color: #ffffff !important;
    }

    /* HTML 테이블 */
    body:has(#pms-d:checked) table { background-color: #252535 !important; border-color: #45475a !important; }
    body:has(#pms-d:checked) th {
        background-color: #0f0f1f !important; color: #ffffff !important;
        border-color: #45475a !important; font-weight: 700 !important;
    }
    body:has(#pms-d:checked) td {
        background-color: #252535 !important; color: #ffffff !important;
        border-color: #313244 !important;
    }
    body:has(#pms-d:checked) tr:nth-child(odd) td { background-color: #1e1e30 !important; }
    body:has(#pms-d:checked) tr:hover td { background-color: #2d2d45 !important; }

    /* 증감 색 — 다크모드에서 더 선명하게 */
    body:has(#pms-d:checked) [style*="color:#E53E3E"] { color: #ff7b7b !important; }
    body:has(#pms-d:checked) [style*="color:#3182CE"] { color: #74b9ff !important; }

    /* 체크박스 / 라디오 */
    body:has(#pms-d:checked) [data-baseweb="checkbox"] div,
    body:has(#pms-d:checked) [data-baseweb="radio"] div { border-color: #ffffff !important; }
    body:has(#pms-d:checked) [data-baseweb="checkbox"] [aria-checked="true"] div,
    body:has(#pms-d:checked) [data-baseweb="radio"] [aria-checked="true"] div {
        background-color: #ffffff !important;
    }

    /* 스크롤바 */
    body:has(#pms-d:checked) ::-webkit-scrollbar { width: 6px; height: 6px; }
    body:has(#pms-d:checked) ::-webkit-scrollbar-track { background: #252535; }
    body:has(#pms-d:checked) ::-webkit-scrollbar-thumb { background: #45475a; border-radius: 3px; }
    body:has(#pms-d:checked) ::-webkit-scrollbar-thumb:hover { background: #585b70; }

    /* 테마 토글 바 */
    body:has(#pms-d:checked) .pms-sw-track { background: #313244; }
    body:has(#pms-d:checked) .pms-btn { color: #ffffff; }

    /* Streamlit canvas-backed tables */
    html[data-pms-theme="d"] [data-testid="stDataEditor"],
    html[data-pms-theme="d"] [data-testid="stDataFrame"],
    body[data-pms-theme="d"] [data-testid="stDataEditor"],
    body[data-pms-theme="d"] [data-testid="stDataFrame"],
    .stApp:has(#pms-d:checked) [data-testid="stDataEditor"],
    .stApp:has(#pms-d:checked) [data-testid="stDataFrame"],
    body:has(#pms-d:checked) [data-testid="stDataEditor"],
    body:has(#pms-d:checked) [data-testid="stDataFrame"] {
        background-color: #1e1e2e !important;
        border-color: #45475a !important;
        color: #ffffff !important;
        border-radius: 8px !important;
        overflow: hidden !important;
    }
    html[data-pms-theme="d"] [data-testid="stDataEditor"] canvas,
    html[data-pms-theme="d"] [data-testid="stDataFrame"] canvas,
    body[data-pms-theme="d"] [data-testid="stDataEditor"] canvas,
    body[data-pms-theme="d"] [data-testid="stDataFrame"] canvas,
    .stApp:has(#pms-d:checked) [data-testid="stDataEditor"] canvas,
    .stApp:has(#pms-d:checked) [data-testid="stDataFrame"] canvas,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] canvas,
    body:has(#pms-d:checked) [data-testid="stDataFrame"] canvas {
        filter: invert(1) hue-rotate(180deg) brightness(0.92) contrast(0.88) saturate(0.9) !important;
    }
    html[data-pms-theme="d"] [data-testid="stDataEditor"] [role="columnheader"],
    html[data-pms-theme="d"] [data-testid="stDataEditor"] [role="gridcell"],
    html[data-pms-theme="d"] [data-testid="stDataFrame"] [role="columnheader"],
    html[data-pms-theme="d"] [data-testid="stDataFrame"] [role="gridcell"],
    .stApp:has(#pms-d:checked) [data-testid="stDataEditor"] [role="columnheader"],
    .stApp:has(#pms-d:checked) [data-testid="stDataEditor"] [role="gridcell"],
    .stApp:has(#pms-d:checked) [data-testid="stDataFrame"] [role="columnheader"],
    .stApp:has(#pms-d:checked) [data-testid="stDataFrame"] [role="gridcell"] {
        background-color: #1e1e2e !important;
        border-color: #45475a !important;
        color: #ffffff !important;
        font-weight: 700 !important;
    }

    @media (prefers-color-scheme: dark) {
        body:has(#pms-s:checked) .stApp,
        body:has(#pms-s:checked) [data-testid="stAppViewContainer"],
        body:has(#pms-s:checked) section.main {
            background-color: #1e1e2e !important; color: #ffffff !important;
        }
        body:has(#pms-s:checked) .main .block-container { background-color: #1e1e2e !important; }
        body:has(#pms-s:checked) h1, body:has(#pms-s:checked) h2, body:has(#pms-s:checked) h3,
        body:has(#pms-s:checked) h4, body:has(#pms-s:checked) h5, body:has(#pms-s:checked) h6 { color: #ffffff !important; }
        body:has(#pms-s:checked) p, body:has(#pms-s:checked) [data-testid="stMarkdownContainer"] p { color: #ffffff !important; }
        body:has(#pms-s:checked) label { color: #ffffff !important; }
        body:has(#pms-s:checked) span:not([style*="color"]) { color: #ffffff !important; }
        body:has(#pms-s:checked) input, body:has(#pms-s:checked) [data-baseweb="input"] input {
            background-color: #252535 !important; color: #ffffff !important; border-color: #45475a !important;
        }
        body:has(#pms-s:checked) [class*="ag-theme"] {
            --ag-background-color: #252535; --ag-foreground-color: #ffffff;
            --ag-secondary-foreground-color: #ffffff;
            --ag-header-background-color: #0f0f1f; --ag-header-foreground-color: #ffffff;
            --ag-border-color: #45475a; --ag-row-border-color: #313244;
            --ag-odd-row-background-color: #1e1e30; --ag-row-hover-color: #2d2d45;
        }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-header,
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-header-row { background-color: #0f0f1f !important; }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-header-cell-text { color: #ffffff !important; font-weight: 700 !important; }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-row-even { background-color: #252535 !important; }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-row-odd  { background-color: #1e1e30 !important; }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-cell,
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-cell-value { color: #ffffff !important; }
        body:has(#pms-s:checked) .js-plotly-plot,
        body:has(#pms-s:checked) .js-plotly-plot .plot-container,
        body:has(#pms-s:checked) .js-plotly-plot .svg-container {
            background: transparent !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .bglayer rect.bg {
            fill: transparent !important;
            stroke: transparent !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .legend rect.bg {
            fill: #252535 !important;
            stroke: #45475a !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot text {
            fill: #ffffff !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .gridlayer path {
            stroke: #45475a !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .zerolinelayer path,
        body:has(#pms-s:checked) .js-plotly-plot .xlines-below path,
        body:has(#pms-s:checked) .js-plotly-plot .ylines-below path {
            stroke: #585b70 !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .modebar-btn path {
            fill: #cdd6f4 !important;
        }
        body:has(#pms-s:checked) [data-testid="stSidebar"],
        body:has(#pms-s:checked) [data-testid="stSidebarContent"] {
            background-color: #171717 !important;
        }
        body:has(#pms-s:checked) .stButton > button,
        body:has(#pms-s:checked) [data-testid="stButton"] button,
        body:has(#pms-s:checked) button[data-testid^="baseButton"] {
            background: #2a2a3e !important;
            background-color: #2a2a3e !important;
            color: #ffffff !important;
            border-color: #45475a !important;
        }
        body:has(#pms-s:checked) .stButton > button:hover,
        body:has(#pms-s:checked) [data-testid="stButton"] button:hover,
        body:has(#pms-s:checked) button[data-testid^="baseButton"]:hover {
            background: #313244 !important;
            background-color: #313244 !important;
            border-color: #ffffff !important;
        }
        body:has(#pms-s:checked) button.pms-home-btn {
            background: transparent !important;
            background-color: transparent !important;
            border: none !important;
            box-shadow: none !important;
        }
        body:has(#pms-s:checked) button.pms-home-btn::before {
            background-color: #cdd6f4;
        }
        body:has(#pms-s:checked) button.pms-home-btn:hover {
            background: transparent !important;
            background-color: transparent !important;
            border: none !important;
        }
        body:has(#pms-s:checked) button.pms-home-btn:hover::before {
            background-color: #818cf8;
        }
        body:has(#pms-s:checked) [data-baseweb="tab"][aria-selected="true"] { color: #ffffff !important; }
        body:has(#pms-s:checked) .pms-sw-track { background: #313244; }
    }
    </style>

    <div class="pms-sw-outer">
        <div class="pms-sw-track">
            <label class="pms-btn" title="라이트">
                <input type="radio" class="pms-theme-radio" name="pms-theme" id="pms-l" checked>
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><circle cx="12" cy="12" r="5"/><line x1="12" y1="1" x2="12" y2="3"/><line x1="12" y1="21" x2="12" y2="23"/><line x1="4.22" y1="4.22" x2="5.64" y2="5.64"/><line x1="18.36" y1="18.36" x2="19.78" y2="19.78"/><line x1="1" y1="12" x2="3" y2="12"/><line x1="21" y1="12" x2="23" y2="12"/><line x1="4.22" y1="19.78" x2="5.64" y2="18.36"/><line x1="18.36" y1="5.64" x2="19.78" y2="4.22"/></svg>
            </label>
            <label class="pms-btn" title="다크">
                <input type="radio" class="pms-theme-radio" name="pms-theme" id="pms-d">
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 12.79A9 9 0 1 1 11.21 3 7 7 0 0 0 21 12.79z"/></svg>
            </label>
        </div>
    </div>

    <!-- localStorage 테마 영속성: onload 인라인 핸들러로 실행 (script 태그는 React innerHTML에서 미실행) -->
    <img src="data:image/gif;base64,R0lGODlhAQABAIAAAAAAAP///yH5BAEAAAAALAAAAAABAAEAAAIBRAA7"
         onload="(function(){
             function readCookie(name){
                 var found = document.cookie.split('; ').find(function(row){ return row.indexOf(name + '=') === 0; });
                 return found ? decodeURIComponent(found.split('=')[1]) : '';
             }
             function normalizeTheme(t){
                 return (t === 'd' || t === 'l') ? t : 'l';
             }
             function saveTheme(t){
                 t = normalizeTheme(t);
                 localStorage.setItem('pms-theme', t);
                 document.cookie = 'pms-theme=' + encodeURIComponent(t) + '; max-age=31536000; path=/; SameSite=Lax';
             }
             function applyTheme(){
                 var t = normalizeTheme(localStorage.getItem('pms-theme') || readCookie('pms-theme') || 'l');
                 saveTheme(t);
                 document.documentElement.setAttribute('data-pms-theme', t);
                 if(document.body){ document.body.setAttribute('data-pms-theme', t); }
                 var r = document.getElementById('pms-' + t);
                 if(r && !r.checked){ r.checked = true; }
             }
             function tagSpecialBtns(){
                 document.querySelectorAll('button').forEach(function(btn){
                     var t = btn.textContent.trim();
                     if(t === '🏠') btn.classList.add('pms-home-btn');
                     if(t === '↻') btn.classList.add('pms-refresh-btn');
                 });
             }
             applyTheme();
             tagSpecialBtns();
             if(!window._pmsThemeReady){
                 window._pmsThemeReady = true;
                 var debounce;
                 window._pmsThemeObs = new MutationObserver(function(){
                     clearTimeout(debounce);
                     debounce = setTimeout(function(){ applyTheme(); tagSpecialBtns(); }, 80);
                 });
                 window._pmsThemeObs.observe(document.body, {childList:true, subtree:true});
                 document.addEventListener('change', function(e){
                     if(e.target && e.target.name === 'pms-theme'){
                         saveTheme(e.target.id.replace('pms-',''));
                         applyTheme();
                     }
                 });
             }
         })()"
         style="display:none" alt="">
    """, unsafe_allow_html=True)

    # ── 모바일 반응형 CSS ──────────────────────────────────────────
    st.markdown("""
    <style>
    @media (max-width: 768px) {
        /* ① 사이드바: 모바일에서 전체 폭 오버레이 */
        [data-testid="stSidebar"] {
            width: 100vw !important;
            min-width: 100vw !important;
            max-width: 100vw !important;
        }
        section[data-testid="stSidebar"] > div:first-child {
            width: 100vw !important;
        }

        /* ② 메인 컨텐츠: 전체 폭, 패딩 최소화 */
        .block-container {
            padding: 0.5rem 0.75rem 2rem !important;
            max-width: 100% !important;
        }
        [data-testid="stAppViewContainer"] > section.main {
            width: 100% !important;
        }

        /* ③ 테마 토글 위치·크기 축소 */
        .pms-sw-outer { top: 6px; right: 6px; }
        .pms-btn { padding: 0 6px; height: 22px; font-size: 10px; }

        /* ④ 표: 가로 스크롤 */
        .pms-report-table {
            overflow-x: auto !important;
            -webkit-overflow-scrolling: touch !important;
            max-width: 100vw !important;
        }
        table { font-size: 11px !important; min-width: max-content; }
        table th, table td { padding: 4px 6px !important; white-space: nowrap !important; }

        /* ⑤ data_editor 가로 스크롤 */
        [data-testid="stDataEditor"] {
            overflow-x: auto !important;
            -webkit-overflow-scrolling: touch !important;
        }

        /* ⑥ 컬럼 레이아웃: 2개 이상이면 가로 스크롤 */
        [data-testid="stHorizontalBlock"] {
            overflow-x: auto !important;
            -webkit-overflow-scrolling: touch !important;
            flex-wrap: nowrap !important;
            gap: 4px !important;
        }
        [data-testid="stHorizontalBlock"] > [data-testid="stColumn"] {
            min-width: 120px !important;
            flex-shrink: 0 !important;
        }

        /* ⑦ 버튼 터치 영역 */
        .stButton > button { min-height: 44px !important; font-size: 14px !important; }

        /* ⑧ 입력 필드 — iOS 자동 확대 방지 */
        input, textarea, select { font-size: 16px !important; }

        /* ⑨ 메트릭 */
        [data-testid="stMetricValue"] { font-size: 1.1rem !important; }
        [data-testid="stMetricLabel"] { font-size: 0.7rem !important; }
        [data-testid="metric-container"] { padding: 8px !important; }

        /* ⑩ 제목 */
        h1 { font-size: 1.3rem !important; }
        h2 { font-size: 1.15rem !important; }
        h3 { font-size: 1.05rem !important; }
        h4 { font-size: 0.95rem !important; }

        /* ⑪ 탭 */
        [data-baseweb="tab"] { padding: 6px 8px !important; font-size: 12px !important; }

        /* ⑫ 사이드바 햄버거 버튼 */
        [data-testid="stSidebarCollapsedControl"] button {
            width: 44px !important; height: 44px !important;
        }
    }
    </style>
    """, unsafe_allow_html=True)


def show_all_staff_summary(staff_names):
    """전체 직원 실적 요약 표시"""
    # 본사 구글시트 로드
    if st.session_state.get("cloud_sheet_df") is None:
        try:
            with st.spinner("본사 구글시트를 불러오는 중..."):
                load_csv_to_state("url_sync", "cloud_sheet_df")
        except Exception as e:
            st.error("⚠️ 본사 구글시트를 불러오지 못했습니다.")
            with st.expander("📋 오류 상세 정보"):
                st.code(str(e))
            st.info("💡 [구글 스트레드시트 연동] 메뉴에서 '본사 구글 시트 CSV URL'을 확인해주세요.")

            if st.button("🔗 구글 스트레드시트 연동 메뉴로 이동", use_container_width=True):
                st.session_state.current_menu = "구글 스트레드시트 연동"
                st.rerun()
            return

    cloud_df = st.session_state.get("cloud_sheet_df")
    if cloud_df is None or cloud_df.empty:
        st.warning("📂 본사 구글시트 데이터가 없습니다.")
        st.info("💡 [구글 스트레드시트 연동] 메뉴에서 URL 설정 및 데이터 새로고침을 해주세요.")

        if st.button("🔗 구글 스트레드시트 연동 메뉴로 이동", use_container_width=True):
            st.session_state.current_menu = "구글 스트레드시트 연동"
            st.rerun()
        return

    # 운영 실적용 데이터: 업로드한 엑셀 파일만 사용 (없으면 None)
    analysis_df = st.session_state.get("admin_uploaded_excel")

    # 은행 구글시트 로드 (하나은행 시트)
    hana_df = None
    if st.session_state.get("hana_sheet_df") is None:
        try:
            load_csv_to_state("url_hana", "hana_sheet_df")
        except:
            pass
    hana_df = st.session_state.get("hana_sheet_df")

    # 데이터 정리
    cloud_df = clean_header_logic(cloud_df.copy())

    # 본사 시트 컬럼 찾기
    owner_col = find_col(cloud_df, ["담당자", "등록자", "성명"])
    open_date_col = find_col(cloud_df, ["개설완료일자", "개설일"])
    erp_date_col = find_col(cloud_df, ["ERP연계일자", "연계일자"])
    biz_num_col = find_col(cloud_df, ["사업자번호", "사업자등록번호"])
    company_col = find_col(cloud_df, ["고객명", "업체명", "상호"])

    if not owner_col or owner_col not in cloud_df.columns:
        st.warning("본사 시트에서 담당자 컬럼을 찾을 수 없습니다.")
        return

    # 은행 시트와 본사 시트 비교 - 누락 데이터 확인
    if hana_df is not None and not hana_df.empty:
        hana_clean = clean_header_logic(hana_df.copy())
        hana_biz_col = find_col(hana_clean, ["사업자번호", "사업자등록번호"])
        hana_company_col = find_col(hana_clean, ["고객명", "업체명", "상호"])
        hana_open_status_col = find_col(hana_clean, ["개설상태"])
        hana_erp_status_col = find_col(hana_clean, ["연계상태", "ERP연계상태"])
        hana_owner_col = find_col(hana_clean, ["담당자", "등록자"])

        cloud_status_col = find_col(cloud_df, ["상태"])

        missing_open = []
        missing_erp = []

        if hana_biz_col and biz_num_col:
            # 개설상태가 "개설완료"인데 본사 시트 상태가 "완료"가 아닌 경우
            if hana_open_status_col:
                hana_open_completed = hana_clean[
                    hana_clean[hana_open_status_col].astype(str).str.contains("개설완료", na=False)
                ].copy()
                for _, row in hana_open_completed.iterrows():
                    biz_num = str(row.get(hana_biz_col, "")).strip()
                    if not biz_num:
                        continue
                    # 본사 시트에서 해당 사업자번호 찾기
                    cloud_match = cloud_df[cloud_df[biz_num_col].astype(str).str.contains(biz_num, na=False, regex=False)]
                    if not cloud_match.empty:
                        # 상태가 "완료"가 아닌 경우
                        if cloud_status_col and cloud_status_col in cloud_match.columns:
                            cloud_status = cloud_match[cloud_status_col].astype(str).str.strip()
                            if not cloud_status.str.contains("완료", na=False).any():
                                missing_open.append({
                                    "사업자번호": biz_num,
                                    "고객명": row.get(hana_company_col, ""),
                                    "담당자": row.get(hana_owner_col, ""),
                                    "은행시트_개설상태": row.get(hana_open_status_col, ""),
                                    "본사시트_상태": cloud_match[cloud_status_col].iloc[0] if not cloud_match.empty else ""
                                })

            # 연계상태가 "연계완료" 또는 "청구완료"인데 본사 시트 ERP연계일자가 없는 경우
            if hana_erp_status_col:
                hana_erp_completed = hana_clean[
                    hana_clean[hana_erp_status_col].astype(str).str.contains("연계완료|청구완료", na=False, regex=True)
                ].copy()
                for _, row in hana_erp_completed.iterrows():
                    biz_num = str(row.get(hana_biz_col, "")).strip()
                    if not biz_num:
                        continue
                    # 본사 시트에서 해당 사업자번호 찾기
                    cloud_match = cloud_df[cloud_df[biz_num_col].astype(str).str.contains(biz_num, na=False, regex=False)]
                    if not cloud_match.empty:
                        # ERP연계일자가 없는 경우
                        if erp_date_col and erp_date_col in cloud_match.columns:
                            if cloud_match[erp_date_col].isna().all():
                                missing_erp.append({
                                    "사업자번호": biz_num,
                                    "고객명": row.get(hana_company_col, ""),
                                    "담당자": row.get(hana_owner_col, ""),
                                    "은행시트_연계상태": row.get(hana_erp_status_col, "")
                                })

        # 누락 데이터 표시
        if missing_open or missing_erp:
            with st.expander("⚠️ 본사 시트 누락 데이터", expanded=False):
                if missing_open:
                    st.markdown("#### 개설 상태 불일치")
                    st.write(f"은행 시트는 개설완료이나 본사 시트 상태가 완료가 아닌 고객사: **{len(missing_open)}건**")
                    st.dataframe(pd.DataFrame(missing_open), use_container_width=True)
                    st.markdown("---")

                if missing_erp:
                    st.markdown("#### ERP연계일자 누락")
                    st.write(f"은행 시트는 연계완료/청구완료이나 본사 시트에 ERP연계일자가 없는 고객사: **{len(missing_erp)}건**")
                    st.dataframe(pd.DataFrame(missing_erp), use_container_width=True)

    # 날짜 컬럼 변환
    if open_date_col and open_date_col in cloud_df.columns:
        cloud_df[open_date_col] = pd.to_datetime(cloud_df[open_date_col], errors="coerce")
    if erp_date_col and erp_date_col in cloud_df.columns:
        cloud_df[erp_date_col] = pd.to_datetime(cloud_df[erp_date_col], errors="coerce")

    # 대상 년월 (2026-05)
    target_year_month = "2026-05"

    # 직원 정보 로드 (직급 정보 가져오기)
    user_db = load_db(DB_FILE, {})

    # 하나지사 실적관리 시트 로드 (방문A 카운트용)
    # 복합 헤더(다중 행)이므로 header=None으로 원본 로딩 후 직접 파싱
    def _load_perf_sheet():
        url = normalize_google_sheet_csv_url(st.session_state.get("url_hana_performance", DEFAULT_URL_HANA_PERFORMANCE))
        raw = None
        _err = ""
        # 방법1: read_csv_cached (캐시 활용)
        try:
            raw = read_csv_cached(url, header=None).copy()
            raw = raw.astype(str)
        except Exception as e1:
            _err += f"cached: {e1} | "
        # 방법2: pandas 직접 로드
        if raw is None:
            try:
                raw = pd.read_csv(url, header=None, dtype=str)
            except Exception as e2:
                _err += f"direct: {e2}"
        if raw is None:
            st.session_state["_perf_debug_err"] = _err
            return None, None, None, None

        # "방문A" 가 들어있는 행을 헤더 행으로 탐색
        header_row_idx = None
        visit_a_col_idx = None
        owner_col_idx = None
        payout_col_idx = None
        for i in range(min(15, len(raw))):
            for j, val in enumerate(raw.iloc[i]):
                if str(val).strip() == "방문A":
                    header_row_idx = i
                    visit_a_col_idx = j
                    break
            if header_row_idx is not None:
                break

        st.session_state["_perf_debug_raw"] = raw.head(5).to_dict()
        if header_row_idx is not None:
            for i in range(max(0, header_row_idx - 3), min(len(raw), header_row_idx + 2)):
                for j, val in enumerate(raw.iloc[i]):
                    if str(val).strip() in ["이름", "담당자", "성명"]:
                        owner_col_idx = j
                    if str(val).strip() in ["배분금액", "지급예상금액", "금액"]:
                        payout_col_idx = j
                if owner_col_idx is not None and payout_col_idx is not None:
                    break
        st.session_state["_perf_debug_hdr"] = (header_row_idx, owner_col_idx, visit_a_col_idx, payout_col_idx)

        if header_row_idx is None or owner_col_idx is None or visit_a_col_idx is None:
            return None, None, None, None

        data = raw.iloc[header_row_idx + 1:].reset_index(drop=True)
        result = pd.DataFrame({
            "_owner":   data.iloc[:, owner_col_idx].astype(str).str.strip(),
            "_visit_a": data.iloc[:, visit_a_col_idx].astype(str).str.strip(),
        })
        if payout_col_idx is not None:
            result["_payout"] = data.iloc[:, payout_col_idx].astype(str).str.strip()
        else:
            result["_payout"] = ""
        result = result[
            result["_owner"].ne("")
            & result["_owner"].ne("nan")
            & result["_owner"].ne("합계")
        ].reset_index(drop=True)
        return result, "_owner", "_visit_a", "_payout"

    _perf_url = normalize_google_sheet_csv_url(st.session_state.get("url_hana_performance", DEFAULT_URL_HANA_PERFORMANCE))
    _perf_parser_version = 3
    _perf_cache = st.session_state.get("_perf_parsed")
    _perf_cache_url = st.session_state.get("_perf_parsed_url")
    _perf_cache_version = st.session_state.get("_perf_parsed_version")
    if _perf_cache is None or _perf_cache_url != _perf_url or _perf_cache_version != _perf_parser_version:
        _perf_df, _perf_owner, _perf_visit, _perf_payout = _load_perf_sheet()
        if _perf_df is not None:
            st.session_state["_perf_parsed"] = (_perf_df, _perf_owner, _perf_visit, _perf_payout)
            st.session_state["_perf_parsed_url"] = _perf_url
            st.session_state["_perf_parsed_version"] = _perf_parser_version
        _perf_cache = (_perf_df, _perf_owner, _perf_visit, _perf_payout)
    else:
        if len(_perf_cache) == 3:
            _perf_df, _perf_owner, _perf_visit = _perf_cache
            _perf_payout = None
        else:
            _perf_df, _perf_owner, _perf_visit, _perf_payout = _perf_cache

    def count_visit_a(staff):
        """담당자의 방문A 값 (실적관리 시트 요약표 기준)"""
        if _perf_df is None or _perf_owner is None or _perf_visit is None:
            return None
        staff_key = str(staff).strip()
        rows = _perf_df[_perf_df[_perf_owner].astype(str).str.strip() == staff_key]
        if rows.empty:
            return None

        def parse_visit_count(value):
            text = str(value).strip().replace(",", "")
            if not text or text.lower() == "nan":
                return 0
            match = re.search(r"-?\d+(?:\.\d+)?", text)
            if not match:
                return 0
            try:
                return int(round(float(match.group())))
            except Exception:
                return 0

        return int(rows[_perf_visit].map(parse_visit_count).sum())

    def perf_sheet_payout_amount(staff):
        """담당자의 배분금액 값 (실적관리 시트 계산 결과 기준)"""
        if _perf_df is None or _perf_owner is None or not _perf_payout:
            return None
        staff_key = str(staff).strip()
        rows = _perf_df[_perf_df[_perf_owner].astype(str).str.strip() == staff_key]
        if rows.empty:
            return None

        def parse_amount(value):
            text = str(value).strip().replace(",", "")
            if not text or text.lower() == "nan" or text == "-":
                return None
            match = re.search(r"-?\d+(?:\.\d+)?", text)
            if not match:
                return None
            try:
                return int(round(float(match.group())))
            except Exception:
                return None

        amounts = [amount for amount in rows[_perf_payout].map(parse_amount).tolist() if amount is not None]
        return amounts[0] if amounts else None

    def apply_perf_sheet_payouts(df):
        if df is None or not isinstance(df, pd.DataFrame) or df.empty or "담당자" not in df.columns:
            return df
        result_df = df.copy()
        for idx, row in result_df.iterrows():
            staff = row.get("담당자", "")
            if str(staff).strip() == "합계":
                continue
            amount = perf_sheet_payout_amount(staff)
            if amount is None:
                continue
            result_df.at[idx, "지급예상금액"] = amount
            result_df.at[idx, "지급포인트"] = int(round(amount / 500))
        return result_df

    # ── 디버그: 실적관리 시트 상태 확인 ──
    with st.expander("🔍 실적관리 시트 디버그 (확인 후 제거 예정)", expanded=False):
        st.write(f"_perf_df: {'로드됨' if _perf_df is not None else 'None'}")
        if st.session_state.get("_perf_debug_err"):
            st.error(f"로드 오류: {st.session_state['_perf_debug_err']}")
        if st.session_state.get("_perf_debug_hdr"):
            st.write(f"헤더행/방문A열 인덱스: {st.session_state['_perf_debug_hdr']}")
        if st.session_state.get("_perf_debug_raw"):
            st.write("원본 상위 5행:")
            st.write(st.session_state["_perf_debug_raw"])
        if _perf_df is not None:
            st.write(f"담당자 목록: {_perf_df['_owner'].tolist()}")
            st.write(f"count_visit_a('이성환') = {count_visit_a('이성환')}")
            st.write(f"perf_sheet_payout_amount('이성환') = {perf_sheet_payout_amount('이성환')}")
            st.dataframe(_perf_df.head(10))

    # 실적 데이터 계산
    performance_data = []

    for staff_name in staff_names:
        # 본사 시트에서 담당자 데이터 필터링
        staff_cloud_df = cloud_df[cloud_df[owner_col] == staff_name].copy()

        # 개설건수: 2026년 5월 개설완료일자 카운트
        open_count = 0
        open_companies = []
        if open_date_col and open_date_col in staff_cloud_df.columns:
            may_2026_data = staff_cloud_df[
                staff_cloud_df[open_date_col].notna() &
                (staff_cloud_df[open_date_col].dt.strftime("%Y-%m") == target_year_month)
            ]
            open_count = len(may_2026_data)
            # 개설고객사 목록
            if company_col and company_col in may_2026_data.columns:
                open_companies = may_2026_data[company_col].astype(str).tolist()

        # 연계건수: 2026년 5월 ERP연계일자 카운트
        erp_count = 0
        erp_companies = []
        if erp_date_col and erp_date_col in staff_cloud_df.columns:
            may_erp_data = staff_cloud_df[
                staff_cloud_df[erp_date_col].notna() &
                (staff_cloud_df[erp_date_col].dt.strftime("%Y-%m") == target_year_month)
            ]
            erp_count = len(may_erp_data)
            # 연계고객사 목록
            if company_col and company_col in may_erp_data.columns:
                erp_companies = may_erp_data[company_col].astype(str).tolist()

        # 운영 실적: 하나지사 실적관리 시트 방문A 카운트 우선, 없으면 업로드 파일 기준
        operation_count = 0
        operation_points = 0

        visit_a_count = count_visit_a(staff_name)
        if visit_a_count is not None:
            operation_count = visit_a_count
            operation_points = operation_count * 30

        if False and analysis_df is not None and not analysis_df.empty:
            analysis_clean = clean_header_logic(analysis_df.copy())
            u_col = find_col(analysis_clean, ["등록자", "담당자", "성명"])
            d_col = find_col(analysis_clean, ["활동상세", "활동내용"])
            category_col = find_col(analysis_clean, ["활동구분", "접수유형"])

            if u_col and u_col in analysis_clean.columns:
                staff_activity = analysis_clean[analysis_clean[u_col] == staff_name].copy()

                if category_col and category_col in staff_activity.columns:
                    # 본사이력 업로드는 활동구분(방문/원격)을 우선 기준으로 운영 실적을 반영한다.
                    category_text = staff_activity[category_col].astype(str)
                    visit_count = int(category_text.str.contains("방문", na=False).sum())
                    remote_count = int(category_text.str.contains("원격", na=False).sum())
                    operation_count_raw = visit_count + remote_count
                    # 운영건수는 최대 60회까지만 카운트
                    operation_count = min(60, operation_count_raw)
                    operation_points = (min(60, visit_count) * 30) + (min(60 - min(60, visit_count), remote_count) * 10)
                elif d_col and d_col in staff_activity.columns:
                    # 활동구분이 없으면 기존처럼 활동상세 기준으로 운영 실적을 계산한다.
                    operation_count_raw = int(
                        staff_activity[d_col].astype(str).str.contains("운영|방문|점검", na=False).sum()
                    )
                    # 운영건수는 최대 60회까지만 카운트
                    operation_count = min(60, operation_count_raw)
                    operation_points = operation_count * 30

        # 포인트 계산
        open_points = open_count * 90  # 개설 1건당 90포인트
        erp_points = erp_count * 120  # 연계 1건당 120포인트

        # 합계포인트는 개설포인트 + 연계포인트 + 운영포인트 합계이며 최대 2800점이다.
        total_points = min(2800, open_points + erp_points + operation_points)

        # 직급 정보 가져오기
        staff_rank = "직원"  # 기본값
        for uid, user_info in user_db.items():
            if isinstance(user_info, dict) and user_info.get("name") == staff_name:
                staff_rank = user_info.get("rank", "직원")
                break

        row_data = {
            "담당자": staff_name,
            "직급": staff_rank,
            "개설건수": open_count,
            "개설고객사": ", ".join(open_companies) if open_companies else "",
            "개설포인트": open_points,
            "연계건수": erp_count,
            "연계고객사": ", ".join(erp_companies) if erp_companies else "",
            "연계포인트": erp_points,
            "운영건수 (실제 활동)": operation_count,
            "운영포인트 (실제 활동)": operation_points,
            "운영포인트 (추가 활동)": 0,
            "합계포인트": total_points,
            "지급포인트": 0,
            "지급예상금액": 0,
        }

        performance_data.append(row_data)

    if not performance_data:
        st.info("집계된 실적이 없습니다.")
        return

    # DataFrame 생성
    perf_df = pd.DataFrame(performance_data)
    perf_df, debug_info = apply_rs_allowance_formula(perf_df, user_db, return_debug=True)
    perf_df = apply_perf_sheet_payouts(perf_df)

    # 디버깅 정보 표시
    st.markdown("---")
    st.markdown("### 🔍 계산 디버깅 정보")

    if debug_info:
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("BU 합산포인트", f"{debug_info['BU합산']:,}")
            st.metric("BU 평균", f"{debug_info['BU평균']:.2f}")
        with col2:
            st.metric("BU 인원", f"{debug_info['BU인원']}명")
            st.metric("외주직원 수", f"{debug_info['외주직원수']}명")
        with col3:
            st.metric("일반직원 수", f"{debug_info['일반직원수']}명")
            st.metric("팀장수당", f"{debug_info['팀장수당']}pt")

        st.divider()
        col4, col5 = st.columns(2)
        with col4:
            st.metric("외주가감 총합", f"{debug_info['외주가감총합']:.2f}pt")
        with col5:
            st.metric("일반직원 1인당 분배", f"{debug_info['일반직원1인당분배']:.2f}pt")
    else:
        st.warning("디버깅 정보가 없습니다.")

    st.markdown("---")

    # 직급 순서 정의 및 정렬용 컬럼 추가
    rank_order = {"부서장": 0, "팀장": 1, "과장": 2, "대리": 3, "주임": 4, "직원": 5}
    perf_df["_rank_order"] = perf_df["직급"].map(lambda x: rank_order.get(x, 99))

    # 필터 및 검색
    col1, col2, col3 = st.columns([2, 2, 1])

    with col1:
        search_name = st.text_input("🔍 담당자 검색", placeholder="이름 입력")

    with col2:
        sort_by = st.selectbox(
            "정렬 기준",
            ["직급", "합계포인트", "지급예상금액", "개설건수", "연계건수", "운영건수 (실제 활동)", "담당자"],
            index=0
        )

    with col3:
        sort_order = st.radio("정렬 순서", ["오름차순 ⬆", "내림차순 ⬇"], horizontal=True)

    # 검색 필터 적용
    if search_name:
        perf_df = perf_df[perf_df["담당자"].astype(str).str.contains(search_name, na=False)]

    if perf_df.empty:
        st.warning("검색 결과가 없습니다.")
        return

    # 정렬 적용
    ascending = sort_order == "오름차순 ⬆"
    if sort_by == "직급":
        perf_df = perf_df.sort_values(by="_rank_order", ascending=ascending)
    else:
        perf_df = perf_df.sort_values(by=sort_by, ascending=ascending)

    # 통계 요약
    st.markdown("#### 📊 전체 실적 요약")
    summary_cols = st.columns(5)

    with summary_cols[0]:
        st.metric("총 인원", f"{len(perf_df)}명")

    with summary_cols[1]:
        total_open = int(perf_df["개설건수"].sum()) if "개설건수" in perf_df.columns else 0
        st.metric("총 개설", f"{total_open}건")

    with summary_cols[2]:
        total_link = int(perf_df["연계건수"].sum()) if "연계건수" in perf_df.columns else 0
        st.metric("총 연계", f"{total_link}건")

    with summary_cols[3]:
        total_points = int(perf_df["합계포인트"].sum()) if "합계포인트" in perf_df.columns else 0
        st.metric("총 포인트", f"{total_points:,}pt")

    with summary_cols[4]:
        total_amount = int(perf_df["지급예상금액"].sum()) if "지급예상금액" in perf_df.columns else 0
        st.metric("총 지급예상", f"{total_amount:,}원")

    st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)

    # 표시할 컬럼 선택
    display_cols = [
        "담당자", "직급", "개설건수", "개설포인트", "연계건수", "연계포인트",
        "운영건수 (실제 활동)", "운영포인트 (실제 활동)", "운영포인트 (추가 활동)",
        "합계포인트", "지급포인트", "지급예상금액", "오류건수", "중복건수"
    ]

    # 존재하는 컬럼만 선택
    display_cols = [col for col in display_cols if col in perf_df.columns]

    # 합계 행 추가
    total_row = {}
    for col in display_cols:
        if col in ["담당자"]:
            total_row[col] = "합계"
        elif col in ["직급"]:
            total_row[col] = ""
        elif col in perf_df.columns:
            if perf_df[col].dtype in ['int64', 'float64']:
                total_row[col] = perf_df[col].sum()
            else:
                total_row[col] = ""
        else:
            total_row[col] = ""

    perf_df_with_total = pd.concat([perf_df[display_cols], pd.DataFrame([total_row])], ignore_index=True)

    # Excel 다운로드 및 새로고침
    st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)

    col_dl1, col_dl2, col_dl3 = st.columns([1, 1, 2])

    with col_dl1:
        from datetime import datetime
        download_time = datetime.now().strftime("%Y%m%d_%H%M%S")
        excel_bytes = BytesIO()
        perf_df[display_cols].to_excel(excel_bytes, index=False, engine="openpyxl")
        excel_bytes.seek(0)

        st.download_button(
            "📥 Excel 다운로드",
            data=excel_bytes.getvalue(),
            file_name=f"전체실적_{download_time}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True
        )

    with col_dl2:
        if st.button("🔄 새로고침", use_container_width=True):
            # 본사 구글시트와 하나은행 시트 데이터 새로고침
            st.session_state.cloud_sheet_df = None
            st.session_state.analysis_lookup_df = None
            st.session_state.hana_sheet_df = None
            st.session_state.admin_uploaded_excel = None
            st.session_state.admin_uploaded_excel_display = None
            st.session_state.admin_history_upload_key = None
            st.session_state.admin_office_upload_key = None
            st.toast("데이터를 새로고침합니다.")
            st.rerun()

    # 엑셀 파일 업로드 기능
    st.markdown("<div style='height:24px'></div>", unsafe_allow_html=True)
    st.markdown("#### 📤 엑셀 파일 업로드")

    # 업로드된 데이터 확인
    if st.session_state.get("admin_uploaded_excel") is not None:
        uploaded_data = st.session_state.get("admin_uploaded_excel")
        with st.expander("🔍 업로드된 데이터 확인", expanded=False):
            st.write(f"**업로드 데이터 건수:** {len(uploaded_data)}")
            if not uploaded_data.empty:
                analysis_clean = clean_header_logic(uploaded_data.copy())
                u_col = find_col(analysis_clean, ["등록자", "담당자", "성명"])
                d_col = find_col(analysis_clean, ["활동상세", "활동내용"])
                st.write(f"**담당자 컬럼명:** {u_col}")
                st.write(f"**활동상세 컬럼명:** {d_col}")
                if u_col and u_col in analysis_clean.columns:
                    unique_users = analysis_clean[u_col].unique()
                    st.write(f"**담당자 목록:** {list(unique_users)}")
                st.write("**데이터 샘플:**")
                st.dataframe(prepare_display_dataframe(uploaded_data).head(10))

        # 중복 이력 및 초과 방문 체크
        with st.expander("⚠️ 중복 이력 및 초과 방문 체크", expanded=False):
            uploaded_clean = clean_header_logic(uploaded_data.copy())
            biz_col = find_col(uploaded_clean, ["사업자번호", "사업자등록번호"])
            u_col = find_col(uploaded_clean, ["등록자", "담당자", "성명"])
            date_col = find_col(uploaded_clean, ["활동일자", "방문일자"])
            d_col = find_col(uploaded_clean, ["활동상세", "활동내용"])
            company_col = find_col(uploaded_clean, ["고객명", "업체명", "상호"])

            # 중복 이력 체크
            dup_df = pd.DataFrame()
            if biz_col and u_col and date_col and d_col:
                _dup_df = uploaded_clean.copy()
                _dup_df["_dup_biz"] = _dup_df[biz_col].astype(str).str.strip()
                _dup_df["_dup_user"] = _dup_df[u_col].astype(str).str.strip()
                _dup_df["_dup_date"] = pd.to_datetime(_dup_df[date_col], errors="coerce").dt.strftime("%Y-%m-%d")
                _dup_df["_dup_detail"] = _dup_df[d_col].astype(str).str.strip()
                _dup_keys = ["_dup_biz", "_dup_user", "_dup_date", "_dup_detail"]
                _dup_df = _dup_df[(_dup_df[_dup_keys] != "").all(axis=1)]
                dup_df = _dup_df[_dup_df.duplicated(subset=_dup_keys, keep=False)].copy()
                if not dup_df.empty:
                    dup_df = dup_df.drop(columns=_dup_keys, errors="ignore")
                    dup_df = dup_df.sort_values(by=[date_col, biz_col, u_col, d_col]) if all(c in dup_df.columns for c in [date_col, biz_col, u_col, d_col]) else dup_df

            # 초과 방문 체크 (하루 5회 이상 방문)
            over_visit_df = pd.DataFrame()
            if biz_col and u_col and date_col and company_col:
                visit_count = uploaded_clean.groupby([biz_col, u_col, date_col]).size().reset_index(name="일방문")
                over_visit = visit_count[visit_count["일방문"] >= 5]
                if not over_visit.empty:
                    # 월 총 방문 횟수 계산
                    month_count = uploaded_clean.groupby([biz_col, u_col]).size().reset_index(name="월총방문")
                    over_visit = over_visit.merge(month_count, on=[biz_col, u_col], how="left")
                    # 업체명 추가
                    company_map = uploaded_clean.groupby(biz_col)[company_col].first().to_dict()
                    over_visit[company_col] = over_visit[biz_col].map(company_map)
                    over_visit_df = over_visit[[company_col, biz_col, u_col, date_col, "일방문", "월총방문"]].rename(columns={date_col: "초과일자"})

            # 운영 60회 초과 체크
            over_operation_df = pd.DataFrame()
            if u_col:
                category_col = find_col(uploaded_clean, ["활동구분", "접수유형"])
                if category_col and category_col in uploaded_clean.columns:
                    # 활동구분 기준
                    operation_count_by_user = uploaded_clean.groupby(u_col).apply(
                        lambda x: int(x[category_col].astype(str).str.contains("방문|원격", na=False).sum())
                    ).reset_index(name="운영건수")
                elif d_col and d_col in uploaded_clean.columns:
                    # 활동상세 기준
                    operation_count_by_user = uploaded_clean.groupby(u_col).apply(
                        lambda x: int(x[d_col].astype(str).str.contains("운영|방문|점검", na=False).sum())
                    ).reset_index(name="운영건수")
                else:
                    operation_count_by_user = pd.DataFrame()

                if not operation_count_by_user.empty:
                    over_operation_df = operation_count_by_user[operation_count_by_user["운영건수"] > 60]

            # 결과 표시
            tabs = st.tabs(["중복 이력", "초과 방문", "운영 60회 초과"])
            with tabs[0]:
                if not dup_df.empty:
                    st.warning(f"⚠️ 중복 이력 {len(dup_df)}건 발견")
                    st.dataframe(dup_df, use_container_width=True, hide_index=True)
                else:
                    st.success("✅ 중복 이력이 없습니다.")

            with tabs[1]:
                if not over_visit_df.empty:
                    st.warning(f"⚠️ 초과 방문 {len(over_visit_df)}건 발견 (하루 5회 이상)")
                    st.dataframe(over_visit_df, use_container_width=True, hide_index=True)
                else:
                    st.success("✅ 초과 방문이 없습니다.")

            with tabs[2]:
                if not over_operation_df.empty:
                    st.warning(f"⚠️ 운영 60회 초과 {len(over_operation_df)}명 발견")
                    st.info("💡 운영건수는 최대 60회까지만 실적에 반영됩니다.")
                    st.dataframe(over_operation_df, use_container_width=True, hide_index=True)
                else:
                    st.success("✅ 운영 60회 초과 담당자가 없습니다.")

    col1, col_convert, col_upload, col_sample, _ = st.columns([1, 1, 1, 1, 2])
    with col1:
        st.markdown("<div style='text-align:center;font-weight:700;margin-bottom:4px;'>은행 이력 업로드</div>", unsafe_allow_html=True)
        history_file = st.file_uploader("은행 이력 업로드", type=["xls", "xlsx"], key="admin_history_upload", label_visibility="collapsed")
        st.caption("은행에서 반출해준 이력파일 그대로 업로드 해주세요.")

    with col_convert:
        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
        if history_file is not None:
            try:
                with st.spinner("은행 이력을 샘플 양식으로 변환 중입니다."):
                    history_df = read_excel_history_file(history_file)
                    # 관리자 모드에서는 담당자명이 없으므로 빈 문자열 사용
                    converted_df, convert_info = convert_history_to_sample_df(history_df, "")
                if converted_df.empty:
                    st.button("변환파일 다운로드", use_container_width=True, disabled=True)
                    st.warning(convert_info.get("error", "변환할 데이터가 없습니다."))
                else:
                    converted_df = normalize_converted_history_df(converted_df)
                    # session_state에 저장하여 운영 카운트에 사용
                    st.session_state.admin_uploaded_excel = converted_df
                    st.session_state.admin_uploaded_excel_display = converted_df.copy()
                    converted_bytes = sample_format_excel_bytes(converted_df)
                    converted_ym = get_uploaded_month(converted_df).replace("-", "") or (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m")
                    st.download_button(
                        "변환파일 다운로드",
                        data=converted_bytes,
                        file_name=f"LMB월간 활동실적_{converted_ym}_관리자.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True,
                    )
                    unmatched = int(convert_info.get("unmatched", 0))
                    if unmatched:
                        st.caption(f"고객번호 매핑 실패 {unmatched}건은 사업자번호가 공란으로 저장됩니다.")
                    st.success("✅ 업로드한 파일이 실적 계산에 반영됩니다.")
                    history_file_key = f"{history_file.name}_{history_file.size}"
                    if st.session_state.get("admin_history_upload_key") != history_file_key:
                        st.session_state.admin_history_upload_key = history_file_key
                        st.rerun()
            except ImportError:
                st.button("변환파일 다운로드", use_container_width=True, disabled=True)
                st.error("xls 파일 변환을 위해 xlrd 패키지가 필요합니다.")
            except Exception as e:
                st.button("변환파일 다운로드", use_container_width=True, disabled=True)
                st.error(f"은행 이력 업로드 실패: {e}")
        else:
            st.button("변환파일 다운로드", use_container_width=True, disabled=True)

    with col_upload:
        st.markdown("<div style='text-align:center;font-weight:700;margin-bottom:4px;'>본사이력 업로드 (선택)</div>", unsafe_allow_html=True)
        u_file = st.file_uploader("본사이력 업로드 (선택)", type=["xlsx"], key="admin_office_upload", label_visibility="collapsed")
        if u_file is not None:
            try:
                office_df = pd.read_excel(u_file)
                st.session_state.admin_uploaded_excel_display = clean_header_logic(office_df.copy())
                office_df = normalize_converted_history_df(office_df)
                # session_state에 저장하여 운영 카운트에 사용
                st.session_state.admin_uploaded_excel = office_df
                st.success("✅ 본사이력 파일이 실적 계산에 반영됩니다.")
                office_file_key = f"{u_file.name}_{u_file.size}"
                if st.session_state.get("admin_office_upload_key") != office_file_key:
                    st.session_state.admin_office_upload_key = office_file_key
                    st.rerun()
            except Exception as e:
                st.error(f"본사이력 업로드 실패: {e}")

    with col_sample:
        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
        if os.path.exists(EXCEL_SAMPLE_FILE):
            with open(EXCEL_SAMPLE_FILE, "rb") as sample_file:
                st.download_button(
                    "샘플파일 다운로드",
                    data=sample_file.read(),
                    file_name="LMB월간 활동실적_000000(샘플).xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )
        else:
            st.button("샘플파일 다운로드", use_container_width=True, disabled=True)

    def load_latest_operation_history(target_perf_df):
        target_counts = {}
        for staff in staff_names:
            visit_count = count_visit_a(staff)
            if visit_count is not None and visit_count > 0:
                target_counts[str(staff).strip()] = min(int(visit_count), 60)

        if not target_counts and isinstance(target_perf_df, pd.DataFrame) and not target_perf_df.empty:
            target_base = target_perf_df[target_perf_df["담당자"].astype(str).str.strip().ne("합계")].copy()
            target_base["_target_count"] = pd.to_numeric(target_base.get("운영건수 (실제 활동)", 0), errors="coerce").fillna(0).astype(int)
            target_base = target_base.groupby("담당자", as_index=False)["_target_count"].max()
            target_counts = {
                str(row["담당자"]).strip(): min(int(row["_target_count"]), 60)
                for _, row in target_base.iterrows()
                if str(row["담당자"]).strip() and int(row["_target_count"]) > 0
            }

        if not target_counts and isinstance(target_perf_df, pd.DataFrame) and not target_perf_df.empty:
            for _, row in target_perf_df.iterrows():
                staff = str(row.get("담당자", "")).strip()
                if not staff or staff == "합계":
                    continue
                try:
                    target_counts[staff] = min(max(0, int(float(row.get("운영건수 (실제 활동)", 0) or 0))), 60)
                except Exception:
                    target_counts[staff] = 0
        target_counts = {staff: count for staff, count in target_counts.items() if count > 0}
        if not target_counts:
            st.warning("가져올 운영건수 대상이 없습니다.")
            return
        target_total = sum(target_counts.values())
        if target_total > (len(target_counts) * 60):
            st.error("목표 운영건수가 비정상적으로 큽니다. 실적관리 시트 새로고침 후 다시 시도해주세요.")
            return

        activity_df = st.session_state.get("analysis_lookup_df")
        if activity_df is None:
            try:
                with st.spinner("하나지사 활동이력 시트를 불러오는 중..."):
                    load_csv_to_state("url_analysis", "analysis_lookup_df")
                    activity_df = st.session_state.get("analysis_lookup_df")
            except Exception as e:
                st.error(f"하나지사 활동이력 구글 시트를 불러오지 못했습니다: {e}")
                return

        if activity_df is None or activity_df.empty:
            st.error("하나지사 활동이력 시트 데이터가 없습니다.")
            return

        with st.spinner("담당자별 운영건수에 맞춰 최신 이력을 가져오는 중..."):
            activity_clean = clean_header_logic(activity_df.copy())
            detail_col = find_col(activity_clean, ["활동상세", "활동내용"])
            owner_col = find_col(activity_clean, ["담당자", "등록자", "성명"])
            date_col = find_col(activity_clean, ["활동일", "활동일자", "일자", "날짜"])
            biz_col = find_col(activity_clean, ["사업자번호", "사업자등록번호"])
            activity_company_col = find_col(activity_clean, ["업체명", "고객명", "상호"])
            activity_title_col = find_col(activity_clean, ["제목"])
            activity_content_col = find_col(activity_clean, ["활동내용", "활동내역"])
            activity_location_col = find_col(activity_clean, ["방문장소", "주소", "지역"])

            if not detail_col or not owner_col:
                st.error("하나지사 활동이력 시트에서 필수 컬럼(활동상세, 담당자)을 찾을 수 없습니다.")
                return
            if not date_col:
                st.error("하나지사 활동이력 시트에서 날짜 컬럼(활동일)을 찾을 수 없습니다.")
                return

            operation_df = activity_clean[
                activity_clean[detail_col].astype(str).str.contains("운영|방문|점검", na=False)
            ].copy()
            if operation_df.empty:
                st.warning("하나지사 활동이력 시트에 운영 활동 데이터가 없습니다.")
                return

            if st.session_state.get("cloud_sheet_df") is None:
                try:
                    load_csv_to_state("url_sync", "cloud_sheet_df")
                except Exception:
                    pass

            if biz_col and biz_col in operation_df.columns:
                op_check_df = attach_cloud_dates(filter_visit_rows(clean_header_logic(operation_df.copy())))
                op_other_errors = build_other_validation_errors(op_check_df)
                if isinstance(op_other_errors, pd.DataFrame) and not op_other_errors.empty and "사업자번호" in op_other_errors.columns:
                    invalid_biz = set(normalize_biz(op_other_errors["사업자번호"]).astype(str))
                    operation_df = operation_df[
                        ~normalize_biz(operation_df[biz_col]).astype(str).isin(invalid_biz)
                    ].copy()
                    if operation_df.empty:
                        st.warning("기타 오류 대상 고객사를 제외한 운영 활동 데이터가 없습니다.")
                        return

            operation_df["_parsed_date"] = pd.to_datetime(operation_df[date_col].map(parse_sheet_date), errors="coerce")
            operation_df = operation_df.sort_values("_parsed_date", ascending=False, na_position="last")
            business_dates = may_2026_business_dates()
            generated = []
            shortages = {}

            def template_from_activity_history(activity_detail, company_name, staff_name="", biz_no=""):
                keyword = "개설|구축|신규" if activity_detail == "개설" else "연계|ERP"
                text_cols = [
                    col for col in [detail_col, activity_title_col, activity_content_col, activity_company_col]
                    if col and col in activity_clean.columns
                ]
                if not text_cols:
                    return "", "", ""
                template_source = activity_clean.copy()
                template_text = template_source[text_cols].fillna("").astype(str).agg(" ".join, axis=1)
                candidates = template_source[template_text.str.contains(keyword, regex=True, na=False)].copy()
                if candidates.empty:
                    return "", "", ""
                company_text = str(company_name or "")
                if "교회" in company_text and activity_company_col and activity_company_col in candidates.columns:
                    church_candidates = candidates[candidates[activity_company_col].astype(str).str.contains("교회", na=False)].copy()
                    if not church_candidates.empty:
                        candidates = church_candidates
                candidates["_template_date"] = pd.to_datetime(candidates[date_col].map(parse_sheet_date), errors="coerce")
                candidates = candidates.sort_values("_template_date", ascending=False, na_position="last")
                seed = f"{activity_detail}|{company_name}|{staff_name}|{biz_no}"
                template_idx = int(hashlib.md5(seed.encode("utf-8")).hexdigest(), 16) % len(candidates)
                template = candidates.iloc[template_idx]
                template_company = str(template.get(activity_company_col, "")).strip() if activity_company_col else ""
                title = str(template.get(activity_title_col, "") or "").strip() if activity_title_col else ""
                content = str(template.get(activity_content_col, "") or "").strip() if activity_content_col else ""
                location = template.get(activity_location_col, "") if activity_location_col else ""

                if template_company and company_text:
                    title = title.replace(template_company, company_text)
                    content = content.replace(template_company, company_text)

                variants = {
                    "개설": [
                        "서버 구축 및 사용자 교육 진행, 초기 사용 환경 점검 완료",
                        "계좌 등록과 권한 설정 안내, 조회/이체 기본 업무 교육 진행",
                        "신규 구축 후 담당자 사용 방법 안내 및 운영 전환 사항 확인",
                        "클라이언트 설치와 기본 메뉴 교육 진행, 추가 요청사항 확인",
                        "초기 세팅 상태 점검 및 개설 이후 사용 일정 협의 완료",
                    ],
                    "연계": [
                        "ERP 연계 대상 업무 확인 및 데이터 내보내기 절차 안내",
                        "연계 설정 상태 점검, 테스트 데이터 확인 후 후속 일정 협의",
                        "ERP 담당자와 연계 범위 확인 및 반영 결과 점검 진행",
                        "연계 오류 가능 항목 사전 확인, 사용자 검증 방법 안내",
                        "연계 완료 후 조회/전송 데이터 확인 및 담당자 교육 진행",
                    ],
                }
                variant_list = variants.get(activity_detail, [])
                if variant_list:
                    variant_idx = int(hashlib.md5(f"{seed}|variant".encode("utf-8")).hexdigest(), 16) % len(variant_list)
                    company_prefix = f"{company_text} " if company_text else ""
                    content = f"{content}\n- {company_prefix}{variant_list[variant_idx]}".strip()
                return title, content, location

            for staff, target_count in target_counts.items():
                source = operation_df[operation_df[owner_col].astype(str).str.strip() == staff].copy()
                if source.empty:
                    source = operation_df.copy()

                source = source.drop(columns=["_parsed_date"], errors="ignore").reset_index(drop=True)
                selected_rows = []
                used_keys = set()
                daily_counts = {}
                source_idx = 0
                attempts = 0
                max_attempts = max(target_count * 20, len(source) * 3, 100)

                while len(selected_rows) < target_count and attempts < max_attempts and not source.empty:
                    base = source.iloc[source_idx % len(source)].copy()
                    source_idx += 1
                    attempts += 1

                    date = business_dates[len(selected_rows) % len(business_dates)]
                    if daily_counts.get(date, 0) >= 5:
                        available_dates = [d for d in business_dates if daily_counts.get(d, 0) < 5]
                        if not available_dates:
                            break
                        date = available_dates[0]

                    biz_value = str(base.get(biz_col, "")).strip() if biz_col else str(source_idx)
                    key = (staff, date, biz_value, str(base.get(detail_col, "")).strip())
                    if key in used_keys:
                        continue

                    base[owner_col] = staff
                    base[date_col] = date
                    selected_rows.append(base)
                    used_keys.add(key)
                    daily_counts[date] = daily_counts.get(date, 0) + 1

                if len(selected_rows) < target_count:
                    shortages[staff] = target_count - len(selected_rows)
                generated.extend(selected_rows)

            if not generated:
                st.warning("가져온 이력이 없습니다.")
                return

            latest_df = normalize_converted_history_df(pd.DataFrame(generated))
            latest_clean = clean_header_logic(latest_df.copy())
            latest_owner_col = find_col(latest_clean, ["등록자", "담당자", "성명"])
            latest_date_col = find_col(latest_clean, ["활동일자", "활동일", "일자"])
            latest_detail_col = find_col(latest_clean, ["활동상세", "활동내용"])
            latest_biz_col = find_col(latest_clean, ["사업자번호"], "사업자번호")

            if latest_owner_col and latest_date_col and latest_detail_col:
                latest_clean[latest_owner_col] = latest_clean[latest_owner_col].astype(str).str.strip()
                latest_clean["_latest_date"] = pd.to_datetime(latest_clean[latest_date_col], errors="coerce").dt.strftime("%Y-%m-%d")
                latest_clean["_latest_detail"] = latest_clean[latest_detail_col].astype(str).str.strip()
                if latest_biz_col and latest_biz_col in latest_clean.columns:
                    latest_clean["_latest_biz"] = normalize_biz(latest_clean[latest_biz_col])
                else:
                    latest_clean["_latest_biz"] = ""

                latest_clean = latest_clean[latest_clean[latest_owner_col].isin(target_counts.keys())].copy()
                latest_clean = latest_clean.drop_duplicates(
                    subset=["_latest_biz", latest_owner_col, "_latest_date", "_latest_detail"],
                    keep="first",
                )
                latest_clean["_daily_seq"] = latest_clean.groupby([latest_owner_col, "_latest_date"]).cumcount()
                latest_clean = latest_clean[latest_clean["_daily_seq"] < 5].copy()
                latest_clean["_staff_seq"] = latest_clean.groupby(latest_owner_col).cumcount()
                latest_clean["_staff_target"] = latest_clean[latest_owner_col].map(target_counts).fillna(0).astype(int)
                latest_clean = latest_clean[latest_clean["_staff_seq"] < latest_clean["_staff_target"]].copy()

                final_check_df = attach_cloud_dates(clean_header_logic(latest_clean.copy()))
                final_other_errors = build_other_validation_errors(final_check_df)
                if isinstance(final_other_errors, pd.DataFrame) and not final_other_errors.empty and "사업자번호" in final_other_errors.columns:
                    invalid_biz = set(normalize_biz(final_other_errors["사업자번호"]).astype(str))
                    latest_clean = latest_clean[
                        ~latest_clean["_latest_biz"].astype(str).isin(invalid_biz)
                    ].copy()

                actual_counts = latest_clean.groupby(latest_owner_col).size().to_dict()
                shortages = {
                    staff: target - int(actual_counts.get(staff, 0))
                    for staff, target in target_counts.items()
                    if int(actual_counts.get(staff, 0)) < target
                }
                latest_df = normalize_converted_history_df(
                    latest_clean.drop(
                        columns=["_latest_date", "_latest_detail", "_latest_biz", "_daily_seq", "_staff_seq", "_staff_target"],
                        errors="ignore",
                    ).reset_index(drop=True)
                )

            cloud_history_rows = []
            cloud_source = st.session_state.get("cloud_sheet_df")
            if isinstance(cloud_source, pd.DataFrame) and not cloud_source.empty:
                cloud_source = clean_header_logic(cloud_source.copy())
                cloud_owner_col = find_col(cloud_source, ["담당자", "등록자", "성명"])
                cloud_open_col = find_col(cloud_source, ["개설완료일자", "개설일"])
                cloud_erp_col = find_col(cloud_source, ["ERP연계일자", "연계일자"])
                cloud_biz_col = find_col(cloud_source, ["사업자번호", "사업자등록번호"])
                cloud_company_col = find_col(cloud_source, ["고객명", "업체명", "상호"])

                def add_cloud_history_row(src_row, activity_detail, activity_date):
                    parsed_date = parse_sheet_date(activity_date)
                    if pd.isna(parsed_date) or parsed_date.strftime("%Y-%m") != target_year_month:
                        return
                    staff = str(src_row.get(cloud_owner_col, "")).strip() if cloud_owner_col else ""
                    if staff not in staff_names:
                        return
                    company = src_row.get(cloud_company_col, "") if cloud_company_col else ""
                    biz_no = normalize_biz(src_row.get(cloud_biz_col, "")) if cloud_biz_col else ""
                    template_title, template_content, template_location = template_from_activity_history(activity_detail, company, staff, biz_no)
                    cloud_history_rows.append({
                        "지사": "HANA지사",
                        "상품": "통합CMS",
                        "업체명": company,
                        "사업자번호": biz_no,
                        "등록자": staff,
                        "활동일자": parsed_date.strftime("%Y-%m-%d"),
                        "방문장소 (시, 군, 구까지)": template_location,
                        "활동구분": "상담",
                        "활동상세": activity_detail,
                        "업무번호": "",
                        "제목": "신규구축" if activity_detail == "개설" else "ERP연계 방문",
                        "활동내역": template_content or f"본사 구글시트 {activity_detail} 실적 반영",
                    })

                if cloud_owner_col and cloud_biz_col:
                    for _, cloud_row in cloud_source.iterrows():
                        if cloud_open_col:
                            add_cloud_history_row(cloud_row, "개설", cloud_row.get(cloud_open_col, ""))
                        if cloud_erp_col:
                            add_cloud_history_row(cloud_row, "연계", cloud_row.get(cloud_erp_col, ""))

            if cloud_history_rows:
                latest_df = pd.concat(
                    [latest_df, normalize_converted_history_df(pd.DataFrame(cloud_history_rows))],
                    ignore_index=True,
                    sort=False,
                )
                latest_all = clean_header_logic(latest_df.copy())
                all_owner_col = find_col(latest_all, ["등록자", "담당자", "성명"])
                all_date_col = find_col(latest_all, ["활동일자", "활동일", "일자"])
                all_detail_col = find_col(latest_all, ["활동상세", "활동내용"])
                all_biz_col = find_col(latest_all, ["사업자번호"], "사업자번호")
                if all_owner_col and all_date_col and all_detail_col and all_biz_col:
                    latest_all["_all_biz"] = normalize_biz(latest_all[all_biz_col])
                    latest_all["_all_date"] = pd.to_datetime(latest_all[all_date_col], errors="coerce").dt.strftime("%Y-%m-%d")
                    latest_all["_all_detail"] = latest_all[all_detail_col].astype(str).str.strip()
                    latest_all[all_owner_col] = latest_all[all_owner_col].astype(str).str.strip()
                    latest_all = latest_all.drop_duplicates(
                        subset=["_all_biz", all_owner_col, "_all_date", "_all_detail"],
                        keep="first",
                    )
                    final_all_check = attach_cloud_dates(clean_header_logic(latest_all.copy()))
                    final_all_errors = build_other_validation_errors(final_all_check)
                    if isinstance(final_all_errors, pd.DataFrame) and not final_all_errors.empty and "사업자번호" in final_all_errors.columns:
                        invalid_biz = set(normalize_biz(final_all_errors["사업자번호"]).astype(str))
                        latest_all = latest_all[~latest_all["_all_biz"].astype(str).isin(invalid_biz)].copy()
                    latest_df = normalize_converted_history_df(
                        latest_all.drop(columns=["_all_biz", "_all_date", "_all_detail"], errors="ignore").reset_index(drop=True)
                    )

            st.session_state.admin_uploaded_excel = latest_df
            st.session_state.admin_uploaded_excel_display = latest_df.copy()
            st.session_state.admin_office_upload_key = f"latest_history_{len(latest_df)}_{int(time.time())}"
            detail_summary = {}
            detail_summary_col = find_col(latest_df, ["활동상세", "활동내용"])
            if detail_summary_col and detail_summary_col in latest_df.columns:
                detail_summary = latest_df[detail_summary_col].astype(str).value_counts().to_dict()
            st.success(
                f"최신 이력 {len(latest_df):,}건을 가져왔습니다. "
                f"(운영 목표 {target_total:,}건 / 개설 {int(detail_summary.get('개설', 0)):,}건 / 연계 {int(detail_summary.get('연계', 0)):,}건)"
            )
            if shortages:
                shortage_text = ", ".join([f"{staff} {count}건" for staff, count in shortages.items()])
                st.warning(f"목표 건수보다 부족한 담당자: {shortage_text}")
            st.rerun()

    latest_col, latest_spacer = st.columns([0.18, 0.82])
    with latest_col:
        if st.button("최신 이력가져오기", use_container_width=True, key="admin_load_latest_history"):
            load_latest_operation_history(perf_df)
    with latest_spacer:
        st.caption("실적관리 시트의 운영건수(실제 활동)에 맞춰 하나지사 활동이력에서 최신 운영 이력을 가져옵니다.")


    if (
        isinstance(st.session_state.get("admin_uploaded_excel"), pd.DataFrame)
        and not st.session_state.admin_uploaded_excel.empty
    ):
        display_source_df = st.session_state.get("admin_uploaded_excel_display")
        if not isinstance(display_source_df, pd.DataFrame) or display_source_df.empty:
            display_source_df = st.session_state.admin_uploaded_excel
        capped_display_df = limit_history_to_total_point_cap(display_source_df, perf_df)
        capped_analysis_df = limit_history_to_total_point_cap(st.session_state.admin_uploaded_excel, perf_df)
        if isinstance(capped_display_df, pd.DataFrame) and not capped_display_df.equals(display_source_df):
            st.session_state.admin_uploaded_excel_display = capped_display_df
            display_source_df = capped_display_df
        if isinstance(capped_analysis_df, pd.DataFrame) and not capped_analysis_df.equals(st.session_state.admin_uploaded_excel):
            st.session_state.admin_uploaded_excel = normalize_converted_history_df(capped_analysis_df)
        search_filters = render_history_search_filters(display_source_df, "admin_uploaded_history_search")
        admin_uploaded_df = apply_history_search_filters(display_source_df.copy(), search_filters)
        filtered_analysis_df = apply_history_search_filters(st.session_state.admin_uploaded_excel.copy(), search_filters)

        analysis_clean = clean_header_logic(filtered_analysis_df.copy())
        u_col = find_col(analysis_clean, ["등록자", "담당자", "성명"])
        d_col = find_col(analysis_clean, ["활동상세", "활동내용"])
        category_col = find_col(analysis_clean, ["활동구분", "접수유형"])
        if u_col and u_col in analysis_clean.columns:
            for idx, row in perf_df.iterrows():
                # 실적관리 시트 방문A 우선 (엑셀 업로드보다 우선)
                va = count_visit_a(row["담당자"])
                if va is not None:
                    operation_count  = va
                    operation_points = va * 30
                else:
                    staff_activity = analysis_clean[analysis_clean[u_col] == row["담당자"]].copy()
                    operation_count = 0
                    operation_points = 0
                    if category_col and category_col in staff_activity.columns:
                        category_text = staff_activity[category_col].astype(str)
                        visit_count = int(category_text.str.contains("방문", na=False).sum())
                        remote_count = int(category_text.str.contains("원격", na=False).sum())
                        operation_count = visit_count + remote_count
                        operation_points = (visit_count * 30) + (remote_count * 10)
                    elif d_col and d_col in staff_activity.columns:
                        operation_count = int(
                            staff_activity[d_col].astype(str).str.contains("운영|방문|점검", na=False).sum()
                        )
                        operation_points = operation_count * 30

                open_points = int(perf_df.at[idx, "개설포인트"]) if "개설포인트" in perf_df.columns else 0
                link_points = int(perf_df.at[idx, "연계포인트"]) if "연계포인트" in perf_df.columns else 0
                total_points = min(2800, open_points + link_points + operation_points)
                perf_df.at[idx, "운영건수 (실제 활동)"] = operation_count
                perf_df.at[idx, "운영포인트 (실제 활동)"] = operation_points
                perf_df.at[idx, "합계포인트"] = total_points

            perf_df = apply_rs_allowance_formula(perf_df, user_db)
            perf_df = apply_perf_sheet_payouts(perf_df)

            total_row = {}
            for col in display_cols:
                if col in ["담당자"]:
                    total_row[col] = "합계"
                elif col in ["직급"]:
                    total_row[col] = ""
                elif col in perf_df.columns:
                    if perf_df[col].dtype in ['int64', 'float64']:
                        total_row[col] = perf_df[col].sum()
                    else:
                        total_row[col] = ""
                else:
                    total_row[col] = ""
            perf_df_with_total = pd.concat([perf_df[display_cols], pd.DataFrame([total_row])], ignore_index=True)

        admin_uploaded_df = prepare_display_dataframe(admin_uploaded_df)
        title_col, add_history_col = st.columns([0.82, 0.18])
        with title_col:
            st.markdown("#### 본사이력 업로드 데이터")
        with add_history_col:
            if False and st.button("추가 이력 가져오기", use_container_width=True, key="admin_add_random_history"):
                # 하나지사 활동이력 구글 시트에서 랜덤 운영 이력 생성 (360건, 담당자별 최대 60건)
                # url_analysis 시트 로드
                activity_df = st.session_state.get("analysis_lookup_df")
                if activity_df is None:
                    try:
                        with st.spinner("하나지사 활동이력 시트를 불러오는 중..."):
                            load_csv_to_state("url_analysis", "analysis_lookup_df")
                            activity_df = st.session_state.get("analysis_lookup_df")
                    except Exception as e:
                        st.error(f"❌ 하나지사 활동이력 구글 시트를 불러오지 못했습니다: {e}")
                        activity_df = None

                if activity_df is None or activity_df.empty:
                    st.error("❌ 하나지사 활동이력 구글 시트를 먼저 불러와주세요.")
                else:
                    with st.spinner("하나지사 활동이력 시트에서 운영 이력을 랜덤으로 생성 중..."):
                        activity_clean = clean_header_logic(activity_df.copy())

                        # 디버깅: 하나지사 활동이력 시트 컬럼명 출력
                        st.write("🔍 하나지사 활동이력 시트 컬럼명:", list(activity_clean.columns))

                        # 컬럼 찾기
                        hana_detail_col = find_col(activity_clean, ["활동상세", "활동내용"])
                        hana_owner_col = find_col(activity_clean, ["담당자", "등록자", "성명"])
                        hana_date_col = find_col(activity_clean, ["활동일", "활동일자", "일자", "날짜"])
                        hana_biz_col = find_col(activity_clean, ["사업자번호", "사업자등록번호"])

                        st.write(f"🔍 찾은 컬럼: detail={hana_detail_col}, owner={hana_owner_col}, date={hana_date_col}, biz={hana_biz_col}")

                        if not hana_detail_col or not hana_owner_col:
                            st.error("❌ 하나지사 활동이력 시트에서 필수 컬럼(활동상세, 담당자)을 찾을 수 없습니다.")
                        elif not hana_date_col:
                            st.error("❌ 하나지사 활동이력 시트에서 날짜 컬럼(활동일)을 찾을 수 없습니다.")
                        else:
                            # 활동상세가 "운영"인 데이터만 필터링
                            operation_df = activity_clean[
                                activity_clean[hana_detail_col].astype(str).str.contains("운영", na=False)
                            ].copy()

                            if operation_df.empty:
                                st.warning("⚠️ 하나지사 활동이력 시트에 '운영' 활동 데이터가 없습니다.")
                            else:
                                # 2026년 5월 영업일 생성 (휴일/공휴일 제외)
                                may_2026_days = []
                                for day in range(1, 32):
                                    try:
                                        date_obj = datetime(2026, 5, day)
                                        # 토요일(5), 일요일(6) 제외
                                        if date_obj.weekday() < 5:
                                            may_2026_days.append(date_obj.strftime("%Y-%m-%d"))
                                    except:
                                        pass

                                # C&S 직원 목록
                                staff_list = staff_names.copy()
                                st.write(f"🔍 디버그: C&S 직원 수 = {len(staff_list)}")
                                st.write(f"🔍 디버그: 직원 목록 = {staff_list}")
                                st.write(f"🔍 디버그: 운영 데이터 건수 = {len(operation_df)}")

                                # 담당자별로 최대 60건씩, 총 360건 생성
                                target_total = 360
                                target_per_staff = 60

                                generated_rows = []
                                staff_count = {}

                                # 랜덤하게 데이터 생성
                                for _ in range(target_total * 2):  # 여유있게 생성
                                    if len(generated_rows) >= target_total:
                                        break

                                    # 랜덤 담당자 선택
                                    staff = np.random.choice(staff_list)

                                    # 담당자별 60건 제한 체크
                                    if staff_count.get(staff, 0) >= target_per_staff:
                                        continue

                                    # 랜덤 데이터 선택
                                    random_row = operation_df.sample(n=1).iloc[0].copy()

                                    # 담당자 변경
                                    random_row[hana_owner_col] = staff

                                    # 랜덤 일자 할당 (2026년 5월 영업일)
                                    random_date = np.random.choice(may_2026_days)
                                    if hana_date_col and hana_date_col in random_row.index:
                                        random_row[hana_date_col] = random_date

                                    generated_rows.append(random_row)
                                    staff_count[staff] = staff_count.get(staff, 0) + 1

                                st.write(f"🔍 디버그: 생성된 행 수 = {len(generated_rows)}")
                                st.write(f"🔍 디버그: 담당자별 건수 = {staff_count}")

                                if generated_rows:
                                    random_df = pd.DataFrame(generated_rows)
                                    st.write(f"🔍 디버그: DataFrame 생성 완료, 행 수 = {len(random_df)}")

                                    # 중복 제거 (담당자 + 일자 + 사업자번호 기준)
                                    if hana_date_col and hana_date_col in random_df.columns and hana_biz_col and hana_biz_col in random_df.columns:
                                        before_dedup = len(random_df)
                                        random_df = random_df.drop_duplicates(
                                            subset=[hana_owner_col, hana_date_col, hana_biz_col],
                                            keep='first'
                                        )
                                        st.write(f"🔍 디버그: 중복 제거 (담당자+일자+사업자번호) 전 {before_dedup}건 → 후 {len(random_df)}건")

                                    # 초과방문 제거 (일자별 최대 5회)
                                    if hana_date_col and hana_date_col in random_df.columns:
                                        before_filter = len(random_df)
                                        # 담당자 + 일자별로 그룹화하여 최대 5개씩만 유지
                                        limited_rows = []
                                        for (staff, date), group in random_df.groupby([hana_owner_col, hana_date_col]):
                                            # 각 그룹에서 최대 5개까지만
                                            limited_rows.append(group.head(5))
                                        random_df = pd.concat(limited_rows, ignore_index=True) if limited_rows else pd.DataFrame()
                                        st.write(f"🔍 디버그: 초과방문 제한 (일별 최대 5회) 전 {before_filter}건 → 후 {len(random_df)}건")

                                    # 세션 스테이트에 저장
                                    st.session_state.admin_uploaded_excel = random_df
                                    st.session_state.admin_uploaded_excel_display = random_df

                                    st.write(f"🔍 디버그: 세션 스테이트 저장 완료")
                                    st.write(f"🔍 디버그: admin_uploaded_excel 건수 = {len(st.session_state.admin_uploaded_excel)}")

                                    st.success(f"✅ 하나지사 활동이력 시트에서 운영 이력 {len(random_df)}건을 랜덤으로 생성했습니다.")
                                    st.info(f"📊 담당자별 최대 60건, 일자는 2026년 5월 영업일")
                                    st.info("⏰ 10초 후 페이지가 자동으로 새로고침됩니다...")
                                    time.sleep(10)
                                    st.rerun()
                                else:
                                    st.warning("⚠️ 생성된 데이터가 없습니다.")
            pass
        st.caption(f"업로드 데이터 건수: {len(admin_uploaded_df):,}건")
        st.dataframe(admin_uploaded_df, use_container_width=True, hide_index=True)

        admin_check_df = attach_cloud_dates(filter_visit_rows(clean_header_logic(filtered_analysis_df.copy())))
        _, admin_err, admin_dup = process_performance_analysis(admin_check_df, st.session_state.get("auto_prev_df"))

        admin_dup_df = pd.DataFrame()
        if admin_dup is not None and not admin_dup.empty:
            admin_dup_df = apply_history_search_filters(admin_dup.copy(), search_filters)

        admin_err_df = pd.DataFrame()
        if admin_err is not None and not admin_err.empty:
            admin_err_df = admin_err[admin_err["일방문"] >= 6].copy() if "일방문" in admin_err.columns else admin_err.copy()
            admin_err_df = apply_history_search_filters(admin_err_df, search_filters)

        admin_missing_open = pd.DataFrame()
        if "본사 개설완료일자" in admin_check_df.columns:
            admin_missing_open = admin_check_df[
                pd.isna(admin_check_df["본사 개설완료일자"]) |
                (admin_check_df["본사 개설완료일자"].astype(str).str.strip() == "")
            ].copy()
            if "본사 신규이행구분" in admin_missing_open.columns:
                admin_missing_open = admin_missing_open[
                    admin_missing_open["본사 신규이행구분"].astype(str).str.strip() != "이행"
                ]
            admin_missing_open = apply_history_search_filters(admin_missing_open, search_filters)

        admin_missing_erp = pd.DataFrame()
        if "본사 ERP연계일자" in admin_check_df.columns:
            if d_col and d_col in admin_check_df.columns:
                admin_erp_target = admin_check_df[admin_check_df[d_col].astype(str).str.contains("연계", na=False)]
            else:
                admin_erp_target = admin_check_df
            admin_missing_erp = admin_erp_target[
                pd.isna(admin_erp_target["본사 ERP연계일자"]) |
                (admin_erp_target["본사 ERP연계일자"].astype(str).str.strip() == "")
            ].copy()
            if "본사 신규이행구분" in admin_missing_erp.columns and "본사 이행추가연계" in admin_missing_erp.columns:
                _is_ihang = admin_missing_erp["본사 신규이행구분"].astype(str).str.strip() == "이행"
                _add_empty = (
                    pd.isna(admin_missing_erp["본사 이행추가연계"]) |
                    (admin_missing_erp["본사 이행추가연계"].astype(str).str.strip() == "")
                )
                admin_missing_erp = admin_missing_erp[~(_is_ihang & _add_empty)]
            admin_missing_erp = apply_history_search_filters(admin_missing_erp, search_filters)

        admin_other_errors = apply_history_search_filters(build_other_validation_errors(admin_check_df), search_filters)

        admin_error_tabs = []
        if not admin_dup_df.empty:
            admin_error_tabs.append("중복 이력")
        if not admin_err_df.empty:
            admin_error_tabs.append("초과 방문")
        if not admin_missing_open.empty:
            admin_error_tabs.append("본사 개설완료일자 누락")
        if not admin_missing_erp.empty:
            admin_error_tabs.append("본사 ERP연계일자 누락")
        if not admin_other_errors.empty:
            admin_error_tabs.append("기타 오류")

        t1, t2, t3, t4, t5 = validation_tabs_with_refresh("refresh_admin_history_validation")
        with t1:
            if not admin_dup_df.empty:
                style_report_logic(admin_dup_df.drop(columns=[c for c in admin_dup_df.columns if c.startswith("_")], errors="ignore"))
            else:
                st.info("중복 이력이 없습니다.")
        with t2:
            if not admin_err_df.empty:
                style_report_logic(admin_err_df.drop(columns=["월총방문"], errors="ignore"))
            else:
                st.info("초과 방문 데이터가 없습니다.")
        with t3:
            if not admin_missing_open.empty:
                style_report_logic(admin_missing_open.drop(columns=["본사 ERP연계일자"], errors="ignore"))
            elif "본사 개설완료일자" not in admin_check_df.columns:
                st.info("본사 구글시트에 개설완료일자 또는 사업자번호 컬럼이 없어 확인할 수 없습니다.")
        with t4:
            if not admin_missing_erp.empty:
                style_report_logic(admin_missing_erp.drop(columns=["본사 개설완료일자"], errors="ignore"))
            elif "본사 ERP연계일자" not in admin_check_df.columns:
                st.info("본사 구글시트에 ERP연계일자 또는 사업자번호 컬럼이 없어 확인할 수 없습니다.")
        with t5:
            style_report_logic(admin_other_errors, align_overrides={"오류 사유": "left"}, default_align="center")

        if admin_error_tabs:
            error_list = ", ".join(admin_error_tabs)
            st.markdown(
                f"<div style='margin-top:12px;margin-bottom:12px;padding:12px 16px;background:#FFF3CD;border-left:4px solid #FFC107;border-radius:4px;'>"
                f"<div style='font-size:14px;color:#856404;'><b>⚠️ 다음 항목을 확인해주세요:</b> {error_list}</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
        else:
            st.markdown(
                f"<div style='margin-top:12px;margin-bottom:12px;padding:12px 16px;background:#D4EDDA;border-left:4px solid #28A745;border-radius:4px;'>"
                f"<div style='font-size:14px;color:#155724;'><b>✓ 잘못된 데이터가 없습니다.</b></div>"
                f"</div>",
                unsafe_allow_html=True,
            )

    detail_title_col, detail_refresh_col = st.columns([0.86, 0.14])
    with detail_title_col:
        st.markdown("#### 담당자별 상세 실적")
    with detail_refresh_col:
        if st.button("🔄 새로고침", key="refresh_staff_detail_perf", use_container_width=True, help="실적관리 시트의 방문A/배분금액을 다시 조회합니다."):
            for key in ["_perf_parsed", "_perf_parsed_url", "_perf_parsed_version", "_perf_debug_err", "_perf_debug_hdr", "_perf_debug_raw"]:
                st.session_state.pop(key, None)
            read_csv_cached.clear()
            st.toast("담당자별 상세 실적을 새로고침합니다.")
            st.rerun()

    st.dataframe(
        perf_df_with_total,
        use_container_width=True,
        hide_index=True
    )

    st.markdown("<div style='height:18px'></div>", unsafe_allow_html=True)
    dl_excel_col, dl_ppt_col, _ = st.columns([1, 1, 2])

    admin_ym = ""
    if isinstance(st.session_state.get("admin_uploaded_excel"), pd.DataFrame):
        admin_ym = get_uploaded_month(st.session_state.admin_uploaded_excel)
    year_month = admin_ym.replace("-", "") if admin_ym else (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m")
    curr_month_label = f"{int(year_month[-2:])}월" if year_month[-2:].isdigit() else "당월"
    adm_sel = st.session_state.get("adm_prev_month", "선택안함")
    prev_month_label = str(int(adm_sel.split("-")[1])) + "월" if adm_sel and adm_sel != "선택안함" else "전월"
    download_time = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m%d_%H%M%S")

    with dl_excel_col:
        if isinstance(st.session_state.get("admin_uploaded_excel"), pd.DataFrame) and not st.session_state.admin_uploaded_excel.empty:
            excel_download_df = st.session_state.get("admin_uploaded_excel_display")
            if not isinstance(excel_download_df, pd.DataFrame) or excel_download_df.empty:
                excel_download_df = st.session_state.admin_uploaded_excel
            excel_download_df = prepare_display_dataframe(excel_download_df)

            # 개설고객사 데이터 (2026년 5월 개설 완료)
            open_companies_df = pd.DataFrame()
            if open_date_col and open_date_col in cloud_df.columns:
                open_companies_df = cloud_df[
                    cloud_df[open_date_col].notna() &
                    (cloud_df[open_date_col].dt.strftime("%Y-%m") == target_year_month)
                ].copy()

            # 연계고객사 데이터 (2026년 5월 ERP 연계)
            erp_companies_df = pd.DataFrame()
            if erp_date_col and erp_date_col in cloud_df.columns:
                erp_companies_df = cloud_df[
                    cloud_df[erp_date_col].notna() &
                    (cloud_df[erp_date_col].dt.strftime("%Y-%m") == target_year_month)
                ].copy()

            # 하나의 시트에 모든 데이터 합치기
            combined_parts = [excel_download_df]

            # 개설고객사 추가
            if not open_companies_df.empty:
                # 빈 행 추가
                empty_row = pd.DataFrame([[""] * len(excel_download_df.columns)], columns=excel_download_df.columns)
                combined_parts.append(empty_row)
                combined_parts.append(empty_row)

                # 제목 행 추가
                title_row = pd.DataFrame([["개설고객사"] + [""] * (len(excel_download_df.columns) - 1)], columns=excel_download_df.columns)
                combined_parts.append(title_row)

                # 개설고객사 데이터 추가 (컬럼 맞추기)
                open_companies_aligned = pd.DataFrame(columns=excel_download_df.columns)
                for col in open_companies_df.columns:
                    if col in open_companies_aligned.columns:
                        open_companies_aligned[col] = open_companies_df[col]
                combined_parts.append(open_companies_df)

            # 연계고객사 추가
            if not erp_companies_df.empty:
                # 빈 행 추가
                empty_row = pd.DataFrame([[""] * len(excel_download_df.columns)], columns=excel_download_df.columns)
                combined_parts.append(empty_row)
                combined_parts.append(empty_row)

                # 제목 행 추가
                title_row = pd.DataFrame([["연계고객사"] + [""] * (len(excel_download_df.columns) - 1)], columns=excel_download_df.columns)
                combined_parts.append(title_row)

                # 연계고객사 데이터 추가
                combined_parts.append(erp_companies_df)

            # 모든 데이터 합치기
            final_df = pd.concat(combined_parts, ignore_index=True, sort=False)

            excel_bytes = dataframe_to_excel_bytes({"실적파일": final_df})
            st.download_button(
                "실적파일 엑셀 다운로드",
                data=excel_bytes,
                file_name=f"LMB월간 활동실적__{year_month}_{download_time}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
            )
        else:
            st.button("실적파일 엑셀 다운로드", use_container_width=True, disabled=True)

    with dl_ppt_col:
        try:
            report_df_for_ppt = perf_df[display_cols].copy()
            upload_df_for_ppt = st.session_state.get("admin_uploaded_excel")
            ppt_bytes = build_report_ppt_bytes(report_df_for_ppt, pd.DataFrame(), curr_month_label, prev_month_label, upload_df_for_ppt)
            st.download_button(
                "실적보고서 PPT 다운로드",
                data=ppt_bytes,
                file_name=f"LMB활동실적보고서_{year_month}_하나지사.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        except Exception as e:
            st.button("실적보고서 PPT 다운로드", use_container_width=True, disabled=True)
            st.caption(f"PPT 생성 준비 중 오류: {e}")


def show_user_history(is_admin_mode=False):
    # 관리자 모드: 담당자 선택
    selected_user = st.session_state.user_name  # 기본값: 본인

    if is_admin_mode:
        st.markdown("### 관리자용 실적 확인")

        # C&S 부서 직원 목록만 가져오기
        user_db = load_db(DB_FILE, {})
        staff_names = []
        for uid, user_info in user_db.items():
            if uid != "1" and isinstance(user_info, dict) and user_info.get("name") and user_info.get("dept_type") == "C&S":
                staff_names.append(user_info.get("name"))

        staff_names = sorted(staff_names)

        if not staff_names:
            st.warning("등록된 직원이 없습니다. [직원 및 권한설정] 메뉴에서 직원을 등록해주세요.")
            return

        # "전체" 옵션 추가
        staff_options = ["📊 전체"] + staff_names

        # 담당자 선택 드롭다운
        col_select1, col_select2 = st.columns([1, 3])
        with col_select1:
            selected_option = st.selectbox(
                "👤 담당자 선택",
                options=staff_options,
                index=0
            )

        # "전체" 선택 시 전체 실적 요약 표시
        if selected_option == "📊 전체":
            show_all_staff_summary(staff_names)
            return

        selected_user = selected_option
        st.markdown(f"**{selected_user}**님의 실적을 확인/수정합니다.")
        st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)

    # 선택된 담당자로 작업 진행
    original_user_name = st.session_state.user_name
    if is_admin_mode:
        st.session_state.user_name = selected_user

    converted_preview_df = None
    analysis_df = None
    converted_ym = ""
    convert_info = {}
    col1, col_convert, col_upload, col_sample, _ = st.columns([1, 1, 1, 1, 2])
    with col1:
        st.markdown("<div style='text-align:center;font-weight:700;margin-bottom:4px;'>은행 이력 업로드</div>", unsafe_allow_html=True)
        history_file = st.file_uploader("은행 이력 업로드", type=["xls", "xlsx"], key="history_convert_upload", label_visibility="collapsed")
        st.caption("은행에서 반출해준 이력파일 그대로 업로드 해주세요.")
    with col_convert:
        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
        if history_file is not None:
            try:
                with st.spinner("은행 이력을 샘플 양식으로 변환 중입니다."):
                    history_df = read_excel_history_file(history_file)
                    converted_df, convert_info = convert_history_to_sample_df(history_df, st.session_state.user_name)
                if converted_df.empty:
                    st.button("변환파일 다운로드", use_container_width=True, disabled=True)
                    st.warning(convert_info.get("error", "변환할 데이터가 없습니다."))
                else:
                    converted_df = normalize_converted_history_df(converted_df)
                    converted_preview_df = converted_df
                    edited_key = "history_convert_preview_editor"
                    data_key = "history_convert_preview_data"
                    # 새 파일 업로드 여부 확인 → 파일이 바뀌면 preview 초기화
                    _file_id = f"{history_file.name}_{history_file.size}"
                    if st.session_state.get("_bank_file_id") != _file_id:
                        st.session_state["_bank_file_id"] = _file_id
                        st.session_state[data_key] = converted_df
                        edited_df = converted_df
                    else:
                        edited_df = st.session_state.get(data_key, converted_df)
                        if not isinstance(edited_df, pd.DataFrame) or list(edited_df.columns) != list(converted_df.columns):
                            edited_df = converted_df
                    edited_df = normalize_converted_history_df(edited_df)
                    st.session_state[data_key] = edited_df
                    st.session_state.user_excel_data = edited_df
                    st.session_state.user_excel_source = "bank"
                    analysis_df = edited_df
                    converted_bytes = sample_format_excel_bytes(edited_df)
                    converted_ym = get_uploaded_month(edited_df).replace("-", "") or (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m")
                    st.download_button(
                        "변환파일 다운로드",
                        data=converted_bytes,
                        file_name=f"LMB월간 활동실적_{converted_ym}_{st.session_state.user_name}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True,
                    )
                    unmatched = int(convert_info.get("unmatched", 0))
                    if unmatched:
                        st.caption(f"고객번호 매핑 실패 {unmatched}건은 사업자번호가 공란으로 저장됩니다.")
            except ImportError:
                st.button("변환파일 다운로드", use_container_width=True, disabled=True)
                st.error("xls 파일 변환을 위해 xlrd 패키지가 필요합니다. 배포 후 requirements.txt 반영을 확인해주세요.")
            except Exception as e:
                st.button("변환파일 다운로드", use_container_width=True, disabled=True)
                st.error(f"은행 이력 업로드 실패: {e}")
        else:
            st.button("변환파일 다운로드", use_container_width=True, disabled=True)
    with col_upload:
        st.markdown("<div style='text-align:center;font-weight:700;margin-bottom:4px;'>본사이력 업로드 (선택)</div>", unsafe_allow_html=True)
        u_file = st.file_uploader("본사이력 업로드 (선택)", type=["xlsx"], label_visibility="collapsed")
    with col_sample:
        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
        if os.path.exists(EXCEL_SAMPLE_FILE):
            with open(EXCEL_SAMPLE_FILE, "rb") as sample_file:
                st.download_button(
                    "샘플파일 다운로드",
                    data=sample_file.read(),
                    file_name="LMB월간 활동실적_000000(샘플).xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )
        else:
            st.button("샘플파일 다운로드", use_container_width=True, disabled=True)

    if False and converted_preview_df is not None and not converted_preview_df.empty:
        st.markdown("#### 변환파일 미리보기")
        st.markdown(
            """
            <style>
            body[data-pms-theme="d"] [data-testid="stDataEditor"],
            body[data-pms-theme="d"] [data-testid="stDataEditor"] > div,
            body[data-pms-theme="d"] [data-testid="stDataEditor"] [role="grid"],
            body[data-pms-theme="d"] [data-testid="stDataEditor"] canvas,
            body:has(#pms-d:checked) [data-testid="stDataEditor"],
            body:has(#pms-d:checked) [data-testid="stDataEditor"] > div,
            body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="grid"],
            body:has(#pms-d:checked) [data-testid="stDataEditor"] canvas {
                background-color: #252535 !important;
                color: #ffffff !important;
            }
            body[data-pms-theme="d"] [data-testid="stDataEditor"] canvas,
            body:has(#pms-d:checked) [data-testid="stDataEditor"] canvas {
                filter: invert(1) hue-rotate(180deg) brightness(0.72) contrast(1.18) saturate(0.85) !important;
            }
            body[data-pms-theme="d"] [data-testid="stDataEditor"] [role="columnheader"],
            body[data-pms-theme="d"] [data-testid="stDataEditor"] [role="gridcell"],
            body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="columnheader"],
            body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="gridcell"] {
                background-color: #252535 !important;
                color: #ffffff !important;
                border-color: #45475a !important;
            }
            body[data-pms-theme="d"] [data-testid="stDataEditor"] input,
            body[data-pms-theme="d"] [data-testid="stDataEditor"] textarea,
            body[data-pms-theme="d"] [data-testid="stDataEditor"] [contenteditable="true"],
            body[data-pms-theme="d"] [data-testid="stDataEditor"] [data-baseweb="input"] input,
            body:has(#pms-d:checked) [data-testid="stDataEditor"] input,
            body:has(#pms-d:checked) [data-testid="stDataEditor"] textarea,
            body:has(#pms-d:checked) [data-testid="stDataEditor"] [contenteditable="true"],
            body:has(#pms-d:checked) [data-testid="stDataEditor"] [data-baseweb="input"] input {
                background-color: #0f0f1f !important;
                color: #ffffff !important;
                caret-color: #ffffff !important;
                border-color: #89b4fa !important;
                -webkit-text-fill-color: #ffffff !important;
            }
            body[data-pms-theme="d"] [data-testid="stDataEditor"] input::selection,
            body[data-pms-theme="d"] [data-testid="stDataEditor"] textarea::selection,
            body:has(#pms-d:checked) [data-testid="stDataEditor"] input::selection,
            body:has(#pms-d:checked) [data-testid="stDataEditor"] textarea::selection {
                background-color: #4f46e5 !important;
                color: #ffffff !important;
            }
            @media (prefers-color-scheme: dark) {
                body:has(#pms-s:checked) [data-testid="stDataEditor"],
                body:has(#pms-s:checked) [data-testid="stDataEditor"] > div,
                body:has(#pms-s:checked) [data-testid="stDataEditor"] [role="grid"],
                body:has(#pms-s:checked) [data-testid="stDataEditor"] canvas {
                    background-color: #252535 !important;
                    color: #ffffff !important;
                }
                body:has(#pms-s:checked) [data-testid="stDataEditor"] canvas {
                    filter: invert(1) hue-rotate(180deg) brightness(0.72) contrast(1.18) saturate(0.85) !important;
                }
                body:has(#pms-s:checked) [data-testid="stDataEditor"] input,
                body:has(#pms-s:checked) [data-testid="stDataEditor"] textarea,
                body:has(#pms-s:checked) [data-testid="stDataEditor"] [contenteditable="true"],
                body:has(#pms-s:checked) [data-testid="stDataEditor"] [data-baseweb="input"] input {
                    background-color: #0f0f1f !important;
                    color: #ffffff !important;
                    caret-color: #ffffff !important;
                    border-color: #89b4fa !important;
                    -webkit-text-fill-color: #ffffff !important;
                }
            }
            </style>
            """,
            unsafe_allow_html=True,
        )
        converted_preview_df = normalize_converted_history_df(converted_preview_df)
        data_key = "history_convert_preview_data"
        editor_source_df = st.session_state.get(data_key, converted_preview_df)
        if not isinstance(editor_source_df, pd.DataFrame) or list(editor_source_df.columns) != list(converted_preview_df.columns):
            editor_source_df = converted_preview_df
        editor_source_df = normalize_converted_history_df(editor_source_df)
        edited_preview_df = st.data_editor(
            editor_source_df,
            key="history_convert_preview_editor",
            use_container_width=True,
            hide_index=True,
            num_rows="dynamic",
            disabled=[col for col in editor_source_df.columns if col not in ["활동일자", "활동구분", "활동상세"]],
            column_config={
                "지사": st.column_config.TextColumn("지사", disabled=True),
                "활동일자": st.column_config.TextColumn("활동일자"),
                "활동구분": st.column_config.SelectboxColumn(
                    "활동구분",
                    options=["방문", "상담", "원격"],
                    required=True,
                ),
                "활동상세": st.column_config.SelectboxColumn(
                    "활동상세",
                    options=["운영", "개설", "연계"],
                    required=True,
                ),
                "활동내역": st.column_config.TextColumn("활동내역"),
            },
        )
        analysis_df = normalize_converted_history_df(edited_preview_df)
        st.session_state[data_key] = analysis_df
        _, add_history_col = st.columns([0.88, 0.12])
        with add_history_col:
            if st.button("이력 추가", use_container_width=True, key="add_history_row"):
                new_row = {col: "" for col in analysis_df.columns}
                new_row.update({
                    "지사": "HANA지사",
                    "상품": "통합CMS",
                    "등록자": st.session_state.user_name,
                    "활동일자": (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d"),
                    "활동구분": "방문",
                    "활동상세": "운영",
                    "제목": "운영방문",
                })
                st.session_state[data_key] = normalize_converted_history_df(
                    pd.concat([analysis_df, pd.DataFrame([new_row])], ignore_index=True)
                )
                st.rerun()
        st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)

    if converted_preview_df is not None and not converted_preview_df.empty:
        _norm_cpdf = normalize_converted_history_df(converted_preview_df)
        _saved_preview = st.session_state.get("history_convert_preview_data")
        if isinstance(_saved_preview, pd.DataFrame) and list(_saved_preview.columns) == list(_norm_cpdf.columns):
            analysis_df = _saved_preview
        else:
            analysis_df = _norm_cpdf
        if st.session_state.get("user_excel_source") in (None, "bank"):
            st.session_state.user_excel_data = analysis_df
            st.session_state.user_excel_source = "bank"

    # 파일이 제거되면 데이터 초기화
    if u_file is None:
        if st.session_state.get("user_excel_source") == "hq" and st.session_state.get("user_excel_data") is not None:
            st.session_state.user_excel_data = None
            st.session_state.user_excel_source = None
            st.session_state.user_excel_file_key = None
            st.rerun()

    # 파일 업로드 시 자동 처리
    if u_file is not None:
        file_key = f"{u_file.name}_{u_file.size}"

        # 파일이 변경되었을 때만 처리
        if st.session_state.get("user_excel_file_key") != file_key:
            st.session_state.user_excel_file_key = file_key
            with st.spinner("본사이력 파일을 불러오는 중입니다..."):
                # cloud_sheet_df가 없을 때만 로드 시도
                if st.session_state.get("cloud_sheet_df") is None:
                    try:
                        load_csv_to_state("url_sync", "cloud_sheet_df")
                    except Exception as e:
                        st.warning(f"⚠️ 본사 구글시트를 불러오는데 실패했습니다: {str(e)}")

                uploaded_df = clean_header_logic(pd.read_excel(u_file, sheet_name=0))

                if not uploaded_df.empty:
                    st.session_state.user_excel_data = uploaded_df
                    st.session_state.user_excel_source = "hq"
                    st.toast("업로드 완료. 분석이 자동으로 시작됩니다.")
                    st.rerun()
                else:
                    st.error("업로드된 데이터가 없습니다. 파일을 확인해주세요.")
                    st.session_state.user_excel_data = None
                    st.session_state.user_excel_source = None

    if (
        st.session_state.get("user_excel_source") == "hq"
        and isinstance(st.session_state.get("user_excel_data"), pd.DataFrame)
        and not st.session_state.user_excel_data.empty
    ):
        hq_uploaded_df = prepare_display_dataframe(st.session_state.user_excel_data)
        st.markdown("#### 본사이력 업로드 데이터")
        st.caption(f"업로드 데이터 건수: {len(hq_uploaded_df):,}건")
        st.dataframe(hq_uploaded_df, use_container_width=True, hide_index=True)
        st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)

    current_df = current_history_analysis_df()
    if not current_df.empty:
        analysis_df = current_df
    elif st.session_state.user_excel_data is not None:
        analysis_df = prepare_history_analysis_df(st.session_state.user_excel_data)

    if analysis_df is None or analysis_df.empty:
        return

    try:
        load_csv_to_state("url_sync", "cloud_sheet_df", force_refresh=True)
    except Exception:
        if st.session_state.get("cloud_sheet_df") is None:
            pass

    df = analysis_df.copy()
    u_col = find_col(df, ["등록자", "담당자", "성명"], "등록자")
    d_col = find_col(df, ["활동상세", "활동내용"], "활동상세")

    current_user_key = re.sub(r"\s+", "", str(st.session_state.user_name).strip())
    if u_col in df.columns:
        df_user = df[df[u_col].apply(lambda value: re.sub(r"\s+", "", str(value).strip())) == current_user_key].copy()
    else:
        df_user = df.iloc[0:0].copy()
    df_user = attach_cloud_dates(df_user)
    df_user_visit = filter_visit_rows(df_user)
    df_visit_all = filter_visit_rows(df)
    df_visit_all_with_cloud = attach_cloud_dates(df_visit_all)

    st.markdown("### 담당자별 활동 수치")
    res, err, dup = process_performance_analysis(df_visit_all, st.session_state.get("auto_prev_df"))

    if isinstance(res, pd.DataFrame) and not res.empty:
        inferred_year_month = get_uploaded_month(analysis_df) or (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m")
        cloud_df = st.session_state.get("cloud_sheet_df")
        if isinstance(cloud_df, pd.DataFrame) and not cloud_df.empty:
            cloud_counts = clean_header_logic(cloud_df.copy())
            cloud_owner_col = find_col(cloud_counts, ["담당자", "등록자", "성명"])
            cloud_open_month_col = find_col(cloud_counts, ["개설년월", "개설월", "개설완료년월", "개설완료월"])
            cloud_open_col = find_col(cloud_counts, ["개설완료일자", "개설일"])
            cloud_erp_month_col = find_col(cloud_counts, ["연계년월", "연계월", "ERP연계년월", "ERP연계월"])
            cloud_erp_col = find_col(cloud_counts, ["ERP연계일자", "연계일자"])
            if cloud_owner_col and cloud_owner_col in cloud_counts.columns:
                if cloud_open_col and cloud_open_col in cloud_counts.columns:
                    cloud_counts[cloud_open_col] = pd.to_datetime(cloud_counts[cloud_open_col], errors="coerce")
                if cloud_erp_col and cloud_erp_col in cloud_counts.columns:
                    cloud_counts[cloud_erp_col] = pd.to_datetime(cloud_counts[cloud_erp_col], errors="coerce")

                def normalize_staff(value):
                    return re.sub(r"\s+", "", str(value or "").strip())

                def value_to_ym(value):
                    if is_blank_value(value):
                        return ""
                    text = str(value).strip()
                    match = re.search(r"(20\d{2})\D*([01]?\d)", text)
                    if match:
                        month = int(match.group(2))
                        if 1 <= month <= 12:
                            return f"{match.group(1)}-{month:02d}"
                    parsed = parse_sheet_date(value)
                    return parsed.strftime("%Y-%m") if pd.notna(parsed) else ""

                def month_mask(series, target_ym):
                    return series.apply(value_to_ym) == target_ym

                available_months = set()
                for month_col, date_col in ((cloud_open_month_col, cloud_open_col), (cloud_erp_month_col, cloud_erp_col)):
                    if month_col and month_col in cloud_counts.columns:
                        available_months.update([ym for ym in cloud_counts[month_col].apply(value_to_ym).tolist() if ym])
                    if date_col and date_col in cloud_counts.columns:
                        available_months.update(
                            cloud_counts[date_col].dropna().dt.strftime("%Y-%m").dropna().astype(str).tolist()
                        )
                month_options = sorted(available_months, reverse=True)
                if inferred_year_month and inferred_year_month not in month_options:
                    month_options.insert(0, inferred_year_month)
                if not month_options:
                    month_options = [inferred_year_month]
                now_kst = datetime.utcnow() + timedelta(hours=9)
                prev_month = (now_kst.replace(day=1) - timedelta(days=1)).strftime("%Y-%m")
                default_target_month = prev_month if prev_month in month_options else inferred_year_month
                if st.session_state.get("upload_perf_target_year_month") not in month_options:
                    st.session_state["upload_perf_target_year_month"] = default_target_month
                elif prev_month in month_options and st.session_state.get("upload_perf_target_year_month") == inferred_year_month:
                    st.session_state["upload_perf_target_year_month"] = prev_month
                default_month_index = month_options.index(st.session_state.get("upload_perf_target_year_month", default_target_month))
                target_year_month = st.selectbox(
                    "개설/연계 집계년월",
                    month_options,
                    index=default_month_index,
                    key="upload_perf_target_year_month",
                )

                for idx, row in res.iterrows():
                    staff_name = str(row.get("담당자", "")).strip()
                    staff_cloud = cloud_counts[
                        cloud_counts[cloud_owner_col].apply(normalize_staff) == normalize_staff(staff_name)
                    ].copy()
                    if cloud_open_month_col and cloud_open_month_col in staff_cloud.columns:
                        open_count = int(month_mask(staff_cloud[cloud_open_month_col], target_year_month).sum())
                    elif cloud_open_col and cloud_open_col in staff_cloud.columns:
                        open_count = int(
                            (
                                staff_cloud[cloud_open_col].notna()
                                & (staff_cloud[cloud_open_col].dt.strftime("%Y-%m") == target_year_month)
                            ).sum()
                        )
                    else:
                        open_count = 0
                    if cloud_erp_month_col and cloud_erp_month_col in staff_cloud.columns:
                        link_count = int(month_mask(staff_cloud[cloud_erp_month_col], target_year_month).sum())
                    elif cloud_erp_col and cloud_erp_col in staff_cloud.columns:
                        link_count = int(
                            (
                                staff_cloud[cloud_erp_col].notna()
                                & (staff_cloud[cloud_erp_col].dt.strftime("%Y-%m") == target_year_month)
                            ).sum()
                        )
                    else:
                        link_count = 0

                    res.at[idx, "개설건수"] = open_count
                    res.at[idx, "개설포인트"] = open_count * 90
                    res.at[idx, "연계건수"] = link_count
                    res.at[idx, "연계포인트"] = link_count * 120
                    operation_points = int(float(row.get("운영포인트 (실제 활동)", 0) or 0))
                    manual_points = int(float(row.get("운영포인트(추가 활동)", row.get("운영포인트 (추가 활동)", 0)) or 0))
                    res.at[idx, "합계포인트"] = min(2800, (open_count * 90) + (link_count * 120) + operation_points + manual_points)
                res = apply_rs_allowance_formula(res, st.session_state.user_db)

        summary_cols = [
            "담당자", "직급",
            "개설건수", "개설포인트",
            "연계건수", "연계포인트",
            "운영건수 (실제 활동)", "운영포인트 (실제 활동)",
            "합계포인트", "지급포인트", "지급예상금액",
        ]
        summary_display = res[[col for col in summary_cols if col in res.columns]].copy()
        total_row = {col: "" for col in summary_display.columns}
        if "담당자" in total_row:
            total_row["담당자"] = "합계"
        for col in summary_display.columns:
            if col in {"담당자", "직급"}:
                continue
            total_row[col] = int(pd.to_numeric(summary_display[col], errors="coerce").fillna(0).sum())
        summary_display = pd.concat([summary_display, pd.DataFrame([total_row])], ignore_index=True)
        style_report_logic(summary_display, compact=True)

        my_res = res[
            res["담당자"].apply(lambda value: re.sub(r"\s+", "", str(value).strip())) == current_user_key
        ].copy()

        # 업로드 전 예상치 계산 (추가 활동 제외)
        if not my_res.empty:
            st.markdown("#### 내 활동 수치")
            before_res = my_res.copy()
            o_p = int(float(before_res.iloc[0].get("개설포인트", 0)))
            l_p = int(float(before_res.iloc[0].get("연계포인트", 0)))
            v_p = int(float(before_res.iloc[0].get("운영포인트 (실제 활동)", 0)))

            # 합계포인트: 개설 + 연계 + 운영(실제)만
            before_total = min(2800, o_p + l_p + v_p)
            before_res.loc[before_res.index[0], "합계포인트"] = before_total

            # 지급포인트
            before_pay_point = max(0, before_total - 1000)
            before_res.loc[before_res.index[0], "지급포인트"] = before_pay_point

            # 지급예상금액
            before_pay = int(before_pay_point * 500)
            before_res.loc[before_res.index[0], "지급예상금액"] = before_pay

            # 전월대비 재계산
            if "전월대비" in before_res.columns:
                prev_df = st.session_state.get("auto_prev_df")
                if prev_df is not None:
                    try:
                        prev_res, _, _ = process_performance_analysis(prev_df)
                        if isinstance(prev_res, pd.DataFrame) and not prev_res.empty:
                            p_map = prev_res.set_index("담당자")["지급예상금액"].to_dict()
                            prev_pay = p_map.get(st.session_state.user_name, 0)
                            before_res.loc[before_res.index[0], "전월대비"] = int(before_pay - prev_pay)
                    except Exception:
                        pass

            hidden_cols = ["운영건수 (추가 활동)", "운영포인트(추가 활동)"] + [c for c in before_res.columns if "전월대비" in c]
            my_res_display = before_res.drop(columns=hidden_cols, errors="ignore")
        else:
            drop_cols = [c for c in my_res.columns if "전월대비" in c]
            my_res_display = my_res.drop(columns=drop_cols, errors="ignore")

        if not my_res_display.empty:
            style_report_logic(my_res_display, compact=True)
    elif isinstance(res, str):
        st.error(res)

    # 탭 데이터 미리 계산 (경고 메시지 표시용)
    # 초과 방문 데이터 계산
    err_filtered = build_upload_over_visit_df(df_visit_all)
    daily_visit_matrix = build_daily_visit_matrix_df(df_visit_all)

    # 기타 오류 데이터 계산
    other_errors_df = build_other_validation_errors(df_visit_all_with_cloud)

    # 중복 이력 데이터
    dup_my = pd.DataFrame()
    if dup is not None and not dup.empty:
        dup_my = dup.copy()

    # 개설완료일자 누락
    missing_open = pd.DataFrame()
    if "본사 개설완료일자" in df_visit_all_with_cloud.columns:
        missing_open = df_visit_all_with_cloud[
            pd.isna(df_visit_all_with_cloud["본사 개설완료일자"]) | (df_visit_all_with_cloud["본사 개설완료일자"].astype(str).str.strip() == "")
        ]
        if "본사 신규이행구분" in missing_open.columns:
            missing_open = missing_open[missing_open["본사 신규이행구분"].astype(str).str.strip() != "이행"]

    # ERP연계일자 누락
    missing_erp = pd.DataFrame()
    if "본사 ERP연계일자" in df_visit_all_with_cloud.columns:
        if d_col and d_col in df_visit_all_with_cloud.columns:
            target = df_visit_all_with_cloud[df_visit_all_with_cloud[d_col].astype(str).str.contains("연계", na=False)]
        else:
            target = df_visit_all_with_cloud
        missing_erp = target[
            pd.isna(target["본사 ERP연계일자"]) | (target["본사 ERP연계일자"].astype(str).str.strip() == "")
        ]
        # 신규/이행구분이 "이행"이고 이행추가연계가 비어있으면 ERP연계 불필요 → 제외
        if "본사 신규이행구분" in missing_erp.columns and "본사 이행추가연계" in missing_erp.columns:
            _is_ihang = missing_erp["본사 신규이행구분"].astype(str).str.strip() == "이행"
            _add_empty = pd.isna(missing_erp["본사 이행추가연계"]) | (missing_erp["본사 이행추가연계"].astype(str).str.strip() == "")
            missing_erp = missing_erp[~(_is_ihang & _add_empty)]

    # 검증 이슈 확인
    has_validation_issues = (
        not dup_my.empty or
        not err_filtered.empty or
        not missing_open.empty or
        not missing_erp.empty or
        not other_errors_df.empty
    )

    # 경고 메시지 표시
    error_tabs = []
    if not dup_my.empty:
        error_tabs.append("중복 이력")
    if not err_filtered.empty:
        error_tabs.append("초과 방문")
    if not missing_open.empty:
        error_tabs.append("본사 개설완료일자 누락")
    if not missing_erp.empty:
        error_tabs.append("본사 ERP연계일자 누락")
    if not other_errors_df.empty:
        error_tabs.append("기타 오류")

    st.divider()
    st.markdown("### 추가 실적 표")

    base = criteria_df()
    saved = load_db(PERF_FILE, {}).get(st.session_state.user_name, {})
    base["입력(건)"] = base["구분"].map(saved).fillna(0).astype(int)

    edited = render_manual_perf_input_table(base)

    edited["입력(건)"] = pd.to_numeric(edited["입력(건)"], errors="coerce").fillna(0).astype(int).clip(lower=0)
    edited["계산점수"] = edited["단위 점수"] * edited["입력(건)"]
    manual_override = st.session_state.setdefault("manual_perf_preview_override", {})
    manual_override[st.session_state.user_name] = edited.set_index("구분")["입력(건)"].to_dict()

    total = calculate_manual_perf_total(edited)
    st.session_state["_manual_perf_total"] = total

    st.metric("추가 실적 합산 점수", f"{total:,} PT")

    st.divider()
    preview_source_df = st.session_state.get("history_convert_preview_data", converted_preview_df)
    if not isinstance(preview_source_df, pd.DataFrame):
        preview_source_df = converted_preview_df if isinstance(converted_preview_df, pd.DataFrame) else df_user_visit
    preview_source_df = normalize_converted_history_df(preview_source_df)
    search_source_df = df_visit_all if isinstance(df_visit_all, pd.DataFrame) and not df_visit_all.empty else preview_source_df
    if search_source_df is None or search_source_df.empty:
        search_source_df = df_user_visit
    search_staff_col = find_col(search_source_df, ["등록자", "담당자", "성명"]) if isinstance(search_source_df, pd.DataFrame) else None
    if search_staff_col and search_staff_col in search_source_df.columns:
        staff_values = {
            re.sub(r"\s+", "", str(value).strip())
            for value in search_source_df[search_staff_col].dropna().tolist()
            if str(value).strip()
        }
        selected_staff_key = re.sub(r"\s+", "", str(st.session_state.get("history_preview_search_staff", "")).strip())
        current_staff_key = re.sub(r"\s+", "", str(st.session_state.user_name).strip())
        if len(staff_values) > 1 and selected_staff_key == current_staff_key:
            st.session_state["history_preview_search_staff"] = "전체"
    search_filters = {
        "company": str(st.session_state.get("history_preview_search_company", "") or "").strip(),
        "staff": st.session_state.get("history_preview_search_staff", "전체"),
        "date": st.session_state.get("history_preview_search_date", "전체"),
        "category": st.session_state.get("history_preview_search_category", "전체"),
        "detail": st.session_state.get("history_preview_search_detail", "전체"),
    }

    _render_preview = None if st.session_state.get("user_excel_source") == "hq" else converted_preview_df
    if (
        st.session_state.get("user_excel_source") != "hq"
        and (_render_preview is None or (isinstance(_render_preview, pd.DataFrame) and _render_preview.empty))
    ):
        _sess_prev = st.session_state.get("history_convert_preview_data")
        if isinstance(_sess_prev, pd.DataFrame) and not _sess_prev.empty:
            _render_preview = _sess_prev
    if (
        st.session_state.get("user_excel_source") != "hq"
        and (_render_preview is None or (isinstance(_render_preview, pd.DataFrame) and _render_preview.empty))
        and isinstance(analysis_df, pd.DataFrame)
    ):
        _render_preview = analysis_df

    if _render_preview is not None and not _render_preview.empty:
        _pre_h = None
        _pd_before = st.session_state.get("history_convert_preview_data")
        if isinstance(_pd_before, pd.DataFrame) and not _pd_before.empty:
            try:
                _pre_h = int(pd.util.hash_pandas_object(_pd_before).sum())
            except Exception:
                pass
        render_converted_preview_editor(_render_preview, search_filters)
        _changed = False
        _pd_after = st.session_state.get("history_convert_preview_data")
        if isinstance(_pd_after, pd.DataFrame) and not _pd_after.empty:
            try:
                if _pre_h is not None and int(pd.util.hash_pandas_object(_pd_after).sum()) != _pre_h:
                    _changed = True
            except Exception:
                pass
        if _changed:
            st.rerun()

        # 추가 방문 운영 이력 등록 안내
        _add_visit_cnt = round(total / 30)
        if _add_visit_cnt > 0:
            st.markdown(
                f"<div style='margin-top:10px;padding:12px 16px;background:#EBF8FF;border-left:4px solid #4299E1;border-radius:6px;font-size:14px;color:#2B6CB0;'>"
                f"<b>📋 추가 이력 등록 안내</b><br>"
                f"추가 실적 합산 점수 <b>{total:,} PT</b> 기준으로 "
                f"<b style='font-size:16px;color:#C05621;'>{_add_visit_cnt}건</b>의 방문 운영 이력을 "
                f"우측 <b>이력 추가</b> 버튼을 통해 등록해주세요."
                f"</div>",
                unsafe_allow_html=True,
            )

    search_filters = render_history_search_filters(search_source_df, "history_preview_search")

    dup_my = apply_history_search_filters(dup_my, search_filters)
    err_filtered = apply_history_search_filters(err_filtered, search_filters)
    missing_open = apply_history_search_filters(missing_open, search_filters)
    missing_erp = apply_history_search_filters(missing_erp, search_filters)
    other_errors_df = apply_history_search_filters(other_errors_df, search_filters)

    t1, t2, t3, t4, t5 = validation_tabs_with_refresh("refresh_user_history_validation")
    with t1:
        if not dup_my.empty:
            _dc = next((c for c in dup_my.columns if "활동일자" in c), None) or next((c for c in dup_my.columns if "활동일" in c), None)
            dup_display = dup_my.copy()
            dup_display.insert(0, "_orig_idx", dup_my.index.astype(int))
            dup_display = dup_display.reset_index(drop=True)
            _ver = st.session_state.get("dup_editor_ver", 0)
            _dup_hidden = {c: None for c in dup_display.columns if c.startswith("_")}
            edited_dup = st.data_editor(
                dup_display,
                key=f"dup_date_editor_{_ver}",
                use_container_width=True,
                hide_index=True,
                disabled=[c for c in dup_display.columns if c != _dc],
                column_config=_dup_hidden,
            )
            if _dc and list(edited_dup[_dc]) != list(dup_display[_dc]):
                _preview = st.session_state.get("history_convert_preview_data")
                if _preview is not None:
                    _preview = _preview.copy()
                    _pd = find_col(_preview, ["활동일자", "활동일"])
                    if _pd:
                        for (_, o_row), (_, e_row) in zip(dup_display.iterrows(), edited_dup.iterrows()):
                            if str(o_row.get(_dc, "")) != str(e_row.get(_dc, "")):
                                try:
                                    _idx = int(o_row["_orig_idx"])
                                    _preview.loc[_idx, _pd] = e_row[_dc]
                                except (KeyError, IndexError, TypeError, ValueError):
                                    pass
                        st.session_state["history_convert_preview_data"] = _preview
                        st.session_state.pop("history_convert_preview_editor", None)
                        st.session_state["dup_editor_ver"] = _ver + 1
                        st.rerun()
        else:
            st.info("중복 이력이 없습니다.")
    with t2:
        if not err_filtered.empty:
            _disp = err_filtered.drop(columns=["월총방문"], errors="ignore").copy()
            _ver2 = st.session_state.get("err_editor_ver", 0)
            _err_hidden = {c: None for c in _disp.columns if c.startswith("_")}
            edited_err = st.data_editor(
                _disp,
                key=f"err_date_editor_{_ver2}",
                use_container_width=True,
                hide_index=True,
                disabled=[c for c in _disp.columns if c != "초과일자"],
                column_config=_err_hidden if _err_hidden else None,
            )
            if "초과일자" in _disp.columns and list(edited_err["초과일자"]) != list(_disp["초과일자"]):
                _preview = st.session_state.get("history_convert_preview_data")
                if _preview is not None:
                    _preview = _preview.copy()
                    _pu = find_col(_preview, ["등록자", "담당자", "성명"])
                    _pb = find_col(_preview, ["사업자번호"])
                    _pd = find_col(_preview, ["활동일자", "활동일"])
                    if _pu and _pb and _pd:
                        for (_, o_row), (_, e_row) in zip(_disp.iterrows(), edited_err.iterrows()):
                            if str(o_row["초과일자"]) != str(e_row["초과일자"]):
                                _mask = (
                                    (_preview[_pb].astype(str) == str(o_row["사업자번호"])) &
                                    (_preview[_pu].astype(str) == str(o_row["담당자"])) &
                                    (_preview[_pd].astype(str) == str(o_row["초과일자"]))
                                )
                                if _mask.any():
                                    _preview.loc[_mask, _pd] = e_row["초과일자"]
                        st.session_state["history_convert_preview_data"] = _preview
                        st.session_state.pop("history_convert_preview_editor", None)
                        st.session_state["err_editor_ver"] = _ver2 + 1
                        st.rerun()
        else:
            st.info("초과 방문 데이터가 없습니다.")
        st.markdown("##### 담당자별 일 방문횟수")
        st.caption("일 방문횟수가 6회 이상인 날짜는 빨간색으로 표시됩니다.")
        render_daily_visit_matrix(daily_visit_matrix)
        visit_change_guide = build_visit_change_guide_df(daily_visit_matrix)
        st.markdown("##### 이력 변경 추천 가이드")
        if visit_change_guide.empty:
            st.info("변경이 필요한 6회 이상 방문일이 없습니다.")
        else:
            st.caption("현재 방문횟수가 6회 이상인 날짜에서 주말/공휴일을 제외한 5회 미만 영업일로 옮기는 것을 기준으로 추천합니다.")
            st.dataframe(visit_change_guide, use_container_width=True, hide_index=True)
    with t3:
        if not missing_open.empty:
            style_report_logic(missing_open.drop(columns=["본사 ERP연계일자"], errors="ignore"))
        elif "본사 개설완료일자" not in df_visit_all_with_cloud.columns:
            st.info("본사 구글시트에 개설완료일자 또는 사업자번호 컬럼이 없어 확인할 수 없습니다.")
    with t4:
        if not missing_erp.empty:
            style_report_logic(missing_erp.drop(columns=["본사 개설완료일자"], errors="ignore"))
        elif "본사 ERP연계일자" not in df_visit_all_with_cloud.columns:
            st.info("본사 구글시트에 ERP연계일자 또는 사업자번호 컬럼이 없어 확인할 수 없습니다.")
    with t5:
        style_report_logic(other_errors_df, align_overrides={"오류 사유": "left"}, default_align="center")

    if error_tabs:
        error_list = ", ".join(error_tabs)
        st.markdown(
            f"<div style='margin-top:12px;margin-bottom:12px;padding:12px 16px;background:#FFF3CD;border-left:4px solid #FFC107;border-radius:4px;'>"
            f"<div style='font-size:14px;color:#856404;'><b>⚠️ 다음 항목을 수정해주세요:</b> {error_list}</div>"
            f"<div style='margin-top:4px;font-size:13px;color:#856404;'>위 탭에서 문제를 해결한 후 엑셀을 다시 업로드해주세요.</div>"
            f"</div>",
            unsafe_allow_html=True,
        )
    else:
        st.markdown(
            f"<div style='margin-top:12px;margin-bottom:12px;padding:12px 16px;background:#D4EDDA;border-left:4px solid #28A745;border-radius:4px;'>"
            f"<div style='font-size:14px;color:#155724;'><b>✓ 잘못된 데이터가 없습니다.</b></div>"
            f"</div>",
            unsafe_allow_html=True,
        )

    st.divider()
    st.markdown("### 최종 실적 확인")
    _ppt_report_df = res.copy() if isinstance(res, pd.DataFrame) else pd.DataFrame()
    _, _final_excel_col, _ppt_download_col = st.columns([0.5, 0.25, 0.25])
    with _final_excel_col:
        render_adjusted_history_download_button(analysis_df, "upload_bottom")
    with _ppt_download_col:
        render_upload_ppt_download_button(_ppt_report_df, analysis_df, "upload_bottom")
    st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)
    show_final_check()

    # 관리자 모드: user_name 복원
    if is_admin_mode:
        st.session_state.user_name = original_user_name


def show_final_check():
    # 저장된 데이터가 있으면 로드
    if current_history_analysis_df().empty and st.session_state.user_excel_data is None:
        saved_db = load_db(SAVED_STATE_FILE, {})
        user_saved = saved_db.get(st.session_state.user_name)
        if user_saved and user_saved.get("user_excel_data"):
            st.session_state.user_excel_data = pd.DataFrame.from_dict(user_saved["user_excel_data"])
            if user_saved.get("user_prev_month_sel"):
                st.session_state.user_prev_month_sel = user_saved["user_prev_month_sel"]

    select_prev_month("auto_prev_df", "user_prev_month_sel")

    original_df = current_history_analysis_df()
    if original_df.empty and st.session_state.user_excel_data is not None:
        original_df = prepare_history_analysis_df(st.session_state.user_excel_data)

    if original_df is None or original_df.empty:
        st.info("먼저 업로드 및 실적 확인 메뉴에서 엑셀을 업로드해주세요.")
        return

    # cloud_sheet_df가 없을 때만 로드 시도
    if st.session_state.get("cloud_sheet_df") is None:
        try:
            load_csv_to_state("url_sync", "cloud_sheet_df")
        except Exception as e:
            st.warning(f"⚠️ 본사 구글시트를 불러오는데 실패했습니다: {str(e)}")

    # Filter user data and attach cloud dates for tabs
    df_for_tabs = original_df.copy()
    u_col = find_col(df_for_tabs, ["등록자", "담당자", "성명"], "등록자")
    d_col = find_col(df_for_tabs, ["활동상세", "활동내용"], "활동상세")
    current_user_key = re.sub(r"\s+", "", str(st.session_state.user_name).strip())
    if u_col in df_for_tabs.columns:
        df_user = df_for_tabs[
            df_for_tabs[u_col].apply(lambda value: re.sub(r"\s+", "", str(value).strip())) == current_user_key
        ].copy()
    else:
        df_user = pd.DataFrame()
    df_user = attach_cloud_dates(df_user)
    df_user_visit = filter_visit_rows(df_user)

    res, err, dup = process_performance_analysis(filter_visit_rows(original_df), st.session_state.get("auto_prev_df"))

    if not isinstance(res, pd.DataFrame) or res.empty:
        st.error(res if isinstance(res, str) else "실적을 계산할 수 없습니다.")
        return

    my_res = res[
        res["담당자"].apply(lambda value: re.sub(r"\s+", "", str(value).strip())) == current_user_key
    ]
    uname = st.session_state.user_name

    # 업로드 전 예상치 계산 (추가 활동 제외)
    if not my_res.empty:
        before_res = my_res.copy()
        o_p = int(float(before_res.iloc[0].get("개설포인트", 0)))
        l_p = int(float(before_res.iloc[0].get("연계포인트", 0)))
        v_p = int(float(before_res.iloc[0].get("운영포인트 (실제 활동)", 0)))

        # 합계포인트: 개설 + 연계 + 운영(실제)만
        before_total = min(2800, o_p + l_p + v_p)
        before_res.loc[before_res.index[0], "합계포인트"] = before_total

        # 지급포인트
        before_pay_point = max(0, before_total - 1000)
        before_res.loc[before_res.index[0], "지급포인트"] = before_pay_point

        # 지급예상금액
        before_pay = int(before_pay_point * 500)
        before_res.loc[before_res.index[0], "지급예상금액"] = before_pay

        # 전월대비 재계산 (필요시)
        if "전월대비" in before_res.columns:
            prev_df = st.session_state.get("auto_prev_df")
            if prev_df is not None:
                try:
                    prev_res, _, _ = process_performance_analysis(prev_df)
                    if isinstance(prev_res, pd.DataFrame) and not prev_res.empty:
                        p_map = prev_res.set_index("담당자")["지급예상금액"].to_dict()
                        prev_pay = p_map.get(name, 0)
                        before_res.loc[before_res.index[0], "전월대비"] = int(before_pay - prev_pay)
                except Exception:
                    pass

    if my_res.empty:
        return

    # 현재 실적 표시
    st.markdown("##### 실적 현황")
    display_res = my_res.drop(columns=["전송시각"], errors="ignore")
    style_report_logic(display_res, compact=True)

    # ── 중복방문/초과방문/누락 확인 탭 ──
    df_user_check = df_user_visit
    _, err_check, dup_check = process_performance_analysis(filter_visit_rows(original_df), st.session_state.get("auto_prev_df"))

    # 탭 데이터 미리 계산 (경고 메시지 표시용 - 재업로드한 경우만)
    # 중복 이력
    dup_my = pd.DataFrame()
    if dup_check is not None and not dup_check.empty:
        u_col_dup = find_col(dup_check, ["등록자", "담당자", "성명"], "담당자")
        if u_col_dup and u_col_dup in dup_check.columns:
            dup_my = dup_check[
                dup_check[u_col_dup].apply(lambda value: re.sub(r"\s+", "", str(value).strip())) == current_user_key
            ]
        else:
            dup_my = dup_check

    # 초과 방문
    err_filtered_final = pd.DataFrame()
    if err_check is not None and not err_check.empty:
        if "담당자" in err_check.columns:
            err_my_final = err_check[
                err_check["담당자"].apply(lambda value: re.sub(r"\s+", "", str(value).strip())) == current_user_key
            ].copy()
        else:
            err_my_final = err_check.copy()
        err_filtered_final = err_my_final[err_my_final["일방문"] >= 6].copy()

    # 개설완료일자 누락
    missing_open_final = pd.DataFrame()
    if "본사 개설완료일자" in df_user_check.columns:
        missing_open_final = df_user_check[
            pd.isna(df_user_check["본사 개설완료일자"]) | (df_user_check["본사 개설완료일자"].astype(str).str.strip() == "")
        ]
        if "본사 신규이행구분" in missing_open_final.columns:
            missing_open_final = missing_open_final[missing_open_final["본사 신규이행구분"].astype(str).str.strip() != "이행"]

    # ERP연계일자 누락
    missing_erp_final = pd.DataFrame()
    if "본사 ERP연계일자" in df_user_check.columns:
        d_col_check = find_col(df_user_check, ["활동상세", "활동내용"], "활동상세")
        if d_col_check and d_col_check in df_user_check.columns:
            target = df_user_check[df_user_check[d_col_check].astype(str).str.contains("연계", na=False)]
        else:
            target = df_user_check
        missing_erp_final = target[
            pd.isna(target["본사 ERP연계일자"]) | (target["본사 ERP연계일자"].astype(str).str.strip() == "")
        ]
        if "본사 신규이행구분" in missing_erp_final.columns and "본사 이행추가연계" in missing_erp_final.columns:
            _is_ihang = missing_erp_final["본사 신규이행구분"].astype(str).str.strip() == "이행"
            _add_empty = pd.isna(missing_erp_final["본사 이행추가연계"]) | (missing_erp_final["본사 이행추가연계"].astype(str).str.strip() == "")
            missing_erp_final = missing_erp_final[~(_is_ihang & _add_empty)]

    # 기타 오류
    other_errors_df_final = build_other_validation_errors(df_user_check)

    # 각 탭의 데이터 존재 여부 확인
    has_dup_data = not dup_my.empty
    has_err_data = not err_filtered_final.empty
    has_missing_open = not missing_open_final.empty
    has_missing_erp = not missing_erp_final.empty
    has_other_errors = not other_errors_df_final.empty

    # 검증 이슈 체크
    has_validation_issues = has_dup_data or has_err_data or has_missing_open or has_missing_erp or has_other_errors

    # 추가 이력 등록건수 불일치 체크
    total = st.session_state.get("_manual_perf_total", 0)
    _req_visit_cnt = round(total / 30) if total > 0 else 0
    _preview_df_check = st.session_state.get("history_convert_preview_data")
    if _req_visit_cnt > 0 and isinstance(_preview_df_check, pd.DataFrame) and not _preview_df_check.empty:
        _is_manual_mask = _preview_df_check.get("_is_manual", pd.Series([False] * len(_preview_df_check))).fillna(False).astype(bool)
        _is_visit_mask = _preview_df_check["활동구분"].astype(str).str.contains("방문", na=False) if "활동구분" in _preview_df_check.columns else pd.Series([False] * len(_preview_df_check))
        _is_oper_mask = _preview_df_check["활동상세"].astype(str).str.contains("운영", na=False) if "활동상세" in _preview_df_check.columns else pd.Series([False] * len(_preview_df_check))
        _actual_visit_cnt = int((_is_manual_mask & _is_visit_mask & _is_oper_mask).sum())
        _visit_count_mismatch = _actual_visit_cnt != _req_visit_cnt
    else:
        _actual_visit_cnt = 0
        _visit_count_mismatch = False

    # 전송 가능 여부
    can_send = not has_validation_issues and not _visit_count_mismatch

    if has_validation_issues:
        st.markdown(
            "<div style='margin-top:8px;padding:10px 16px;background:#FFF5F5;border:1px solid #FC8181;border-radius:8px;font-size:13px;color:#C53030;font-weight:700;'>"
            "❌ 검증 오류가 있습니다. 위 탭에서 문제를 해결한 후 실적을 전송해주세요.</div>",
            unsafe_allow_html=True,
        )
    if _visit_count_mismatch:
        st.markdown(
            f"<div style='margin-top:8px;padding:10px 16px;background:#FFFAF0;border:1px solid #F6AD55;border-radius:8px;font-size:13px;color:#C05621;font-weight:700;'>"
            f"⚠️ 추가 이력등록건수 불일치 — 안내 기준 <b>{_req_visit_cnt}건</b> 필요, 현재 등록 <b>{_actual_visit_cnt}건</b>.<br>"
            f"방문 운영 이력을 <b>{_req_visit_cnt - _actual_visit_cnt}건</b> 추가 등록해주세요.</div>",
            unsafe_allow_html=True,
        )
    if not has_validation_issues and not my_res.empty:
        개설건수 = int(my_res.iloc[0].get("개설건수", 0))
        연계건수 = int(my_res.iloc[0].get("연계건수", 0))
        운영건수_실제 = int(my_res.iloc[0].get("운영건수 (실제 활동)", 0))
        운영건수_추가 = int(my_res.iloc[0].get("운영건수 (추가 활동)", 0))
        운영건수 = min(60, 운영건수_실제 + 운영건수_추가)
        지급예상금액 = int(my_res.iloc[0].get("지급예상금액", 0))

        전월대비_text = ""
        전월대비_col = next((c for c in my_res.columns if "전월대비" in c), None)
        if 전월대비_col:
            전월대비 = int(my_res.iloc[0].get(전월대비_col, 0))
            user_sel = st.session_state.get("user_prev_month_sel", "선택안함")
            if user_sel and user_sel != "선택안함":
                prev_m = str(int(user_sel.split("-")[1])) + "월"
                비교월_label = f"{prev_m} 대비"
            else:
                비교월_label = "전월 대비"
            if 전월대비 > 0:
                전월대비_text = f" {비교월_label} <b>{전월대비:,}</b>원 증가하였습니다."
            elif 전월대비 < 0:
                전월대비_text = f" {비교월_label} <b>{abs(전월대비):,}</b>원 감소하였습니다."

        st.markdown(
            f"<div style='margin-top:8px;padding:10px 16px;background:#EBF8FF;border-radius:8px;font-size:13px;color:#2B6CB0;'>"
            f"<b>{html.escape(uname)}</b>님의 최종 실적은 개설 <b>{개설건수}</b>개, 연계 <b>{연계건수}</b>개, 운영 <b>{운영건수}</b>개 에 금액은 <b>{지급예상금액:,}</b>원 입니다."
            f"{전월대비_text}</div>",
            unsafe_allow_html=True,
        )
        if not _visit_count_mismatch:
            st.markdown(
                "<div style='margin-top:8px;padding:10px 16px;background:#F0FFF4;border:1px solid #9AE6B4;border-radius:8px;font-size:13px;color:#276749;font-weight:700;'>"
                "✅ 실적 결과를 전송할 수 있습니다.</div>",
                unsafe_allow_html=True,
            )

    st.markdown("<div style='height:32px'></div>", unsafe_allow_html=True)

    _, ppt_col, save_col, send_col = st.columns([0.45, 0.25, 0.15, 0.15])

    with ppt_col:
        render_upload_ppt_download_button(display_res.copy(), original_df, "final_user")

    with save_col:
        st.markdown('<div class="action-btn">', unsafe_allow_html=True)
        do_save = st.button("저장", use_container_width=True, type="secondary")
        st.markdown('</div>', unsafe_allow_html=True)

    with send_col:
        st.markdown('<div class="action-btn">', unsafe_allow_html=True)
        do_send = st.button("실적 결과 전송", use_container_width=True, type="primary", disabled=not can_send)
        st.markdown('</div>', unsafe_allow_html=True)

    if do_save:
        save_manual_perf_override_for_current_user()
        saved_db = load_db(SAVED_STATE_FILE, {})
        user_saved = {
            "user_excel_data": st.session_state.user_excel_data.to_dict() if st.session_state.user_excel_data is not None else None,
            "user_prev_month_sel": st.session_state.get("user_prev_month_sel", "선택안함"),
            "saved_at": (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S")
        }
        saved_db[st.session_state.user_name] = user_saved
        save_db(SAVED_STATE_FILE, saved_db)
        st.success("저장 완료")
        time.sleep(0.5)
        st.rerun()

    if do_send:
        save_manual_perf_override_for_current_user()
        sent_db = load_db(SENT_FILE, {})
        sent_uploads_db = load_db(SENT_UPLOADS_FILE, {})
        row_data = my_res.iloc[0].to_dict()
        row_data["전송시각"] = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S")

        ym = get_uploaded_month(original_df)
        if ym:
            row_data["이력작성월"] = ym

        user_sel = st.session_state.get("user_prev_month_sel", "선택안함")
        if "전월대비" in row_data and user_sel and user_sel != "선택안함":
            prev_m = str(int(user_sel.split("-")[1])) + "월"
            row_data[f"전월대비({prev_m})"] = row_data.pop("전월대비")

        sent_db[st.session_state.user_name] = row_data

        # 본인 이름에 맞는 이력만 저장
        u_col_send = find_col(original_df, ["등록자", "담당자", "성명"], "등록자")
        if u_col_send and u_col_send in original_df.columns:
            user_only_df = original_df[
                original_df[u_col_send].apply(lambda value: re.sub(r"\s+", "", str(value).strip())) == current_user_key
            ].copy()
            sent_uploads_db[st.session_state.user_name] = dataframe_to_upload_payload(user_only_df)
        else:
            sent_uploads_db[st.session_state.user_name] = dataframe_to_upload_payload(original_df)

        save_db(SENT_FILE, sent_db)
        save_db(SENT_UPLOADS_FILE, sent_uploads_db)

        st.success("전송 완료")
        time.sleep(3)
        st.rerun()


def sent_results_df():
    sent_db = load_db(SENT_FILE, {})
    if not sent_db:
        return pd.DataFrame()

    staff_names = get_staff_names()
    if staff_names:
        sent_db = {k: v for k, v in sent_db.items() if k in staff_names}
    if not sent_db:
        return pd.DataFrame()

    sent_df = pd.DataFrame(list(sent_db.values()))
    sent_df = sent_df.drop(columns=[c for c in sent_df.columns if "관리자전월대비" in c], errors="ignore")

    if "이력작성월" in sent_df.columns:
        idx = sent_df.columns.tolist().index(next((c for c in sent_df.columns if "전월대비" in c), sent_df.columns[-1]))
        sent_df.insert(idx, "등록월", sent_df["이력작성월"])
        sent_df = sent_df.drop(columns=["이력작성월"], errors="ignore")
    elif "전송시각" in sent_df.columns:
        idx = sent_df.columns.tolist().index(next((c for c in sent_df.columns if "전월대비" in c), sent_df.columns[-1]))
        sent_df.insert(idx, "등록월", sent_df["전송시각"].astype(str).str[:7])

    return sort_by_rank_name(apply_rank_from_user_db(sent_df))


def apply_admin_prev_diff(sent_df):
    if sent_df.empty:
        return sent_df

    result_df = sent_df.copy()
    result_df = result_df.drop(columns=[c for c in result_df.columns if "전월대비" in str(c)], errors="ignore")

    selected_month = st.session_state.get("adm_prev_month", "선택안함")
    prev_df = st.session_state.get("auto_prev_df")

    if selected_month == "선택안함" or prev_df is None or prev_df.empty:
        return result_df

    prev_res, _, _ = process_performance_analysis(prev_df)
    if not isinstance(prev_res, pd.DataFrame) or prev_res.empty:
        return result_df

    prev_pay_map = prev_res.set_index("담당자")["지급예상금액"].to_dict()

    if "담당자" in result_df.columns and "지급예상금액" in result_df.columns:
        diff_col = str(int(selected_month.split("-")[1])) + "월 대비"
        result_df[diff_col] = result_df.apply(
            lambda r: int(float(r.get("지급예상금액", 0))) - int(float(prev_pay_map.get(r.get("담당자"), 0))),
            axis=1,
        )

        cols = result_df.columns.tolist()
        cols.remove(diff_col)
        insert_at = cols.index("전송시각") if "전송시각" in cols else len(cols)
        cols.insert(insert_at, diff_col)
        result_df = result_df[cols]

    return result_df


def dataframe_to_excel_bytes(sheets):
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        for sheet_name, df in sheets.items():
            safe_name = str(sheet_name)[:31] or "Sheet"
            df.to_excel(writer, index=False, sheet_name=safe_name)
    output.seek(0)
    return output.getvalue()


def sent_uploaded_files_excel_bytes():
    sent_uploads_db = load_db(SENT_UPLOADS_FILE, {})
    if not sent_uploads_db:
        return dataframe_to_excel_bytes({
            "안내": pd.DataFrame([{"내용": "아직 전송된 업로드 엑셀 파일이 없습니다. 직원이 실적 결과 전송을 다시 진행하면 이 파일에 포함됩니다."}])
        })

    staff_names = get_staff_names()
    if staff_names:
        sent_uploads_db = {k: v for k, v in sent_uploads_db.items() if k in staff_names}
    if not sent_uploads_db:
        return dataframe_to_excel_bytes({
            "안내": pd.DataFrame([{"내용": "직원 목록에 해당하는 전송 데이터가 없습니다."}])
        })

    analysis_result = st.session_state.get("analysis_result")
    if analysis_result is not None and not analysis_result.empty:
        if "담당자" in analysis_result.columns:
            report_names = set(analysis_result["담당자"].unique())
        else:
            report_names = set(sent_uploads_db.keys())
    else:
        report_names = set(sent_uploads_db.keys())

    all_data = []
    for name, payload in sent_uploads_db.items():
        if name in report_names:
            df = upload_payload_to_dataframe(payload)
            if not df.empty:
                all_data.append(df)

    if not all_data:
        return dataframe_to_excel_bytes({
            "안내": pd.DataFrame([{"내용": "저장된 업로드 파일 데이터가 비어 있습니다."}])
        })

    # 원본 업로드 양식과 동일하게 저장
    combined_df = pd.concat(all_data, ignore_index=True)
    return dataframe_to_excel_bytes({"전체 활동 이력": combined_df})


def report_month_info(report_df):
    ym = ""
    if "등록월" in report_df.columns:
        months = report_df["등록월"].dropna().astype(str)
        months = months[months.str.match(r"^\d{4}-\d{2}$", na=False)]
        if not months.empty:
            ym = months.value_counts().idxmax()
    if not ym:
        ym = get_uploaded_month(st.session_state.user_excel_data) if st.session_state.user_excel_data is not None else ""
    if not ym:
        ym = datetime.utcnow().strftime("%Y-%m")
    year = int(ym.split("-")[0])
    month = int(ym.split("-")[1])
    return ym, year, month


def as_date_series(series):
    return pd.to_datetime(series, errors="coerce")


def cloud_customer_counts(name=None):
    cloud = st.session_state.get("cloud_sheet_df")
    if cloud is None:
        try:
            load_csv_to_state("url_sync", "cloud_sheet_df")
            cloud = st.session_state.get("cloud_sheet_df")
        except Exception:
            cloud = None

    analysis_source_df = st.session_state.get("analysis_result")
    if not isinstance(analysis_source_df, pd.DataFrame):
        analysis_source_df = pd.DataFrame()
    ym, year, _ = report_month_info(analysis_source_df)
    empty_counts = {
        "manage_total": 0, "manage_link": 0,
        "year_open": 0, "year_link": 0,
        "month_open": 0, "month_link": 0,
    }
    if cloud is None or cloud.empty:
        return empty_counts

    df = clean_header_logic(cloud.copy())
    owner_col = find_col(df, ["담당자", "등록자", "성명"])
    if name and owner_col and owner_col in df.columns:
        df = df[df[owner_col].astype(str).str.strip() == str(name).strip()].copy()

    status_col = find_col(df, ["상태항목", "상태"])
    open_col = find_col(df, ["개설완료일자", "개설일"])
    erp_col = find_col(df, ["ERP연계일자", "ERP연계", "연계일자"])

    if status_col and status_col in df.columns:
        status = df[status_col].astype(str).str.strip()
        manage_total = int(status.isin(["대기", "완료"]).sum())
        manage_link = int(((status == "완료") & as_date_series(df[erp_col]).notna()).sum()) if erp_col and erp_col in df.columns else 0
    else:
        manage_total = 0
        manage_link = int(as_date_series(df[erp_col]).notna().sum()) if erp_col and erp_col in df.columns else 0

    open_dates = as_date_series(df[open_col]) if open_col and open_col in df.columns else pd.Series(dtype="datetime64[ns]")
    erp_dates = as_date_series(df[erp_col]) if erp_col and erp_col in df.columns else pd.Series(dtype="datetime64[ns]")

    return {
        "manage_total": manage_total,
        "manage_link": manage_link,
        "year_open": int((open_dates.dt.year == year).sum()) if not open_dates.empty else 0,
        "year_link": int((erp_dates.dt.year == year).sum()) if not erp_dates.empty else 0,
        "month_open": int((open_dates.dt.strftime("%Y-%m") == ym).sum()) if not open_dates.empty else 0,
        "month_link": int((erp_dates.dt.strftime("%Y-%m") == ym).sum()) if not erp_dates.empty else 0,
    }


def sent_activity_counts(report_df, name=None, upload_df=None):
    if isinstance(upload_df, pd.DataFrame) and not upload_df.empty:
        upload = clean_header_logic(upload_df.copy())
        u_col = find_col(upload, ["등록자", "담당자", "성명"])
        d_col = find_col(upload, ["활동상세", "활동내용"])
        biz_col = find_col(upload, ["사업자번호"])
        if name and u_col and u_col in upload.columns:
            upload = upload[upload[u_col].astype(str).str.strip() == str(name).strip()].copy()
        if d_col and d_col in upload.columns:
            detail = upload[d_col].astype(str)
            open_count = int(detail.str.contains("개설", na=False).sum())
            link_count = int(detail.str.contains("연계", na=False).sum())
            operation_count = int(detail.str.contains("개설|연계|운영|방문|점검", na=False).sum())
            total_count = operation_count
            unique_customers = 0
            if biz_col and biz_col in upload.columns:
                unique_customers = len({v for v in normalize_biz(upload[biz_col]).tolist() if v})
            return {
                "open": open_count,
                "link": link_count,
                "operation": operation_count,
                "total": total_count,
                "unique_customers": unique_customers or total_count,
            }

    df = report_df.copy()
    if name and "담당자" in df.columns:
        df = df[df["담당자"].astype(str).str.strip() == str(name).strip()]

    def sum_col(col):
        return int(pd.to_numeric(df[col], errors="coerce").fillna(0).sum()) if col in df.columns else 0

    open_count = sum_col("개설건수")
    link_count = sum_col("연계건수")
    operation_col = "운영건수 (실제 활동)" if "운영건수 (실제 활동)" in df.columns else "운영건수"
    operation_count = sum_col(operation_col)
    total_count = open_count + link_count + operation_count

    unique_customers = uploaded_activity_customer_count(name)
    if unique_customers == 0:
        unique_customers = total_count

    return {
        "open": open_count,
        "link": link_count,
        "operation": operation_count,
        "total": total_count,
        "unique_customers": unique_customers,
    }


def uploaded_activity_customer_count(name=None):
    sent_uploads_db = load_db(SENT_UPLOADS_FILE, {})
    biz_keys = set()

    for saved_name, payload in sent_uploads_db.items():
        df = upload_payload_to_dataframe(payload)
        if df.empty:
            continue
        u_col = find_col(df, ["등록자", "담당자", "성명"])
        d_col = find_col(df, ["활동상세", "활동내용"])
        biz_col = find_col(df, ["사업자번호"])
        if name and str(saved_name).strip() != str(name).strip():
            if not u_col or u_col not in df.columns:
                continue
            df = df[df[u_col].astype(str).str.strip() == str(name).strip()].copy()
        if d_col and d_col in df.columns:
            df = df[df[d_col].astype(str).str.contains("개설|연계|운영|방문|점검", na=False)]
        if biz_col and biz_col in df.columns:
            biz_keys.update([v for v in normalize_biz(df[biz_col]).tolist() if v])

    return len(biz_keys)


def summarize_activity_note(value, max_lines=2, max_chars_per_line=34):
    text = "" if value is None else str(value)
    text = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text or text.lower() == "nan":
        return ""

    parts = [
        re.sub(r"\s+", " ", part).strip(" -ㆍ*•●○◦·.,;:/\\|")
        for part in re.split(r"(?:\n+|[.;。]|[.!?]\s+)", text)
        if part.strip(" -ㆍ*•●○◦·.,;:/\\|")
    ]
    if not parts:
        parts = [re.sub(r"^[^\w가-힣]+", "", text)]

    lines = []
    for part in parts:
        part = re.sub(r"^[^\w가-힣]+", "", part).strip()
        if not part:
            continue
        while len(part) > max_chars_per_line and len(lines) < max_lines:
            cut_at = part.rfind(" ", 0, max_chars_per_line + 1)
            if cut_at < max_chars_per_line // 2:
                cut_at = max_chars_per_line
            lines.append(part[:cut_at].strip())
            part = part[cut_at:].strip()
        if part and len(lines) < max_lines:
            lines.append(part[:max_chars_per_line].strip())
        if len(lines) >= max_lines:
            break

    if not lines:
        return ""
    lines[0] = f"● {lines[0]}"
    return "\n".join(lines[:max_lines])


def is_major_note_col(table, col_idx):
    try:
        header = table.cell(0, col_idx).text.replace(" ", "")
        return "비고" in header
    except Exception:
        return False


def uploaded_major_rows(name, keyword, row_type, upload_df=None):
    if isinstance(upload_df, pd.DataFrame) and not upload_df.empty:
        df = clean_header_logic(upload_df.copy())
    else:
        payload = load_db(SENT_UPLOADS_FILE, {}).get(name)
        df = upload_payload_to_dataframe(payload)
    if df.empty:
        return []

    u_col = find_col(df, ["등록자", "담당자", "성명"])
    d_col = find_col(df, ["활동상세"]) or find_col(df, ["활동내용"])
    comp_col = find_col(df, ["업체명", "상호", "고객명"])
    product_col = find_col(df, ["상품"])
    system_col = find_col(df, ["내부시스템", "ERP", "시스템"])
    report_col = find_col(df, ["구축보고서"])
    note_col = find_col(df, ["활동내용"]) or find_col(df, ["비고", "제목"])

    if u_col and u_col in df.columns:
        df = df[df[u_col].astype(str).str.strip() == str(name).strip()]
    if d_col and d_col in df.columns:
        df = df[df[d_col].astype(str).str.contains(keyword, na=False)]

    rows = []
    for _, row in df.iterrows():
        rows.append([
            row.get(comp_col, "") if comp_col else "",
            row_type,
            row.get(product_col, "통합CMS") if product_col else "통합CMS",
            row.get(system_col, "") if system_col else "",
            row.get(report_col, "X") if report_col else "X",
            summarize_activity_note(row.get(note_col, "")) if note_col else "",
        ])
    return rows


def build_report_ppt_bytes(report_df, compare_df, curr_month_label, prev_month_label, upload_df=None):
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    from copy import deepcopy

    if not os.path.exists(PPT_TEMPLATE_FILE):
        raise FileNotFoundError(f"PPT 템플릿 파일을 찾을 수 없습니다: {PPT_TEMPLATE_FILE}")

    if report_df.empty:
        raise ValueError("실적 데이터가 비어있습니다. 마감 후 다시 시도해주세요.")

    prs = Presentation(PPT_TEMPLATE_FILE)
    ym, _, _ = report_month_info(report_df)

    def fmt(value):
        if value == "":
            return ""
        if value is None:
            return "-"
        try:
            if pd.isna(value):
                return ""
        except Exception:
            pass
        try:
            if isinstance(value, str) and value.strip() == "-":
                return "-"
            if isinstance(value, (int, float, np.integer, np.floating)):
                return f"{int(value):,}"
            if str(value).replace(",", "").replace("-", "").isdigit():
                return f"{int(str(value).replace(',', '')):,}"
        except Exception:
            pass
        return str(value)

    def set_cell_text(cell, value, align=None, font_size=14):
        cell.text = fmt(value)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = align if align is not None else PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = "맑은 고딕"
                run.font.size = Pt(font_size)

    def slide_tables(slide):
        return [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]

    def ensure_rows(table, required_rows):
        if len(table.rows) == 0 or len(table._tbl.tr_lst) == 0:
            return
        while len(table.rows) < required_rows:
            try:
                new_tr = deepcopy(table._tbl.tr_lst[-1])
                table._tbl.append(new_tr)
                for cell in table.rows[-1].cells:
                    cell.text = ""
            except (IndexError, AttributeError):
                break

    def fill_table_row(table, row_idx, values):
        ensure_rows(table, row_idx + 1)
        if row_idx >= len(table.rows):
            return
        for col_idx, value in enumerate(values):
            if col_idx < len(table.columns):
                set_cell_text(table.cell(row_idx, col_idx), value, PP_ALIGN.CENTER)

    def fill_stats_slide(slide, counts, cloud_counts):
        tables = slide_tables(slide)
        if len(tables) < 4:
            return

        fill_table_row(tables[0], 1, ["-", counts["open"], counts["link"], counts["operation"]])
        fill_table_row(tables[2], 1, [counts["unique_customers"], counts["total"], counts["total"], "-", "-", "-"])
        fill_table_row(tables[3], 2, [
            cloud_counts["manage_total"],
            cloud_counts["manage_link"],
            "-",
            cloud_counts["year_open"],
            cloud_counts["year_link"],
            "-",
            cloud_counts["month_open"],
            cloud_counts["month_link"],
        ])

    def fill_major_slide(slide, rows):
        tables = slide_tables(slide)
        if not tables:
            return
        table = tables[0]
        if len(table.rows) == 0 or len(table.columns) == 0:
            return

        ensure_rows(table, len(rows) + 1)

        # 기존 행 초기화 (헤더 제외)
        for r in range(1, len(table.rows)):
            for c in range(len(table.columns)):
                try:
                    set_cell_text(table.cell(r, c), "", font_size=10)
                except:
                    pass

        # 데이터 채우기
        for r_idx, row_values in enumerate(rows, start=1):
            if r_idx >= len(table.rows):
                break
            if not isinstance(row_values, (list, tuple)):
                continue
            for c_idx, value in enumerate(row_values[:len(table.columns)]):
                try:
                    align = PP_ALIGN.LEFT if is_major_note_col(table, c_idx) else PP_ALIGN.CENTER
                    set_cell_text(table.cell(r_idx, c_idx), value, align, font_size=10)
                except:
                    pass

    def delete_slide(index):
        slide_id_list = prs.slides._sldIdLst
        slide_id = slide_id_list[index]
        prs.part.drop_rel(slide_id.rId)
        del slide_id_list[index]

    # 전체 브랜치 서비스(BS성과)_운영: 4페이지
    try:
        if len(prs.slides) > 3:
            fill_stats_slide(prs.slides[3], sent_activity_counts(report_df, upload_df=upload_df), cloud_customer_counts())
    except Exception as e:
        raise Exception(f"전체 브랜치 슬라이드(4페이지) 처리 중 오류: {e}")

    sections = {
        "이성환": [5, 6, 7, 8],
        "임인지": [9, 10, 11, 12],
        "전준수": [13, 14, 15, 16],
        "이수현": [17, 18, 19, 20],
        "하성춘": [21, 22, 23, 24],
        "길민종": [25, 26, 27],
    }
    sent_names = set(report_df["담당자"].dropna().astype(str)) if "담당자" in report_df.columns else set()

    for name, idxs in sections.items():
        if name not in sent_names:
            continue
        stats_idx = idxs[1]
        open_idx = idxs[2] if len(idxs) > 2 else None
        link_idx = idxs[3] if len(idxs) > 3 else None
        try:
            if stats_idx < len(prs.slides):
                fill_stats_slide(prs.slides[stats_idx], sent_activity_counts(report_df, name, upload_df), cloud_customer_counts(name))
            if open_idx is not None and open_idx < len(prs.slides):
                fill_major_slide(prs.slides[open_idx], uploaded_major_rows(name, "개설", "일반 개설", upload_df))
            if link_idx is not None and link_idx < len(prs.slides):
                fill_major_slide(prs.slides[link_idx], uploaded_major_rows(name, "연계", "추가 연계", upload_df))
        except Exception as e:
            raise Exception(f"{name} 슬라이드 처리 중 오류: {e}")

    # 전송된 담당자 외 개인성과 섹션은 제거
    remove_indices = []
    for name, idxs in sections.items():
        if name not in sent_names:
            remove_indices.extend(idxs)
    for idx in sorted(remove_indices, reverse=True):
        if idx < len(prs.slides):
            delete_slide(idx)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def render_report_action_buttons(report_df, compare_df, curr_month_label, prev_month_label):
    st.markdown("<div style='height:18px'></div>", unsafe_allow_html=True)

    is_closed = bool(st.session_state.get("report_closed"))

    dc1, dc2 = st.columns([0.85, 0.15])
    with dc2:
        if not is_closed:
            if st.button("마감", use_container_width=True, type="primary"):
                close_time = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M")
                st.session_state.report_closed = close_time
                saved_db = load_db(SAVED_STATE_FILE, {})
                saved_db["report_closed"] = {"time": close_time, "by": st.session_state.get("user_name", "")}
                save_db(SAVED_STATE_FILE, saved_db)
                st.success("실적 보고서 마감 완료")
                time.sleep(0.5)
                st.rerun()
        else:
            st.info(f"마감 완료: {st.session_state.report_closed}")

    c1, c2, c3 = st.columns(3)

    with c1:
        st.button("실적파일 구글시트 전송", use_container_width=True, disabled=not is_closed)

    ym = get_uploaded_month(st.session_state.user_excel_data) if st.session_state.user_excel_data is not None else ""
    if ym:
        year_month = ym.replace("-", "")
    else:
        year_month = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m")

    download_time = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m%d_%H%M%S")
    file_name = f"LMB월간 활동실적__{year_month}_{download_time}.xlsx"

    with c2:
        if is_closed:
            excel_bytes = sent_uploaded_files_excel_bytes()
            st.download_button(
                "실적파일 엑셀 다운로드",
                data=excel_bytes,
                file_name=file_name,
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
            )
        else:
            st.button("실적파일 엑셀 다운로드", use_container_width=True, disabled=True)

    with c3:
        _is_admin = st.session_state.get("user_role") == "관리자"
        if is_closed or _is_admin:
            try:
                upload_df_for_ppt = None
                if isinstance(st.session_state.get("admin_uploaded_excel"), pd.DataFrame):
                    upload_df_for_ppt = st.session_state.admin_uploaded_excel
                ppt_bytes = build_report_ppt_bytes(report_df, compare_df, curr_month_label, prev_month_label, upload_df_for_ppt)
                st.download_button(
                    "실적보고서 PPT 다운로드",
                    data=ppt_bytes,
                    file_name=f"LMB활동실적보고서_{year_month}_하나지사.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )
            except Exception as e:
                st.button("실적보고서 PPT 다운로드", use_container_width=True, disabled=True)
                st.caption(f"PPT 생성 준비 중 오류: {e}")
        else:
            st.button("실적보고서 PPT 다운로드", use_container_width=True, disabled=True)


def show_admin_performance():
    # 관리자 모드로 사용자 실적 확인 페이지 호출
    show_user_history(is_admin_mode=True)


def show_admin_analysis():
    select_prev_month("auto_prev_df", "adm_prev_month")

    st.markdown("### 직원 전송 실적 내역")
    sent_df = filter_by_staff(apply_admin_prev_diff(sent_results_df()))

    if sent_df.empty:
        st.info("아직 전송된 실적이 없습니다.")
        return

    style_report_logic(sent_df)

    c1, c2, c3 = st.columns([0.8, 0.1, 0.1])

    with c1:
        if st.button("전송 내역 초기화"):
            save_db(SENT_FILE, {})
            save_db(SENT_UPLOADS_FILE, {})
            st.success("초기화 완료")
            time.sleep(0.5)
            st.rerun()

    with c2:
        st.markdown('<div class="action-btn">', unsafe_allow_html=True)
        if st.button("저장", use_container_width=True, type="secondary"):
            admin_saved = {
                "sent_df": sent_df.to_dict(),
                "sent_uploads_db": load_db(SENT_UPLOADS_FILE, {}),
                "adm_prev_month": st.session_state.get("adm_prev_month", "선택안함"),
                "saved_at": (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S")
            }
            saved_db = load_db(SAVED_STATE_FILE, {})
            saved_db["admin_analysis"] = admin_saved
            save_db(SAVED_STATE_FILE, saved_db)
            st.success("저장 완료")
            time.sleep(0.5)
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)

    with c3:
        st.markdown('<div class="action-btn">', unsafe_allow_html=True)
        if st.button("마감", use_container_width=True, type="primary"):
            deadline_time = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M")
            st.session_state.deadline_time = deadline_time
            st.session_state.analysis_result = sent_df.copy()
            saved_db = load_db(SAVED_STATE_FILE, {})
            saved_db["deadline"] = {"time": deadline_time, "by": st.session_state.get("user_name", "")}
            save_db(SAVED_STATE_FILE, saved_db)
            send_kakao_notify(f"[LMB 실적관리] 실적 분석/계산 마감 완료\n마감시각: {deadline_time}\n마감자: {st.session_state.get('user_name', '')}")
            st.success("마감 처리 완료")
            time.sleep(0.5)
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)

    if st.session_state.deadline_time:
        st.info(f"마감 완료: {st.session_state.deadline_time}")


def show_report():
    if st.session_state.analysis_result is None:
        st.info("실적 분석/계산 메뉴에서 마감 후 확인 가능합니다.")
        return

    drop_cols = ["전송시각"] + [c for c in st.session_state.analysis_result.columns if "관리자전월대비" in c]
    report_df = filter_by_staff(st.session_state.analysis_result.drop(columns=drop_cols, errors="ignore"))

    curr_month_label = "당월"
    prev_month_label = "전월"

    ym = get_uploaded_month(st.session_state.user_excel_data) if st.session_state.user_excel_data is not None else ""
    if ym:
        curr_month_label = str(int(ym.split("-")[1])) + "월"

    adm_sel = st.session_state.get("adm_prev_month", "선택안함")
    if adm_sel and adm_sel != "선택안함":
        prev_month_label = str(int(adm_sel.split("-")[1])) + "월"

    st.markdown(f"### {curr_month_label} 실적보고서")
    style_report_logic(report_df)

    prev_df = st.session_state.get("auto_prev_df")

    if prev_df is None or prev_df.empty:
        st.info("실적 분석/계산 메뉴에서 비교할 전월을 선택하면 비교 리포트가 표시됩니다.")
        return

    prev_res, _, _ = process_performance_analysis(prev_df)

    if not isinstance(prev_res, pd.DataFrame) or prev_res.empty:
        st.info("전월 데이터를 분석할 수 없습니다.")
        return

    prev_res = filter_by_staff(prev_res)

    st.markdown("### 비교 리포트")

    compare_cols = [
        "개설건수",
        "개설포인트",
        "연계건수",
        "연계포인트",
        "운영건수 (실제 활동)",
        "운영포인트 (실제 활동)",
        "운영건수 (추가 활동)",
        "운영포인트(추가 활동)",
        "합계포인트",
        "지급포인트",
        "지급예상금액",
    ]

    shared_cols = [c for c in compare_cols if c in report_df.columns and c in prev_res.columns]
    compare_all_rows = []

    names = report_df["담당자"].tolist()
    for start in range(0, len(names), 2):
        person_cols = st.columns(2)
        for person_col, name in zip(person_cols, names[start:start + 2]):
            curr_row = report_df[report_df["담당자"] == name]
            prev_row = prev_res[prev_res["담당자"] == name]

            if curr_row.empty:
                continue

            rank = curr_row.iloc[0].get("직급", "")
            rows = []
            for col in shared_cols:
                c_val = int(float(curr_row.iloc[0][col]))
                p_val = int(float(prev_row.iloc[0][col])) if not prev_row.empty else 0
                rows.append({"항목": col, prev_month_label: p_val, curr_month_label: c_val, "증감": c_val - p_val})
                compare_all_rows.append({
                    "담당자": name,
                    "직급": rank,
                    "항목": col,
                    prev_month_label: p_val,
                    curr_month_label: c_val,
                    "증감": c_val - p_val,
                })

            with person_col:
                st.markdown(f"#### {name} ({rank})")
                style_report_logic(pd.DataFrame(rows), compact=True)

    render_report_action_buttons(report_df, pd.DataFrame(compare_all_rows), curr_month_label, prev_month_label)


def show_staff_admin():
    st.markdown("### 직원 목록")

    # 파일 기준으로 재로드 (_deleted 목록 포함 반영)
    st.session_state.user_db = load_user_db()

    if st.session_state.pop("reset_staff_edit_sel", False):
        st.session_state.staff_edit_sel = "선택안함"

    # ── 직원 직접 추가 ──
    with st.expander("직원 직접 추가", expanded=False):
        with st.form("add_staff_form"):
            ac1, ac2 = st.columns(2)
            with ac1:
                new_uid = st.text_input("아이디", key="add_staff_uid")
                new_name = st.text_input("성명", key="add_staff_name")
                new_email = st.text_input("메일주소", key="add_staff_email")
            with ac2:
                new_pw = st.text_input("비밀번호", type="password", key="add_staff_pw")
                new_access = st.selectbox("로그인 허용 여부", ["허용", "불가"], key="add_staff_access")
                new_role = st.selectbox("메뉴 접근 권한", ["사용자 메뉴", "관리자 메뉴"], key="add_staff_role")
            add_submitted = st.form_submit_button("추가", use_container_width=True, type="primary")

        if add_submitted:
            uid_str = new_uid.strip()
            name_str = new_name.strip()
            pw_str = new_pw.strip()
            if not uid_str or not name_str or not pw_str:
                st.error("아이디, 성명, 비밀번호는 필수 입력입니다.")
            elif uid_str in st.session_state.user_db:
                st.error(f"이미 존재하는 아이디입니다: {uid_str}")
            else:
                st.session_state.user_db[uid_str] = {
                    "pw": pw_str,
                    "name": name_str,
                    "email": new_email.strip(),
                    "access": new_access,
                    "role": "관리자" if new_role == "관리자 메뉴" else "사용자",
                    "dept_type": "사업부",
                    "staff_type": "정규직",
                    "outsource": "아니오",
                    "outsource_period": "해당없음",
                }
                # _deleted 목록 유지
                _cur_file = {}
                if os.path.exists(DB_FILE):
                    try:
                        with open(DB_FILE, "r", encoding="utf-8") as _f:
                            _cur_file = json.load(_f)
                    except Exception:
                        pass
                save_data = dict(st.session_state.user_db)
                if "_deleted" in _cur_file:
                    save_data["_deleted"] = _cur_file["_deleted"]
                save_db(DB_FILE, save_data)
                st.success(f"직원 '{name_str}({uid_str})'을(를) 추가했습니다.")
                st.rerun()

    staff_rows = []
    for uid, info in st.session_state.user_db.items():
        if uid == "1":
            continue
        staff_rows.append({
            "ID": uid,
            "성명": info.get("name", ""),
            "직급": info.get("rank", "직원"),
            "메일주소": info.get("email", ""),
            "부서구분": info.get("dept_type", "사업부"),
            "직원구분": info.get("staff_type", "정규직"),
            "외주여부": info.get("outsource", "아니오"),
            "외주 근무기간": info.get("outsource_period", "해당없음"),
            "로그인 허용 여부": info.get("access", "불가"),
            "메뉴 접근 권한": "관리자 메뉴" if info.get("role") == "관리자" else "사용자 메뉴",
        })

    if not staff_rows:
        st.info("등록된 직원이 없습니다.")
        return

    # 직급 순서로 정렬
    rank_order = {"부서장": 0, "팀장": 1, "과장": 2, "대리": 3, "주임": 4, "직원": 5}
    staff_rows.sort(key=lambda x: rank_order.get(x["직급"], 99))

    # ── 직원 목록 HTML 테이블 표시 (다크모드 호환, 가운데 정렬) ──
    render_plain_html_table(pd.DataFrame(staff_rows), center_align=True)

    st.markdown("---")
    st.markdown("#### 직원 정보 수정")

    uid_options = [f"{r['ID']} — {r['성명']}" for r in staff_rows]
    sel_col, _ = st.columns([0.32, 0.68])
    with sel_col:
        sel = st.selectbox("수정할 직원 선택", ["선택안함"] + uid_options, key="staff_edit_sel")

    if sel == "선택안함":
        return

    sel_uid = sel.split(" — ")[0].strip()
    info = st.session_state.user_db.get(sel_uid, {})

    st.markdown(f"**메일주소:** {info.get('email', '—')}")

    c1, c2, c3, c4, c5, c6, c7 = st.columns([1, 1, 1, 1, 1.15, 1.15, 1.2])
    with c1:
        new_rank = st.selectbox("직급", ["부서장", "팀장", "과장", "대리", "주임", "직원"],
                                index=["부서장", "팀장", "과장", "대리", "주임", "직원"].index(info.get("rank", "직원")),
                                key="edit_rank")
    with c2:
        dept_opts = ["사업부", "C&S"]
        new_dept_type = st.selectbox("부서구분", dept_opts,
                                     index=dept_opts.index(info.get("dept_type", "사업부")) if info.get("dept_type", "사업부") in dept_opts else 0,
                                     key="edit_dept_type")
    with c3:
        new_staff_type = st.selectbox("직원구분", ["정규직", "계약직", "파견직", "외주"],
                                      index=["정규직", "계약직", "파견직", "외주"].index(info.get("staff_type", "정규직")),
                                      key="edit_staff_type")
    with c4:
        new_outsource = st.selectbox("외주여부", ["아니오", "예"],
                                     index=["아니오", "예"].index(info.get("outsource", "아니오")),
                                     key="edit_outsource")
    with c5:
        period_opts = ["해당없음", "1년 미만", "1년 이상", "2년 이상"]
        new_period = st.selectbox("외주 근무기간", period_opts,
                                  index=period_opts.index(info.get("outsource_period", "해당없음")),
                                  key="edit_period")
    with c6:
        new_access = st.selectbox("로그인 허용 여부", ["허용", "불가"],
                                  index=["허용", "불가"].index(info.get("access", "불가")),
                                  key="edit_access")
    with c7:
        new_role = st.selectbox("메뉴 접근 권한", ["사용자 메뉴", "관리자 메뉴"],
                                index=0 if info.get("role") != "관리자" else 1,
                                key="edit_role")

    bc1, bc2, _ = st.columns([0.15, 0.15, 0.7])
    with bc1:
        if st.button("저장", type="primary", use_container_width=True):
            st.session_state.user_db[sel_uid]["rank"] = new_rank
            st.session_state.user_db[sel_uid]["dept_type"] = new_dept_type
            st.session_state.user_db[sel_uid]["staff_type"] = new_staff_type
            st.session_state.user_db[sel_uid]["outsource"] = new_outsource
            st.session_state.user_db[sel_uid]["outsource_period"] = new_period
            st.session_state.user_db[sel_uid]["access"] = new_access
            st.session_state.user_db[sel_uid]["role"] = "관리자" if new_role == "관리자 메뉴" else "사용자"
            # _deleted 목록 유지
            _cur_file = {}
            if os.path.exists(DB_FILE):
                try:
                    with open(DB_FILE, "r", encoding="utf-8") as _f:
                        _cur_file = json.load(_f)
                except Exception:
                    pass
            save_data = dict(st.session_state.user_db)
            if "_deleted" in _cur_file:
                save_data["_deleted"] = _cur_file["_deleted"]
            save_db(DB_FILE, save_data)
            st.session_state.reset_staff_edit_sel = True
            st.success("저장 완료")
            time.sleep(0.5)
            st.rerun()
    with bc2:
        if st.button("삭제", type="secondary", use_container_width=True):
            del st.session_state.user_db[sel_uid]
            # _deleted 목록을 파일에 함께 저장해 기본 계정에서도 복원되지 않도록 함
            _cur_file = {}
            if os.path.exists(DB_FILE):
                try:
                    with open(DB_FILE, "r", encoding="utf-8") as _f:
                        _cur_file = json.load(_f)
                except Exception:
                    pass
            _deleted_set = set(_cur_file.get("_deleted", []))
            _deleted_set.add(sel_uid)
            save_data = dict(st.session_state.user_db)
            save_data["_deleted"] = list(_deleted_set)
            save_db(DB_FILE, save_data, allow_shrink=True)
            st.session_state.reset_staff_edit_sel = True
            st.success(f"{sel} 삭제 완료")
            time.sleep(0.5)
            st.rerun()


KPI_MONTHLY_TARGETS = {
    "개설": 20,
    "연계": 12,
    "운영 순증": 13,
    "유통 활동": 40,
}


def _kpi_number(value):
    if is_blank_value(value):
        return 0
    parsed = pd.to_numeric(str(value).replace(",", "").strip(), errors="coerce")
    return int(parsed) if pd.notna(parsed) else 0


def build_kpi_activity_recommendations(hana_sheet, billing_sheet, user_name=None, use_dart=False):
    if hana_sheet is None or billing_sheet is None or hana_sheet.empty or billing_sheet.empty:
        return pd.DataFrame()

    hana = hana_sheet.copy()
    billing = billing_sheet.copy()
    hana.columns = [str(c).strip() for c in hana.columns]
    billing.columns = [str(c).strip() for c in billing.columns]

    hana_customer_col = find_col(hana, ["고객번호"])
    hana_owner_col = find_col(hana, ["담당자"])
    hana_company_col = find_col(hana, ["고객명", "업체명", "상호"])
    hana_biz_col = find_col(hana, ["사업자번호"])
    hana_build_type_col = find_col(hana, ["구축형"])
    hana_open_status_col = find_col(hana, ["개설상태"])
    hana_link_status_col = find_col(hana, ["연계상태", "ERP연계상태"])
    hana_manage_col = find_col(hana, ["관리구분"])
    hana_open_date_col = find_col(hana, ["개설/이행일", "개설일", "이행일"])

    billing_customer_col = find_col(billing, ["고객번호"])
    billing_company_col = find_col(billing, ["고객명", "업체명", "상호"])
    billing_biz_col = find_col(billing, ["사업자번호", "사업자등록번호"])
    billing_new_col = find_col(billing, ["신규일자", "최초신규일자"])
    billing_login_col = find_col(billing, ["최종로그인일자", "최종로그인", "로그인일자"])
    billing_transfer_col = find_col(billing, ["최종이체일자", "최종이체", "이체일자"])
    billing_end_col = find_col(billing, ["해지일자", "해지일", "해약일"])
    billing_login_count_col = find_col(billing, ["로그인건수", "로그인횟수"])
    billing_menu_col = exact_col(billing, ["메뉴사용"]) or find_col(billing, ["메뉴사용", "메뉴클릭수"])
    billing_service_detail_col = find_col(billing, ["서비스상세", "서비스 구분"])

    if not hana_customer_col or not hana_owner_col or not billing_customer_col:
        return pd.DataFrame()

    hana["_고객번호"] = hana[hana_customer_col].apply(normalize_billing_customer_no)
    hana = hana[hana["_고객번호"].ne("")]
    if user_name:
        hana = hana[hana[hana_owner_col].astype(str).str.strip() == str(user_name).strip()].copy()
    if hana_manage_col and hana_manage_col in hana.columns:
        hana = hana[~hana[hana_manage_col].astype(str).str.strip().str.contains("해지|취소", case=False, na=False)].copy()
    if hana.empty:
        return pd.DataFrame()

    base = pd.DataFrame()
    base["_고객번호"] = hana["_고객번호"]
    base["담당자"] = hana[hana_owner_col].fillna("").astype(str)
    base["고객명"] = hana[hana_company_col].fillna("").astype(str) if hana_company_col else ""
    base["사업자번호"] = hana[hana_biz_col].apply(normalize_biz_no) if hana_biz_col else ""
    base["구축형"] = hana[hana_build_type_col].fillna("").astype(str) if hana_build_type_col else ""
    base["개설상태"] = hana[hana_open_status_col].fillna("").astype(str) if hana_open_status_col else ""
    base["연계상태"] = hana[hana_link_status_col].fillna("").astype(str) if hana_link_status_col else ""
    base["개설/이행일"] = hana[hana_open_date_col].apply(parse_sheet_date) if hana_open_date_col else pd.NaT
    base = base.drop_duplicates("_고객번호", keep="first")

    billing["_고객번호"] = billing[billing_customer_col].apply(normalize_billing_customer_no)
    billing = billing[billing["_고객번호"].ne("")].copy()
    lookup = pd.DataFrame()
    lookup["_고객번호"] = billing["_고객번호"]
    lookup["청구고객명"] = billing[billing_company_col].fillna("").astype(str) if billing_company_col else ""
    lookup["청구사업자번호"] = billing[billing_biz_col].apply(normalize_biz_no) if billing_biz_col else ""
    lookup["신규일자"] = billing[billing_new_col].apply(parse_sheet_date) if billing_new_col else pd.NaT
    lookup["최종로그인일자"] = billing[billing_login_col].apply(parse_sheet_date) if billing_login_col else pd.NaT
    lookup["최종이체일자"] = billing[billing_transfer_col].apply(parse_sheet_date) if billing_transfer_col else pd.NaT
    lookup["해지일자"] = billing[billing_end_col].apply(parse_sheet_date) if billing_end_col else pd.NaT
    lookup["로그인건수"] = billing[billing_login_count_col].apply(_kpi_number) if billing_login_count_col else 0
    lookup["메뉴사용"] = billing[billing_menu_col].apply(_kpi_number) if billing_menu_col else 0
    lookup["서비스상세"] = billing[billing_service_detail_col].fillna("").astype(str) if billing_service_detail_col else ""
    lookup = lookup.drop_duplicates("_고객번호", keep="first")

    merged = base.merge(lookup, on="_고객번호", how="inner")
    merged = merged[merged["해지일자"].isna()].copy()
    if merged.empty:
        return pd.DataFrame()

    today = pd.Timestamp(datetime.utcnow() + timedelta(hours=9)).normalize()

    # DART 공시 기업코드 맵 로드 (use_dart=True 이고 API 키 있을 때만)
    dart_key = _get_dart_api_key() if use_dart else ""
    dart_corp_map = _load_dart_corp_map(dart_key) if dart_key else {}

    rows = []
    for _, row in merged.iterrows():
        score = 0
        areas = []
        reasons = []
        guide = []
        action = "해피콜"

        last_login = row.get("최종로그인일자")
        last_transfer = row.get("최종이체일자")
        new_date = row.get("신규일자")
        login_count = int(row.get("로그인건수", 0) or 0)
        menu_count = int(row.get("메뉴사용", 0) or 0)
        build_type = str(row.get("구축형", ""))
        link_status = str(row.get("연계상태", ""))
        service_detail = str(row.get("서비스상세", ""))

        days_from_login = None
        if pd.isna(last_login):
            score += 55
            areas.append("운영/MAU")
            reasons.append("청구 시트상 로그인 이력 없음")
            guide.append("담당자 변경, 사용 의사, 로그인 장애 여부 확인 후 매뉴얼/방문교육 제안")
        else:
            days_from_login = int((today - pd.Timestamp(last_login).normalize()).days)
            if days_from_login >= 90:
                score += 45
                areas.append("해지방어")
                reasons.append(f"최종 로그인 {days_from_login}일 경과")
                guide.append("미사용 사유를 유형화하고 해지 가능성 또는 잠재 활성화 여부를 비고에 기록")
            elif days_from_login >= 30:
                score += 25
                areas.append("MAU")
                reasons.append(f"최종 로그인 {days_from_login}일 경과")
                guide.append("최근 업무 일정 확인 후 이번 달 1회 사용 목표로 안내")

        if pd.isna(last_transfer):
            if login_count >= 100:
                score += 45
                areas.append("이체 활성화")
                reasons.append(f"로그인 {login_count}회 이상이나 이체 이력 없음")
                guide.append("자금이체/지급거래 흐름을 확인하고 지급거래 교육 또는 RM 연계 제안")
                action = "방문/원격교육"
            elif login_count > 0:
                score += 30
                areas.append("이체 활성화")
                reasons.append("로그인 이력은 있으나 이체 이력 없음")
                guide.append("조회만 사용하는 사유 확인 후 첫 이체까지 동행 지원")
                action = "해피콜+교육"

        if pd.notna(new_date) and (today - pd.Timestamp(new_date).normalize()).days <= 60:
            score += 25
            areas.append("초기정착")
            reasons.append("최근 2개월 내 신규 고객")
            guide.append("초기 2개월 집중관리 대상으로 방문 또는 사용자 교육 일정 확정")

        open_status = str(row.get("개설상태", ""))
        # 구축형이 연계형 + 개설상태가 이행완료인 경우 이미 연계 완료 → 연계전환 추천 불가
        already_linked = ("연계" in build_type) and ("이행완료" in open_status)
        is_not_linked = (not already_linked) and (("연계" not in build_type) or ("완료" not in link_status))
        if is_not_linked and (login_count >= 30 or menu_count >= 50 or "단독" in service_detail):
            score += 25
            areas.append("연계 전환")
            reasons.append("사용 흔적이 있어 ERP연계 전환 제안 가능")
            guide.append("사용 ERP, 고도화 여부, 연계 항목, 접속/DB 정보 사전 확인")
            action = "연계 니즈확인"

        if login_count >= 100 or menu_count >= 100:
            score += 20
            areas.append("유통")
            reasons.append(f"사용량 높음(로그인 {login_count}, 메뉴 {menu_count})")
            guide.append("대시보드/IHB/We-Hub/이음택스 중 재무업무 니즈 확인")

        # DART 공시 기반 유통활동 가산점
        if dart_corp_map:
            company_name = str(row.get("고객명", "") or row.get("청구고객명", "")).strip()
            dart_bonus, dart_reasons = _dart_enrich(dart_key, dart_corp_map, company_name)
            if dart_bonus > 0:
                score += dart_bonus
                if "유통" not in areas:
                    areas.append("유통")
                reasons.extend(dart_reasons)
                guide.append("DART 공시 내용 확인 후 재무 니즈에 맞춰 서비스 제안")

        if score < 35:
            continue

        priority = "상" if score >= 85 else "중" if score >= 60 else "하"
        due = "이번 주" if priority == "상" else "이번 달 2주 내" if priority == "중" else "이번 달 내"
        customer_name = str(row.get("고객명", "")).strip() or str(row.get("청구고객명", "")).strip()
        biz_no = str(row.get("사업자번호", "")).strip() or str(row.get("청구사업자번호", "")).strip()

        rows.append({
            "우선순위": priority,
            "추천점수": score,
            "KPI영역": " / ".join(dict.fromkeys(areas)),
            "담당자": str(row.get("담당자", "")).strip(),
            "고객명": customer_name,
            "사업자번호": biz_no,
            "구축형": build_type,
            "개설상태": row.get("개설상태", ""),
            "연계상태": link_status,
            "최종로그인일자": last_login.strftime("%Y-%m-%d") if pd.notna(last_login) else "없음",
            "최종이체일자": last_transfer.strftime("%Y-%m-%d") if pd.notna(last_transfer) else "없음",
            "로그인건수": login_count,
            "메뉴사용": menu_count,
            "추천활동": action,
            "활동사유": " / ".join(dict.fromkeys(reasons)),
            "활동가이드": " / ".join(dict.fromkeys(guide)),
            "권장기한": due,
        })

    if not rows:
        return pd.DataFrame()
    return pd.DataFrame(rows).sort_values(["추천점수", "로그인건수", "메뉴사용"], ascending=[False, False, False]).reset_index(drop=True)


def render_kpi_activity_recommendations(hana_sheet, billing_sheet, user_name=None, key_prefix="kpi_rec"):
    st.markdown("#### KPI 집중 활동 추천")
    remaining_months = 6
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("개설 잔여 목표", f"{KPI_MONTHLY_TARGETS['개설'] * remaining_months:,}건", "월 20건")
    c2.metric("연계 잔여 목표", f"{KPI_MONTHLY_TARGETS['연계'] * remaining_months:,}건", "월 12건")
    c3.metric("운영 순증 잔여 목표", f"{KPI_MONTHLY_TARGETS['운영 순증'] * remaining_months:,}개", "월 13개")
    c4.metric("유통 활동 잔여 목표", f"{KPI_MONTHLY_TARGETS['유통 활동'] * remaining_months:,}건", "월 40건")

    st.caption("2026년 KPI 계획 기준: 개설 20건/월, 연계 12건/월, 운영 순증 13개/월, 유통 활동 40건/월. 잔여 목표는 남은 6개월 기준입니다.")

    # DART 공시 연동 토글 (DART_API_KEY 설정된 경우에만 표시)
    _dart_available = bool(_get_dart_api_key())
    use_dart = False
    if _dart_available:
        use_dart = st.checkbox(
            "📰 DART 전자공시 기반 유통활동 추천 보정 (공시 있는 고객 점수 상향)",
            key=f"{key_prefix}_dart",
            value=False,
        )
        if use_dart:
            st.caption(
                "체크 시 고객명으로 DART 전자공시 시스템을 조회해 최근 6개월 내 아래 공시가 있는 고객의 "
                "유통활동 추천 점수를 자동으로 높입니다. 공시 내용은 활동사유 컬럼에 `[DART공시]` 형태로 표시됩니다.\n\n"
                "- **유상증자**: 회사가 새 주식을 발행해 외부에서 자금을 조달하는 것 → 현금 유입으로 이체·지급 거래 활성화 기대\n"
                "- **합병**: 두 회사가 하나로 합쳐지는 것 → 계좌 통합·정리 니즈 발생, 접점 확대 기회\n"
                "- **수주**: 기업이 대규모 계약·프로젝트를 따내는 것 → 매출 증가 예상, 자금 흐름 활성화\n"
                "- **신규사업**: 기존 사업 외 새로운 사업에 진출하는 것 → 향후 거래 규모 확대 가능성"
            )
        else:
            st.caption("체크하면 DART 전자공시(유상증자·합병·수주·신규사업 등) 기반으로 유통활동 대상 고객의 추천 점수를 자동 보정합니다.")
    else:
        st.caption("💡 DART 전자공시 연동을 활성화하려면 [설정 → 구글 스프레드시트 연동] 메뉴에서 DART API 키를 입력하세요.")

    rec_df = build_kpi_activity_recommendations(hana_sheet, billing_sheet, user_name, use_dart=use_dart)
    if rec_df.empty:
        st.info("청구 시트와 하나은행 시트 기준으로 추천 가능한 KPI 활동 고객이 없습니다.")
        return

    f1, f2, f3, f4 = st.columns([2, 2, 2, 4])
    with f1:
        area_options = ["전체"] + sorted({area for text in rec_df["KPI영역"].dropna() for area in str(text).split(" / ") if area})
        selected_area = st.selectbox("KPI영역", area_options, key=f"{key_prefix}_area")
    with f2:
        selected_priority = st.selectbox("우선순위", ["전체", "상", "중", "하"], key=f"{key_prefix}_priority")
    with f3:
        owner_options = ["전체"] + sorted(v for v in rec_df["담당자"].dropna().astype(str).str.strip().unique().tolist() if v)
        selected_owner = st.selectbox("담당자", owner_options, key=f"{key_prefix}_owner")

    filtered = rec_df.copy()
    if selected_area != "전체":
        filtered = filtered[filtered["KPI영역"].astype(str).str.contains(selected_area, regex=False, na=False)]
    if selected_priority != "전체":
        filtered = filtered[filtered["우선순위"] == selected_priority]
    if selected_owner != "전체":
        filtered = filtered[filtered["담당자"] == selected_owner]
    filtered = filtered.head(100).reset_index(drop=True)
    filtered.insert(0, "순번", range(1, len(filtered) + 1))

    st.metric("추천 고객", f"{len(filtered):,}건")
    render_plain_html_table(filtered, max_rows=100, center_align=False)
    today_str = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m%d")
    st.download_button(
        "KPI 추천 고객 다운로드",
        data=filtered.to_csv(index=False, encoding="utf-8-sig").encode("utf-8-sig"),
        file_name=f"KPI추천고객_{selected_area}_{selected_priority}_{today_str}.csv",
        mime="text/csv",
        use_container_width=True,
    )


def show_target_customers():
    st.markdown("### 이번달 활동 대상고객 추천")

    # 하나은행 구글시트 로드 (개설 미완료용)
    hana_sheet = None
    try:
        hana_raw = read_google_csv(st.session_state.get("url_hana", DEFAULT_URL_HANA), header=2)
        hana_sheet = hana_raw
    except Exception as e:
        st.error(f"하나은행 구글시트를 불러오지 못했습니다: {e}")

    # 하나은행 청구 시트 로드 (운영 관리 대상용)
    hana_billing_sheet = None
    try:
        billing_raw = read_google_csv(st.session_state.get("url_hana_billing", DEFAULT_URL_HANA_BILLING))
        billing_raw = billing_raw.dropna(how="all").reset_index(drop=True)
        hana_billing_sheet = billing_raw
    except Exception as e:
        st.error(f"하나은행 청구 시트를 불러오지 못했습니다: {e}")

    # 본사 구글시트 로드 (ERP연계 미완료용)
    if st.session_state.get("cloud_sheet_df") is None:
        try:
            load_csv_to_state("url_sync", "cloud_sheet_df")
        except Exception as e:
            st.error(f"본사 구글시트를 불러오지 못했습니다: {e}")

    cloud = st.session_state.get("cloud_sheet_df")

    # 활동 이력 시트 로드 (최근 활동 정보용)
    if st.session_state.get("analysis_lookup_df") is None:
        try:
            load_csv_to_state("url_analysis", "analysis_lookup_df")
        except Exception:
            pass

    hana_df = st.session_state.get("analysis_lookup_df")

    user_name = st.session_state.user_name

    def _only_current_user_rows(result_df):
        if result_df is None or result_df.empty or "담당자" not in result_df.columns:
            return result_df
        return result_df[result_df["담당자"].astype(str).str.strip() == str(user_name).strip()].reset_index(drop=True)

    render_kpi_activity_recommendations(hana_sheet, hana_billing_sheet, user_name, key_prefix="target_kpi_rec")
    st.divider()

    cloud = clean_header_logic(cloud.copy())

    # 본사 시트 컬럼 탐색
    owner_col  = find_col(cloud, ["담당자", "등록자", "성명"])
    comp_col   = find_col(cloud, ["업체명", "고객명", "상호"])
    biz_col    = find_col(cloud, ["사업자번호"])
    open_col   = find_col(cloud, ["개설완료일자", "개설일"])
    erp_col    = find_col(cloud, ["ERP연계일자", "연계일자"])
    div_col    = find_col(cloud, ["신규/이행구분", "이행구분", "신규이행"])
    add_col    = find_col(cloud, ["이행추가연계"])
    status_col = find_col(cloud, ["상태항목", "상태"])
    open_cancel_col = find_col(cloud, ["개설취소"])
    erp_cancel_col = find_col(cloud, ["ERP연계취소", "연계취소"])
    manage_div_col = find_col(cloud, ["관리구분", "관리 구분"])
    open_status_col = find_col(cloud, ["개설상태"])
    erp_status_col = find_col(cloud, ["연계상태", "ERP연계상태"])

    if not owner_col:
        st.warning("구글시트에 담당자 컬럼을 찾을 수 없습니다.")
        return

    # 본인 담당 고객만 필터
    df = cloud[cloud[owner_col].astype(str).str.strip() == str(user_name).strip()].copy()

    # 관리구분이 "해지" 또는 "취소"를 포함하는 고객 제외
    if manage_div_col and manage_div_col in df.columns:
        manage_values = df[manage_div_col].astype(str).str.strip()
        # "해지" 또는 "취소"가 포함된 경우 모두 제외
        is_terminated = manage_values.str.contains("해지|취소", case=False, na=False)
        df = df[~is_terminated].copy()

    if df.empty:
        st.info(f"구글시트에 {user_name}님 담당 고객이 없습니다.")
        return

    # 날짜 정규화
    for col in [open_col, erp_col]:
        if col and col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce")

    # 하나은행 시트에서 활동 이력 가져오기
    def get_recent_activity(biz_no):
        if hana_df is None or hana_df.empty or not biz_no:
            return {"last_date": "", "last_detail": "", "activity_count": 0}

        hana_clean = clean_header_logic(hana_df.copy())
        hana_biz_col = find_col(hana_clean, ["사업자번호"])
        hana_date_col = find_col(hana_clean, ["활동일", "일자"])
        hana_detail_col = find_col(hana_clean, ["활동상세", "활동내용"])

        if not hana_biz_col or hana_biz_col not in hana_clean.columns:
            return {"last_date": "", "last_detail": "", "activity_count": 0}

        # 사업자번호 정규화하여 매칭
        biz_normalized = normalize_biz(pd.Series([biz_no])).iloc[0]
        hana_clean[hana_biz_col] = normalize_biz(hana_clean[hana_biz_col])

        matched = hana_clean[hana_clean[hana_biz_col] == biz_normalized].copy()

        if matched.empty:
            return {"last_date": "", "last_detail": "", "activity_count": 0}

        # 날짜 정렬
        if hana_date_col and hana_date_col in matched.columns:
            matched[hana_date_col] = pd.to_datetime(matched[hana_date_col], errors="coerce")
            matched = matched.sort_values(hana_date_col, ascending=False)
            last_date = matched.iloc[0][hana_date_col]
            last_date_str = last_date.strftime("%Y-%m-%d") if pd.notna(last_date) else ""
        else:
            last_date_str = ""

        last_detail = matched.iloc[0].get(hana_detail_col, "") if hana_detail_col else ""
        activity_count = len(matched)

        return {
            "last_date": last_date_str,
            "last_detail": str(last_detail),
            "activity_count": activity_count
        }

    def _empty(col):
        if not col or col not in df.columns:
            return pd.Series(True, index=df.index)
        return df[col].isna()

    def _is_ihang_no_add(row):
        if div_col and div_col in df.columns and add_col and add_col in df.columns:
            return (str(row.get(div_col, "")).strip() == "이행" and
                    str(row.get(add_col, "")).strip() == "")
        return False

    # 표시 컬럼 선택 및 활동 이력 추가
    base_show_cols = [c for c in [comp_col, biz_col, div_col, open_col, erp_col, add_col] if c and c in df.columns]

    def fmt_df(sub):
        if sub.empty:
            return sub
        d = sub[base_show_cols].copy().reset_index(drop=True)

        # 날짜 포맷팅
        for col in [open_col, erp_col]:
            if col and col in d.columns:
                d[col] = pd.to_datetime(d[col], errors="coerce").dt.strftime("%Y-%m-%d").fillna("")

        # 활동 이력 추가
        if biz_col and biz_col in d.columns:
            activities = []
            for _, row in d.iterrows():
                biz_no = row.get(biz_col, "")
                activity = get_recent_activity(biz_no)
                activities.append(activity)

            d["최근활동일"] = [a["last_date"] for a in activities]
            d["최근활동내용"] = [a["last_detail"] for a in activities]
            d["총활동건수"] = [a["activity_count"] for a in activities]

        return d

    # ── 개설 미완료 (하나은행 시트 기반) ──────────────────────
    st.markdown("#### 🟠 개설 미완료 고객")

    if hana_sheet is not None and not hana_sheet.empty:
        hana_df_open = hana_sheet.copy()

        # 컬럼명 확인
        hana_owner_col = "담당자"
        hana_guchuk_col = "구축구분"
        hana_manage_col = "관리구분"
        hana_open_status_col = "개설상태"
        hana_open_date_col = "개설/이행일"
        hana_comp_col = find_col(hana_df_open, ["업체명", "고객명", "상호"])
        hana_biz_col = find_col(hana_df_open, ["사업자번호"])

        # 필터링 조건 적용
        # 1. 본인 담당자
        if hana_owner_col in hana_df_open.columns:
            hana_df_open = hana_df_open[hana_df_open[hana_owner_col].astype(str).str.strip() == str(user_name).strip()]

        # 2. 구축구분 = 신규
        if hana_guchuk_col in hana_df_open.columns:
            hana_df_open = hana_df_open[hana_df_open[hana_guchuk_col].astype(str).str.strip() == "신규"]

        # 3. 관리구분 = 정상
        if hana_manage_col in hana_df_open.columns:
            hana_df_open = hana_df_open[hana_df_open[hana_manage_col].astype(str).str.strip() == "정상"]

        # 4. 개설상태 = 개설대기
        if hana_open_status_col in hana_df_open.columns:
            hana_df_open = hana_df_open[hana_df_open[hana_open_status_col].astype(str).str.strip() == "개설대기"]

        # 5. 개설/이행일 = 공백
        if hana_open_date_col in hana_df_open.columns:
            hana_df_open[hana_open_date_col] = pd.to_datetime(hana_df_open[hana_open_date_col], errors="coerce")
            hana_df_open = hana_df_open[hana_df_open[hana_open_date_col].isna()]

        # 결과 표시
        if not hana_df_open.empty:
            # 표시할 컬럼 선택
            display_cols = [c for c in [hana_comp_col, hana_biz_col, hana_guchuk_col, hana_manage_col, hana_open_status_col, hana_open_date_col] if c and c in hana_df_open.columns]

            # 활동 이력 추가
            result_df = hana_df_open[display_cols].copy().reset_index(drop=True)

            if hana_biz_col and hana_biz_col in result_df.columns:
                activities = []
                for _, row in result_df.iterrows():
                    biz_no = row.get(hana_biz_col, "")
                    activity = get_recent_activity(biz_no)
                    activities.append(activity)

                result_df["최근활동일"] = [a["last_date"] for a in activities]
                result_df["최근활동내용"] = [a["last_detail"] for a in activities]
                result_df["총활동건수"] = [a["activity_count"] for a in activities]

            st.caption(f"{len(result_df)}건")
            st.dataframe(result_df, use_container_width=True, hide_index=True)
        else:
            st.success("개설 미완료 고객 없음")
    else:
        st.warning("하나은행 시트 데이터를 불러올 수 없습니다.")

    st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)

    # ── ERP연계 미완료 (하나은행 시트 기반) ───────────────────
    st.markdown("#### 🔵 ERP연계 미완료 고객")

    if hana_sheet is not None and not hana_sheet.empty:
        hana_df_erp = hana_sheet.copy()

        # 컬럼명 확인
        hana_owner_col = "담당자"
        hana_manage_col = "관리구분"
        hana_guchuk_type_col = find_col(hana_df_erp, ["구축형", "구축 형"])
        hana_erp_status_col = find_col(hana_df_erp, ["연계상태", "ERP연계상태"])
        hana_comp_col = find_col(hana_df_erp, ["업체명", "고객명", "상호"])
        hana_biz_col = find_col(hana_df_erp, ["사업자번호"])

        # 필터링 조건 적용
        # 1. 본인 담당자
        if hana_owner_col in hana_df_erp.columns:
            hana_df_erp = hana_df_erp[hana_df_erp[hana_owner_col].astype(str).str.strip() == str(user_name).strip()]

        # 2. 관리구분 = 정상
        if hana_manage_col in hana_df_erp.columns:
            hana_df_erp = hana_df_erp[hana_df_erp[hana_manage_col].astype(str).str.strip() == "정상"]

        # 3. 구축형 = 연계형
        if hana_guchuk_type_col and hana_guchuk_type_col in hana_df_erp.columns:
            hana_df_erp = hana_df_erp[hana_df_erp[hana_guchuk_type_col].astype(str).str.strip() == "연계형"]

        # 4. 연계상태 = ERP연계대기 또는 ERP연계진행
        if hana_erp_status_col and hana_erp_status_col in hana_df_erp.columns:
            hana_df_erp = hana_df_erp[hana_df_erp[hana_erp_status_col].astype(str).str.strip().isin(["ERP연계대기", "ERP연계진행"])]

        # 결과 표시
        if not hana_df_erp.empty:
            # 표시할 컬럼 선택
            display_cols = [c for c in [hana_comp_col, hana_biz_col, hana_guchuk_type_col, hana_manage_col, hana_erp_status_col] if c and c in hana_df_erp.columns]

            # 활동 이력 추가
            result_df = hana_df_erp[display_cols].copy().reset_index(drop=True)

            if hana_biz_col and hana_biz_col in result_df.columns:
                activities = []
                for _, row in result_df.iterrows():
                    biz_no = row.get(hana_biz_col, "")
                    activity = get_recent_activity(biz_no)
                    activities.append(activity)

                result_df["최근활동일"] = [a["last_date"] for a in activities]
                result_df["최근활동내용"] = [a["last_detail"] for a in activities]
                result_df["총활동건수"] = [a["activity_count"] for a in activities]

            st.caption(f"{len(result_df)}건")
            st.dataframe(result_df, use_container_width=True, hide_index=True)
        else:
            st.success("ERP연계 미완료 고객 없음")
    else:
        st.warning("하나은행 시트 데이터를 불러올 수 없습니다.")

    # ── 미로그인 고객 현황 ──────────────────────────────────────────
    _tg_force = st.session_state.pop("_nl_rf_tg_nl_refresh", False)
    if _tg_force:
        try:
            hana_sheet = read_google_csv(
                st.session_state.get("url_hana", DEFAULT_URL_HANA), header=2, force_refresh=True
            )
            billing_raw = read_google_csv(
                st.session_state.get("url_hana_billing", DEFAULT_URL_HANA_BILLING), force_refresh=True
            )
            hana_billing_sheet = billing_raw.dropna(how="all").reset_index(drop=True)
        except Exception as e:
            st.error(f"새로고침 실패: {e}")

    if hana_sheet is not None and not hana_sheet.empty and hana_billing_sheet is not None and not hana_billing_sheet.empty:
        _hana = hana_sheet.dropna(how="all").reset_index(drop=True)
        _billing = hana_billing_sheet

        st.divider()
        st.markdown("#### 개설/이행일 이후 미로그인 고객")
        st.caption("개설/이행일 기준으로 청구시트 최종로그인일자가 없거나 개설/이행일 이전인 고객")
        _render_no_login_section(
            _only_current_user_rows(build_no_login_after_open(_hana, _billing)),
            year_key="target_open_year", owner_key="target_open_owner",
            label="미로그인 고객", download_prefix="미로그인고객_개설이행일",
            refresh_key="tg_nl_refresh", exclude_key="tg_nl_open_exclude",
        )

        st.divider()
        st.markdown("#### 연계청구일자 이후 미로그인 고객")
        st.caption("연계청구일자 기준으로 청구시트 최종로그인일자가 없거나 연계청구일자 이전인 고객")
        _render_no_login_section(
            _only_current_user_rows(build_no_login_after_link_billing(_hana, _billing)),
            year_key="target_link_year", owner_key="target_link_owner",
            label="연계청구 미로그인 고객", download_prefix="미로그인고객_연계청구일자",
            exclude_key="tg_nl_link_exclude",
        )

        st.divider()
        st.markdown("#### 로그인 100회 이상 · 미이체 고객")
        st.caption("청구시트 기준 로그인 100회 이상이지만 이체 이력이 없는 고객 (해지 제외)")
        _render_high_login_no_transfer(
            _only_current_user_rows(build_high_login_no_transfer(_hana, _billing)),
            owner_key="tg_high_login_owner",
            download_prefix="미이체고객",
        )
    else:
        st.divider()
        st.warning("하나은행 구글시트 또는 청구시트 데이터를 불러올 수 없습니다.")


VISIT_HISTORY_FILE = "visit_history.json"


def show_visit_history():
    st.markdown("### 방문이력 작성")
    user_name = st.session_state.get("user_name", "")
    db = load_db(VISIT_HISTORY_FILE, {})
    my_entries = db.get(user_name, [])

    # 고객명 자동검색용 목록 — 하나은행 구글시트에서 추출
    customer_options = []
    hana_df = st.session_state.get("hana_sheet_df")
    if hana_df is not None and not hana_df.empty:
        comp_col = find_col(hana_df, ["고객명", "업체명", "상호"])
        owner_col = find_col(hana_df, ["담당자"])
        if comp_col:
            filtered_hana = hana_df
            if owner_col:
                filtered_hana = hana_df[
                    hana_df[owner_col].astype(str).str.strip() == str(user_name).strip()
                ]
            customer_options = sorted(
                filtered_hana[comp_col].dropna().astype(str).str.strip()
                .replace("", pd.NA).dropna().unique().tolist()
            )

    # ── 등록 폼 ──
    with st.expander("방문이력 등록", expanded=True):
        # 고객명 자동검색 (폼 바깥 — 실시간 필터링)
        if customer_options:
            v_customer_sel = st.selectbox(
                "고객명 검색",
                ["직접 입력"] + customer_options,
                key="visit_customer_sel",
            )
        else:
            v_customer_sel = "직접 입력"

        with st.form("visit_form"):
            vc1, vc2 = st.columns(2)
            with vc1:
                v_date = st.date_input("방문일자")
            with vc2:
                v_purpose = st.selectbox("방문목적", ["개설", "운영", "연계"])
            # 고객명: 자동검색 선택값 또는 직접 입력
            if v_customer_sel and v_customer_sel != "직접 입력":
                v_customer = st.text_input("고객명", value=v_customer_sel)
            else:
                v_customer = st.text_input("고객명", placeholder="고객명을 입력하세요")
            v_content = st.text_area("내용", height=120)
            v_submitted = st.form_submit_button("등록", use_container_width=True, type="primary")

        if v_submitted:
            if not v_customer.strip():
                st.error("고객명을 입력해주세요.")
            else:
                now_text = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S")
                entry = {
                    "id": hashlib.md5(f"{user_name}-{now_text}-{v_customer}".encode()).hexdigest(),
                    "방문일자": str(v_date),
                    "고객명": v_customer.strip(),
                    "방문목적": v_purpose,
                    "내용": v_content.strip(),
                    "담당자": user_name,
                    "작성시각": now_text,
                }
                db.setdefault(user_name, []).append(entry)
                save_db(VISIT_HISTORY_FILE, db)
                st.success("방문이력을 등록했습니다.")
                st.rerun()

    # ── 이력 목록 ──
    st.markdown("#### 내 방문이력")
    if not my_entries:
        st.info("등록된 방문이력이 없습니다.")
        return

    sorted_entries = sorted(my_entries, key=lambda x: x.get("방문일자", ""), reverse=True)
    show_cols = ["방문일자", "고객명", "방문목적", "내용", "작성시각"]
    display_df = pd.DataFrame(sorted_entries)
    display_df = display_df[[c for c in show_cols if c in display_df.columns]].copy()
    display_df.insert(0, "선택", False)

    edited_df = st.data_editor(
        display_df,
        use_container_width=True,
        hide_index=True,
        column_config={"선택": st.column_config.CheckboxColumn("선택", default=False, width="small")},
        disabled=[c for c in display_df.columns if c != "선택"],
        key="visit_history_editor",
    )

    selected_ids = [sorted_entries[i].get("id") for i, v in enumerate(edited_df["선택"]) if v]
    _, del_col, _ = st.columns([15, 15, 70])
    with del_col:
        if st.button("삭제", use_container_width=True, disabled=not selected_ids, key="visit_del_btn"):
            db[user_name] = [r for r in my_entries if r.get("id") not in selected_ids]
            save_db(VISIT_HISTORY_FILE, db)
            st.success(f"{len(selected_ids)}건 삭제했습니다.")
            st.rerun()

    today_str = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m%d")
    dl_df = pd.DataFrame(sorted_entries)[[c for c in show_cols if c in pd.DataFrame(sorted_entries).columns]]
    st.download_button(
        "방문이력 다운로드",
        data=dl_df.to_csv(index=False, encoding="utf-8-sig").encode("utf-8-sig"),
        file_name=f"방문이력_{user_name}_{today_str}.csv",
        mime="text/csv",
        use_container_width=True,
    )


def show_weekly_report_user():
    opening_categories = ["개설대기", "개설진행", "개설완료"]
    link_categories = ["연계대기", "연계진행", "연계완료"]
    report_categories = opening_categories + link_categories
    categories = ["전체"] + report_categories + ["운영부문"]
    user_name = st.session_state.get("user_name", "")

    def load_hana_sheet_for_weekly(force_refresh=False):
        try:
            url = st.session_state.get("url_hana", DEFAULT_URL_HANA)
            df = read_google_csv(url, header=2, force_refresh=force_refresh)
            st.session_state.hana_sheet_df = df
            st.session_state.weekly_hana_loaded_at = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S")
        except Exception as e:
            df = st.session_state.get("hana_sheet_df")
            if df is None or df.empty:
                st.error(f"하나은행 구글 시트를 불러오지 못했습니다: {e}")
                return pd.DataFrame()
            st.warning("하나은행 구글 시트 최신 데이터를 불러오지 못해 세션에 저장된 데이터를 사용합니다.")
        df = df.copy()
        df.columns = [str(c).strip() for c in df.columns]
        return df.dropna(how="all").reset_index(drop=True)

    def weekly_category_mask(df, category, start_date=None, end_date=None):
        open_status_col = find_col(df, ["개설상태"])
        open_date_col = find_col(df, ["개설/이행일", "개설일", "이행일"])
        link_status_col = find_col(df, ["연계상태"])
        link_date_col = find_col(df, ["연계일자", "연계일"])
        build_type_col = find_col(df, ["구축형"])
        build_category_col = find_col(df, ["구축구분"])
        end_col = find_col(df, ["해지일자"])

        open_status = df[open_status_col].astype(str).str.strip() if open_status_col else pd.Series("", index=df.index)
        link_status = df[link_status_col].astype(str).str.strip() if link_status_col else pd.Series("", index=df.index)
        build_type = df[build_type_col].astype(str).str.strip() if build_type_col else pd.Series("", index=df.index)
        build_category = df[build_category_col].astype(str).str.strip() if build_category_col else pd.Series("", index=df.index)
        open_dates = pd.to_datetime(df[open_date_col].map(parse_sheet_date), errors="coerce") if open_date_col else pd.Series(pd.NaT, index=df.index)
        link_dates = pd.to_datetime(df[link_date_col].map(parse_sheet_date), errors="coerce") if link_date_col else pd.Series(pd.NaT, index=df.index)
        open_date = open_dates.notna()
        link_date = link_dates.notna()
        active = df[end_col].map(is_blank_value) if end_col else pd.Series(True, index=df.index)

        open_done = active & (open_status.str.contains("완료|이행완료", na=False) | open_date)
        link_done = active & (link_status.str.contains("완료", na=False) | link_date)
        open_progress = active & ~open_done & open_status.str.contains("진행|구축중|처리중", na=False)
        link_progress = active & ~link_done & link_status.str.contains("진행|연계중|처리중", na=False)
        link_target = active & (build_type.str.contains("연계|이행", na=False) | build_category.str.contains("연계|이행", na=False) | link_status.ne(""))

        if category == "개설완료":
            if start_date is not None and end_date is not None:
                return active & open_dates.between(pd.Timestamp(start_date), pd.Timestamp(end_date), inclusive="both")
            return active & open_date
        if category == "개설진행":
            return open_progress
        if category == "개설대기":
            return active & open_status.eq("개설대기")
        if category == "연계완료":
            link_complete = link_status.str.replace(r"\s+", "", regex=True).eq("ERP연계완료")
            if start_date is not None and end_date is not None:
                return active & link_complete & link_dates.between(pd.Timestamp(start_date), pd.Timestamp(end_date), inclusive="both")
            return active & link_complete & link_date
        if category == "연계진행":
            return link_progress
        if category == "연계대기":
            return active & link_status.str.contains("ERP연계대기", na=False)
        if category == "운영부문":
            return active & open_done & (~link_target | link_done)
        return pd.Series(True, index=df.index)

    def display_customer_df(df):
        wanted_keys = ["고객번호", "고객명", "구축구분", "구축형", "신규접수일", "추가연계접수일", "개설상태", "개설/이행일", "연계상태", "연계일자", "구축예정일", "담당자"]
        cols = []
        for key in wanted_keys:
            col = find_col(df, [key])
            if col and col in df.columns and col not in cols:
                cols.append(col)
        if not cols:
            cols = list(df.columns[:8])
        show = df[cols].copy()
        for col in show.columns:
            if any(k in str(col) for k in ["일자", "일", "예정"]):
                parsed = pd.to_datetime(show[col].map(parse_sheet_date), errors="coerce")
                show[col] = show[col].where(parsed.isna(), parsed.dt.strftime("%Y-%m-%d"))
        return show.reset_index(drop=True)

    def format_weekly_date(value):
        parsed = parse_sheet_date(value)
        if pd.isna(parsed):
            return "" if is_blank_value(value) else str(value).strip()
        return pd.Timestamp(parsed).strftime("%Y-%m-%d")

    def date_input_value(value, fallback):
        parsed = parse_sheet_date(value)
        if pd.isna(parsed):
            return fallback
        return pd.Timestamp(parsed).date()

    def render_weekly_detail_header(title, key, level=4):
        title_col, refresh_col = st.columns([0.82, 0.18])
        with title_col:
            st.markdown(f"{'#' * level} {title}")
        with refresh_col:
            if st.button("새로고침", key=f"weekly_refresh_{key}", use_container_width=True):
                st.session_state.weekly_force_hana_refresh = True
                st.rerun()

    title_col, refresh_col = st.columns([0.82, 0.18])
    with title_col:
        st.markdown("### 주간보고 작성")
    force_refresh = bool(st.session_state.pop("weekly_force_hana_refresh", False))
    with refresh_col:
        if st.button("새로고침", key="weekly_user_top_refresh", use_container_width=True):
            force_refresh = True
    st.caption("하나은행 구글 시트의 본인 담당 고객을 기준으로 카테고리별 현황을 선택해 주간보고 이력을 저장합니다.")

    hana = load_hana_sheet_for_weekly(force_refresh=force_refresh)
    if hana.empty:
        return

    staff_col = find_col(hana, ["담당자"])
    company_col = find_col(hana, ["고객명", "업체명", "상호"])
    branch_col = find_col(hana, ["신청점", "담당 부서", "담당부서"])
    service_col = find_col(hana, ["구축구분", "서비스"])
    build_type_col = find_col(hana, ["구축형"])
    customer_no_col = find_col(hana, ["고객번호"])
    open_status_col = find_col(hana, ["개설상태"])
    link_status_col = find_col(hana, ["연계상태"])
    link_date_col = find_col(hana, ["연계일자", "연계일"])

    if not staff_col or not company_col:
        st.error("하나은행 시트에서 담당자 또는 고객명 컬럼을 찾을 수 없습니다.")
        return

    mine = hana[hana[staff_col].astype(str).str.strip() == str(user_name).strip()].copy()
    if mine.empty:
        st.info(f"하나은행 시트에 {user_name}님 담당 고객이 없습니다.")
        return

    today = datetime.utcnow() + timedelta(hours=9)
    week_start_default = (today - timedelta(days=today.weekday() + 7)).date()
    week_end_default = week_start_default + timedelta(days=4)

    col_a, col_b, col_c = st.columns([1.1, 1.1, 1.4])
    with col_a:
        week_start = st.date_input("보고 시작일", value=week_start_default, key="weekly_report_start_v2")
    with col_b:
        week_end = st.date_input("보고 종료일", value=week_end_default, key="weekly_report_end_v2")
    with col_c:
        category = st.selectbox("카테고리", categories, key="weekly_report_category")
    st.caption(f"선택 기간: {week_start} ~ {week_end}")
    if st.session_state.get("weekly_hana_loaded_at"):
        st.caption(f"하나은행 구글 시트 새로고침: {st.session_state.weekly_hana_loaded_at}")

    if category == "전체":
        categorized_map = {
            cat: mine[weekly_category_mask(mine, cat, week_start, week_end)].copy()
            for cat in report_categories
        }
        st.markdown("#### 전체 현황")
        st.caption("개설 부문과 ERP연계 부문을 PPT 주간보고 구성에 맞춰 표시합니다.")
        st.markdown("##### 개설 부문")
        for cat in opening_categories:
            render_weekly_detail_header(cat, f"all_{cat}", level=5)
            st.caption(f"{len(categorized_map[cat])}건 · 하나은행 구글 시트 기준")
            render_plain_html_table(display_customer_df(categorized_map[cat]), max_rows=200)
        st.markdown("##### ERP연계 부문")
        for cat in link_categories:
            render_weekly_detail_header(cat, f"all_{cat}", level=5)
            st.caption(f"{len(categorized_map[cat])}건 · 하나은행 구글 시트 기준")
            render_plain_html_table(display_customer_df(categorized_map[cat]), max_rows=200)
    else:
        categorized = mine[weekly_category_mask(mine, category, week_start, week_end)].copy()
        render_weekly_detail_header(f"{category} 현황", f"single_{category}", level=4)
        st.caption(f"{len(categorized)}건 · 하나은행 구글 시트 기준")
        render_plain_html_table(display_customer_df(categorized), max_rows=200)

        if categorized.empty and category != "운영부문":
            st.info("선택한 카테고리에 작성할 고객이 없습니다.")
            if category == "연계완료" and link_status_col and link_date_col:
                status_clean = mine[link_status_col].astype(str).str.strip().str.replace(r"\s+", "", regex=True)
                candidates = mine[status_clean.eq("ERP연계완료")].copy()
                if not candidates.empty:
                    candidates["_연계일자변환"] = pd.to_datetime(candidates[link_date_col].map(parse_sheet_date), errors="coerce")
                    candidates["_선택기간포함"] = candidates["_연계일자변환"].between(pd.Timestamp(week_start), pd.Timestamp(week_end), inclusive="both")
                    with st.expander("ERP연계완료 후보 확인", expanded=False):
                        cols = [company_col, link_status_col, link_date_col, "_연계일자변환", "_선택기간포함"]
                        cols = [c for c in cols if c and c in candidates.columns]
                        show_candidates = candidates[cols].copy()
                        if "_연계일자변환" in show_candidates.columns:
                            show_candidates["_연계일자변환"] = show_candidates["_연계일자변환"].dt.strftime("%Y-%m-%d")
                        st.dataframe(show_candidates.reset_index(drop=True), use_container_width=True, hide_index=True)

    st.markdown("---")
    st.markdown("#### 주간보고 이력 입력")
    if category == "전체":
        input_part = st.radio("입력 부문", ["개설 부문", "ERP연계 부문"], horizontal=True, key="weekly_input_part")
        input_category_options = opening_categories if input_part == "개설 부문" else link_categories
        input_category = st.selectbox("세부 카테고리", input_category_options, key="weekly_input_category")
    else:
        input_category = category
        input_part = "ERP연계 부문" if input_category in link_categories else "개설 부문"
        st.caption(f"{input_part} · {input_category}")

    if st.session_state.get("_weekly_prev_input_category") != input_category:
        st.session_state.pop("weekly_customer_combo", None)
        for _cat in ["개설대기", "개설진행", "연계대기", "연계진행"]:
            st.session_state.pop(f"weekly_plan_date_{_cat}", None)
        st.session_state["_weekly_prev_input_category"] = input_category

    categorized = mine[weekly_category_mask(mine, input_category, week_start, week_end)].copy()

    tmp = categorized.copy()
    customer_options = []
    if not tmp.empty:
        tmp["_company_label"] = tmp.apply(
            lambda r: f"{str(r.get(company_col, '')).strip()} | {str(r.get(customer_no_col, '')).strip() if customer_no_col else ''}".strip(" |"),
            axis=1,
        )
        customer_options = [v for v in tmp["_company_label"].tolist() if v]
    label_to_row = {str(row["_company_label"]): row for _, row in tmp.iterrows()} if not tmp.empty else {}
    if not customer_options:
        customer_options = ["선택 가능한 고객 없음"]

    if st.session_state.get("weekly_customer_combo") not in customer_options:
        st.session_state.pop("weekly_customer_combo", None)
    selected_company_label = st.selectbox("고객명", customer_options, key="weekly_customer_combo")
    selected_row = label_to_row.get(selected_company_label)
    default_company = "" if selected_row is None else str(selected_row.get(company_col, "")).strip()
    default_customer_no = "" if selected_row is None or not customer_no_col else str(selected_row.get(customer_no_col, "")).strip()
    default_service = "" if selected_row is None or not service_col else str(selected_row.get(service_col, "")).strip()
    default_build = "" if selected_row is None or not build_type_col else str(selected_row.get(build_type_col, "")).strip()
    default_branch = "" if selected_row is None or not branch_col else str(selected_row.get(branch_col, "")).strip()
    open_receipt_col = find_col(mine, ["신규접수일"])
    link_receipt_col = find_col(mine, ["추가연계접수일"])
    schedule_col = find_col(mine, ["구축예정일"])
    receipt_col = link_receipt_col if input_category in link_categories else open_receipt_col
    default_receipt_date = "" if selected_row is None or not receipt_col else format_weekly_date(selected_row.get(receipt_col, ""))
    default_plan_date = "" if selected_row is None or not schedule_col else format_weekly_date(selected_row.get(schedule_col, ""))
    plan_date_label = "구축/피드백 예정일"
    if input_category == "개설완료":
        plan_date_label = "개설완료일자"
        open_done_col = find_col(mine, ["개설/이행일", "개설일", "이행일"])
        default_plan_date = "" if selected_row is None or not open_done_col else format_weekly_date(selected_row.get(open_done_col, ""))
    elif input_category == "연계완료":
        plan_date_label = "연계완료일자"
        link_done_col = find_col(mine, ["연계일자", "연계일"])
        default_plan_date = "" if selected_row is None or not link_done_col else format_weekly_date(selected_row.get(link_done_col, ""))
    status_default = ""
    if selected_row is not None:
        if input_category in link_categories and link_status_col:
            status_default = str(selected_row.get(link_status_col, "")).strip()
        elif open_status_col:
            status_default = str(selected_row.get(open_status_col, "")).strip()

    with st.form("weekly_report_form", border=True):
        st.markdown(f"##### {input_part} - {input_category}")
        company = default_company
        f1, f2 = st.columns(2)
        with f1:
            service = st.text_input("서비스", value=default_service)
        with f2:
            build_type = st.text_input("구분", value=default_build)

        f4, f5, f6 = st.columns([1.2, 1.0, 1.0])
        with f4:
            branch = st.text_input("신청점/부서", value=default_branch)
        with f5:
            customer_no = st.text_input("고객번호", value=default_customer_no)
        with f6:
            status = st.text_input("상태", value=status_default)

        f7, f8 = st.columns(2)
        with f7:
            receipt_date = st.text_input("접수일자", value=default_receipt_date)
        with f8:
            if input_category in ["개설대기", "개설진행", "연계대기", "연계진행"]:
                plan_date = st.date_input(
                    plan_date_label,
                    value=date_input_value(default_plan_date, today.date()),
                    key=f"weekly_plan_date_{input_category}",
                ).strftime("%Y-%m-%d")
            else:
                plan_date = st.text_input(plan_date_label, value=default_plan_date)

        issue = st.text_area("이슈 / 진행내용", height=100, placeholder="PPT 주간보고의 이슈 항목처럼 고객별 진행상황을 작성")
        special_note = st.text_area("특이사항", height=80, placeholder="예정 작업, 고객 회신 필요사항, 은행 전달사항 등")

        submitted = st.form_submit_button("주간보고 저장", use_container_width=True, type="primary")

    if submitted:
        if not company.strip():
            st.warning("고객명을 입력해주세요.")
        else:
            db = load_db(WEEKLY_REPORT_FILE, {})
            entries = db.setdefault(user_name, [])
            now_text = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S")
            entries.append({
                "id": hashlib.md5(f"{user_name}-{now_text}-{company}-{input_category}".encode("utf-8")).hexdigest(),
                "보고시작일": str(week_start),
                "보고종료일": str(week_end),
                "부문": input_part,
                "카테고리": input_category,
                "서비스": service.strip(),
                "구분": build_type.strip(),
                "고객명": company.strip(),
                "고객번호": customer_no.strip(),
                "신청점": branch.strip(),
                "접수일자": receipt_date.strip(),
                "일자구분": plan_date_label,
                "일자": plan_date.strip(),
                "이슈": issue.strip(),
                "특이사항": special_note.strip(),
                "상태": status.strip(),
                "담당자": user_name,
                "작성시각": now_text,
            })
            save_db(WEEKLY_REPORT_FILE, db)
            st.success("주간보고 이력을 저장했습니다.")
            st.rerun()

    db = load_db(WEEKLY_REPORT_FILE, {})
    my_entries = db.get(user_name, [])
    st.markdown("#### 내 작성 이력")
    if not my_entries:
        st.info("아직 저장된 주간보고 이력이 없습니다.")
        return

    sorted_entries = sorted(my_entries, key=lambda x: x.get("작성시각", ""), reverse=True)
    history_df = pd.DataFrame(sorted_entries)
    show_cols = ["보고시작일", "보고종료일", "부문", "카테고리", "고객명", "접수일자", "일자구분", "일자", "이슈", "특이사항", "상태", "작성시각"]
    show_cols = [c for c in show_cols if c in history_df.columns]

    display_df = history_df[show_cols].copy()
    display_df.insert(0, "선택", False)

    edited_df = st.data_editor(
        display_df,
        use_container_width=True,
        hide_index=True,
        column_config={"선택": st.column_config.CheckboxColumn("선택", default=False, width="small")},
        disabled=show_cols,
        key="weekly_history_editor",
    )

    selected_indices = [i for i, v in enumerate(edited_df["선택"]) if v]
    selected_ids = [sorted_entries[i].get("id") for i in selected_indices]

    btn_col1, btn_col2, _ = st.columns([15, 15, 70])
    with btn_col1:
        edit_clicked = st.button("수정", use_container_width=True, disabled=(len(selected_ids) != 1), key="weekly_edit_btn")
    with btn_col2:
        del_clicked = st.button("삭제", use_container_width=True, disabled=(len(selected_ids) == 0), key="weekly_del_btn")

    if edit_clicked and len(selected_ids) == 1:
        st.session_state["weekly_edit_id"] = selected_ids[0]
        st.session_state["weekly_edit_mode"] = True
        st.rerun()

    if del_clicked and selected_ids:
        db[user_name] = [row for row in my_entries if row.get("id") not in selected_ids]
        save_db(WEEKLY_REPORT_FILE, db)
        st.success(f"{len(selected_ids)}개 이력을 삭제했습니다.")
        st.rerun()

    if st.session_state.get("weekly_edit_mode") and st.session_state.get("weekly_edit_id"):
        edit_id = st.session_state["weekly_edit_id"]
        edit_entry = next((row for row in my_entries if row.get("id") == edit_id), None)
        if edit_entry:
            st.markdown("---")
            st.markdown("##### 선택 이력 수정")
            with st.form("weekly_edit_form"):
                ec1, ec2 = st.columns(2)
                with ec1:
                    e_start = st.text_input("보고시작일", value=edit_entry.get("보고시작일", ""))
                    e_category = st.text_input("카테고리", value=edit_entry.get("카테고리", ""))
                    e_receipt = st.text_input("접수일자", value=edit_entry.get("접수일자", ""))
                    e_date = st.text_input("일자", value=edit_entry.get("일자", ""))
                with ec2:
                    e_end = st.text_input("보고종료일", value=edit_entry.get("보고종료일", ""))
                    e_status = st.text_input("상태", value=edit_entry.get("상태", ""))
                    e_date_type = st.text_input("일자구분", value=edit_entry.get("일자구분", ""))
                e_issue = st.text_area("이슈 / 진행내용", value=edit_entry.get("이슈", ""), height=100)
                e_special = st.text_area("특이사항", value=edit_entry.get("특이사항", ""), height=80)
                sc1, sc2 = st.columns(2)
                with sc1:
                    save_edit = st.form_submit_button("저장", use_container_width=True, type="primary")
                with sc2:
                    cancel_edit = st.form_submit_button("취소", use_container_width=True)

            if save_edit:
                for i, row in enumerate(db.get(user_name, [])):
                    if row.get("id") == edit_id:
                        db[user_name][i].update({
                            "보고시작일": e_start.strip(),
                            "보고종료일": e_end.strip(),
                            "카테고리": e_category.strip(),
                            "접수일자": e_receipt.strip(),
                            "일자구분": e_date_type.strip(),
                            "일자": e_date.strip(),
                            "상태": e_status.strip(),
                            "이슈": e_issue.strip(),
                            "특이사항": e_special.strip(),
                        })
                        break
                save_db(WEEKLY_REPORT_FILE, db)
                st.session_state["weekly_edit_mode"] = False
                st.session_state["weekly_edit_id"] = None
                st.success("수정되었습니다.")
                st.rerun()

            if cancel_edit:
                st.session_state["weekly_edit_mode"] = False
                st.session_state["weekly_edit_id"] = None
                st.rerun()


def load_weekly_hana_for_status():
    try:
        df = read_google_csv(st.session_state.get("url_hana", DEFAULT_URL_HANA), header=2)
    except Exception:
        df = st.session_state.get("hana_sheet_df")
    if df is None:
        return pd.DataFrame()
    df = df.copy()
    df.columns = [str(c).strip() for c in df.columns]
    return df.dropna(how="all").reset_index(drop=True)


def weekly_customer_status_tables(hana_df, year=2026, snapshot_end_date=None, save_snapshot=False):
    if hana_df is None or hana_df.empty:
        return {}

    df = hana_df.copy()
    open_receipt_col = find_col(df, ["신규접수일"])
    link_receipt_col = find_col(df, ["추가연계접수일"])
    open_date_col = find_col(df, ["개설/이행일", "개설일", "이행일"])
    link_date_col = find_col(df, ["연계일자", "연계일"])
    link_billing_date_col = find_col(df, ["연계청구일자", "연계청구일", "청구일자"])
    open_status_col = find_col(df, ["개설상태"])
    link_status_col = find_col(df, ["연계상태"])
    build_type_col = find_col(df, ["구축형"])
    manage_col = find_col(df, ["관리구분"])
    end_col = find_col(df, ["해지일자"])

    if not open_status_col:
        return {}

    def dates(col):
        if not col or col not in df.columns:
            return pd.Series(pd.NaT, index=df.index)
        return pd.to_datetime(df[col].map(parse_sheet_date), errors="coerce")

    def fmt(value):
        try:
            value = int(value)
        except Exception:
            value = 0
        return "-" if value == 0 else f"{value:,}"

    def count(mask):
        return int(mask.fillna(False).sum())

    open_receipt_dates = dates(open_receipt_col)
    link_receipt_dates = dates(link_receipt_col).fillna(open_receipt_dates)
    open_done_dates = dates(open_date_col)
    link_done_dates = dates(link_date_col)
    link_billing_dates = dates(link_billing_date_col)
    end_dates = dates(end_col)
    as_of_date = pd.to_datetime(snapshot_end_date, errors="coerce")
    if pd.isna(as_of_date):
        as_of_date = pd.Timestamp.max

    year_start = pd.Timestamp(f"{year}-01-01")
    active = end_dates.isna() | (end_dates > as_of_date)
    # 월별 표용: 전체 기간(연도 무관)
    open_received_by_asof = open_receipt_dates.notna() & (open_receipt_dates <= as_of_date)
    link_received_by_asof = link_receipt_dates.notna() & (link_receipt_dates <= as_of_date)
    # 상태 표용: year 기준(2026.01.01~) 필터
    open_received_in_year = open_receipt_dates.notna() & (open_receipt_dates >= year_start) & (open_receipt_dates <= as_of_date)
    link_received_in_year = link_receipt_dates.notna() & (link_receipt_dates >= year_start) & (link_receipt_dates <= as_of_date)
    open_done_by_asof = open_done_dates.notna() & (open_done_dates <= as_of_date)
    link_done_by_asof = link_done_dates.notna() & (link_done_dates <= as_of_date)
    link_billing_by_asof = link_billing_dates.notna() & (link_billing_dates <= as_of_date)
    open_status = df[open_status_col].astype(str).str.strip()
    link_status = df[link_status_col].astype(str).str.strip() if link_status_col else pd.Series("", index=df.index)
    link_clean = link_status.str.replace(r"\s+", "", regex=True)
    build_type = df[build_type_col].astype(str).str.strip() if build_type_col else pd.Series("", index=df.index)
    manage = df[manage_col].astype(str).str.strip() if manage_col else pd.Series("", index=df.index)

    basic = build_type.str.contains("기본", na=False)
    linked = build_type.str.contains("연계|이행", na=False)
    transfer = manage.str.contains("이관", na=False)

    # 개설상태 명시값 기준 분류 (소거법 사용 X, NaN 제외)
    open_cancel = open_received_by_asof & open_status.str.contains("취소|DROP|드랍", case=False, na=False)
    link_cancel = link_received_by_asof & link_status.str.contains("취소|DROP|드랍", case=False, na=False)
    open_done = active & open_received_by_asof & ~open_cancel & (open_status.str.contains("완료|이행완료", na=False) & open_done_by_asof)
    open_progress = active & open_received_by_asof & ~open_cancel & ~open_done & open_status.str.contains("진행|구축중|처리중", na=False)
    open_wait = active & open_received_by_asof & ~open_cancel & ~open_done & ~open_progress & open_status.str.contains("대기", na=False)
    # 연계완료: 연계일자 유무와 관계없이 상태값 "완료" 기준으로 판단
    link_done = active & link_received_by_asof & ~link_cancel & link_status.str.contains("완료", na=False)
    # 연계진행/대기: 패턴을 넓혀 "진행", "대기" 포함 변형 모두 매칭
    link_progress = active & link_received_by_asof & ~link_cancel & ~link_done & link_status.str.contains("진행", na=False)
    link_wait = active & link_received_by_asof & ~link_cancel & ~link_done & ~link_progress & link_status.str.contains("대기", na=False)
    link_receipt = active & (link_received_by_asof | link_wait | link_progress | link_done)
    terminated = end_dates.notna() & (end_dates <= as_of_date)

    # 상태 표: 구글시트 기준과 동일하게 연계상태 컬럼값만으로 단순 집계
    # active/link_received_by_asof 등 부가 조건 없이 연계상태 텍스트만 매칭
    st_link_wait     = link_status.str.contains("대기", na=False)
    st_link_progress = link_status.str.contains("진행", na=False)
    st_link_done     = link_status.str.contains("완료", na=False)
    status_rows = [
        ["전체", "", count(open_wait), count(open_progress), count(open_done), count(st_link_wait), count(st_link_progress), count(st_link_done), count(terminated)],
        ["", "기본형", count(open_wait & basic), count(open_progress & basic), count(open_done & basic), "-", "-", "-", count(terminated & basic)],
        ["", "연계형", count(open_wait & linked), count(open_progress & linked), count(open_done & linked), count(st_link_wait), count(st_link_progress), count(st_link_done), count(terminated & linked)],
    ]
    status_df = pd.DataFrame(status_rows, columns=["구분", "유형", "구축대기", "구축진행", "구축완료", "연계대기", "연계진행", "연계완료", "해지"])

    month_cols = ["2025년"] + [f"{year}{m:02d}" for m in range(1, 13)] + [f"{str(year)[-2:]}년합계", "합계"]

    def month_count_values(mask, event_date):
        values = []
        bounded_date = event_date.notna() & (event_date <= as_of_date)
        prev_year_value = count(mask & bounded_date & (event_date.dt.year == year - 1))
        values.append(prev_year_value)
        year_total = 0
        for month in range(1, 13):
            value = count(mask & bounded_date & (event_date.dt.year == year) & (event_date.dt.month == month))
            values.append(value)
            year_total += value
        values.append(year_total)                    # 26년합계 (2026년 월별 합산)
        values.append(prev_year_value + year_total)  # 합계 (2025년 + 2026년)
        return values

    def fmt_values(values):
        return [fmt(v) for v in values]

    def summed_values(*rows):
        return [sum(values) for values in zip(*rows)]

    def debug_rows(mask, event_date, reason):
        debug_cols = [
            find_col(df, ["업체명", "고객명", "상호", "회사명"]),
            find_col(df, ["담당자"]),
            find_col(df, ["구축형"]),
            find_col(df, ["관리구분"]),
            open_receipt_col,
            open_status_col,
            open_date_col,
            end_col,
        ]
        labels = ["고객명", "담당자", "구축형", "관리구분", "신규접수일", "개설상태", "개설/이행일", "해지일자"]
        rows = []
        for idx in df[mask].index:
            row = {}
            for label, col in zip(labels, debug_cols):
                row[label] = df.at[idx, col] if col and col in df.columns else ""
            row["신규접수일(파싱)"] = event_date.at[idx].strftime("%Y-%m-%d") if pd.notna(event_date.at[idx]) else ""
            row["산출사유"] = reason
            rows.append(row)
        return pd.DataFrame(rows, columns=labels + ["신규접수일(파싱)", "산출사유"])

    open_wait_values = month_count_values(open_wait, open_receipt_dates)
    open_progress_values = month_count_values(open_progress, open_receipt_dates)
    open_done_values = month_count_values(open_done, open_receipt_dates)
    open_cancel_values = month_count_values(open_cancel, open_receipt_dates)
    debug_open_wait_prev_year = debug_rows(
        open_wait & open_receipt_dates.notna() & (open_receipt_dates.dt.year == year - 1) & (open_receipt_dates <= as_of_date),
        open_receipt_dates,
        f"{year - 1}년 신규접수 + 기준일 현재 구축대기",
    )

    # 구축 접수: 신규접수일 기준 직접 카운트 (개설상태 패턴 불일치로 인한 누락 방지)
    open_receipt_count_values = month_count_values(open_received_by_asof, open_receipt_dates)
    open_month_rows = [
        ["구축", "접수"] + fmt_values(open_receipt_count_values),
        ["", "대기"] + fmt_values(open_wait_values),
        ["", "진행"] + fmt_values(open_progress_values),
        ["", "완료"] + fmt_values(open_done_values),
        ["", "취소"] + fmt_values(open_cancel_values),
    ]

    # 연계 월별 표: 연계상태 텍스트만 매칭 (구축형 필터 없음 → 전체 행 기준과 일치)
    _lr = link_receipt_dates.notna() & (link_receipt_dates <= as_of_date)
    link_received_mask   = _lr
    mt_link_wait         = _lr & link_status.str.contains("대기", na=False)
    mt_link_progress     = _lr & link_status.str.contains("진행", na=False)
    mt_link_done         = _lr & link_status.str.contains("완료", na=False)
    mt_link_cancel       = _lr & link_status.str.contains("취소|DROP|드랍", case=False, na=False)
    link_receipt_count_values = month_count_values(link_received_mask, link_receipt_dates)
    link_wait_values     = month_count_values(mt_link_wait,     link_receipt_dates)
    link_progress_values = month_count_values(mt_link_progress, link_receipt_dates)
    link_done_values     = month_count_values(mt_link_done,     link_receipt_dates)
    link_cancel_values   = month_count_values(mt_link_cancel,   link_receipt_dates)

    link_month_rows = [
        ["연계", "접수"] + fmt_values(summed_values(link_wait_values, link_progress_values, link_done_values, link_cancel_values)),
        ["", "대기"] + fmt_values(link_wait_values),
        ["", "진행"] + fmt_values(link_progress_values),
        ["", "완료"] + fmt_values(link_done_values),
        ["", "취소"] + fmt_values(link_cancel_values),
    ]
    monthly_df = pd.DataFrame(open_month_rows + link_month_rows, columns=["구분", "상태"] + month_cols)

    current_group = ""
    monthly_counts = {}
    for row in open_month_rows + link_month_rows:
        if row[0]:
            current_group = row[0]
        monthly_counts[f"{current_group}|{row[1]}"] = [int(str(v).replace(",", "").replace("-", "0")) for v in row[2:]]

    snapshot_key = str(snapshot_end_date or "").strip()
    snapshots = load_db(WEEKLY_STATUS_SNAPSHOT_FILE, {})
    prev_snapshot = None
    if snapshot_key and snapshots:
        prev_keys = sorted(k for k in snapshots.keys() if k < snapshot_key)
        if prev_keys:
            prev_snapshot = snapshots.get(prev_keys[-1], {}).get("monthly_counts", {})

    delta_rows = []
    current_group_for_delta = ""
    for _, row in monthly_df.iterrows():
        if str(row["구분"]).strip():
            current_group_for_delta = str(row["구분"]).strip()
        status = str(row["상태"]).strip()
        key = f"{current_group_for_delta}|{status}"
        current_values = monthly_counts.get(key, [0] * len(month_cols))
        previous_values = prev_snapshot.get(key, [0] * len(month_cols)) if prev_snapshot else None
        delta_values = []
        for idx, current_value in enumerate(current_values):
            if previous_values is None:
                delta_values.append("-")
                continue
            prev_value = previous_values[idx] if idx < len(previous_values) else 0
            delta = current_value - int(prev_value)
            if delta > 0:
                delta_values.append(f"▲{delta:,}")
            elif delta < 0:
                delta_values.append(f"▼{abs(delta):,}")
            else:
                delta_values.append("-")
        delta_rows.append(["", "전주대비"] + delta_values)
    expanded_rows = []
    for idx, row in monthly_df.iterrows():
        expanded_rows.append(row.tolist())
        expanded_rows.append(delta_rows[idx])
    monthly_display_df = pd.DataFrame(expanded_rows, columns=monthly_df.columns)

    if save_snapshot and snapshot_key:
        snapshots[snapshot_key] = {
            "saved_at": (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S"),
            "monthly_counts": monthly_counts,
        }
        save_db(WEEKLY_STATUS_SNAPSHOT_FILE, snapshots)

    complete_cols = ["구분", "-"] + [f"{year}{m:02d}" for m in range(1, 13)] + [f"{str(year)[-2:]}' 합계"]
    open_complete_values = []
    link_complete_values = []
    open_total = 0
    link_total = 0
    link_billing_done = link_clean.eq("ERP청구완료")
    for month in range(1, 13):
        ov = count(open_done_by_asof & (open_done_dates.dt.year == year) & (open_done_dates.dt.month == month))
        lv = count(link_billing_done & link_billing_by_asof & (link_billing_dates.dt.year == year) & (link_billing_dates.dt.month == month))
        open_complete_values.append(fmt(ov))
        link_complete_values.append(fmt(lv))
        open_total += ov
        link_total += lv
    complete_df = pd.DataFrame([
        ["구축", "-"] + open_complete_values + [fmt(open_total)],
        ["연계", "-"] + link_complete_values + [fmt(link_total)],
    ], columns=complete_cols)

    return {
        "status": status_df,
        "monthly": monthly_display_df,
        "monthly_base": monthly_df,
        "complete": complete_df,
        "debug_open_wait_prev_year": debug_open_wait_prev_year,
    }


def weekly_prev_friday():
    today = datetime.utcnow() + timedelta(hours=9)
    return (today - timedelta(days=(today.weekday() - 4) % 7 or 7)).date()


def save_weekly_status_snapshot(hana_df, snapshot_end_date):
    weekly_customer_status_tables(hana_df, snapshot_end_date=snapshot_end_date, save_snapshot=True)


def render_weekly_front_status_tables(hana_df, snapshot_end_date=None):
    tables = weekly_customer_status_tables(hana_df, snapshot_end_date=snapshot_end_date)
    if not tables:
        st.warning("하나은행 구글 시트 기준 고객 현황을 계산할 수 없습니다.")
        return

    prev_friday = pd.to_datetime(snapshot_end_date, errors="coerce")
    if pd.isna(prev_friday):
        prev_friday = pd.Timestamp(weekly_prev_friday())
    st.markdown(f"#### 2026년 기준 (2026.01.01 ~ {prev_friday.strftime('%Y.%m.%d')})")
    render_plain_html_table(tables["status"], max_rows=30, center_align=True, merge_cols=["구분"])
    st.markdown(f"#### ■ 26년 월별 접수고객 진행 현황 관리 (2026.01.01 ~ {prev_friday.strftime('%Y.%m.%d')})")
    render_plain_html_table(tables["monthly"], max_rows=30, center_align=True, merge_cols=["구분"])
    debug_open_wait_prev_year = tables.get("debug_open_wait_prev_year")
    if debug_open_wait_prev_year is not None and not debug_open_wait_prev_year.empty:
        with st.expander(f"2025년 구축대기 산출내역 ({len(debug_open_wait_prev_year)}건)"):
            render_plain_html_table(debug_open_wait_prev_year, max_rows=100)
    st.markdown("#### ■ 26년 월별 개설완료 현황 관리")
    render_plain_html_table(tables["complete"], max_rows=10, center_align=True, merge_cols=["구분"])


def build_weekly_report_ppt_bytes(report_df, week_start="", week_end="", hana_df=None, save_snapshot=False):
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    from copy import deepcopy

    if not os.path.exists(WEEKLY_PPT_TEMPLATE_FILE):
        raise FileNotFoundError(f"주간보고 PPT 템플릿 파일을 찾을 수 없습니다: {WEEKLY_PPT_TEMPLATE_FILE}")

    prs = Presentation(WEEKLY_PPT_TEMPLATE_FILE)
    df = report_df.copy() if report_df is not None else pd.DataFrame()
    if df.empty:
        raise ValueError("다운로드할 주간보고 이력이 없습니다.")

    def clean(value):
        if value is None:
            return ""
        try:
            if pd.isna(value):
                return ""
        except Exception:
            pass
        text = str(value).strip()
        return "" if text.lower() in ["nan", "nat", "none"] else text

    def set_cell_text(cell, value, font_size=8, align=PP_ALIGN.CENTER):
        cell.text = clean(value)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = "맑은 고딕"
                run.font.size = Pt(font_size)

    def slide_tables(slide):
        return [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]

    def ensure_rows(table, required_rows):
        while len(table.rows) < required_rows:
            new_tr = deepcopy(table._tbl.tr_lst[-1])
            table._tbl.append(new_tr)
            for cell in table.rows[-1].cells:
                cell.text = ""

    def clear_body(table):
        for r in range(1, len(table.rows)):
            for c in range(len(table.columns)):
                set_cell_text(table.cell(r, c), "", font_size=8)

    def fill_table(table, rows, font_size=8, left_cols=None):
        left_cols = set(left_cols or [])
        required_rows = max(2, len(rows) + 1)
        ensure_rows(table, required_rows)
        clear_body(table)
        for r_idx, values in enumerate(rows, start=1):
            if r_idx >= len(table.rows):
                break
            for c_idx, value in enumerate(values[:len(table.columns)]):
                align = PP_ALIGN.LEFT if c_idx in left_cols else PP_ALIGN.CENTER
                set_cell_text(table.cell(r_idx, c_idx), value, font_size=font_size, align=align)

    def normalize_hana(df_hana):
        if df_hana is None or df_hana.empty:
            try:
                url = st.session_state.get("url_hana", DEFAULT_URL_HANA)
            except Exception:
                url = DEFAULT_URL_HANA
            df_hana = read_google_csv(url, header=2)
        df_hana = df_hana.copy()
        df_hana.columns = [str(c).strip() for c in df_hana.columns]
        return df_hana.dropna(how="all").reset_index(drop=True)

    def event_dates(series):
        return pd.to_datetime(series.map(parse_sheet_date), errors="coerce")

    def month_columns_for(table, start_col=3, end_col=None):
        end_col = end_col if end_col is not None else len(table.columns) - 1
        cols = []
        for c in range(start_col, end_col):
            label = clean(table.cell(0, c).text).replace("\n", "").replace(" ", "")
            m = re.search(r"(\d{2})\.(\d{2})", label)
            if m:
                cols.append((c, int(m.group(1)), int(m.group(2))))
        return cols

    def fmt_count(value):
        try:
            value = int(value)
        except Exception:
            value = 0
        return "-" if value == 0 else f"{value:,}"

    def fmt_delta(value):
        try:
            value = int(value)
        except Exception:
            value = 0
        if value > 0:
            return f"▲{value:,}"
        if value < 0:
            return f"▼{abs(value):,}"
        return "-"

    def count_by_month(mask, dates, year, month):
        return int((mask & dates.notna() & (dates.dt.year == year) & (dates.dt.month == month)).sum())

    def count_before_year(mask, dates, year):
        return int((mask & dates.notna() & (dates.dt.year < year)).sum())

    def week_mask(dates):
        if not week_start or not week_end:
            return pd.Series(False, index=dates.index)
        return dates.between(pd.Timestamp(week_start), pd.Timestamp(week_end), inclusive="both")

    def fill_monthly_table(table, metric_rows, year=2026, has_prev_year=True):
        start_col = 3 if has_prev_year else 2
        month_cols = month_columns_for(table, start_col=start_col, end_col=len(table.columns) - 1)
        for row_idx, metric in metric_rows.items():
            mask = metric["mask"]
            dates = metric["dates"]
            total = 0
            if has_prev_year and len(table.columns) > 2:
                prev_count = count_before_year(mask, dates, year)
                total += prev_count
                set_cell_text(table.cell(row_idx, 2), fmt_count(prev_count), font_size=8)
            for col_idx, yy, mm in month_cols:
                full_year = 2000 + yy
                count = count_by_month(mask, dates, full_year, mm)
                total += count
                set_cell_text(table.cell(row_idx, col_idx), fmt_count(count), font_size=8)
                delta_row = row_idx + 1
                if delta_row < len(table.rows):
                    delta = int((mask & week_mask(dates) & (dates.dt.year == full_year) & (dates.dt.month == mm)).sum())
                    set_cell_text(table.cell(delta_row, col_idx), fmt_delta(delta), font_size=8)
            set_cell_text(table.cell(row_idx, len(table.columns) - 1), fmt_count(total), font_size=8)

    def fill_overall_table(table, metrics):
        build_type = metrics["build_type"]
        rows = [
            ("전체", pd.Series(True, index=metrics["base"].index)),
            ("기본형", build_type.str.contains("기본", na=False)),
            ("연계형", build_type.str.contains("연계|이행", na=False)),
        ]
        keys = ["open_wait", "open_progress", "open_done", "link_wait", "link_progress", "link_done"]
        for r_idx, (_, row_mask) in enumerate(rows, start=1):
            if r_idx >= len(table.rows):
                break
            for c_idx, key in enumerate(keys, start=1):
                value = int((metrics[key] & row_mask).sum())
                set_cell_text(table.cell(r_idx, c_idx), fmt_count(value), font_size=8)

    def fill_customer_status_slides():
        try:
            hana = normalize_hana(hana_df)
        except Exception:
            return
        if hana.empty or len(prs.slides) < 4:
            return
        front_tables = weekly_customer_status_tables(hana, snapshot_end_date=week_end, save_snapshot=save_snapshot)

        open_receipt_col = find_col(hana, ["신규접수일"])
        link_receipt_col = find_col(hana, ["추가연계접수일"])
        open_date_col = find_col(hana, ["개설/이행일", "개설일", "이행일"])
        link_date_col = find_col(hana, ["연계일자", "연계일"])
        link_billing_date_col = find_col(hana, ["연계청구일자", "연계청구일", "청구일자"])
        open_status_col = find_col(hana, ["개설상태"])
        link_status_col = find_col(hana, ["연계상태"])
        build_type_col = find_col(hana, ["구축형"])
        end_col = find_col(hana, ["해지일자"])

        if not open_receipt_col or not open_status_col:
            return

        base = hana.copy()
        active = base[end_col].map(is_blank_value) if end_col else pd.Series(True, index=base.index)
        open_status = base[open_status_col].astype(str).str.strip() if open_status_col else pd.Series("", index=base.index)
        link_status = base[link_status_col].astype(str).str.strip() if link_status_col else pd.Series("", index=base.index)
        link_clean = link_status.str.replace(r"\s+", "", regex=True)
        build_type = base[build_type_col].astype(str).str.strip() if build_type_col else pd.Series("", index=base.index)

        open_receipt_dates = event_dates(base[open_receipt_col]) if open_receipt_col else pd.Series(pd.NaT, index=base.index)
        link_receipt_dates = event_dates(base[link_receipt_col]) if link_receipt_col else open_receipt_dates
        link_receipt_dates = link_receipt_dates.fillna(open_receipt_dates)
        open_done_dates = event_dates(base[open_date_col]) if open_date_col else pd.Series(pd.NaT, index=base.index)
        link_done_dates = event_dates(base[link_date_col]) if link_date_col else pd.Series(pd.NaT, index=base.index)
        link_billing_dates = event_dates(base[link_billing_date_col]) if link_billing_date_col else pd.Series(pd.NaT, index=base.index)

        open_done = active & (open_status.str.contains("완료|이행완료", na=False) | open_done_dates.notna())
        open_progress = active & ~open_done & open_status.str.contains("진행|구축중|처리중", na=False)
        open_wait = active & ~open_done & ~open_progress
        open_drop = open_status.str.contains("취소|DROP|드랍", case=False, na=False)

        link_done = active & link_status.str.contains("완료", na=False)
        link_billing_done = active & link_clean.eq("ERP청구완료")
        link_progress = active & ~link_done & link_status.str.contains("ERP연계진행|연계진행|진행", na=False)
        link_wait = active & link_status.str.contains("ERP연계대기", na=False)
        link_drop = link_status.str.contains("취소|DROP|드랍", case=False, na=False)
        link_receipt = active & (link_receipt_dates.notna() | link_wait | link_progress | link_done)

        metrics = {
            "base": base,
            "build_type": build_type,
            "type_col": build_type_col,
            "open_wait": open_wait,
            "open_progress": open_progress,
            "open_done": open_done,
            "link_wait": link_wait,
            "link_progress": link_progress,
            "link_done": link_done,
        }

        slide2_tables = slide_tables(prs.slides[1])
        if slide2_tables:
            fill_monthly_table(slide2_tables[0], {
                1: {"mask": active & open_receipt_dates.notna(), "dates": open_receipt_dates},
                3: {"mask": open_wait, "dates": open_receipt_dates},
                5: {"mask": open_progress, "dates": open_receipt_dates},
                7: {"mask": open_done, "dates": open_receipt_dates},
                9: {"mask": open_drop, "dates": open_receipt_dates},
            }, year=2026, has_prev_year=True)

        slide3_tables = slide_tables(prs.slides[2])
        if slide3_tables:
            fill_monthly_table(slide3_tables[0], {
                1: {"mask": link_receipt, "dates": link_receipt_dates},
                3: {"mask": link_wait, "dates": link_receipt_dates},
                5: {"mask": link_progress, "dates": link_receipt_dates},
                7: {"mask": link_done, "dates": link_receipt_dates},
                9: {"mask": link_drop, "dates": link_receipt_dates},
            }, year=2026, has_prev_year=True)

        slide4_tables = slide_tables(prs.slides[3])
        if slide4_tables:
            fill_monthly_table(slide4_tables[0], {
                1: {"mask": open_done, "dates": open_done_dates},
                3: {"mask": link_billing_done, "dates": link_billing_dates},
            }, year=2026, has_prev_year=False)
            if len(slide4_tables) > 1:
                fill_overall_table(slide4_tables[1], metrics)

    def update_title_dates():
        date_text = ""
        if week_start and week_end:
            date_text = f"{week_start} - {week_end}"
        elif "보고시작일" in df.columns and "보고종료일" in df.columns:
            starts = [clean(v) for v in df["보고시작일"].tolist() if clean(v)]
            ends = [clean(v) for v in df["보고종료일"].tolist() if clean(v)]
            if starts and ends:
                date_text = f"{min(starts)} - {max(ends)}"
        if not date_text or not prs.slides:
            return
        for shape in prs.slides[0].shapes:
            if getattr(shape, "has_text_frame", False) and "2026." in shape.text:
                shape.text = date_text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "맑은 고딕"
                        run.font.size = Pt(18)
                break

    def rows_for(category, no_start=1):
        part = df[df["카테고리"].astype(str) == category].copy() if "카테고리" in df.columns else pd.DataFrame()
        rows = []
        for no, (_, row) in enumerate(part.iterrows(), start=no_start):
            service = clean(row.get("서비스", "")) or "일반"
            build_type = clean(row.get("구분", ""))
            company = clean(row.get("고객명", ""))
            branch = clean(row.get("신청점", ""))
            receipt = clean(row.get("접수일자", ""))
            issue = clean(row.get("이슈", ""))
            note = clean(row.get("특이사항", ""))
            detail = issue if not note else f"{issue}\n{note}" if issue else note
            staff = clean(row.get("담당자", ""))
            date_value = clean(row.get("일자", ""))
            status = clean(row.get("상태", ""))
            rows.append({
                "no": no,
                "service": service,
                "type": build_type,
                "company": company,
                "branch": branch,
                "receipt": receipt,
                "issue": detail,
                "staff": staff,
                "date": date_value,
                "status": status,
            })
        return rows

    def opening_wait_rows(rows):
        return [[r["no"], r["service"], r["type"], r["company"], r["branch"], r["receipt"], r["date"], r["issue"], r["staff"]] for r in rows]

    def opening_progress_rows(rows):
        return [[r["no"], r["type"], r["company"], r["branch"], r["receipt"], r["issue"], r["staff"], r["date"], r["status"]] for r in rows]

    def link_wait_rows(rows):
        return [[r["no"], r["company"], r["branch"], r["receipt"], r["date"], r["issue"], r["staff"]] for r in rows]

    def link_progress_rows(rows):
        return [[r["no"], r["company"], r["branch"], r["receipt"], "", r["status"], r["issue"], r["staff"], r["date"]] for r in rows]

    def operation_rows(rows):
        return [[r["no"], r["type"], r["company"], r["branch"], r["receipt"], r["issue"], r["staff"], r["date"], r["status"]] for r in rows]

    def fill_slide_table(slide_idx, values, left_cols):
        if slide_idx >= len(prs.slides):
            return
        tables = slide_tables(prs.slides[slide_idx])
        if tables:
            fill_table(tables[0], values, font_size=7, left_cols=left_cols)

    update_title_dates()
    for slide_idx in range(4, min(24, len(prs.slides))):
        for table in slide_tables(prs.slides[slide_idx]):
            clear_body(table)

    open_wait = rows_for("개설대기")
    open_progress = rows_for("개설진행")
    open_done = rows_for("개설완료")
    link_wait = rows_for("연계대기")
    link_progress = rows_for("연계진행")
    link_done = rows_for("연계완료")
    operation = rows_for("운영부문")

    fill_slide_table(4, opening_wait_rows(open_wait), left_cols={6, 7})
    fill_slide_table(7, opening_progress_rows(open_done), left_cols={5})
    fill_slide_table(9, opening_progress_rows(open_progress), left_cols={5})
    fill_slide_table(10, link_wait_rows(link_wait), left_cols={5})
    fill_slide_table(14, link_progress_rows(link_progress), left_cols={6})
    fill_slide_table(19, link_progress_rows(link_done), left_cols={6})
    fill_slide_table(21, operation_rows(operation), left_cols={5})

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


OPERATION_PLAN_FILE = "operation_plan.json"


def show_operation_plan():
    st.markdown("### 운영계획")
    tab_kpi_rec, tab_no_login = st.tabs(["KPI 집중 추천", "미로그인 고객 현황"])

    # ── 탭1: KPI 집중 추천 ──
    with tab_kpi_rec:
        _force_kpi = st.session_state.pop("_op_kpi_refresh", False)
        refresh_col, _ = st.columns([1, 5])
        with refresh_col:
            if st.button("새로고침", key="op_kpi_refresh_btn", use_container_width=True):
                st.session_state["_op_kpi_refresh"] = True
                st.rerun()
        try:
            hana_df_kpi = read_google_csv(
                st.session_state.get("url_hana", DEFAULT_URL_HANA),
                header=2, force_refresh=_force_kpi,
            )
            billing_df_kpi = load_hana_billing_df(force_refresh=_force_kpi)
        except Exception as e:
            st.error(f"데이터를 불러오지 못했습니다: {e}")
            hana_df_kpi = pd.DataFrame()
            billing_df_kpi = pd.DataFrame()

        if hana_df_kpi is None or hana_df_kpi.empty or billing_df_kpi is None or billing_df_kpi.empty:
            st.info("KPI 추천에 사용할 하나은행 구글시트 또는 청구시트 데이터가 없습니다.")
        else:
            hana_df_kpi = hana_df_kpi.dropna(how="all").reset_index(drop=True)
            billing_df_kpi = billing_df_kpi.dropna(how="all").reset_index(drop=True)
            render_kpi_activity_recommendations(
                hana_df_kpi,
                billing_df_kpi,
                user_name=None,
                key_prefix="op_kpi_rec",
            )

    # ── 탭2: 미로그인 고객 현황 ──
    with tab_no_login:
        _force_nl = st.session_state.pop("_nl_rf_op_nl_refresh", False)
        try:
            hana_df_nl    = read_google_csv(
                st.session_state.get("url_hana", DEFAULT_URL_HANA),
                header=2, force_refresh=_force_nl,
            )
            billing_df_nl = load_hana_billing_df(force_refresh=_force_nl)
        except Exception as e:
            st.error(f"데이터를 불러오지 못했습니다: {e}")
            return
        if hana_df_nl is None or hana_df_nl.empty or billing_df_nl is None or billing_df_nl.empty:
            st.info("구글시트 또는 청구시트 데이터가 없습니다.")
        else:
            hana_df_nl    = hana_df_nl.dropna(how="all").reset_index(drop=True)
            billing_df_nl = billing_df_nl.dropna(how="all").reset_index(drop=True)

            st.markdown("#### 개설/이행일 이후 미로그인 고객")
            st.caption("개설/이행일 기준으로 청구시트 최종로그인일자가 없거나 개설/이행일 이전인 고객")
            _render_no_login_section(
                build_no_login_after_open(hana_df_nl, billing_df_nl),
                year_key="no_login_open_year", owner_key="no_login_open_owner",
                label="미로그인 고객", download_prefix="미로그인고객_개설이행일",
                refresh_key="op_nl_refresh", exclude_key="op_nl_open_exclude",
            )

            st.divider()
            st.markdown("#### 연계청구일자 이후 미로그인 고객")
            st.caption("연계청구일자 기준으로 청구시트 최종로그인일자가 없거나 연계청구일자 이전인 고객")
            _render_no_login_section(
                build_no_login_after_link_billing(hana_df_nl, billing_df_nl),
                year_key="no_login_link_year", owner_key="no_login_link_owner",
                label="연계청구 미로그인 고객", download_prefix="미로그인고객_연계청구일자",
                exclude_key="op_nl_link_exclude",
            )

            st.divider()
            st.markdown("#### 로그인 100회 이상 · 미이체 고객")
            st.caption("청구시트 기준 로그인 100회 이상이지만 이체 이력이 없는 고객 (해지 제외)")
            _render_high_login_no_transfer(
                build_high_login_no_transfer(hana_df_nl, billing_df_nl),
                owner_key="op_high_login_owner",
                download_prefix="미이체고객",
            )


def show_weekly_report_admin():
    db = load_db(WEEKLY_REPORT_FILE, {})
    rows = []
    for entries in db.values():
        rows.extend(entries if isinstance(entries, list) else [])
    title_col, refresh_col = st.columns([5, 1])
    with title_col:
        st.markdown("### 주간보고 취합")
    with refresh_col:
        if st.button("새로고침", key="weekly_admin_refresh", use_container_width=True):
            try:
                st.session_state.hana_sheet_df = read_google_csv(st.session_state.get("url_hana", DEFAULT_URL_HANA), header=2, force_refresh=True)
                st.session_state.weekly_hana_loaded_at = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d %H:%M:%S")
                st.success("새로고침 완료")
            except Exception as e:
                st.warning(f"새로고침 실패: {e}")
            st.rerun()
    raw_df_for_period = pd.DataFrame(rows) if rows else pd.DataFrame()
    default_snapshot_end = str(weekly_prev_friday())
    if "보고종료일" in raw_df_for_period.columns and not raw_df_for_period.empty:
        ends = [str(v) for v in raw_df_for_period["보고종료일"].dropna().tolist() if str(v).strip()]
        if ends:
            default_snapshot_end = max(ends)
    hana_status_df = load_weekly_hana_for_status()
    render_weekly_front_status_tables(hana_status_df, default_snapshot_end)
    st.markdown("---")

    if not rows:
        st.info("취합할 주간보고 이력이 없습니다.")
        return
    df = pd.DataFrame(rows)
    categories = ["전체", "개설대기", "개설진행", "개설완료", "연계대기", "연계진행", "연계완료", "운영부문"]
    c1, c2 = st.columns(2)
    with c1:
        category = st.selectbox("카테고리", categories, key="weekly_admin_category")
    with c2:
        staff_options = ["전체"] + sorted(df["담당자"].dropna().astype(str).unique().tolist()) if "담당자" in df.columns else ["전체"]
        staff = st.selectbox("담당자", staff_options, key="weekly_admin_staff")
    if category != "전체" and "카테고리" in df.columns:
        df = df[df["카테고리"] == category]
    if staff != "전체" and "담당자" in df.columns:
        df = df[df["담당자"] == staff]
    display_df = df.copy()
    if "일자" in display_df.columns:
        display_df["구축/피드백 예정일"] = display_df["일자"]
    show_cols = ["보고시작일", "보고종료일", "카테고리", "구분", "고객명", "접수일자", "이슈", "상태", "담당자", "구축/피드백 예정일", "작성시각"]
    show_cols = [c for c in show_cols if c in display_df.columns]
    render_plain_html_table(
        display_df[show_cols].sort_values("작성시각", ascending=False),
        max_rows=1000,
        center_align=False,
        merge_cols=["보고시작일", "보고종료일", "카테고리", "구분"],
    )

    st.markdown("<div style='height:10px'></div>", unsafe_allow_html=True)
    try:
        hana_for_ppt = None
        try:
            hana_for_ppt = read_google_csv(st.session_state.get("url_hana", DEFAULT_URL_HANA), header=2)
        except Exception:
            hana_for_ppt = st.session_state.get("hana_sheet_df")
        week_start = ""
        week_end = ""
        if "보고시작일" in df.columns and "보고종료일" in df.columns and not df.empty:
            starts = [str(v) for v in df["보고시작일"].dropna().tolist() if str(v).strip()]
            ends = [str(v) for v in df["보고종료일"].dropna().tolist() if str(v).strip()]
            week_start = min(starts) if starts else ""
            week_end = max(ends) if ends else ""
        ppt_bytes = build_weekly_report_ppt_bytes(df, week_start, week_end, hana_for_ppt)
        file_period = f"{week_start}_{week_end}".replace("-", "") if week_start and week_end else (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m%d")
        st.download_button(
            "주간보고 PPT 다운로드",
            data=ppt_bytes,
            file_name=f"주간보고_통합CMS고객_개설운영_주간보고_{file_period}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            on_click=save_weekly_status_snapshot,
            args=(hana_for_ppt, week_end or default_snapshot_end),
        )
    except Exception as e:
        st.button("주간보고 PPT 다운로드", use_container_width=True, disabled=True)
        st.caption(f"PPT 생성 준비 중 오류: {e}")


def show_dashboard():
    import plotly.graph_objects as go

    # ── 대시보드 진입 시 캐시된 구글시트 데이터를 재사용 ─────────────
    try:
        with st.spinner("데이터 불러오는 중..."):
            hana_raw = read_google_csv(st.session_state.get("url_hana", DEFAULT_URL_HANA), header=2)
            st.session_state.hana_sheet_df = hana_raw
            raw_act = read_google_csv(st.session_state.get("url_analysis", DEFAULT_URL_ANALYSIS))
            st.session_state.analysis_lookup_df = raw_act
            billing_raw = read_google_csv(st.session_state.get("url_hana_billing", DEFAULT_URL_HANA_BILLING))
            billing_raw = billing_raw.dropna(how="all").reset_index(drop=True)
            st.session_state.hana_billing_df = billing_raw
    except Exception as e:
        st.warning(f"데이터 로드 실패: {e}")

    df_hana = st.session_state.get("hana_sheet_df")
    user_name = st.session_state.user_name

    if df_hana is None or df_hana.empty:
        st.info("📂 데이터를 불러올 수 없습니다. [구글 스트레드시트 연동] 메뉴에서 URL을 확인해주세요.")
        return

    # ── 하나은행 데이터 전처리 (당월용) ───────────────
    df_hana = df_hana.copy()
    u_col = "담당자"
    gaeseol_date_col = "개설/이행일"
    yeonge_date_col = "연계일자"
    gubun_col = "구축구분"

    if u_col not in df_hana.columns or gaeseol_date_col not in df_hana.columns:
        st.error("데이터 컬럼을 찾을 수 없습니다. (담당자, 개설/이행일 필요)")
        return

    df_hana[gaeseol_date_col] = pd.to_datetime(
        df_hana[gaeseol_date_col].astype(str).str.strip().str[:8], format="%Y%m%d", errors="coerce"
    )
    if yeonge_date_col in df_hana.columns:
        df_hana[yeonge_date_col] = pd.to_datetime(
            df_hana[yeonge_date_col].astype(str).str.strip().str[:8], format="%Y%m%d", errors="coerce"
        )
    df_user_hana = df_hana[df_hana[u_col].astype(str).str.strip() == user_name].copy()

    # ── 하나지사 활동이력 전처리 (전월/전년동월용) ────
    df_user_act = None
    df_act = None
    act_u_col = None
    act_date_col = "활동일"
    act_d_col = "활동상세"
    df_activity = st.session_state.get("analysis_lookup_df")
    if df_activity is not None and not df_activity.empty:
        df_act = df_activity.copy()
        df_act.columns = [str(c).strip() for c in df_act.columns]
        act_u_col = next((c for c in df_act.columns if c.strip() in ["등록자", "담당자", "성명"]), None)
        act_date_col = next((c for c in df_act.columns if "활동일" in c.replace(" ", "") or "일자" in c.replace(" ", "")), None)
        act_d_col = next((c for c in df_act.columns if c.strip() in ["활동상세", "활동내용"]), None)
        if act_u_col and act_date_col and act_d_col:
            df_act[act_date_col] = pd.to_datetime(df_act[act_date_col], errors="coerce")
            df_user_act = df_act[df_act[act_u_col].astype(str).str.strip() == user_name].dropna(subset=[act_date_col]).copy()

    now = datetime.utcnow() + timedelta(hours=9)
    curr_ym = now.strftime("%Y-%m")
    prev_dt = (now.replace(day=1) - timedelta(days=1))
    prev_ym = prev_dt.strftime("%Y-%m")
    prev_year_ym = f"{now.year - 1}-{now.month:02d}"

    # ── 공통 계산 헬퍼 (유저 이름 인자) ──────────────────
    def calc_hana_for(uname, ym):
        df_u = df_hana[df_hana[u_col].astype(str).str.strip() == uname]
        gdf  = df_u[df_u[gaeseol_date_col].dt.strftime("%Y-%m") == ym]
        o = int((gdf[gubun_col].astype(str).str.strip() == "신규").sum()) if gubun_col in gdf.columns else len(gdf)
        ldf = df_u[df_u[yeonge_date_col].dt.strftime("%Y-%m") == ym] if yeonge_date_col in df_u.columns else df_u.iloc[0:0]
        l   = len(ldf)
        v = int((gdf[gubun_col].astype(str).str.strip() == "이행").sum()) if gubun_col in gdf.columns else 0
        o_p, l_p, v_p = o * 90, l * 120, v * 30
        total = min(2800, o_p + l_p + v_p)
        return {"개설건수": o, "연계건수": l, "운영건수": v, "개설포인트": o_p, "연계포인트": l_p, "운영포인트": v_p, "합계포인트": total}

    def calc_act_for(uname, ym):
        empty = {"개설건수": 0, "연계건수": 0, "운영건수": 0, "합계포인트": 0}
        try:
            if df_activity is None or act_date_col is None or act_d_col is None or act_u_col is None:
                return empty
            df_u = df_act[df_act[act_u_col].astype(str).str.strip() == uname]
            df_m = df_u[df_u[act_date_col].dt.strftime("%Y-%m") == ym]
            if df_m.empty:
                return empty
            o = int(df_m[act_d_col].astype(str).str.contains("개설", na=False).sum())
            l = int(df_m[act_d_col].astype(str).str.contains("연계", na=False).sum())
            v = int(df_m[act_d_col].astype(str).str.contains("운영|방문|점검", na=False).sum())
            o_p, l_p, v_p = o * 90, l * 120, v * 30
            total = min(2800, o_p + l_p + v_p)
            return {"개설건수": o, "연계건수": l, "운영건수": v, "합계포인트": total}
        except Exception:
            return empty

    def filter_month_gaeseol(df, ym):
        return df[df[gaeseol_date_col].dt.strftime("%Y-%m") == ym].copy()

    def filter_month_yeonge(df, ym):
        if yeonge_date_col not in df.columns:
            return df.iloc[0:0]
        return df[df[yeonge_date_col].dt.strftime("%Y-%m") == ym].copy()

    def calc_points_hana(ym):
        return calc_hana_for(user_name, ym)

    def calc_points_activity(ym):
        empty = {"개설건수": 0, "연계건수": 0, "운영건수": 0, "개설포인트": 0, "연계포인트": 0, "운영포인트": 0, "합계포인트": 0}
        try:
            if df_user_act is None or df_user_act.empty or act_date_col is None or act_d_col is None:
                return empty
            df_m = df_user_act[df_user_act[act_date_col].dt.strftime("%Y-%m") == ym].copy()
            if df_m.empty:
                return empty
            o = int(df_m[act_d_col].astype(str).str.contains("개설", na=False).sum())
            l = int(df_m[act_d_col].astype(str).str.contains("연계", na=False).sum())
            v = int(df_m[act_d_col].astype(str).str.contains("운영|방문|점검", na=False).sum())
            o_p, l_p, v_p = o * 90, l * 120, v * 30
            total = min(2800, o_p + l_p + v_p)
            return {"개설건수": o, "연계건수": l, "운영건수": v, "개설포인트": o_p, "연계포인트": l_p, "운영포인트": v_p, "합계포인트": total}
        except Exception:
            return empty

    # ── 관리자: 전체 직원 현황 2열 그리드 ────────────────
    if st.session_state.user_role == "관리자":
        st.markdown(f"### 전체 직원 {curr_ym} 실적 현황")
        staff_items = [
            (info.get("name", ""), info.get("rank", "직원"))
            for uid, info in st.session_state.user_db.items()
            if uid != "1"
            and info.get("name")
            and info.get("access") == "허용"
            and info.get("rank") != "부서장"
            and info.get("dept_type", "사업부") == "C&S"
            and info.get("staff_type", "정규직") != "파견직"
        ]
        staff_items = sorted(staff_items, key=lambda item: (_RANK_ORDER.get(item[1], 9), item[0]))
        cards = []
        for uname, _rank in staff_items:
            c = calc_hana_for(uname, curr_ym)
            p = calc_act_for(uname, prev_ym)
            delta = c["합계포인트"] - p["합계포인트"]
            cards.append((uname, c, p, delta))

        for i in range(0, len(cards), 2):
            cols = st.columns(2)
            for j, col in enumerate(cols):
                if i + j >= len(cards):
                    break
                uname, c, p, delta = cards[i + j]
                with col:
                    delta_cls = "pms-delta-up" if delta >= 0 else "pms-delta-down"
                    delta_sym = "▲" if delta >= 0 else "▼"
                    st.markdown(f"""
                    <div class="pms-staff-card">
                        <div class="pms-card-name">👤 {uname}</div>
                        <div class="pms-card-stats">
                            <span>개설 <b>{c['개설건수']}건</b></span>
                            <span>연계 <b>{c['연계건수']}건</b></span>
                            <span>이행 <b>{c['운영건수']}건</b></span>
                        </div>
                        <div class="pms-card-points">
                            {c['합계포인트']:,} pt
                            <span class="pms-delta {delta_cls}">{delta_sym} {abs(delta):,}pt (전월비)</span>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
        st.markdown("---")
        st.markdown(f"### {user_name}님 상세 현황")

    curr = calc_points_hana(curr_ym)
    prev = calc_points_activity(prev_ym)
    py   = calc_points_activity(prev_year_ym)

    diff_prev = curr["합계포인트"] - prev["합계포인트"]
    diff_year = curr["합계포인트"] - py["합계포인트"]
    max_add = max(0, 2800 - curr["합계포인트"])

    # ── 요약 카드 ──────────────────────────────────────
    st.markdown(f"### {user_name}님의 {curr_ym} 실적 현황")
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("이번달 추정 포인트", f"{curr['합계포인트']:,} pt")
    c2.metric("전월 대비", f"{prev['합계포인트']:,} pt", delta=f"{diff_prev:+,} pt")
    c3.metric("전년 동월 대비", f"{py['합계포인트']:,} pt", delta=f"{diff_year:+,} pt")
    c4.metric("최대 추가 가능", f"{max_add:,} pt")

    # ── 이번달 고객사 개설/운영 현황 ──────────────────────
    curr_g = filter_month_gaeseol(df_user_hana, curr_ym)
    curr_y = filter_month_yeonge(df_user_hana, curr_ym)
    if not curr_g.empty or not curr_y.empty:
        st.markdown(f"**{curr_ym} 고객사 진행 현황**")
        _comp_col = next((c for c in df_user_hana.columns if any(k in c for k in ["업체명", "고객명", "상호"])), None)
        _dg, _dy = st.columns(2)
        with _dg:
            st.caption(f"개설/이행 ({len(curr_g)}건)")
            if not curr_g.empty:
                _dcols = [c for c in [_comp_col, gubun_col, gaeseol_date_col] if c and c in curr_g.columns]
                _show = curr_g[_dcols].copy() if _dcols else curr_g.iloc[:, :3].copy()
                _show[gaeseol_date_col] = _show[gaeseol_date_col].dt.strftime("%Y-%m-%d") if gaeseol_date_col in _show.columns else _show.get(gaeseol_date_col, "")
                st.dataframe(_show.reset_index(drop=True), use_container_width=True, hide_index=True)
            else:
                st.info("이번달 개설/이행 없음")
        with _dy:
            st.caption(f"연계 ({len(curr_y)}건)")
            if not curr_y.empty:
                _dcols2 = [c for c in [_comp_col, yeonge_date_col] if c and c in curr_y.columns]
                _show2 = curr_y[_dcols2].copy() if _dcols2 else curr_y.iloc[:, :2].copy()
                if yeonge_date_col in _show2.columns:
                    _show2[yeonge_date_col] = _show2[yeonge_date_col].dt.strftime("%Y-%m-%d")
                st.dataframe(_show2.reset_index(drop=True), use_container_width=True, hide_index=True)
            else:
                st.info("이번달 연계 없음")

    st.markdown("---")

    # ── 활동 유형별 포인트 차트 ────────────────────────
    col_l, col_r = st.columns(2)

    with col_l:
        st.markdown("**활동 유형별 포인트 비교 (당월 / 전월 / 전년 동월)**")
        categories = ["개설포인트", "연계포인트", "운영포인트"]
        labels = ["개설", "연계", "운영"]
        fig = go.Figure(data=[
            go.Bar(name=f"당월 ({curr_ym})",   x=labels, y=[curr[c] for c in categories], marker_color="#6366f1"),
            go.Bar(name=f"전월 ({prev_ym})",   x=labels, y=[prev[c] for c in categories], marker_color="#a78bfa"),
            go.Bar(name=f"전년 동월 ({prev_year_ym})", x=labels, y=[py[c]   for c in categories], marker_color="#94a3b8"),
        ])
        _txt0 = "#1e293b"
        fig.update_layout(**_chart_layout(height=300, barmode="group",
                          legend=dict(orientation="h", y=-0.3, font=dict(color=_txt0))))
        st.plotly_chart(fig, use_container_width=True, theme=None)

    with col_r:
        st.markdown("**이번달 포인트 구성**")
        pie_labels = ["개설포인트", "연계포인트", "운영포인트"]
        pie_values = [curr["개설포인트"], curr["연계포인트"], curr["운영포인트"]]
        if sum(pie_values) > 0:
            _txt  = "#1e293b"
            fig2 = go.Figure(go.Pie(
                labels=["개설", "연계", "운영"],
                values=pie_values,
                hole=0.4,
                marker_colors=["#6366f1", "#7C3AED", "#a78bfa"],
                textfont=dict(color=_txt),
            ))
            fig2.update_layout(**_chart_layout(height=300,
                legend=dict(font=dict(color=_txt))))
            st.plotly_chart(fig2, use_container_width=True, theme=None)
        else:
            st.info("이번달 집계된 활동이 없습니다.")

    st.markdown("---")

    # ── 일별 개설 추이 ─────────────────────────────────
    curr_df = filter_month_gaeseol(df_user_hana, curr_ym)
    if not curr_df.empty:
        st.markdown("**이번달 일별 개설 건수 추이**")
        daily = curr_df.groupby(curr_df[gaeseol_date_col].dt.strftime("%Y-%m-%d")).size().reset_index()
        daily.columns = ["날짜", "건수"]
        fig3 = go.Figure(go.Scatter(x=daily["날짜"], y=daily["건수"], mode="lines+markers",
                                    line=dict(color="#6366f1", width=2), marker=dict(size=7, color="#a78bfa")))
        fig3.update_layout(**_chart_layout(height=220, xaxis_title="날짜", yaxis_title="건수"))
        st.plotly_chart(fig3, use_container_width=True, theme=None)

    st.markdown("---")

    # ── 전년도 / 올해 업무별 월평균 ──────────────────────
    if df_user_act is not None and not df_user_act.empty:
        st.markdown("**📊 업무별 월평균 비교**")

        this_year  = str(now.year)
        last_year  = str(now.year - 1)

        def monthly_avg(year_str):
            df_y = df_user_act[df_user_act[act_date_col].dt.strftime("%Y") == year_str].copy()
            if df_y.empty:
                return {"개설": 0.0, "연계": 0.0, "운영": 0.0}
            months = df_y[act_date_col].dt.strftime("%Y-%m").unique()
            totals = {"개설": 0, "연계": 0, "운영": 0}
            for ym in months:
                dm = df_y[df_y[act_date_col].dt.strftime("%Y-%m") == ym]
                totals["개설"] += int(dm[act_d_col].astype(str).str.contains("개설", na=False).sum())
                totals["연계"] += int(dm[act_d_col].astype(str).str.contains("연계", na=False).sum())
                totals["운영"] += int(dm[act_d_col].astype(str).str.contains("운영|방문|점검", na=False).sum())
            n = len(months)
            return {k: round(v / n, 1) for k, v in totals.items()}

        avg_last = monthly_avg(last_year)
        avg_this = monthly_avg(this_year)

        avg_labels = ["개설", "연계", "운영"]
        _txt2  = "#1e293b"
        fig_avg = go.Figure(data=[
            go.Bar(name=f"{last_year}년 월평균", x=avg_labels,
                   y=[avg_last[k] for k in avg_labels], marker_color="#a78bfa"),
            go.Bar(name=f"{this_year}년 월평균", x=avg_labels,
                   y=[avg_this[k] for k in avg_labels], marker_color="#6366f1"),
        ])
        fig_avg.update_layout(**_chart_layout(height=260, barmode="group",
                              legend=dict(orientation="h", y=-0.3, font=dict(color=_txt2))))
        st.plotly_chart(fig_avg, use_container_width=True, theme=None)

        ca, cb, cc = st.columns(3)
        for col_ui, key in zip([ca, cb, cc], avg_labels):
            col_ui.metric(
                f"{key} 월평균",
                f"{avg_this[key]}건 (올해)",
                delta=f"{avg_this[key] - avg_last[key]:+.1f}건 vs {last_year}년",
            )

    st.markdown("---")

    # ── 실적 보완 가이드 ───────────────────────────────
    st.markdown("**💡 실적 보완 가이드**")
    ol_used = min(1000, curr["개설포인트"] + curr["연계포인트"])
    ol_remain = max(0, 1000 - (curr["개설포인트"] + curr["연계포인트"]))
    v_remain = max(0, 1800 - curr["운영포인트"])

    guides = []
    if ol_remain > 0:
        add_o = ol_remain // 90
        add_l = ol_remain // 120
        guides.append(f"📌 개설 추가 시 최대 **{add_o}건** ({ol_remain:,}pt 확보 가능, 건당 90pt)")
        guides.append(f"📌 연계 추가 시 최대 **{add_l}건** ({ol_remain:,}pt 확보 가능, 건당 120pt)")
    else:
        guides.append("✅ 개설·연계 포인트 한도(1,000pt) 달성!")

    if v_remain > 0:
        add_v = v_remain // 30
        guides.append(f"📌 운영·방문 추가 시 최대 **{add_v}건** ({v_remain:,}pt 확보 가능, 건당 30pt)")
    else:
        guides.append("✅ 운영 포인트 한도(1,800pt) 달성!")

    if max_add == 0:
        guides.append("🎉 최대 포인트(2,800pt) 달성! 수고하셨습니다.")

    for g in guides:
        st.markdown(f"- {g}")

    if ol_remain > 0:
        st.markdown("")
        st.info(
            "💡 개설·연계 포인트 확보가 어렵다면 **교차판매·활성화·VOC 활동**을 통해 부족한 포인트를 채울 수 있습니다.\n\n"
            "- **교차판매**: 타겟고객 선별·메일발송·방문설명회 등 영업활동 — 건당 2\~60pt, 월 최대 100\~제한없음\n"
            "- **활성화**: 이체·집금 활성화, 조회업무 활성화 — 건당 10\~30pt, 월 최대 100\~150pt\n"
            "- **VOC**: 고객 개선 아이디어 제출 — 건당 10pt, 월 최대 50pt\n\n"
            "위 활동들은 개설·연계 실적 없이도 단독으로 포인트를 적립할 수 있는 가장 현실적인 방법입니다."
        )
        st.info("💡 위 활동들은 **[업로드 및 실적 확인] 메뉴 → 추가 실적 입력** 표에서 건수를 직접 입력하여 포인트를 적립할 수 있습니다.")

def load_hana_billing_df(force_refresh=False):
    billing_raw = read_google_csv(
        st.session_state.get("url_hana_billing", DEFAULT_URL_HANA_BILLING),
        force_refresh=force_refresh,
    )
    billing_raw = billing_raw.dropna(how="all").reset_index(drop=True)
    st.session_state.hana_billing_df = billing_raw
    return billing_raw


def normalize_billing_customer_no(value):
    if is_blank_value(value):
        return ""
    text = str(value).strip()
    if re.match(r"^\d+\.0$", text):
        text = text[:-2]
    digits = re.sub(r"\D", "", text)
    if not digits:
        return ""
    return digits.zfill(9)


def normalize_biz_no(value):
    if is_blank_value(value):
        return ""
    text = str(value).strip()
    if re.match(r"^\d+\.0$", text):
        text = text[:-2]
    digits = re.sub(r"\D", "", text)
    return digits.zfill(10) if digits else ""


def format_yyyymmdd(value):
    parsed = parse_sheet_date(value)
    if pd.isna(parsed):
        return ""
    return parsed.strftime("%Y%m%d")


def series_yyyymm(series):
    parsed = pd.to_datetime(series.apply(parse_sheet_date), errors="coerce")
    return parsed.dt.strftime("%Y-%m")


def exact_col(df, names):
    normalized = {str(c).replace(" ", "").replace("\n", "").replace("\r", ""): c for c in df.columns}
    for name in names:
        found = normalized.get(str(name).replace(" ", "").replace("\n", "").replace("\r", ""))
        if found is not None:
            return found
    return None


def build_billing_lookup(billing_df):
    lookup = billing_df.copy()
    lookup.columns = [str(c).strip() for c in lookup.columns]

    customer_col = find_col(lookup, ["고객번호", "고객No", "고객no"])
    company_col = find_col(lookup, ["고객명", "업체명", "상호"])
    biz_col = find_col(lookup, ["사업자번호", "사업자등록번호"])
    first_login_col = exact_col(lookup, ["신규일자"]) or find_col(lookup, ["신규일자", "최초신규일자"])
    # "최종로그인" 단독 컬럼명도 커버 (일자 없는 변형 포함)
    last_login_col = find_col(lookup, ["최종로그인일자", "최근로그인일자", "최종로그인", "최근로그인", "로그인일자"])
    login_count_col  = find_col(lookup, ["로그인건수", "로그인횟수"])
    menu_click_col   = exact_col(lookup, ["메뉴사용"]) or find_col(lookup, ["메뉴사용", "메뉴클릭수"])
    transfer_col     = find_col(lookup, ["최종이체일자", "이체일자", "최종이체"])
    transfer_cnt_col = find_col(lookup, ["이체건수", "이체횟수"])

    if not customer_col:
        return pd.DataFrame(columns=["_고객번호", "_은행고객명", "_은행사업자번호", "_최초로그인", "_최종로그인", "_로그인횟수", "_메뉴클릭수", "_최종이체일자", "_이체건수"])

    rows = pd.DataFrame()
    rows["_고객번호"]     = lookup[customer_col].apply(normalize_billing_customer_no)
    rows["_은행고객명"]   = lookup[company_col].astype(str).str.strip() if company_col else ""
    rows["_은행사업자번호"] = lookup[biz_col].apply(normalize_biz_no) if biz_col else ""
    rows["_최초로그인"]   = lookup[first_login_col].apply(format_yyyymmdd) if first_login_col else ""
    rows["_최종로그인"]   = lookup[last_login_col].apply(format_yyyymmdd) if last_login_col else ""
    rows["_로그인횟수"]   = lookup[login_count_col].fillna("").astype(str).str.replace(r"\.0$", "", regex=True) if login_count_col else ""
    rows["_메뉴클릭수"]   = lookup[menu_click_col].fillna("").astype(str).str.replace(r"\.0$", "", regex=True) if menu_click_col else ""
    rows["_최종이체일자"] = lookup[transfer_col].apply(format_yyyymmdd) if transfer_col else ""
    rows["_이체건수"]     = lookup[transfer_cnt_col].fillna("").astype(str).str.replace(r"\.0$", "", regex=True) if transfer_cnt_col else ""
    rows = rows[rows["_고객번호"] != ""].drop_duplicates("_고객번호", keep="first")
    return rows


def add_bank_compare_columns(source_df, billing_lookup, source_company_col):
    result = source_df.merge(billing_lookup, on="_고객번호", how="left")
    result["청구시트 여부"] = np.where(result.get("_은행고객명", "").fillna("").astype(str).str.strip().ne(""), "있음", "없음")
    result["최초로그인"] = result.get("_최초로그인", "").fillna("")
    result["사업자번호_은행"] = result.get("_은행사업자번호", "").fillna("")
    result["최종로그인"] = result.get("_최종로그인", "").fillna("")
    result["로그인횟수"] = result.get("_로그인횟수", "").fillna("")
    result["메뉴클릭수"] = result.get("_메뉴클릭수", "").fillna("")
    result["청구원본 고객명"] = result[source_company_col].fillna("").astype(str) if source_company_col in result.columns else ""
    result["실적파일 고객명"] = result.get("_은행고객명", "").fillna("")
    return result


def build_open_billing_status(hana_df, billing_lookup, selected_month):
    hana = hana_df.copy()
    hana.columns = [str(c).strip() for c in hana.columns]

    customer_col = find_col(hana, ["고객번호"])
    biz_col = find_col(hana, ["사업자번호"])
    company_col = find_col(hana, ["고객명", "업체명", "상호"])
    build_type_col = find_col(hana, ["구축형"])
    receipt_col = find_col(hana, ["신규접수일"])
    open_date_col = find_col(hana, ["개설/이행일", "개설일", "이행일"])
    owner_col = find_col(hana, ["담당자"])
    manage_col = find_col(hana, ["관리구분"])
    open_status_col = find_col(hana, ["개설상태"])

    if not customer_col or not open_date_col:
        return pd.DataFrame()

    hana["_고객번호"] = hana[customer_col].apply(normalize_billing_customer_no)
    hana["_기준월"] = series_yyyymm(hana[open_date_col])
    mask = hana["_고객번호"].ne("") & hana["_기준월"].eq(selected_month)
    if manage_col and manage_col in hana.columns:
        mask &= ~hana[manage_col].astype(str).str.strip().isin(["해지", "취소"])
    if open_status_col and open_status_col in hana.columns:
        mask &= ~hana[open_status_col].astype(str).str.contains("취소|반려", na=False)

    base = hana[mask].copy().reset_index(drop=True)
    if base.empty:
        return pd.DataFrame()

    out = pd.DataFrame()
    out["순번"] = range(1, len(base) + 1)
    out["고객번호"] = base["_고객번호"]
    out["사업자번호"] = base[biz_col].apply(normalize_biz_no) if biz_col else ""
    out["업체명"] = base[company_col].fillna("").astype(str) if company_col else ""
    out["ERP연계 여부"] = base[build_type_col].fillna("").astype(str) if build_type_col else ""
    out["접수일자"] = base[receipt_col].apply(format_yyyymmdd) if receipt_col else ""
    out["구축일자"] = base[open_date_col].apply(format_yyyymmdd)
    out["방문일자"] = base[open_date_col].apply(format_yyyymmdd)
    out["담당자"] = base[owner_col].fillna("").astype(str) if owner_col else ""
    out["비고"] = ""
    out["_고객번호"] = base["_고객번호"]

    out = add_bank_compare_columns(out, billing_lookup, "업체명")
    out = out.rename(columns={"최종로그인": "최종로그인일자"})
    return out[["순번", "고객번호", "사업자번호", "업체명", "ERP연계 여부", "접수일자", "구축일자", "방문일자", "담당자", "비고", "최초로그인", "최종로그인일자", "로그인횟수", "메뉴클릭수", "청구시트 여부"]]


def build_link_billing_status(hana_df, billing_lookup, selected_month):
    hana = hana_df.copy()
    hana.columns = [str(c).strip() for c in hana.columns]

    customer_col = find_col(hana, ["고객번호"])
    biz_col = find_col(hana, ["사업자번호"])
    company_col = find_col(hana, ["고객명", "업체명", "상호"])
    add_receipt_col = find_col(hana, ["추가연계접수일"])
    owner_col = find_col(hana, ["담당자"])
    open_date_col = find_col(hana, ["개설/이행일", "개설일", "이행일"])
    link_date_col = find_col(hana, ["연계일자"])
    link_billing_col = find_col(hana, ["연계청구일자", "연계청구일", "청구일자"])
    link_status_col = find_col(hana, ["연계상태"])

    if not customer_col:
        return pd.DataFrame()

    hana["_고객번호"] = hana[customer_col].apply(normalize_billing_customer_no)
    mask = hana["_고객번호"].ne("")
    if link_status_col and link_status_col in hana.columns:
        mask &= hana[link_status_col].astype(str).str.strip().str.contains("연계완료", na=False)
    else:
        return pd.DataFrame()

    base = hana[mask].copy().reset_index(drop=True)
    if base.empty:
        return pd.DataFrame()

    link_type = pd.Series("신규", index=base.index)
    if add_receipt_col and add_receipt_col in base.columns:
        link_type = np.where(base[add_receipt_col].apply(format_yyyymmdd).astype(str).ne(""), "추가", "신규")

    out = pd.DataFrame()
    out["순서"] = range(1, len(base) + 1)
    out["고객번호"] = base["_고객번호"]
    out["사업자번호"] = base[biz_col].apply(normalize_biz_no) if biz_col else ""
    out["업체명"] = base[company_col].fillna("").astype(str) if company_col else ""
    out["구분"] = link_type
    out["추가연계신청일자"] = base[add_receipt_col].apply(format_yyyymmdd) if add_receipt_col else ""
    out["담당자"] = base[owner_col].fillna("").astype(str) if owner_col else ""
    out["구축일"] = base[open_date_col].apply(format_yyyymmdd) if open_date_col else ""
    out["연계시작일자"] = base[link_date_col].apply(format_yyyymmdd) if link_date_col else ""
    out["은행연계완료일자"] = base[link_billing_col].apply(format_yyyymmdd) if link_billing_col else ""
    out["수령여부"] = ""
    out["비고"] = ""
    out["_고객번호"] = base["_고객번호"]

    out = add_bank_compare_columns(out, billing_lookup, "업체명")
    out = out.rename(columns={"사업자번호_은행": "은행 사업자번호"})
    return out[["순서", "고객번호", "사업자번호", "업체명", "구분", "추가연계신청일자", "담당자", "구축일", "연계시작일자", "은행연계완료일자", "수령여부", "비고", "청구시트 여부", "최초로그인", "은행 사업자번호", "최종로그인", "로그인횟수", "메뉴클릭수", "청구원본 고객명", "실적파일 고객명"]]


def build_billing_status_excel_bytes(open_df, link_df):
    return dataframe_to_excel_bytes({
        "개설현황": open_df,
        "연계현황": link_df,
    })


def _merge_with_billing_lookup(hana_df, billing_df, ref_date_col_keys, ref_date_label, extra_hana_col_keys=None):
    """
    기존 build_billing_lookup을 재사용해 hana_df와 billing_df를 고객번호로 대사한 뒤
    ref_date_label 이후 최종로그인일자가 없는 고객을 반환하는 공통 헬퍼.
    """
    hana = hana_df.copy()
    hana.columns = [str(c).strip() for c in hana.columns]

    hana_cust_col  = find_col(hana, ["고객번호"])
    ref_date_col   = find_col(hana, ref_date_col_keys)
    company_col    = find_col(hana, ["고객명", "업체명", "상호"])
    owner_col      = find_col(hana, ["담당자"])
    build_type_col = find_col(hana, ["구축형"])
    extra_col      = find_col(hana, extra_hana_col_keys) if extra_hana_col_keys else None
    end_col        = find_col(hana, ["해지일자", "해지일", "해약일"])

    if not all([hana_cust_col, ref_date_col]):
        return pd.DataFrame()

    hana["_고객번호"] = hana[hana_cust_col].apply(normalize_billing_customer_no)
    hana["_ref_dt"]   = pd.to_datetime(hana[ref_date_col].map(parse_sheet_date), errors="coerce")
    # 해지일자 있는 고객 제외
    if end_col:
        hana["_end_dt"] = pd.to_datetime(hana[end_col].map(parse_sheet_date), errors="coerce")
        hana = hana[hana["_end_dt"].isna()].copy()
    hana_valid = hana[hana["_고객번호"].ne("") & hana["_ref_dt"].notna()].copy()

    # 기존 검증된 build_billing_lookup 사용 (고객번호 매칭 동일 보장)
    bill_lkp = build_billing_lookup(billing_df)
    merged = hana_valid.merge(bill_lkp, on="_고객번호", how="left")

    # _최종로그인은 format_yyyymmdd 결과("YYYYMMDD" 또는 "") → datetime 변환
    merged["_last_dt"] = merged["_최종로그인"].apply(
        lambda v: pd.NaT if (v == "" or pd.isna(v))
        else pd.to_datetime(str(v).strip(), format="%Y%m%d", errors="coerce")
    )

    mask = merged["_last_dt"].isna() | (merged["_last_dt"] <= merged["_ref_dt"])
    result = merged[mask].copy().reset_index(drop=True)
    if result.empty:
        return pd.DataFrame()

    # 직전 3개월 구간 계산 (당월 제외 → 이전 3개월)
    today_kst   = datetime.utcnow() + timedelta(hours=9)
    month_start = today_kst.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
    three_ago   = month_start - pd.DateOffset(months=3)
    m_labels    = [(three_ago + pd.DateOffset(months=i)).strftime("%m월") for i in range(3)]
    prev3_label = f"직전3개월({m_labels[0]}~{m_labels[2]})미로그인"

    out = pd.DataFrame()
    out["순번"]          = range(1, len(result) + 1)
    out["고객번호"]      = result["_고객번호"]
    out["고객명"]        = result[company_col].fillna("").astype(str)    if company_col    else ""
    out["담당자"]        = result[owner_col].fillna("").astype(str)      if owner_col      else ""
    out["구축형"]        = result[build_type_col].fillna("").astype(str) if build_type_col else ""
    out[ref_date_label]  = result["_ref_dt"].dt.strftime("%Y-%m-%d")
    if extra_col:
        out[extra_hana_col_keys[0]] = result[extra_col].fillna("").astype(str)
    out["최종로그인일자"] = result["_last_dt"].dt.strftime("%Y-%m-%d").fillna("없음")
    # 직전 3개월 내 로그인 이력 없음 여부 (개설/이행일 이후이면서 직전 3개월 내 미로그인)
    three_ago_ts = pd.Timestamp(three_ago)
    out[prev3_label] = result.apply(
        lambda r: "미로그인" if (
            pd.isna(r["_last_dt"]) or r["_last_dt"] < three_ago_ts
        ) and (pd.notna(r["_ref_dt"]) and r["_ref_dt"] < three_ago_ts)
        else "",
        axis=1,
    )
    return out


def build_no_login_after_open(hana_df, billing_df):
    """구글시트 개설/이행일 이후 청구시트에 최종로그인일자가 없는 고객 목록 반환."""
    return _merge_with_billing_lookup(
        hana_df, billing_df,
        ref_date_col_keys=["개설/이행일", "개설일", "이행일"],
        ref_date_label="개설/이행일",
        extra_hana_col_keys=["개설상태"],
    )


def build_no_login_after_link_billing(hana_df, billing_df):
    """구글시트 연계청구일자 이후 청구시트에 최종로그인일자가 없는 고객 목록 반환."""
    return _merge_with_billing_lookup(
        hana_df, billing_df,
        ref_date_col_keys=["연계청구일자", "연계청구일", "청구일자"],
        ref_date_label="연계청구일자",
        extra_hana_col_keys=["연계상태"],
    )


def build_high_login_no_transfer(hana_df, billing_df, login_threshold=100):
    """로그인 100회 이상이지만 이체 이력이 없는 고객 목록 반환."""
    hana = hana_df.copy()
    hana.columns = [str(c).strip() for c in hana.columns]

    hana_cust_col  = find_col(hana, ["고객번호"])
    company_col    = find_col(hana, ["고객명", "업체명", "상호"])
    owner_col      = find_col(hana, ["담당자"])
    build_type_col = find_col(hana, ["구축형"])
    open_date_col  = find_col(hana, ["개설/이행일", "개설일", "이행일"])
    open_status_col= find_col(hana, ["개설상태"])
    end_col        = find_col(hana, ["해지일자", "해지일", "해약일"])

    if not hana_cust_col:
        return pd.DataFrame()

    hana["_고객번호"] = hana[hana_cust_col].apply(normalize_billing_customer_no)

    # 해지 고객 제외
    if end_col:
        hana["_end_dt"] = pd.to_datetime(hana[end_col].map(parse_sheet_date), errors="coerce")
        hana = hana[hana["_end_dt"].isna()].copy()

    hana_valid = hana[hana["_고객번호"].ne("")].copy()

    bill_lkp = build_billing_lookup(billing_df)
    merged = hana_valid.merge(bill_lkp, on="_고객번호", how="inner")  # 청구시트에 있는 고객만

    # 로그인횟수 숫자 변환
    merged["_login_cnt"] = pd.to_numeric(
        merged["_로그인횟수"].replace("", "0"), errors="coerce"
    ).fillna(0)

    # 조건: 로그인 >= threshold AND 최종이체일자 없음
    mask = (merged["_login_cnt"] >= login_threshold) & (merged["_최종이체일자"].eq("") | merged["_최종이체일자"].isna())
    result = merged[mask].copy().reset_index(drop=True)
    if result.empty:
        return pd.DataFrame()

    out = pd.DataFrame()
    out["순번"]        = range(1, len(result) + 1)
    out["고객번호"]    = result["_고객번호"]
    out["고객명"]      = result[company_col].fillna("").astype(str)    if company_col     else ""
    out["담당자"]      = result[owner_col].fillna("").astype(str)      if owner_col       else ""
    out["구축형"]      = result[build_type_col].fillna("").astype(str) if build_type_col  else ""
    out["개설상태"]    = result[open_status_col].fillna("").astype(str) if open_status_col else ""
    out["개설/이행일"] = result[open_date_col].apply(parse_sheet_date).apply(
        lambda d: d.strftime("%Y-%m-%d") if pd.notna(d) else ""
    ) if open_date_col else ""
    out["로그인횟수"]   = result["_로그인횟수"]
    out["최종로그인일자"] = result["_최종로그인"].apply(
        lambda v: pd.to_datetime(v, format="%Y%m%d", errors="coerce").strftime("%Y-%m-%d")
        if v and not pd.isna(v) else "없음"
    )
    out["최종이체일자"] = result["_최종이체일자"].apply(
        lambda v: pd.to_datetime(v, format="%Y%m%d", errors="coerce").strftime("%Y-%m-%d")
        if v and not pd.isna(v) else "없음"
    )
    return out


def _render_high_login_no_transfer(df, owner_key, download_prefix, year_key=None):
    """로그인 100회+ 미이체 고객 렌더링 헬퍼."""
    if df.empty:
        st.info("해당 조건의 고객이 없습니다.")
        return
    date_col = "개설/이행일" if "개설/이행일" in df.columns else None
    years  = sorted(df[date_col].str[:4].dropna().unique().tolist(), reverse=True) if date_col else []
    owners = sorted(df["담당자"].dropna().unique().tolist())
    fc1, fc2, _ = st.columns([2, 2, 6])
    with fc1:
        sel_year = st.selectbox("연도 조회", ["전체"] + years, key=year_key or f"{owner_key}_year")
    with fc2:
        sel_owner = st.selectbox("담당자 조회", ["전체"] + owners, key=owner_key)
    filtered = df.copy()
    if sel_year != "전체" and date_col:
        filtered = filtered[filtered[date_col].str.startswith(sel_year)]
    if sel_owner != "전체":
        filtered = filtered[filtered["담당자"] == sel_owner]
    filtered = filtered.reset_index(drop=True)
    filtered["순번"] = range(1, len(filtered) + 1)
    st.metric("미이체 고객 수", f"{len(filtered):,}건")
    render_plain_html_table(filtered, max_rows=500, center_align=False)
    today_str = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m%d")
    st.download_button(
        "미이체 고객 다운로드",
        data=filtered.to_csv(index=False, encoding="utf-8-sig").encode("utf-8-sig"),
        file_name=f"{download_prefix}_{sel_year}_{today_str}.csv",
        mime="text/csv",
        use_container_width=True,
    )


def _render_no_login_section(df, year_key, owner_key, label, download_prefix,
                             refresh_key=None, exclude_key=None):
    """미로그인 고객 공통 렌더링 헬퍼 (연도·담당자 필터 + 표 + 다운로드)."""
    if df.empty:
        st.info("해당 조건의 고객이 없습니다.")
        return
    date_col = "개설/이행일" if "개설/이행일" in df.columns else "연계청구일자"
    years  = sorted(df[date_col].str[:4].dropna().unique().tolist(), reverse=True)
    owners = sorted(df["담당자"].dropna().unique().tolist())

    fc1, fc2, fc3, fc4 = st.columns([2, 2, 1, 5])
    with fc1:
        sel_year  = st.selectbox("연도 조회",  ["전체"] + years,  key=year_key)
    with fc2:
        sel_owner = st.selectbox("담당자 조회", ["전체"] + owners, key=owner_key)
    with fc3:
        if refresh_key and st.button("새로고침", key=refresh_key, use_container_width=True):
            st.session_state[f"_nl_rf_{refresh_key}"] = True
            st.rerun()
    exclude_current = False
    with fc4:
        if exclude_key:
            current_ym = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m")
            exclude_current = st.checkbox(f"당월({current_ym}) 개설 제외", key=exclude_key)

    filtered = df.copy()
    if sel_year  != "전체":
        filtered = filtered[filtered[date_col].str.startswith(sel_year)]
    if sel_owner != "전체":
        filtered = filtered[filtered["담당자"] == sel_owner]
    if exclude_current:
        filtered = filtered[~filtered[date_col].str.startswith(current_ym)]

    filtered = filtered.reset_index(drop=True)
    filtered["순번"] = range(1, len(filtered) + 1)
    st.metric(f"{label} 수", f"{len(filtered):,}건")
    render_plain_html_table(filtered, max_rows=500, center_align=False)
    today_str = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y%m%d")
    st.download_button(
        f"{label} 다운로드",
        data=filtered.to_csv(index=False, encoding="utf-8-sig").encode("utf-8-sig"),
        file_name=f"{download_prefix}_{sel_year}_{today_str}.csv",
        mime="text/csv",
        use_container_width=True,
    )


def show_billing_materials():
    title_col, refresh_col = st.columns([5, 1])
    with title_col:
        st.markdown("### 청구자료 작성")
    with refresh_col:
        refresh = st.button("새로고침", key="billing_refresh", use_container_width=True)

    try:
        hana_df = read_google_csv(
            st.session_state.get("url_hana", DEFAULT_URL_HANA),
            header=2,
            force_refresh=refresh,
        )
        billing_df = load_hana_billing_df(force_refresh=refresh)
        if refresh:
            st.success("새로고침 완료")
    except Exception as e:
        st.error(f"청구자료 생성용 구글 시트를 불러오지 못했습니다: {e}")
        st.info("[구글 스트레드시트 연동] 메뉴에서 하나은행 구글 시트 CSV URL과 하나은행 청구 시트 CSV URL을 확인해주세요.")
        return

    if hana_df is None or hana_df.empty or billing_df is None or billing_df.empty:
        st.info("청구자료로 작성할 데이터가 없습니다.")
        return

    hana_df = hana_df.dropna(how="all").reset_index(drop=True)
    billing_df = billing_df.dropna(how="all").reset_index(drop=True)

    hana_cols = [str(c).strip() for c in hana_df.columns]
    hana_df.columns = hana_cols
    open_date_col = find_col(hana_df, ["개설/이행일", "개설일", "이행일"])
    link_billing_col = find_col(hana_df, ["연계청구일자", "연계청구일", "청구일자"])
    link_date_col = find_col(hana_df, ["연계일자"])

    month_candidates = []
    for date_col in [open_date_col, link_billing_col, link_date_col]:
        if date_col and date_col in hana_df.columns:
            month_candidates.extend(series_yyyymm(hana_df[date_col]).dropna().astype(str).tolist())
    month_values = sorted([m for m in set(month_candidates) if m and m != "NaT"], reverse=True)
    month_options = ["전체"] + month_values
    current_month = (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m")
    default_month_index = month_options.index(current_month) if current_month in month_options else 0

    c1, c2 = st.columns([1, 2])
    with c1:
        selected_month = st.selectbox("청구 기준월", month_options, index=default_month_index, key="billing_month")
    with c2:
        st.caption("하나은행 구글 시트와 하나은행 청구 시트를 고객번호 기준으로 대사해 개설/연계현황을 생성합니다.")

    if selected_month == "전체":
        st.info("다운로드할 청구 기준월을 선택해주세요.")
        return

    billing_lookup = build_billing_lookup(billing_df)
    open_df = build_open_billing_status(hana_df, billing_lookup, selected_month)
    link_df = build_link_billing_status(hana_df, billing_lookup, selected_month)

    open_missing_billing = 0
    if "청구시트 여부" in open_df.columns:
        open_missing_billing = int(open_df["청구시트 여부"].astype(str).eq("없음").sum())

    m1, m2, m3, m4 = st.columns(4)
    m1.metric("개설현황", f"{len(open_df):,}건")
    m2.metric("연계현황", f"{len(link_df):,}건")
    m3.metric("개설 청구시트 없음", f"{open_missing_billing:,}건")
    m4.metric("기준월", selected_month)

    tab_open, tab_link = st.tabs(["개설현황", "연계현황"])
    with tab_open:
        selected_open_df = open_df.copy()
        if open_df.empty:
            render_plain_html_table(open_df, max_rows=1000, center_align=False)
        else:
            open_select_df = open_df.copy()
            open_select_df["다운로드"] = False
            edited_open_df = st.data_editor(
                open_select_df,
                use_container_width=True,
                hide_index=True,
                num_rows="fixed",
                disabled=[c for c in open_select_df.columns if c != "다운로드"],
                column_config={"다운로드": st.column_config.CheckboxColumn("다운로드", default=False, width="small")},
                key=f"billing_open_select_{selected_month}",
            )
            selected_open_df = edited_open_df[edited_open_df["다운로드"]].drop(columns=["다운로드"], errors="ignore")
            st.caption(f"다운로드 선택: {len(selected_open_df):,}건")
    with tab_link:
        selected_link_df = link_df.copy()
        if link_df.empty:
            render_plain_html_table(link_df, max_rows=1000, center_align=False)
        else:
            link_select_df = link_df.copy()
            link_select_df["다운로드"] = False
            edited_link_df = st.data_editor(
                link_select_df,
                use_container_width=True,
                hide_index=True,
                num_rows="fixed",
                disabled=[c for c in link_select_df.columns if c != "다운로드"],
                column_config={"다운로드": st.column_config.CheckboxColumn("다운로드", default=False, width="small")},
                key=f"billing_link_select_{selected_month}",
            )
            selected_link_df = edited_link_df[edited_link_df["다운로드"]].drop(columns=["다운로드"], errors="ignore")
            st.caption(f"다운로드 선택: {len(selected_link_df):,}건")

    file_month = selected_month.replace("-", "")
    st.download_button(
        "개설/연계현황 다운로드",
        data=build_billing_status_excel_bytes(selected_open_df, selected_link_df),
        file_name=f"개설_연계현황_{file_month}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True,
        disabled=selected_open_df.empty and selected_link_df.empty,
    )


def show_google_sync():
    st.session_state.url_sync = st.text_input("본사 구글 시트 CSV URL", value=st.session_state.url_sync)
    st.session_state.url_analysis = st.text_input("하나지사 활동이력 구글 시트 CSV URL", value=st.session_state.url_analysis)
    st.session_state.url_hana = st.text_input("하나은행 구글 시트 CSV URL", value=st.session_state.url_hana)
    st.session_state.url_hana_billing = st.text_input("하나은행 청구 시트 CSV URL", value=st.session_state.url_hana_billing)
    st.session_state.url_hana_performance = st.text_input("하나지사 실적관리 CSV URL", value=st.session_state.get("url_hana_performance", DEFAULT_URL_HANA_PERFORMANCE))

    st.divider()
    st.markdown("##### DART 전자공시 API 키")
    st.caption("opendart.fss.or.kr 에서 무료 발급 — 유통활동 추천 보정에 사용됩니다.")
    dart_input = st.text_input(
        "DART API 키",
        value=st.session_state.get("dart_api_key", ""),
        type="password",
        placeholder="발급받은 40자리 API 키를 입력하세요",
        key="dart_api_key_input",
    )
    if dart_input != st.session_state.get("dart_api_key", ""):
        st.session_state.dart_api_key = dart_input

    if st.button("데이터 저장", type="primary"):
        errors = []
        try:
            load_csv_to_state("url_sync", "temp_cloud_df", force_refresh=True)
            st.session_state.cloud_sheet_df = st.session_state.temp_cloud_df
        except Exception:
            errors.append("본사 구글 시트")
        try:
            load_csv_to_state("url_analysis", "analysis_lookup_df", force_refresh=True)
        except Exception:
            errors.append("하나지사 활동이력 구글 시트")
        try:
            hana_raw = read_google_csv(st.session_state.url_hana, header=2, force_refresh=True)
            hana_raw = hana_raw.dropna(how="all").reset_index(drop=True)
            st.session_state.hana_sheet_df = hana_raw
        except Exception:
            errors.append("하나은행 구글 시트")
        try:
            billing_raw = read_google_csv(st.session_state.url_hana_billing, force_refresh=True)
            billing_raw = billing_raw.dropna(how="all").reset_index(drop=True)
            st.session_state.hana_billing_df = billing_raw
        except Exception:
            errors.append("하나은행 청구 시트")
        try:
            perf_raw = read_google_csv(st.session_state.get("url_hana_performance", DEFAULT_URL_HANA_PERFORMANCE), force_refresh=True)
            perf_raw = perf_raw.dropna(how="all").reset_index(drop=True)
            st.session_state.hana_performance_df = perf_raw
        except Exception:
            errors.append("하나지사 실적관리 시트")
        if errors:
            st.error(f"불러오기 실패: {', '.join(errors)} — URL을 확인해주세요.")
        else:
            st.success("불러오기 및 저장 완료")

    if st.session_state.temp_cloud_df is not None:
        st.markdown("**본사 구글 시트 데이터** (상위 20행)")
        render_plain_html_table(strip_activity_time_columns(st.session_state.temp_cloud_df), max_rows=20)

    if st.session_state.analysis_lookup_df is not None:
        st.markdown("**하나지사 활동이력 구글 시트 데이터** (상위 20행)")
        render_plain_html_table(strip_activity_time_columns(st.session_state.analysis_lookup_df), max_rows=20)

    if st.session_state.hana_sheet_df is not None:
        st.markdown("**하나은행 구글 시트 데이터** (상위 20행)")
        render_plain_html_table(strip_activity_time_columns(st.session_state.hana_sheet_df), max_rows=20)

    if st.session_state.hana_billing_df is not None:
        st.markdown("**하나은행 청구 시트 데이터** (상위 20행)")
        render_plain_html_table(strip_activity_time_columns(st.session_state.hana_billing_df), max_rows=20)

    if st.session_state.get("hana_performance_df") is not None:
        st.markdown("**하나지사 실적관리 데이터** (상위 20행)")
        render_plain_html_table(st.session_state.hana_performance_df, max_rows=20)


def inject_theme_toggle():
    # JavaScript는 st.markdown innerHTML으로 실행 불가 → CSS :has() + radio 버튼 방식 사용
    st.markdown("""
    <style>
    .pms-sw-outer {
        position: fixed; top: 14px; right: 18px; z-index: 999999;
    }
    .pms-sw-track {
        display: inline-flex; align-items: center;
        background: #e5e7eb; border-radius: 22px; padding: 3px;
        box-shadow: 0 1px 6px rgba(0,0,0,0.12); gap: 2px;
    }
    .pms-btn {
        position: relative; cursor: pointer;
        padding: 0 10px; height: 28px; border-radius: 18px;
        line-height: 28px; opacity: 0.45; color: #374151;
        display: flex; align-items: center; justify-content: center;
        transition: opacity 0.2s, background 0.15s, color 0.15s;
        user-select: none;
    }
    .pms-btn:hover { opacity: 0.75; }
    /* radio를 label 전체에 투명하게 덮어서 한 번 클릭에 즉시 체크 */
    .pms-theme-radio {
        position: absolute; inset: 0;
        opacity: 0; cursor: pointer; margin: 0;
    }
    .pms-btn:has(input:checked) {
        opacity: 1; background: #4F46E5; color: white;
    }

    /* ══ 다크 모드 컬러 팔레트 ══
       배경:    #1e1e2e  표면:   #252535  카드:    #2a2a3e
       테두리:  #45475a  텍스트: #cdd6f4  보조:    #a6adc8
       강조:    #89b4fa  성공:   #a6e3a1  경고:    #f9e2af
       오류:    #f38ba8  정보:   #89dceb  사이드:  #16162a  */

    /* 기본 배경·텍스트 — 모든 글씨 흰색 */
    body:has(#pms-d:checked) .stApp,
    body:has(#pms-d:checked) [data-testid="stAppViewContainer"],
    body:has(#pms-d:checked) [data-testid="stMain"],
    body:has(#pms-d:checked) section.main {
        background-color: #1e1e2e !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) .main .block-container { background-color: #1e1e2e !important; }

    /* 텍스트 — 전체 흰색 */
    body:has(#pms-d:checked) h1, body:has(#pms-d:checked) h2, body:has(#pms-d:checked) h3,
    body:has(#pms-d:checked) h4, body:has(#pms-d:checked) h5, body:has(#pms-d:checked) h6 {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) p, body:has(#pms-d:checked) li,
    body:has(#pms-d:checked) [data-testid="stMarkdownContainer"],
    body:has(#pms-d:checked) [data-testid="stMarkdownContainer"] p {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) label { color: #ffffff !important; }
    /* span: 인라인 style 색상이 있는 셀(대비/증감)은 제외하고 흰색 적용 */
    body:has(#pms-d:checked) span:not([style*="color"]) { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-testid="stMarkdownContainer"] span:not([style*="color"]) { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-testid="stSidebar"] span { color: #ffffff !important; }
    body:has(#pms-d:checked) .stButton span { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-baseweb="tab"] span { color: inherit !important; }
    body:has(#pms-d:checked) [data-testid="stMetricValue"] span { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-testid="stMetricLabel"] span { color: #ffffff !important; }

    /* 파일 업로더 다크모드 — 모든 자식 요소 포함 */
    body:has(#pms-d:checked) [data-testid="stFileUploader"],
    body:has(#pms-d:checked) [data-testid="stFileUploader"] > div,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] > div > div {
        background-color: #252535 !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderDropzone"],
    body:has(#pms-d:checked) [data-testid="stFileUploader"] section,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] [class*="uploadDropzone"],
    body:has(#pms-d:checked) [data-testid="stFileUploader"] [class*="drop"] {
        background-color: #1e1e30 !important;
        border: 2px dashed #45475a !important;
        border-radius: 8px !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderDropzone"] *,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] section * {
        color: #ffffff !important;
        background-color: transparent !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderDropzone"] button,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] section button {
        background-color: #313244 !important; color: #ffffff !important;
        border: 1px solid #45475a !important; border-radius: 6px !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderDropzone"] button:hover,
    body:has(#pms-d:checked) [data-testid="stFileUploader"] section button:hover {
        background-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-testid="stFileUploaderFileData"],
    body:has(#pms-d:checked) [data-testid="stFileUploader"] [data-testid="stFileUploaderFileData"] {
        background-color: #252535 !important; border-color: #45475a !important;
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) small, body:has(#pms-d:checked) [data-testid="stCaptionContainer"] {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) hr { border-color: #45475a !important; }

    /* 사이드바 */
    body:has(#pms-d:checked) [data-testid="stSidebar"],
    body:has(#pms-d:checked) [data-testid="stSidebarContent"] {
        background-color: #171717 !important;
    }
    body:has(#pms-d:checked) [data-testid="stSidebar"] p,
    body:has(#pms-d:checked) [data-testid="stSidebar"] span,
    body:has(#pms-d:checked) [data-testid="stSidebar"] label {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stSidebar"] .stButton > button {
        background-color: transparent !important; color: #ffffff !important;
        border-color: transparent !important; box-shadow: none !important;
    }
    body:has(#pms-d:checked) [data-testid="stSidebar"] .stButton > button:hover {
        background-color: #2f2f2f !important;
    }

    /* 입력 컨트롤 */
    body:has(#pms-d:checked) input,
    body:has(#pms-d:checked) textarea,
    body:has(#pms-d:checked) [data-baseweb="input"] input,
    body:has(#pms-d:checked) [data-baseweb="textarea"] textarea {
        background-color: #252535 !important; color: #ffffff !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="input"],
    body:has(#pms-d:checked) [data-baseweb="textarea"] {
        background-color: #252535 !important; border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="select"] > div,
    body:has(#pms-d:checked) [data-baseweb="select"] div[role="combobox"] {
        background-color: #252535 !important; border-color: #45475a !important; color: #ffffff !important;
    }
    /* 셀렉트박스 드롭다운 팝업 — 전체 컨테이너 및 모든 하위 요소 */
    body:has(#pms-d:checked) [data-baseweb="popover"],
    body:has(#pms-d:checked) [data-baseweb="popover"] > div,
    body:has(#pms-d:checked) [data-baseweb="popover"] > div > div,
    body:has(#pms-d:checked) [data-baseweb="popover"] > div > div > div {
        background-color: #252535 !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="menu"],
    body:has(#pms-d:checked) [data-baseweb="menu"] ul,
    body:has(#pms-d:checked) [role="listbox"] {
        background-color: #252535 !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="menu"] li,
    body:has(#pms-d:checked) [data-baseweb="menu"] [role="option"],
    body:has(#pms-d:checked) [role="option"],
    body:has(#pms-d:checked) [role="listbox"] > div {
        background-color: #252535 !important;
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-baseweb="menu"] li span,
    body:has(#pms-d:checked) [data-baseweb="menu"] [role="option"] span,
    body:has(#pms-d:checked) [role="option"] span {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-baseweb="menu"] li:hover,
    body:has(#pms-d:checked) [data-baseweb="menu"] [aria-selected="true"],
    body:has(#pms-d:checked) [role="option"]:hover,
    body:has(#pms-d:checked) [role="option"][aria-selected="true"] {
        background-color: #3a3a5e !important;
    }
    body:has(#pms-d:checked) [data-testid="stNumberInput"] button {
        background-color: #313244 !important; color: #ffffff !important;
        border-color: #45475a !important;
    }

    /* 버튼 */
    body:has(#pms-d:checked) .stButton > button,
    body:has(#pms-d:checked) [data-testid="stButton"] button,
    body:has(#pms-d:checked) button[data-testid^="baseButton"] {
        background: #2a2a3e !important;
        background-color: #2a2a3e !important; color: #ffffff !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) .stButton > button:hover,
    body:has(#pms-d:checked) [data-testid="stButton"] button:hover,
    body:has(#pms-d:checked) button[data-testid^="baseButton"]:hover {
        background: #313244 !important;
        background-color: #313244 !important; border-color: #ffffff !important;
    }
    body:has(#pms-d:checked) .stButton > button[kind="primary"],
    body:has(#pms-d:checked) .stButton > button[data-testid="baseButton-primary"] {
        background-color: #4F46E5 !important; color: #ffffff !important;
        border-color: #6366f1 !important;
    }
    body:has(#pms-d:checked) .stButton > button[kind="primary"]:hover {
        background-color: #6366f1 !important;
    }

    /* 홈 버튼 — 아이콘 전용: 배경·테두리 없음 */
    button.pms-home-btn,
    button.pms-home-btn:hover {
        display: inline-flex !important;
        align-items: center !important;
        justify-content: center !important;
        width: 82px !important;
        height: 82px !important;
        padding: 0 !important;
        background: transparent !important;
        border: none !important;
        box-shadow: none !important;
        color: transparent !important;
        font-size: 0 !important;
    }
    button.pms-home-btn *,
    button.pms-home-btn:hover * {
        color: transparent !important;
        display: none !important;
        font-size: 0 !important;
    }
    button.pms-home-btn::before {
        content: "";
        width: 75px;
        height: 75px;
        display: block;
        flex: 0 0 auto;
        background: url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 96 96'%3E%3Cpath fill='%23000' d='M32 8h56v56c0 13.255-10.745 24-24 24H8V32C8 18.745 18.745 8 32 8Z'/%3E%3Cpath fill='%23fff' d='M20 52 48 26l28 26h-10v26H30V52H20Z'/%3E%3Cpath fill='%23fff' d='M58 34h12v18H58z'/%3E%3Cpath fill='%23000' d='M41 52h8v8h-8zM53 52h8v8h-8zM41 64h8v8h-8zM53 64h8v8h-8z'/%3E%3C/svg%3E") center / contain no-repeat;
    }
    button.pms-home-btn:hover::before {
        filter: opacity(0.86);
    }
    body:has(#pms-d:checked) button.pms-home-btn {
        background: transparent !important;
        background-color: transparent !important;
        border: none !important;
        box-shadow: none !important;
    }
    body:has(#pms-d:checked) button.pms-home-btn::before {
        filter: none;
    }
    body:has(#pms-d:checked) button.pms-home-btn:hover {
        background: transparent !important;
        background-color: transparent !important;
        border: none !important;
    }
    body:has(#pms-d:checked) button.pms-home-btn:hover::before {
        filter: opacity(0.86);
    }

    /* 새로고침 버튼 다크모드 */
    body:has(#pms-d:checked) button.pms-refresh-btn {
        background: #2a2a3e !important;
        background-color: #2a2a3e !important;
        color: #cdd6f4 !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) button.pms-refresh-btn:hover {
        background: #313244 !important;
        background-color: #313244 !important;
        color: #ffffff !important;
    }

    /* 탭 */
    body:has(#pms-d:checked) [data-baseweb="tab-list"] {
        background-color: #1e1e2e !important; border-bottom: 2px solid #45475a !important;
    }
    body:has(#pms-d:checked) [data-baseweb="tab"] {
        background-color: transparent !important; color: rgba(255,255,255,0.55) !important;
    }
    body:has(#pms-d:checked) [data-baseweb="tab"][aria-selected="true"] {
        color: #ffffff !important; font-weight: 700;
    }
    body:has(#pms-d:checked) [data-baseweb="tab-highlight"] {
        background-color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-baseweb="tab-border"] {
        background-color: #45475a !important;
    }

    /* 메트릭 카드 */
    body:has(#pms-d:checked) [data-testid="metric-container"] {
        background: #252535 !important; border: 1px solid #45475a !important; border-radius: 8px;
    }
    body:has(#pms-d:checked) [data-testid="stMetricValue"] {
        color: #ffffff !important; font-weight: 700;
    }
    body:has(#pms-d:checked) [data-testid="stMetricLabel"] { color: #ffffff !important; }
    body:has(#pms-d:checked) [data-testid="stMetricDeltaIcon-Up"]   { color: #a6e3a1 !important; }
    body:has(#pms-d:checked) [data-testid="stMetricDeltaIcon-Down"] { color: #f38ba8 !important; }

    /* Plotly 차트 */
    body:has(#pms-d:checked) .js-plotly-plot,
    body:has(#pms-d:checked) .js-plotly-plot .plot-container,
    body:has(#pms-d:checked) .js-plotly-plot .svg-container {
        background: transparent !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .bglayer rect.bg {
        fill: transparent !important;
        stroke: transparent !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .legend rect.bg {
        fill: #252535 !important;
        stroke: #45475a !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot text {
        fill: #ffffff !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .gridlayer path {
        stroke: #45475a !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .zerolinelayer path,
    body:has(#pms-d:checked) .js-plotly-plot .xlines-below path,
    body:has(#pms-d:checked) .js-plotly-plot .ylines-below path {
        stroke: #585b70 !important;
    }
    body:has(#pms-d:checked) .js-plotly-plot .modebar-btn path {
        fill: #cdd6f4 !important;
    }

    /* 알림 박스 */
    body:has(#pms-d:checked) [data-testid="stNotification"],
    body:has(#pms-d:checked) .stAlert { border-radius: 6px !important; }
    body:has(#pms-d:checked) [data-testid="stNotification"][kind="info"],
    body:has(#pms-d:checked) .stAlert.stInfo {
        background-color: #1a2a3a !important; border-left: 4px solid #89dceb !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stNotification"][kind="success"],
    body:has(#pms-d:checked) .stAlert.stSuccess {
        background-color: #1a2e1a !important; border-left: 4px solid #a6e3a1 !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stNotification"][kind="warning"],
    body:has(#pms-d:checked) .stAlert.stWarning {
        background-color: #2e2a1a !important; border-left: 4px solid #f9e2af !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stNotification"][kind="error"],
    body:has(#pms-d:checked) .stAlert.stError {
        background-color: #2e1a1a !important; border-left: 4px solid #f38ba8 !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stNotification"] p,
    body:has(#pms-d:checked) .stAlert p { color: #ffffff !important; }

    /* 익스팬더 */
    body:has(#pms-d:checked) [data-testid="stExpander"] {
        background-color: #252535 !important; border: 1px solid #45475a !important; border-radius: 6px;
    }
    body:has(#pms-d:checked) [data-testid="stExpander"] summary {
        background-color: #252535 !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stExpander"] summary:hover {
        background-color: #313244 !important;
    }

    /* ─── AG Grid 완전 다크모드 ─── */

    /* 1. CSS 변수 (AG Grid가 읽는 경우 적용) */
    body:has(#pms-d:checked) [class*="ag-theme"] {
        --ag-background-color: #252535;
        --ag-foreground-color: #ffffff;
        --ag-secondary-foreground-color: #ffffff;
        --ag-border-color: #45475a;
        --ag-secondary-border-color: #313244;
        --ag-row-border-color: #313244;
        --ag-header-background-color: #0f0f1f;
        --ag-header-foreground-color: #ffffff;
        --ag-header-column-separator-color: #45475a;
        --ag-odd-row-background-color: #1e1e30;
        --ag-row-hover-color: #2d2d45;
        --ag-selected-row-background-color: #3a3a5e;
        --ag-range-selection-background-color: rgba(255,255,255,0.1);
        --ag-input-focus-border-color: #ffffff;
        --ag-cell-horizontal-border: solid #313244;
        --ag-font-size: 13px;
        --ag-data-color: #ffffff;
        --ag-alpine-active-color: #4F46E5;
    }

    /* 2. 최외곽 래퍼 + 내부 모든 div 배경 강제 */
    body:has(#pms-d:checked) [data-testid="stDataFrame"],
    body:has(#pms-d:checked) [data-testid="stDataEditor"],
    body:has(#pms-d:checked) [data-testid="stDataFrame"] > div,
    body:has(#pms-d:checked) [data-testid="stDataFrame"] > div > div,
    body:has(#pms-d:checked) [data-testid="stDataFrame"] > div > div > div,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] > div,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] > div > div,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="grid"] {
        background-color: #13131f !important;
        border-color: #2d2d4a !important;
        color: #e2e8f0 !important;
    }
    /* stDataEditor canvas — invert 필터로 라이트 렌더링을 다크로 반전 */
    body:has(#pms-d:checked) [data-testid="stDataEditor"] canvas {
        filter: invert(1) hue-rotate(180deg) brightness(0.92) contrast(0.88) saturate(0.9) !important;
    }
    /* stDataEditor 인라인 편집 입력창 */
    body:has(#pms-d:checked) [data-testid="stDataEditor"] input,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] textarea,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [contenteditable="true"],
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [data-baseweb="input"] input {
        background-color: #1e1e34 !important;
        color: #ffffff !important;
        caret-color: #ffffff !important;
        border-color: #6366f1 !important;
        -webkit-text-fill-color: #ffffff !important;
    }

    /* 3. ag-theme 컨테이너 전체 */
    body:has(#pms-d:checked) [class*="ag-theme"],
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-root-wrapper,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-root-wrapper-body,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-root,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-body,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-body-viewport,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-body-horizontal-scroll,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-body-horizontal-scroll-viewport,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-center-cols-clipper,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-center-cols-viewport,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-center-cols-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-pinned-left-cols-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-pinned-right-cols-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-full-width-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-floating-top,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-floating-bottom,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-overlay {
        background-color: #252535 !important;
        border-color: #45475a !important;
    }

    /* 4. 헤더 */
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-row,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-viewport,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-container,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-pinned-left-header,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-pinned-right-header {
        background-color: #0f0f1f !important;
        border-bottom: 2px solid #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-cell,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-cell-comp-wrapper,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-group-cell {
        background-color: #0f0f1f !important;
        color: #ffffff !important;
        border-right: 1px solid #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-cell-text,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-cell-label {
        color: #ffffff !important;
        font-weight: 700 !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-sort-indicator-icon,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-header-icon {
        color: #ffffff !important;
    }

    /* 5. 행 — .ag-row 자체에 배경 지정 (ag-row-even 없을 때 대비) */
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row {
        background-color: #252535 !important;
        border-bottom-color: #313244 !important;
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row-odd {
        background-color: #1e1e30 !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row:hover,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row-hover {
        background-color: #2d2d45 !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-row-selected {
        background-color: #3a3a5e !important;
    }

    /* 6. 셀 */
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-cell {
        border-right-color: #313244 !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-cell:not([style*="color"]),
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-cell-value:not([style*="color"]) {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-cell span:not([style*="color"]) {
        color: #ffffff !important;
    }

    /* 7. role 속성 기반 (columnheader / gridcell) */
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="columnheader"],
    body:has(#pms-d:checked) [data-testid="stDataFrame"]  [role="columnheader"] {
        background-color: #0f0f1f !important;
        color: #ffffff !important;
        border-bottom-color: #45475a !important;
        border-right-color: #45475a !important;
    }
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="columnheader"] *,
    body:has(#pms-d:checked) [data-testid="stDataFrame"]  [role="columnheader"] * {
        color: #ffffff !important;
    }
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="gridcell"]:not([style*="color"]),
    body:has(#pms-d:checked) [data-testid="stDataFrame"]  [role="gridcell"]:not([style*="color"]),
    body:has(#pms-d:checked) [data-testid="stDataEditor"] [role="gridcell"]:not([style*="color"]) *,
    body:has(#pms-d:checked) [data-testid="stDataFrame"]  [role="gridcell"]:not([style*="color"]) * {
        color: #ffffff !important;
    }

    /* 8. 셀 편집 팝업 (SelectboxColumn 등) */
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-popup,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-popup-editor,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-popup-child {
        background-color: #252535 !important; border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select-list {
        background-color: #252535 !important; border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select-row,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-virtual-list-item {
        background-color: #252535 !important; color: #ffffff !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select-row:hover,
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-rich-select-row-selected {
        background-color: #3a3a5e !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-text-field-input {
        background-color: #252535 !important; color: #ffffff !important;
        border-color: #45475a !important;
    }
    body:has(#pms-d:checked) [class*="ag-theme"] .ag-checkbox-input-wrapper::after {
        color: #ffffff !important;
    }

    /* HTML 테이블 */
    body:has(#pms-d:checked) table { background-color: #252535 !important; border-color: #45475a !important; }
    body:has(#pms-d:checked) th {
        background-color: #0f0f1f !important; color: #ffffff !important;
        border-color: #45475a !important; font-weight: 700 !important;
    }
    body:has(#pms-d:checked) td {
        background-color: #252535 !important; color: #ffffff !important;
        border-color: #313244 !important;
    }
    body:has(#pms-d:checked) tr:nth-child(odd) td { background-color: #1e1e30 !important; }
    body:has(#pms-d:checked) tr:hover td { background-color: #2d2d45 !important; }

    /* 증감 색 — 다크모드에서 더 선명하게 */
    body:has(#pms-d:checked) [style*="color:#E53E3E"] { color: #ff7b7b !important; }
    body:has(#pms-d:checked) [style*="color:#3182CE"] { color: #74b9ff !important; }

    /* 체크박스 / 라디오 */
    body:has(#pms-d:checked) [data-baseweb="checkbox"] div,
    body:has(#pms-d:checked) [data-baseweb="radio"] div { border-color: #ffffff !important; }
    body:has(#pms-d:checked) [data-baseweb="checkbox"] [aria-checked="true"] div,
    body:has(#pms-d:checked) [data-baseweb="radio"] [aria-checked="true"] div {
        background-color: #ffffff !important;
    }

    /* 스크롤바 */
    body:has(#pms-d:checked) ::-webkit-scrollbar { width: 6px; height: 6px; }
    body:has(#pms-d:checked) ::-webkit-scrollbar-track { background: #252535; }
    body:has(#pms-d:checked) ::-webkit-scrollbar-thumb { background: #45475a; border-radius: 3px; }
    body:has(#pms-d:checked) ::-webkit-scrollbar-thumb:hover { background: #585b70; }

    /* 테마 토글 바 */
    body:has(#pms-d:checked) .pms-sw-track { background: #313244; }
    body:has(#pms-d:checked) .pms-btn { color: #ffffff; }

    /* Streamlit canvas-backed tables */
    html[data-pms-theme="d"] [data-testid="stDataEditor"],
    html[data-pms-theme="d"] [data-testid="stDataFrame"],
    body[data-pms-theme="d"] [data-testid="stDataEditor"],
    body[data-pms-theme="d"] [data-testid="stDataFrame"],
    .stApp:has(#pms-d:checked) [data-testid="stDataEditor"],
    .stApp:has(#pms-d:checked) [data-testid="stDataFrame"],
    body:has(#pms-d:checked) [data-testid="stDataEditor"],
    body:has(#pms-d:checked) [data-testid="stDataFrame"] {
        background-color: #1e1e2e !important;
        border-color: #45475a !important;
        color: #ffffff !important;
        border-radius: 8px !important;
        overflow: hidden !important;
    }
    html[data-pms-theme="d"] [data-testid="stDataEditor"] canvas,
    html[data-pms-theme="d"] [data-testid="stDataFrame"] canvas,
    body[data-pms-theme="d"] [data-testid="stDataEditor"] canvas,
    body[data-pms-theme="d"] [data-testid="stDataFrame"] canvas,
    .stApp:has(#pms-d:checked) [data-testid="stDataEditor"] canvas,
    .stApp:has(#pms-d:checked) [data-testid="stDataFrame"] canvas,
    body:has(#pms-d:checked) [data-testid="stDataEditor"] canvas,
    body:has(#pms-d:checked) [data-testid="stDataFrame"] canvas {
        filter: invert(1) hue-rotate(180deg) brightness(0.92) contrast(0.88) saturate(0.9) !important;
    }
    html[data-pms-theme="d"] [data-testid="stDataEditor"] [role="columnheader"],
    html[data-pms-theme="d"] [data-testid="stDataEditor"] [role="gridcell"],
    html[data-pms-theme="d"] [data-testid="stDataFrame"] [role="columnheader"],
    html[data-pms-theme="d"] [data-testid="stDataFrame"] [role="gridcell"],
    .stApp:has(#pms-d:checked) [data-testid="stDataEditor"] [role="columnheader"],
    .stApp:has(#pms-d:checked) [data-testid="stDataEditor"] [role="gridcell"],
    .stApp:has(#pms-d:checked) [data-testid="stDataFrame"] [role="columnheader"],
    .stApp:has(#pms-d:checked) [data-testid="stDataFrame"] [role="gridcell"] {
        background-color: #1e1e2e !important;
        border-color: #45475a !important;
        color: #ffffff !important;
        font-weight: 700 !important;
    }

    @media (prefers-color-scheme: dark) {
        body:has(#pms-s:checked) .stApp,
        body:has(#pms-s:checked) [data-testid="stAppViewContainer"],
        body:has(#pms-s:checked) section.main {
            background-color: #1e1e2e !important; color: #ffffff !important;
        }
        body:has(#pms-s:checked) .main .block-container { background-color: #1e1e2e !important; }
        body:has(#pms-s:checked) h1, body:has(#pms-s:checked) h2, body:has(#pms-s:checked) h3,
        body:has(#pms-s:checked) h4, body:has(#pms-s:checked) h5, body:has(#pms-s:checked) h6 { color: #ffffff !important; }
        body:has(#pms-s:checked) p, body:has(#pms-s:checked) [data-testid="stMarkdownContainer"] p { color: #ffffff !important; }
        body:has(#pms-s:checked) label { color: #ffffff !important; }
        body:has(#pms-s:checked) span:not([style*="color"]) { color: #ffffff !important; }
        body:has(#pms-s:checked) input, body:has(#pms-s:checked) [data-baseweb="input"] input {
            background-color: #252535 !important; color: #ffffff !important; border-color: #45475a !important;
        }
        body:has(#pms-s:checked) [class*="ag-theme"] {
            --ag-background-color: #252535; --ag-foreground-color: #ffffff;
            --ag-secondary-foreground-color: #ffffff;
            --ag-header-background-color: #0f0f1f; --ag-header-foreground-color: #ffffff;
            --ag-border-color: #45475a; --ag-row-border-color: #313244;
            --ag-odd-row-background-color: #1e1e30; --ag-row-hover-color: #2d2d45;
        }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-header,
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-header-row { background-color: #0f0f1f !important; }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-header-cell-text { color: #ffffff !important; font-weight: 700 !important; }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-row-even { background-color: #252535 !important; }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-row-odd  { background-color: #1e1e30 !important; }
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-cell,
        body:has(#pms-s:checked) [class*="ag-theme"] .ag-cell-value { color: #ffffff !important; }
        body:has(#pms-s:checked) .js-plotly-plot,
        body:has(#pms-s:checked) .js-plotly-plot .plot-container,
        body:has(#pms-s:checked) .js-plotly-plot .svg-container {
            background: transparent !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .bglayer rect.bg {
            fill: transparent !important;
            stroke: transparent !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .legend rect.bg {
            fill: #252535 !important;
            stroke: #45475a !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot text {
            fill: #ffffff !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .gridlayer path {
            stroke: #45475a !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .zerolinelayer path,
        body:has(#pms-s:checked) .js-plotly-plot .xlines-below path,
        body:has(#pms-s:checked) .js-plotly-plot .ylines-below path {
            stroke: #585b70 !important;
        }
        body:has(#pms-s:checked) .js-plotly-plot .modebar-btn path {
            fill: #cdd6f4 !important;
        }
        body:has(#pms-s:checked) [data-testid="stSidebar"],
        body:has(#pms-s:checked) [data-testid="stSidebarContent"] {
            background-color: #171717 !important;
        }
        body:has(#pms-s:checked) .stButton > button,
        body:has(#pms-s:checked) [data-testid="stButton"] button,
        body:has(#pms-s:checked) button[data-testid^="baseButton"] {
            background: #2a2a3e !important;
            background-color: #2a2a3e !important;
            color: #ffffff !important;
            border-color: #45475a !important;
        }
        body:has(#pms-s:checked) .stButton > button:hover,
        body:has(#pms-s:checked) [data-testid="stButton"] button:hover,
        body:has(#pms-s:checked) button[data-testid^="baseButton"]:hover {
            background: #313244 !important;
            background-color: #313244 !important;
            border-color: #ffffff !important;
        }
        body:has(#pms-s:checked) button.pms-home-btn {
            background: transparent !important;
            background-color: transparent !important;
            border: none !important;
            box-shadow: none !important;
        }
        body:has(#pms-s:checked) button.pms-home-btn::before {
            background-color: #cdd6f4;
        }
        body:has(#pms-s:checked) button.pms-home-btn:hover {
            background: transparent !important;
            background-color: transparent !important;
            border: none !important;
        }
        body:has(#pms-s:checked) button.pms-home-btn:hover::before {
            background-color: #818cf8;
        }
        body:has(#pms-s:checked) [data-baseweb="tab"][aria-selected="true"] { color: #ffffff !important; }
        body:has(#pms-s:checked) .pms-sw-track { background: #313244; }
    }
    </style>

    <div class="pms-sw-outer">
        <div class="pms-sw-track">
            <label class="pms-btn" title="라이트">
                <input type="radio" class="pms-theme-radio" name="pms-theme" id="pms-l" checked>
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><circle cx="12" cy="12" r="5"/><line x1="12" y1="1" x2="12" y2="3"/><line x1="12" y1="21" x2="12" y2="23"/><line x1="4.22" y1="4.22" x2="5.64" y2="5.64"/><line x1="18.36" y1="18.36" x2="19.78" y2="19.78"/><line x1="1" y1="12" x2="3" y2="12"/><line x1="21" y1="12" x2="23" y2="12"/><line x1="4.22" y1="19.78" x2="5.64" y2="18.36"/><line x1="18.36" y1="5.64" x2="19.78" y2="4.22"/></svg>
            </label>
            <label class="pms-btn" title="다크">
                <input type="radio" class="pms-theme-radio" name="pms-theme" id="pms-d">
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 12.79A9 9 0 1 1 11.21 3 7 7 0 0 0 21 12.79z"/></svg>
            </label>
        </div>
    </div>

    <!-- localStorage 테마 영속성: onload 인라인 핸들러로 실행 (script 태그는 React innerHTML에서 미실행) -->
    <img src="data:image/gif;base64,R0lGODlhAQABAIAAAAAAAP///yH5BAEAAAAALAAAAAABAAEAAAIBRAA7"
         onload="(function(){
             function readCookie(name){
                 var found = document.cookie.split('; ').find(function(row){ return row.indexOf(name + '=') === 0; });
                 return found ? decodeURIComponent(found.split('=')[1]) : '';
             }
             function normalizeTheme(t){
                 return (t === 'd' || t === 'l') ? t : 'l';
             }
             function saveTheme(t){
                 t = normalizeTheme(t);
                 localStorage.setItem('pms-theme', t);
                 document.cookie = 'pms-theme=' + encodeURIComponent(t) + '; max-age=31536000; path=/; SameSite=Lax';
             }
             function applyTheme(){
                 var t = normalizeTheme(localStorage.getItem('pms-theme') || readCookie('pms-theme') || 'l');
                 saveTheme(t);
                 document.documentElement.setAttribute('data-pms-theme', t);
                 if(document.body){ document.body.setAttribute('data-pms-theme', t); }
                 var r = document.getElementById('pms-' + t);
                 if(r && !r.checked){ r.checked = true; }
             }
             function tagSpecialBtns(){
                 document.querySelectorAll('button').forEach(function(btn){
                     var t = btn.textContent.trim();
                     if(t === '🏠') btn.classList.add('pms-home-btn');
                     if(t === '↻') btn.classList.add('pms-refresh-btn');
                 });
             }
             applyTheme();
             tagSpecialBtns();
             if(!window._pmsThemeReady){
                 window._pmsThemeReady = true;
                 var debounce;
                 window._pmsThemeObs = new MutationObserver(function(){
                     clearTimeout(debounce);
                     debounce = setTimeout(function(){ applyTheme(); tagSpecialBtns(); }, 80);
                 });
                 window._pmsThemeObs.observe(document.body, {childList:true, subtree:true});
                 document.addEventListener('change', function(e){
                     if(e.target && e.target.name === 'pms-theme'){
                         saveTheme(e.target.id.replace('pms-',''));
                         applyTheme();
                     }
                 });
             }
         })()"
         style="display:none" alt="">
    """, unsafe_allow_html=True)

    # ── 모바일 반응형 CSS ──────────────────────────────────────────
    st.markdown("""
    <style>
    @media (max-width: 768px) {
        /* ① 사이드바: 모바일에서 전체 폭 오버레이 */
        [data-testid="stSidebar"] {
            width: 100vw !important;
            min-width: 100vw !important;
            max-width: 100vw !important;
        }
        section[data-testid="stSidebar"] > div:first-child {
            width: 100vw !important;
        }

        /* ② 메인 컨텐츠: 전체 폭, 패딩 최소화 */
        .block-container {
            padding: 0.5rem 0.75rem 2rem !important;
            max-width: 100% !important;
        }
        [data-testid="stAppViewContainer"] > section.main {
            width: 100% !important;
        }

        /* ③ 테마 토글 위치·크기 축소 */
        .pms-sw-outer { top: 6px; right: 6px; }
        .pms-btn { padding: 0 6px; height: 22px; font-size: 10px; }

        /* ④ 표: 가로 스크롤 */
        .pms-report-table {
            overflow-x: auto !important;
            -webkit-overflow-scrolling: touch !important;
            max-width: 100vw !important;
        }
        table { font-size: 11px !important; min-width: max-content; }
        table th, table td { padding: 4px 6px !important; white-space: nowrap !important; }

        /* ⑤ data_editor 가로 스크롤 */
        [data-testid="stDataEditor"] {
            overflow-x: auto !important;
            -webkit-overflow-scrolling: touch !important;
        }

        /* ⑥ 컬럼 레이아웃: 2개 이상이면 가로 스크롤 */
        [data-testid="stHorizontalBlock"] {
            overflow-x: auto !important;
            -webkit-overflow-scrolling: touch !important;
            flex-wrap: nowrap !important;
            gap: 4px !important;
        }
        [data-testid="stHorizontalBlock"] > [data-testid="stColumn"] {
            min-width: 120px !important;
            flex-shrink: 0 !important;
        }

        /* ⑦ 버튼 터치 영역 */
        .stButton > button { min-height: 44px !important; font-size: 14px !important; }

        /* ⑧ 입력 필드 — iOS 자동 확대 방지 */
        input, textarea, select { font-size: 16px !important; }

        /* ⑨ 메트릭 */
        [data-testid="stMetricValue"] { font-size: 1.1rem !important; }
        [data-testid="stMetricLabel"] { font-size: 0.7rem !important; }
        [data-testid="metric-container"] { padding: 8px !important; }

        /* ⑩ 제목 */
        h1 { font-size: 1.3rem !important; }
        h2 { font-size: 1.15rem !important; }
        h3 { font-size: 1.05rem !important; }
        h4 { font-size: 0.95rem !important; }

        /* ⑪ 탭 */
        [data-baseweb="tab"] { padding: 6px 8px !important; font-size: 12px !important; }

        /* ⑫ 사이드바 햄버거 버튼 */
        [data-testid="stSidebarCollapsedControl"] button {
            width: 44px !important; height: 44px !important;
        }
    }
    </style>
    """, unsafe_allow_html=True)

def read_billing_login_upload(uploaded_file):
    file_name = (uploaded_file.name or "").lower()
    raw = uploaded_file.getvalue()
    if file_name.endswith(".csv"):
        try:
            return pd.read_csv(BytesIO(raw), dtype=str).fillna("")
        except UnicodeDecodeError:
            return pd.read_csv(BytesIO(raw), dtype=str, encoding="cp949").fillna("")
    return pd.read_excel(BytesIO(raw), dtype=str).fillna("")


def read_billing_login_upload_with_header_scan(uploaded_file):
    df = read_billing_login_upload(uploaded_file)
    if find_exact_col(df, ["고객번호"]) or find_col(df, ["고객번호"]):
        return df

    file_name = (uploaded_file.name or "").lower()
    raw = uploaded_file.getvalue()
    if file_name.endswith(".csv"):
        return df

    raw_df = pd.read_excel(BytesIO(raw), dtype=str, header=None).fillna("")
    for idx, row in raw_df.iterrows():
        values = [str(value).strip() for value in row.tolist()]
        if any(value == "고객번호" for value in values) and any(value in ["고객명", "고객사명"] for value in values):
            headers = [value or f"빈컬럼{pos + 1}" for pos, value in enumerate(values)]
            data = raw_df.iloc[idx + 1 :].copy()
            data.columns = headers
            return data.replace({np.nan: ""}).reset_index(drop=True)
    return df


def find_exact_col(df, names):
    normalized_names = {str(name).replace(" ", "").replace("　", "").lower() for name in names}
    for col in df.columns:
        col_normalized = str(col).replace(" ", "").replace("　", "").lower()
        if col_normalized in normalized_names:
            return col
    return None


def read_billing_source_upload(uploaded_file):
    file_name = (uploaded_file.name or "").lower()
    raw = uploaded_file.getvalue()
    if file_name.endswith(".csv"):
        try:
            return {"업로드 파일": pd.read_csv(BytesIO(raw), dtype=str, header=None).fillna("")}
        except UnicodeDecodeError:
            return {"업로드 파일": pd.read_csv(BytesIO(raw), dtype=str, encoding="cp949", header=None).fillna("")}
    sheets = pd.read_excel(BytesIO(raw), sheet_name=None, dtype=str, header=None)
    return {str(name): df.fillna("") for name, df in sheets.items()}


def normalize_billing_login_df(df):
    source = clean_header_logic(df.copy()).replace({np.nan: ""})
    customer_col = find_exact_col(source, ["고객번호"]) or find_col(source, ["고객번호"])
    company_col = find_exact_col(source, ["고객명", "고객사명"]) or find_col(source, ["고객명", "고객사명"])
    latest_login_col = find_exact_col(source, ["최근로그인", "최종로그인일자", "최종로그인"]) or find_col(
        source,
        ["최근로그인", "최종로그인일자", "최종로그인"],
    )
    first_login_col = find_exact_col(source, ["신규일자"]) or find_col(source, ["신규일자"])
    login_count_col = find_exact_col(source, ["로그인", "로그인횟수", "로그인수"]) or find_col(
        source,
        ["로그인횟수", "로그인수"],
    )
    is_reference_only = not login_count_col

    missing = [
        label
        for label, col in {
            "고객번호": customer_col,
            "고객명": company_col,
        }.items()
        if not col or col not in source.columns
    ]
    if not is_reference_only:
        missing.extend(
            [
                label
                for label, col in {
                    "최근로그인": latest_login_col,
                    "로그인": login_count_col,
                }.items()
                if not col or col not in source.columns
            ]
        )
    if missing:
        return pd.DataFrame(), missing

    result = pd.DataFrame(
        {
            "고객번호": source[customer_col].apply(normalize_customer_no),
            "고객명": source[company_col].astype(str).str.strip(),
            "최초로그인": source[first_login_col].astype(str).str.strip() if first_login_col else "",
            "최근로그인": source[latest_login_col].astype(str).str.strip() if latest_login_col and login_count_col else "",
            "로그인": pd.to_numeric(source[login_count_col].astype(str).str.replace(",", "", regex=False), errors="coerce")
            .fillna(0)
            .astype(int)
            if login_count_col
            else 0,
        }
    )
    result = result[result["고객번호"].astype(str).str.strip().ne("")]
    return result.reset_index(drop=True), []


def merge_billing_login_dfs(dfs):
    if not dfs:
        return pd.DataFrame(columns=["고객번호", "고객명", "최초로그인", "최근로그인", "로그인"])

    merged = pd.concat(dfs, ignore_index=True).replace({np.nan: ""})
    if merged.empty:
        return merged

    merged["_login_num"] = pd.to_numeric(merged["로그인"], errors="coerce").fillna(0).astype(int)
    merged["_login_date"] = pd.to_datetime(merged["최근로그인"], errors="coerce")
    merged["_first_login_date"] = pd.to_datetime(merged["최초로그인"], errors="coerce")
    rows = []
    for customer_no, group in merged.groupby("고객번호", sort=False):
        with_login = group[group["_login_num"] > 0]
        base = with_login.iloc[-1] if not with_login.empty else group.iloc[-1]
        latest = group["_login_date"].max()
        latest_text = latest.strftime("%Y%m%d") if pd.notna(latest) else str(base.get("최근로그인", "")).strip()
        first = group["_first_login_date"].max()
        first_values = [str(value).strip() for value in group["최초로그인"].tolist() if str(value).strip()]
        first_text = first.strftime("%Y%m%d") if pd.notna(first) else (first_values[-1] if first_values else "")
        name_values = [str(value).strip() for value in group["고객명"].tolist() if str(value).strip()]
        rows.append(
            {
                "고객번호": customer_no,
                "고객명": name_values[-1] if name_values else "",
                "최초로그인": first_text,
                "최근로그인": latest_text,
                "로그인": int(group["_login_num"].sum()),
            }
        )
    return pd.DataFrame(rows, columns=["고객번호", "고객명", "최초로그인", "최근로그인", "로그인"])


def normalize_billing_source_sheet(df):
    source = df.copy().replace({np.nan: ""})
    header_index = None
    for idx, row in source.iterrows():
        values = {str(value).strip() for value in row.tolist() if str(value).strip()}
        if is_billing_source_header(values):
            header_index = idx
            break
    if header_index is None:
        return clean_header_logic(source)

    headers = [str(value).strip() or f"빈컬럼{pos + 1}" for pos, value in enumerate(source.loc[header_index].tolist())]
    data = source.iloc[header_index + 1 :].copy()
    data.columns = headers
    data = data.replace({np.nan: ""})
    data = data.loc[~data.apply(lambda row: all(str(value).strip() == "" for value in row), axis=1)]
    return clean_header_logic(data.reset_index(drop=True))


def is_billing_source_header(values):
    return any("고객번호" in value for value in values) and any(
        "사업자번호" in value or "사업자등록번호" in value for value in values
    )


def parse_billing_source_sections(df):
    source = df.copy().replace({np.nan: ""})
    sections = []
    current_title = "업로드 데이터"
    row_count = len(source)
    idx = 0
    while idx < row_count:
        row_values = [str(value).strip() for value in source.iloc[idx].tolist()]
        nonempty_values = [value for value in row_values if value]
        is_header = is_billing_source_header(nonempty_values)
        if nonempty_values and not is_header:
            current_title = nonempty_values[0]
            idx += 1
            continue
        if not is_header:
            idx += 1
            continue

        headers = [value or f"빈컬럼{pos + 1}" for pos, value in enumerate(row_values)]
        data_rows = []
        idx += 1
        while idx < row_count:
            next_values = [str(value).strip() for value in source.iloc[idx].tolist()]
            next_nonempty = [value for value in next_values if value]
            next_is_header = is_billing_source_header(next_nonempty)
            if next_is_header:
                break
            if len(next_nonempty) == 1 and not next_nonempty[0].isdigit():
                break
            if next_nonempty:
                data_rows.append(next_values)
            idx += 1

        table = pd.DataFrame(data_rows, columns=headers).replace({np.nan: ""})
        sections.append((current_title, clean_header_logic(table.reset_index(drop=True))))
    return sections


def pick_billing_source_sheet(source_sheets, kind):
    normalized = {name: normalize_billing_source_sheet(df) for name, df in source_sheets.items()}
    if kind == "open":
        for name, df in normalized.items():
            if any(token in name for token in ["구축", "5월", "6월"]) or "ERP연계 여부" in df.columns:
                return name, df
    if kind == "erp":
        for name, df in normalized.items():
            if any(token in name for token in ["연계", "수령"]) or "연계시작일자" in df.columns:
                return name, df
    fallback_name = next(iter(source_sheets.keys()), "")
    return fallback_name, normalized.get(fallback_name, pd.DataFrame())


def login_lookup_from_df(login_df):
    if login_df is None or login_df.empty:
        return {}
    lookup = {}
    for _, row in login_df.iterrows():
        customer_keys = billing_customer_keys(row.get("고객번호", ""))
        if not customer_keys:
            continue
        login_info = {
            "고객명": str(row.get("고객명", "")).strip(),
            "최초로그인": str(row.get("최초로그인", "")).strip(),
            "최근로그인": str(row.get("최근로그인", "")).strip(),
            "로그인": row.get("로그인", 0),
        }
        for customer_key in customer_keys:
            lookup.setdefault(customer_key, login_info)
    return lookup


def has_billing_login_history(login_lookup, customer_no):
    for customer_key in billing_customer_keys(customer_no):
        info = login_lookup.get(customer_key)
        if not info:
            continue
        try:
            if int(float(str(info.get("로그인", 0)).replace(",", "") or 0)) > 0:
                return True
        except Exception:
            pass
        if str(info.get("최근로그인", "")).strip() or str(info.get("최초로그인", "")).strip():
            return True
    return False


def slice_billing_sections_by_count(sections, target_count):
    selected = []
    remaining = []
    left = max(0, int(target_count or 0))
    for title, df in sections:
        if df is None or df.empty:
            selected.append((title, pd.DataFrame(columns=getattr(df, "columns", []))))
            continue
        if left <= 0:
            selected.append((title, df.iloc[0:0].copy()))
            remaining.append((title, df.copy()))
            continue
        take_count = min(left, len(df))
        selected.append((title, df.iloc[:take_count].copy()))
        if take_count < len(df):
            remaining.append((title, df.iloc[take_count:].copy()))
        left -= take_count
    return selected, remaining


def append_billing_sections_excel(writer, sheet_name, sections):
    safe_name = str(sheet_name)[:31] or "Sheet"
    startrow = 0
    for title, df in sections:
        sheet_df = df.copy() if isinstance(df, pd.DataFrame) else pd.DataFrame()
        pd.DataFrame([[title]]).to_excel(writer, index=False, header=False, sheet_name=safe_name, startrow=startrow)
        startrow += 1
        sheet_df.to_excel(writer, index=False, sheet_name=safe_name, startrow=startrow)
        startrow += len(sheet_df) + 3


def billing_sections_excel_bytes(open_sections, erp_sections):
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        append_billing_sections_excel(writer, "청구원본(개설업로드)", open_sections)
        append_billing_sections_excel(writer, "청구원본(연계업로드)", erp_sections)
    output.seek(0)
    return output.getvalue()


def build_billing_download_sections(parsed_open_sections, parsed_erp_sections, open_count, link_count, login_df):
    open_main_sections = [(title, df) for title, df in parsed_open_sections if "사용자교육" not in title]
    education_sections = [(title, df) for title, df in parsed_open_sections if "사용자교육" in title]

    selected_open, remaining_open = slice_billing_sections_by_count(open_main_sections, open_count)
    selected_erp, _ = slice_billing_sections_by_count(parsed_erp_sections, link_count)

    login_lookup = login_lookup_from_df(login_df)
    moved_rows = []
    edu_columns = None
    for _, edu_df in education_sections:
        if isinstance(edu_df, pd.DataFrame) and edu_columns is None:
            edu_columns = list(edu_df.columns)
    for _, remain_df in remaining_open:
        if remain_df is None or remain_df.empty:
            continue
        if edu_columns is None:
            edu_columns = list(remain_df.columns)
        customer_col = find_col(remain_df, ["고객번호", "고객번호(당월)"])
        if not customer_col or customer_col not in remain_df.columns:
            continue
        for _, row in remain_df.iterrows():
            if has_billing_login_history(login_lookup, row.get(customer_col, "")):
                moved_rows.append(row.to_dict())

    moved_df = pd.DataFrame(moved_rows)
    if edu_columns:
        moved_df = moved_df.reindex(columns=edu_columns, fill_value="")
    education_title = education_sections[0][0] if education_sections else "사용자교육(방문)대기 고객사"
    final_open_sections = selected_open + [(education_title, moved_df)]
    return final_open_sections, selected_erp


def billing_display_value(value):
    if is_blank_value(value):
        return ""
    if isinstance(value, (datetime, pd.Timestamp)):
        return pd.to_datetime(value).strftime("%Y-%m-%d")
    if isinstance(value, (int, float, np.integer, np.floating)):
        num = float(value)
        if 20000 <= num <= 60000:
            return (datetime(1899, 12, 30) + timedelta(days=int(num))).strftime("%Y-%m-%d")
        if num.is_integer():
            return str(int(num))
    return str(value).strip()


def billing_format_date(value):
    if is_blank_value(value):
        return ""
    parsed = parse_sheet_date(value)
    if pd.notna(parsed):
        return parsed.strftime("%Y-%m-%d")
    return str(value).strip()


_ERP_TYPE_MAP = {"연계형": "O", "기본형": "X"}


def _map_erp_type(value):
    return _ERP_TYPE_MAP.get(str(value).strip(), value)


def find_col_by_priority(df, keys):
    for key in keys:
        col = find_col(df, [key])
        if col and col in df.columns:
            return col
    return None


def billing_customer_keys(value):
    normalized = normalize_customer_no(value)
    keys = [normalized] if normalized else []
    if normalized:
        stripped = normalized.lstrip("0")
        if stripped and stripped not in keys:
            keys.append(stripped)
        if len(normalized) == 8:
            padded = normalized.zfill(9)
            if padded not in keys:
                keys.append(padded)
    return keys


def load_billing_customer_reference():
    try:
        ref_df = clean_header_logic(read_google_csv(st.session_state.get("url_hana", DEFAULT_URL_HANA), header=2))
    except Exception as exc:
        return {}, f"고객원장 참조 시트를 불러올 수 없습니다: {exc}"

    ref_df = clean_header_logic(ref_df.copy()).replace({np.nan: ""})
    customer_col = find_col(ref_df, ["고객번호", "고개번호", "고객NO", "고객 No", "고객"])
    if not customer_col or customer_col not in ref_df.columns:
        return {}, "고객원장 참조 시트에서 고객번호 컬럼을 찾을 수 없습니다."

    col_map = {
        "사업자번호": find_col_by_priority(ref_df, ["사업자번호"]),
        "업체명": find_col_by_priority(ref_df, ["고객명", "고객사명", "업체명"]),
        "ERP연계 여부": find_col_by_priority(ref_df, ["구축형"]),
        "관리구분": find_col_by_priority(ref_df, ["관리구분", "관리 구분"]),
        "접수일자": find_col_by_priority(ref_df, ["신규접수일", "접수일자"]),
        "구축일자": find_col_by_priority(ref_df, ["개설/이행일", "구축일자", "구축일"]),
        "방문일자": find_col_by_priority(ref_df, ["개설/이행일", "방문일자", "방문일", "구축일자", "구축일"]),
        "해지일자": find_col_by_priority(ref_df, ["해지일자", "해지일", "해약일"]),
        "담당자": find_col_by_priority(ref_df, ["담당자"]),
        "청구원본 고객명": find_col_by_priority(ref_df, ["고객명", "고객사명", "업체명"]),
    }

    lookup = {}
    for _, row in ref_df.iterrows():
        customer_keys = billing_customer_keys(row.get(customer_col, ""))
        if not customer_keys:
            continue
        reference_info = {
            target_col: billing_display_value(row.get(source_col, ""))
            for target_col, source_col in col_map.items()
            if source_col and source_col in ref_df.columns
        }
        for customer_key in customer_keys:
            lookup.setdefault(customer_key, reference_info)
        biz_key = normalize_biz_no(reference_info.get("사업자번호", ""))
        if biz_key:
            lookup.setdefault(f"biz:{biz_key}", reference_info)
    return lookup, ""


def billing_value(row, candidates):
    for candidate in candidates:
        if candidate in row.index:
            return str(row.get(candidate, "")).strip()
    return ""


def billing_fill_blank(value, reference_info, ref_key):
    text = billing_display_value(value)
    return text if text else reference_info.get(ref_key, "")


def billing_dates_differ(left, right):
    left_date = parse_sheet_date(left)
    right_date = parse_sheet_date(right)
    if pd.isna(left_date) or pd.isna(right_date):
        return False
    return left_date.normalize() != right_date.normalize()


def billing_prelogin_note(build_date, visit_date):
    build_dt = parse_sheet_date(build_date)
    visit_dt = parse_sheet_date(visit_date)
    if pd.isna(build_dt) or pd.isna(visit_dt):
        return ""
    build_month = build_dt.to_period("M")
    visit_month = visit_dt.to_period("M")
    if build_month >= visit_month:
        return ""
    return f"{build_dt.month:02d}월 사전 로그인 후 {visit_dt.month:02d}월 사용자교육 진행"


def append_billing_note(remark, note):
    remark_text = billing_display_value(remark)
    if not note or note in remark_text:
        return remark_text
    return f"{remark_text} / {note}" if remark_text else note


def billing_reference_info(reference_lookup, customer_no):
    for customer_key in billing_customer_keys(customer_no):
        if customer_key in reference_lookup:
            return reference_lookup[customer_key]
    return {}


def billing_reference_info_from_row(reference_lookup, row):
    customer_no = billing_value(row, ["고객번호", "고객번호(당월)"])
    reference_info = billing_reference_info(reference_lookup, customer_no)
    if reference_info:
        return reference_info
    biz_no = billing_value(row, ["사업자번호", "사업자등록번호", "사업자번호(당월)", "사업자등록번호(당월)"])
    biz_key = normalize_biz_no(biz_no)
    if biz_key:
        return reference_lookup.get(f"biz:{biz_key}", {})
    return {}


def billing_termination_text(reference_info):
    if not reference_info:
        return ""
    manage_text = str(reference_info.get("관리구분", "")).strip()
    end_text = billing_format_date(reference_info.get("해지일자", ""))
    has_terminated_status = bool(re.search("해지|취소", manage_text, flags=re.IGNORECASE))
    if end_text:
        return f"해지 이력 있음({end_text})"
    if has_terminated_status:
        return f"해지 이력 있음({manage_text})" if manage_text else "해지 이력 있음"
    return ""


def build_open_billing_table(source_df, login_df=None, reference_lookup=None):
    open_columns = [
        "순번",
        "고객번호",
        "사업자번호",
        "업체명",
        "ERP연계 여부",
        "접수일자",
        "최종로그인일자",
        "방문일자",
        "담당자",
        "비고",
        "최초로그인",
        "최근로그인일자",
        "로그인횟수",
        "청구원본 고객명",
        "실적파일 고객명",
        "해지체크",
    ]
    if source_df is None or source_df.empty:
        return pd.DataFrame(columns=open_columns)
    login_lookup = login_lookup_from_df(login_df)
    reference_lookup = reference_lookup or {}
    rows = []
    for idx, row in source_df.iterrows():
        customer_no = billing_value(row, ["고객번호", "고객번호(당월)"])
        if not customer_no and not billing_value(row, ["업체명", "고객명"]):
            continue
        login_info = billing_reference_info(login_lookup, customer_no)
        reference_info = billing_reference_info_from_row(reference_lookup, row)
        first_login = login_info.get("최초로그인", "")
        build_date = billing_fill_blank(billing_value(row, ["구축일자", "구축일"]), reference_info, "구축일자")
        if first_login and billing_dates_differ(build_date, first_login):
            build_date = first_login
        visit_date = billing_fill_blank(billing_value(row, ["방문일자", "방문일"]), reference_info, "방문일자")
        remark = append_billing_note(billing_value(row, ["비고"]), billing_prelogin_note(build_date, visit_date))
        rows.append(
            {
                "순번": billing_value(row, ["순번"]) or str(idx + 1),
                "고객번호": customer_no,
                "사업자번호": billing_fill_blank(
                    billing_value(row, ["사업자번호", "사업자등록번호", "사업자번호(당월)", "사업자등록번호(당월)"]),
                    reference_info,
                    "사업자번호",
                ),
                "업체명": billing_fill_blank(
                    billing_value(row, ["업체명", "고객명", "업체명(당월)", "고객명(당월)"]),
                    reference_info,
                    "업체명",
                ),
                "ERP연계 여부": _map_erp_type(billing_fill_blank(billing_value(row, ["ERP연계 여부", "ERP연계여부"]), reference_info, "ERP연계 여부")),
                "해지체크": billing_termination_text(reference_info),
                "접수일자": billing_fill_blank(billing_value(row, ["접수일자"]), reference_info, "접수일자"),
                "최종로그인일자": build_date,
                "방문일자": visit_date,
                "담당자": billing_fill_blank(billing_value(row, ["담당자", "담당자(당월)"]), reference_info, "담당자"),
                "비고": remark,
                "최초로그인": first_login,
                "최근로그인일자": login_info.get("최근로그인", ""),
                "로그인횟수": login_info.get("로그인", ""),
                "청구원본 고객명": billing_fill_blank(
                    billing_value(row, ["업체명", "고객명", "업체명(당월)", "고객명(당월)"]),
                    reference_info,
                    "청구원본 고객명",
                ),
                "실적파일 고객명": login_info.get("고객명", ""),
            }
        )
    df = pd.DataFrame(rows, columns=open_columns)
    for col in ("접수일자", "최종로그인일자", "방문일자", "최초로그인", "최근로그인일자"):
        if col in df.columns:
            df[col] = df[col].apply(billing_format_date)
    return df


def build_erp_billing_table(source_df, login_df=None, reference_lookup=None):
    erp_columns = [
        "순서",
        "고객번호",
        "사업자번호",
        "업체명",
        "구분",
        "추가연계신청일자",
        "담당자",
        "구축일",
        "연계시작일자",
        "은행연계완료일자",
        "수령여부",
        "비고",
        "최초로그인",
        "최종로그인일자",
        "로그인횟수",
        "청구원본 고객명",
        "실적파일 고객명",
        "해지체크",
    ]
    if source_df is None or source_df.empty:
        return pd.DataFrame(columns=erp_columns)
    login_lookup = login_lookup_from_df(login_df)
    reference_lookup = reference_lookup or {}
    rows = []
    for idx, row in source_df.iterrows():
        customer_no = billing_value(row, ["고객번호", "고객번호(당월)"])
        if not customer_no and not billing_value(row, ["업체명", "고객명"]):
            continue
        login_info = billing_reference_info(login_lookup, customer_no)
        reference_info = billing_reference_info_from_row(reference_lookup, row)
        rows.append(
            {
                "순서": billing_value(row, ["순서", "순번"]) or str(idx + 1),
                "고객번호": customer_no,
                "사업자번호": billing_value(row, ["사업자번호", "사업자등록번호"]),
                "업체명": billing_value(row, ["업체명", "고객명"]),
                "해지체크": billing_termination_text(reference_info),
                "구분": billing_value(row, ["구분"]),
                "추가연계신청일자": billing_value(row, ["추가연계신청일자"]),
                "담당자": billing_value(row, ["담당자"]),
                "구축일": billing_value(row, ["구축일", "구축일자"]),
                "연계시작일자": billing_value(row, ["연계시작일자"]),
                "은행연계완료일자": billing_value(row, ["은행연계완료일자"]),
                "수령여부": billing_value(row, ["수령여부"]),
                "비고": billing_value(row, ["비고"]),
                "최초로그인": login_info.get("최초로그인", ""),
                "최종로그인일자": login_info.get("최근로그인", ""),
                "로그인횟수": login_info.get("로그인", ""),
                "청구원본 고객명": billing_value(row, ["업체명", "고객명"]),
                "실적파일 고객명": login_info.get("고객명", ""),
            }
        )
    df = pd.DataFrame(rows, columns=erp_columns)
    for col in ("추가연계신청일자", "구축일", "연계시작일자", "은행연계완료일자", "최초로그인", "최종로그인일자"):
        if col in df.columns:
            df[col] = df[col].apply(billing_format_date)
    return df


def load_education_waiting_section(login_df=None, reference_lookup=None):
    try:
        raw_df = clean_header_logic(pd.read_csv(URL_EDUCATION_WAITING).replace({np.nan: ""}))
    except Exception as exc:
        return "사용자교육(방문)대기 고객사", pd.DataFrame(), str(exc)
    status_col = find_col(raw_df, ["개설상태"])
    if status_col:
        mask = raw_df[status_col].astype(str).str.strip().isin(["개설대기", "개설진행"])
        raw_df = raw_df[mask].reset_index(drop=True)
    table_df = build_open_billing_table(raw_df, login_df, reference_lookup or {})
    return "사용자교육(방문)대기 고객사", table_df, None


def billing_preview_style(df):
    if df is None or df.empty or "해지체크" not in df.columns:
        return df

    def highlight_termination(value):
        return "background-color: #fde2e2; color: #991b1b; font-weight: 600;" if str(value).strip() else ""

    styled = df.style
    if hasattr(styled, "map"):
        return styled.map(highlight_termination, subset=["해지체크"])
    return styled.applymap(highlight_termination, subset=["해지체크"])


def render_billing_source_tables(source_upload=None, login_df=None):
    st.markdown("#### 청구 원본 표")
    source_file = source_upload or st.file_uploader(
        "구축 및 연계 리스트 업로드",
        type=["xlsx", "xls", "csv"],
        key="billing_source_upload",
        help="2026년 6월 구축 및 연계 리스트_최종본.xlsx 파일을 업로드해주세요.",
    )
    open_sections = [("2026년 6월 구축 실적", build_open_billing_table(pd.DataFrame(), login_df))]
    erp_sections = [("당월 ERP연계 청구 고객사", build_erp_billing_table(pd.DataFrame(), login_df))]
    education_sections = [("사용자교육(방문)대기 고객사", pd.DataFrame())]
    parsed_open_sections = []
    parsed_erp_sections = []
    reference_lookup = {}
    source_loaded = False

    if source_file:
        try:
            source_sheets = read_billing_source_upload(source_file)
            open_sheet_name, open_source = pick_billing_source_sheet(source_sheets, "open")
            erp_sheet_name, erp_source = pick_billing_source_sheet(source_sheets, "erp")
            parsed_open_sections = parse_billing_source_sections(source_sheets.get(open_sheet_name, open_source))
            parsed_erp_sections = parse_billing_source_sections(source_sheets.get(erp_sheet_name, erp_source))
            if not parsed_open_sections:
                parsed_open_sections = [("2026년 6월 구축 실적", open_source)]
            if not parsed_erp_sections:
                parsed_erp_sections = [("당월 ERP연계 청구 고객사", erp_source)]
            reference_lookup, reference_error = load_billing_customer_reference()
            if reference_error:
                st.warning(reference_error)
            education_sections = [
                (title, build_open_billing_table(section_df, login_df, reference_lookup))
                for title, section_df in parsed_open_sections
                if "사용자교육" in title
            ] or education_sections
            open_sections = [
                (title, build_open_billing_table(section_df, login_df, reference_lookup))
                for title, section_df in parsed_open_sections
                if "사용자교육" not in title
            ] or [("2026년 6월 구축 실적", build_open_billing_table(open_source, login_df, reference_lookup))]
            erp_sections = [
                (title, build_erp_billing_table(section_df, login_df, reference_lookup))
                for title, section_df in parsed_erp_sections
            ] or [("당월 ERP연계 청구 고객사", build_erp_billing_table(erp_source, login_df, reference_lookup))]
            source_loaded = True
        except Exception as exc:
            st.error(f"구축 및 연계 리스트 파일을 읽을 수 없습니다: {exc}")
    else:
        st.warning("구축 및 연계 리스트를 업로드하면 청구원본 표가 Google Sheet 업로드 탭 화면 구조로 채워집니다.")

    open_sections = [(t, d) for t, d in open_sections if "사용자교육" not in t]
    open_sections.extend(education_sections)

    open_tab, erp_tab = st.tabs(["청구원본(개설업로드)", "청구원본(연계업로드)"])
    with open_tab:
        st.caption("개설 청구자료 생성 화면 구성입니다.")
        section_tabs = st.tabs([title for title, _ in open_sections])
        for tab, (section_title, section_df) in zip(section_tabs, open_sections):
            with tab:
                st.markdown(f"##### {section_title}")
                if "구축 실적" in section_title or "사용자교육" in section_title:
                    section_df = section_df.drop(columns=["최종로그인일자"], errors="ignore")
                if section_df.empty:
                    st.info("표시할 데이터가 없습니다.")
                else:
                    st.dataframe(billing_preview_style(section_df), use_container_width=True, hide_index=True)
    with erp_tab:
        st.caption("연계 청구자료 생성 화면 구성입니다.")
        for section_title, section_df in erp_sections:
            st.markdown(f"##### {section_title}")
            st.dataframe(billing_preview_style(section_df), use_container_width=True, hide_index=True)

    st.divider()
    st.markdown("#### 청구자료 엑셀 다운로드")
    input_col1, input_col2, download_col = st.columns([0.22, 0.22, 0.24])
    with input_col1:
        open_count = st.number_input("개설청구 갯수 입력", min_value=0, step=1, value=0, key="billing_open_count")
    with input_col2:
        link_count = st.number_input("연계청구 갯수 입력", min_value=0, step=1, value=0, key="billing_link_count")
    with download_col:
        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
        if source_loaded:
            download_open_sections, download_erp_sections = build_billing_download_sections(
                parsed_open_sections,
                parsed_erp_sections,
                open_count,
                link_count,
                login_df,
            )
            st.download_button(
                "청구자료 엑셀 다운로드",
                data=billing_sections_excel_bytes(download_open_sections, download_erp_sections),
                file_name=f"청구자료_{(datetime.utcnow() + timedelta(hours=9)).strftime('%Y%m%d_%H%M%S')}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
                key="billing_source_final_download",
            )
        else:
            st.button("청구자료 엑셀 다운로드", use_container_width=True, disabled=True)
            st.caption("구축 및 연계 리스트를 업로드하면 다운로드할 수 있습니다.")


def read_excel_sheet_with_header_scan(file, sheet_name, required_keys=None, max_scan_rows=15):
    required_keys = required_keys or []
    raw = pd.read_excel(file, sheet_name=sheet_name, header=None)
    header_idx = 0
    for idx in range(min(len(raw), max_scan_rows)):
        row_values = [str(v).strip() for v in raw.iloc[idx].tolist()]
        compact_values = [v.replace(" ", "").replace("\n", "") for v in row_values]
        if all(any(key.replace(" ", "").replace("\n", "") in value for value in compact_values) for key in required_keys):
            header_idx = idx
            break
    df = raw.iloc[header_idx + 1:].copy()
    df.columns = _dedupe_columns([str(v).strip() for v in raw.iloc[header_idx].tolist()])
    df = df.loc[:, ~pd.Series(df.columns).astype(str).str.contains("^nan$|^Unnamed", case=False, na=False).values]
    return df.dropna(how="all").reset_index(drop=True)


def _dedupe_columns(columns):
    seen = {}
    result = []
    for col in columns:
        name = str(col).strip()
        count = seen.get(name, 0)
        seen[name] = count + 1
        result.append(name if count == 0 else f"{name}.{count}")
    return result


def _infer_month_from_name(file_name):
    match = re.search(r"(20\d{2})[-_ .]?(0[1-9]|1[0-2])", str(file_name or ""))
    if not match:
        return None
    return f"{match.group(1)}-{match.group(2)}"


def _infer_month_from_bank_df(bank_df):
    if bank_df is None or bank_df.empty:
        return None
    login_col = find_col(bank_df, ["최종로그인일자", "최근로그인일자", "최종로그인", "최근로그인"])
    if not login_col:
        return None
    dates = bank_df[login_col].apply(parse_sheet_date).dropna()
    if dates.empty:
        return None
    return dates.max().strftime("%Y-%m")


def _month_bounds(year_month):
    start = pd.to_datetime(f"{year_month}-01", errors="coerce")
    if pd.isna(start):
        return pd.NaT, pd.NaT
    return start, start + pd.offsets.MonthEnd(0)


def _format_date_for_display(value):
    parsed = parse_sheet_date(value)
    if pd.isna(parsed):
        return ""
    return parsed.strftime("%Y-%m-%d")


def _operation_contact_reason(row):
    reasons = []
    if bool(row.get("_is_this_year_build")):
        reasons.append("올해 구축")
    elif bool(row.get("_is_gap_watch_build")):
        reasons.append("2025년 이후 구축")
    if bool(row.get("_is_this_year_link")):
        reasons.append("올해 연계")
    elif bool(row.get("_is_gap_watch_link")):
        reasons.append("2025년 이후 연계")
    if bool(row.get("_login_stopped")):
        reasons.append("최근 3개월 로그인 없음")
    if bool(row.get("_is_billable")):
        reasons.append("최근 3개월 로그인 있음")
    if bool(row.get("_has_recent_menu")):
        reasons.append("메뉴사용 이력 있음")
    return " / ".join(reasons)


def build_operation_activity_targets(
    manage_df,
    bank_df,
    reference_month,
    contact_only=True,
    billable_only=True,
    max_targets=None,
    sort_mode="recent_build",
    exclude_current_open_month=True,
):
    if manage_df is None or manage_df.empty:
        return pd.DataFrame(), "관리 엑셀의 고객원장 데이터를 찾을 수 없습니다."
    if bank_df is None or bank_df.empty:
        return pd.DataFrame(), "은행 엑셀의 명단 데이터를 찾을 수 없습니다."

    manage = clean_header_logic(manage_df.copy())
    bank = bank_df.copy()
    manage.columns = _dedupe_columns([str(c).strip() for c in manage.columns])
    bank.columns = _dedupe_columns([str(c).strip() for c in bank.columns])

    m_customer_col = find_col(manage, ["고객번호"])
    m_biz_col = find_col(manage, ["사업자번호", "사업자등록번호"])
    m_company_col = find_col(manage, ["고객명", "고객사명", "업체명", "상호"])
    m_owner_col = find_col(manage, ["담당자"])
    m_manage_col = find_col(manage, ["관리구분"])
    m_open_status_col = find_col(manage, ["개설상태"])
    m_open_date_col = find_col(manage, ["개설/이행일", "개설일", "이행일"])
    m_link_status_col = find_col(manage, ["연계상태", "ERP연계상태"])
    m_link_date_col = find_col(manage, ["연계청구일자", "연계일자", "ERP연계일자"])
    m_cancel_col = find_col(manage, ["해지일자"])

    b_customer_col = find_col(bank, ["고객번호"])
    b_biz_col = find_col(bank, ["사업자번호", "사업자등록번호"])
    b_company_col = find_col(bank, ["고객명", "고객사명", "업체명", "상호"])
    b_last_login_col = find_col(bank, ["최종로그인일자", "최근로그인일자", "최종로그인", "최근로그인"])
    b_last_transfer_col = find_col(bank, ["최종이체일자", "최근이체일자", "최종이체", "최근이체"])
    b_login_count_col = find_col(bank, ["로그인건수", "로그인횟수"])
    b_menu_col = exact_col(bank, ["메뉴사용"]) or find_col(bank, ["메뉴사용", "메뉴클릭수"])
    b_bill_col = find_col(bank, ["청구구분"])

    if not m_customer_col and not m_biz_col:
        return pd.DataFrame(), "관리 엑셀에서 고객번호 또는 사업자번호 컬럼을 찾을 수 없습니다."
    if not b_customer_col and not b_biz_col:
        return pd.DataFrame(), "은행 엑셀에서 고객번호 또는 사업자번호 컬럼을 찾을 수 없습니다."

    start_month, end_month = _month_bounds(reference_month)
    if pd.isna(start_month):
        return pd.DataFrame(), "기준월을 확인할 수 없습니다."
    recent_start = start_month - pd.DateOffset(months=2)
    year_start = pd.Timestamp(year=start_month.year, month=1, day=1)
    year_end = pd.Timestamp(year=start_month.year, month=12, day=31)
    gap_watch_start = pd.Timestamp(year=max(2025, start_month.year - 1), month=1, day=1)

    manage_key = pd.Series([""] * len(manage), index=manage.index, dtype="object")
    if m_customer_col:
        manage_key = manage[m_customer_col].apply(normalize_billing_customer_no)
    if m_biz_col:
        biz_key = manage[m_biz_col].apply(normalize_biz_no)
        manage_key = manage_key.where(manage_key.astype(str) != "", biz_key)

    bank_key = pd.Series([""] * len(bank), index=bank.index, dtype="object")
    if b_customer_col:
        bank_key = bank[b_customer_col].apply(normalize_billing_customer_no)
    if b_biz_col:
        bank_biz_key = bank[b_biz_col].apply(normalize_biz_no)
        bank_key = bank_key.where(bank_key.astype(str) != "", bank_biz_key)

    bank_norm = pd.DataFrame({
        "_match_key": bank_key,
        "은행고객명": bank[b_company_col].fillna("").astype(str).str.strip() if b_company_col else "",
        "은행사업자번호": bank[b_biz_col].apply(normalize_biz_no) if b_biz_col else "",
        "은행최종로그인일자": bank[b_last_login_col].apply(_format_date_for_display) if b_last_login_col else "",
        "은행최종이체일자": bank[b_last_transfer_col].apply(_format_date_for_display) if b_last_transfer_col else "",
        "은행로그인건수": pd.to_numeric(bank[b_login_count_col], errors="coerce").fillna(0).astype(int) if b_login_count_col else 0,
        "은행메뉴사용": pd.to_numeric(bank[b_menu_col], errors="coerce").fillna(0).astype(int) if b_menu_col else 0,
        "청구구분": bank[b_bill_col].fillna("").astype(str).str.strip() if b_bill_col else "",
    })
    bank_norm["_last_login_dt"] = pd.to_datetime(bank_norm["은행최종로그인일자"], errors="coerce")
    bank_norm = bank_norm[bank_norm["_match_key"].astype(str) != ""]
    bank_norm = bank_norm.sort_values(["_last_login_dt", "은행로그인건수", "은행메뉴사용"], ascending=[False, False, False])
    bank_norm = bank_norm.drop_duplicates("_match_key", keep="first")

    open_dates = manage[m_open_date_col].apply(parse_sheet_date) if m_open_date_col else pd.Series(pd.NaT, index=manage.index)
    link_dates = manage[m_link_date_col].apply(parse_sheet_date) if m_link_date_col else pd.Series(pd.NaT, index=manage.index)
    cancel_dates = manage[m_cancel_col].apply(parse_sheet_date) if m_cancel_col else pd.Series(pd.NaT, index=manage.index)

    manage_norm = pd.DataFrame({
        "_match_key": manage_key,
        "고객번호": manage[m_customer_col].apply(normalize_billing_customer_no) if m_customer_col else "",
        "사업자번호": manage[m_biz_col].apply(normalize_biz_no) if m_biz_col else "",
        "고객사명": manage[m_company_col].fillna("").astype(str).str.strip() if m_company_col else "",
        "담당자": manage[m_owner_col].fillna("").astype(str).str.strip() if m_owner_col else "",
        "관리구분": manage[m_manage_col].fillna("").astype(str).str.strip() if m_manage_col else "",
        "개설상태": manage[m_open_status_col].fillna("").astype(str).str.strip() if m_open_status_col else "",
        "개설/이행일": open_dates.dt.strftime("%Y-%m-%d").fillna(""),
        "연계상태": manage[m_link_status_col].fillna("").astype(str).str.strip() if m_link_status_col else "",
        "연계일자": link_dates.dt.strftime("%Y-%m-%d").fillna(""),
        "해지일자": cancel_dates.dt.strftime("%Y-%m-%d").fillna(""),
        "_open_dt": open_dates,
        "_link_dt": link_dates,
        "_cancel_dt": cancel_dates,
    })
    manage_norm = manage_norm[manage_norm["_match_key"].astype(str) != ""]
    manage_norm = manage_norm.drop_duplicates("_match_key", keep="first")

    merged = manage_norm.merge(bank_norm, on="_match_key", how="left")
    last_login = pd.to_datetime(merged.get("은행최종로그인일자", ""), errors="coerce")
    login_count = pd.to_numeric(merged.get("은행로그인건수", 0), errors="coerce").fillna(0)
    menu_count = pd.to_numeric(merged.get("은행메뉴사용", 0), errors="coerce").fillna(0)

    valid_manage = ~merged["관리구분"].astype(str).str.contains("해지|취소", na=False)
    valid_manage &= merged["_cancel_dt"].isna()
    if exclude_current_open_month:
        valid_manage &= merged["_open_dt"].dt.strftime("%Y-%m").fillna("") != reference_month
    this_year_build = merged["_open_dt"].between(year_start, year_end, inclusive="both")
    this_year_link = merged["_link_dt"].between(year_start, year_end, inclusive="both")
    gap_watch_build = merged["_open_dt"].between(gap_watch_start, year_end, inclusive="both")
    gap_watch_link = merged["_link_dt"].between(gap_watch_start, year_end, inclusive="both")
    recent_login = last_login.between(recent_start, end_month, inclusive="both")

    merged["_is_this_year_build"] = this_year_build
    merged["_is_this_year_link"] = this_year_link
    merged["_is_gap_watch_build"] = gap_watch_build
    merged["_is_gap_watch_link"] = gap_watch_link
    merged["_is_billable"] = recent_login
    merged["_has_recent_menu"] = menu_count.gt(0)
    merged["_login_stopped"] = (gap_watch_build | gap_watch_link) & ~recent_login
    merged["_recent_build_dt"] = pd.concat([merged["_open_dt"], merged["_link_dt"]], axis=1).max(axis=1)
    merged["_recent_build_rank_dt"] = merged["_recent_build_dt"].fillna(pd.Timestamp.min)
    merged["_priority_score"] = (
        merged["_login_stopped"].astype(int) * 100
        + merged["_is_billable"].astype(int) * 20
        + menu_count.clip(upper=50)
        + login_count.clip(upper=30)
    )

    billable_mask = valid_manage & merged["_is_billable"] & merged["_has_recent_menu"]
    stopped_mask = valid_manage & merged["_login_stopped"]
    result_billable = merged[billable_mask].copy()
    result_stopped = merged[stopped_mask].copy() if contact_only else pd.DataFrame()

    if sort_mode == "recent_login":
        sort_cols = ["_last_login_dt", "_recent_build_rank_dt", "담당자", "고객사명"]
        sort_ascending = [False, False, True, True]
    elif sort_mode == "usage":
        sort_cols = ["은행메뉴사용", "은행로그인건수", "_recent_build_rank_dt", "담당자", "고객사명"]
        sort_ascending = [False, False, False, True, True]
    else:
        sort_cols = ["_recent_build_rank_dt", "_last_login_dt", "담당자", "고객사명"]
        sort_ascending = [False, False, True, True]

    result_billable = result_billable.sort_values(sort_cols, ascending=sort_ascending)
    if max_targets:
        result_billable = result_billable.head(int(max_targets))
    if not result_stopped.empty:
        result_stopped = result_stopped.sort_values(["_recent_build_rank_dt", "담당자", "고객사명"], ascending=[False, True, True])
        result = pd.concat([result_stopped, result_billable], ignore_index=True)
        result = result.drop_duplicates("_match_key", keep="first")
    else:
        result = result_billable

    if result.empty:
        return pd.DataFrame(), None

    result["구분"] = np.where(result["_login_stopped"], "로그인 끊김 점검", "청구 가능 운영관리")
    result["선정사유"] = result.apply(_operation_contact_reason, axis=1)
    result["최근3개월기준"] = f"{recent_start.strftime('%Y-%m-%d')} ~ {end_month.strftime('%Y-%m-%d')}"
    result["최근구축/연계일"] = result["_recent_build_dt"].dt.strftime("%Y-%m-%d").fillna("")
    result["활동권장"] = np.where(
        result["_login_stopped"],
        "올해 구축/연계 후 로그인 공백 확인 및 사용 재유도",
        "최근 사용 고객 청구 유지 및 추가 사용 메뉴 점검",
    )
    for col in ["은행로그인건수", "은행메뉴사용"]:
        if col in result.columns:
            result[col] = pd.to_numeric(result[col], errors="coerce").fillna(0).astype(int)
    display_cols = [
        "구분", "담당자", "고객사명", "고객번호", "사업자번호", "관리구분",
        "개설상태", "개설/이행일", "연계상태", "연계일자", "최근구축/연계일",
        "은행최종로그인일자", "은행최종이체일자", "은행로그인건수", "은행메뉴사용",
        "청구구분", "최근3개월기준", "선정사유", "활동권장",
    ]
    result["_type_order"] = np.where(result["_login_stopped"], 0, 1)
    result = result.sort_values(["_type_order"] + sort_cols, ascending=[True] + sort_ascending)
    return result[[col for col in display_cols if col in result.columns]].reset_index(drop=True), None


def show_operation_activity_targets():
    st.markdown("### 운영관리 활동고객 선정")
    st.caption("HANA사업부 고객관리 엑셀과 은행 통합CMS 고객명단을 비교해 청구 가능 고객과 로그인 공백 점검 고객을 추립니다.")

    default_manage_path = r"C:\Users\이성환\Downloads\25년_HANA사업부_고객관리 (1).xlsx"
    default_bank_path = r"C:\Users\이성환\Downloads\00.(명단)통합CMS고객명단_202606.xlsx"
    default_files_available = os.path.exists(default_manage_path) and os.path.exists(default_bank_path)

    use_default = st.checkbox(
        "로컬 Downloads의 기본 엑셀 파일 사용",
        value=default_files_available,
        disabled=not default_files_available,
        key="operation_targets_use_default_files",
    )

    col_a, col_b = st.columns(2)
    with col_a:
        manage_upload = st.file_uploader("HANA사업부 고객관리 엑셀", type=["xlsx", "xls"], key="operation_manage_upload")
    with col_b:
        bank_upload = st.file_uploader("은행 통합CMS 고객명단 엑셀", type=["xlsx", "xls"], key="operation_bank_upload")

    if use_default:
        manage_source = default_manage_path
        bank_source = default_bank_path
    else:
        manage_source = manage_upload
        bank_source = bank_upload

    if not manage_source or not bank_source:
        st.info("두 엑셀 파일을 업로드하면 활동고객사를 선정합니다.")
        return

    try:
        manage_df = read_excel_sheet_with_header_scan(manage_source, "고객원장", required_keys=["고객번호", "사업자번호", "고객명"])
        bank_df = pd.read_excel(bank_source, sheet_name="명단")
    except Exception as exc:
        st.error(f"엑셀을 읽는 중 오류가 발생했습니다: {exc}")
        return

    bank_file_name = os.path.basename(default_bank_path) if use_default else getattr(bank_upload, "name", "")
    inferred_month = (
        _infer_month_from_bank_df(bank_df)
        or _infer_month_from_name(bank_file_name)
        or (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m")
    )
    if "operation_reference_month" not in st.session_state:
        st.session_state["operation_reference_month"] = inferred_month
    ctrl_a, ctrl_b, ctrl_c, ctrl_d = st.columns([0.18, 0.18, 0.24, 0.4])
    reference_month = ctrl_a.text_input("기준월", key="operation_reference_month", help="예: 2026-07")
    max_targets = ctrl_b.number_input("활동 고객수", min_value=1, max_value=500, value=50, step=10, key="operation_target_limit")
    sort_label = ctrl_c.selectbox(
        "우선순위",
        ["최근 구축/연계순", "최근 로그인순", "사용량순"],
        key="operation_target_sort",
    )
    billable_only = ctrl_d.checkbox(
        "비용 청구 가능 고객만",
        value=True,
        key="operation_billable_only",
        help="직전 3개월 최종로그인 이력이 있는 고객만 선정합니다.",
    )
    include_stopped = st.checkbox(
        "2025년 이후 구축/연계 후 로그인 끊긴 고객도 별도 점검 대상으로 포함",
        value=True,
        key="operation_include_stopped",
    )
    exclude_current_open_month = st.checkbox(
        "당월 개설/이행 고객 제외",
        value=True,
        key="operation_exclude_current_open_month",
        help="기준월과 개설/이행일의 월이 같은 고객은 활동 리스트에서 제외합니다.",
    )
    sort_mode = {
        "최근 구축/연계순": "recent_build",
        "최근 로그인순": "recent_login",
        "사용량순": "usage",
    }.get(sort_label, "recent_build")

    targets, error = build_operation_activity_targets(
        manage_df,
        bank_df,
        reference_month,
        contact_only=include_stopped,
        billable_only=billable_only,
        max_targets=max_targets,
        sort_mode=sort_mode,
        exclude_current_open_month=exclude_current_open_month,
    )
    if error:
        st.error(error)
        return
    if targets.empty:
        st.info("조건에 맞는 활동고객사가 없습니다.")
        return

    stopped_count = int(targets["구분"].astype(str).eq("로그인 끊김 점검").sum()) if "구분" in targets.columns else 0
    billable_count = int(targets["구분"].astype(str).eq("청구 가능 운영관리").sum()) if "구분" in targets.columns else 0
    c1, c2, c3 = st.columns(3)
    c1.metric("전체 표시", f"{len(targets):,}개")
    c2.metric("청구 가능 활동", f"{billable_count:,}개")
    c3.metric("로그인 공백 점검", f"{stopped_count:,}개")

    owner_options = ["전체"] + sorted(v for v in targets.get("담당자", pd.Series(dtype=str)).astype(str).str.strip().unique() if v)
    f1, f2 = st.columns([0.25, 0.75])
    selected_owner = f1.selectbox("담당자", owner_options, key="operation_target_owner")
    keyword = f2.text_input("고객사 검색", key="operation_target_keyword")

    view = targets.copy()
    if selected_owner != "전체" and "담당자" in view.columns:
        view = view[view["담당자"].astype(str).str.strip() == selected_owner]
    if keyword.strip() and "고객사명" in view.columns:
        view = view[view["고객사명"].astype(str).str.contains(keyword.strip(), case=False, na=False)]

    gap_view = view[view["구분"].astype(str).eq("로그인 끊김 점검")].reset_index(drop=True) if "구분" in view.columns else pd.DataFrame()
    billable_view = view[view["구분"].astype(str).eq("청구 가능 운영관리")].reset_index(drop=True) if "구분" in view.columns else view

    gap_tab, billable_tab, all_tab = st.tabs([
        f"로그인 공백 점검 {len(gap_view):,}",
        f"청구 가능 활동 {len(billable_view):,}",
        f"전체 {len(view):,}",
    ])
    with gap_tab:
        st.caption("2025년 이후 개설/연계 고객 중 직전 3개월 로그인 이력이 없는 고객입니다.")
        st.dataframe(gap_view, use_container_width=True, hide_index=True)
    with billable_tab:
        st.caption("직전 3개월 로그인 이력이 있어 비용 청구 가능한 고객 중 활동 고객수만큼 선정한 목록입니다.")
        st.dataframe(billable_view, use_container_width=True, hide_index=True)
    with all_tab:
        st.dataframe(view, use_container_width=True, hide_index=True)

    st.download_button(
        "활동고객사 엑셀 다운로드",
        data=dataframe_to_excel_bytes({
            "로그인공백점검": gap_view,
            "청구가능활동": billable_view,
            "전체": view,
        }),
        file_name=f"운영관리_활동고객사_{reference_month.replace('-', '')}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True,
    )


def show_billing_generation():
    st.markdown("### 청구자료 생성")
    st.caption("청구 원본과 은행 로그인 실적파일을 대사해 청구자료 생성 전 확인 목록을 만듭니다.")
    st.link_button(
        "원본 Google Sheet 열기",
        "https://docs.google.com/spreadsheets/d/12BeCTDegUWD-jomaG3WS75Jx1dJ9Lqjxn1hVs3FrpE4/edit?gid=1244892381#gid=1244892381",
    )
    st.info(
        "FastAPI/React 화면에서는 해당 시트의 개설·연계 청구원본과 은행 실적파일을 대사합니다. "
        "현재 Google Sheet가 비공개이면 CSV export 권한 안내가 표시됩니다."
    )
    uploaded_login = st.file_uploader(
        "은행로그인실적파일(은행) 엑셀 업로드",
        type=["xlsx", "xls", "csv"],
        accept_multiple_files=True,
        help="메뉴사용현황 파일과 통합CMS 가입명세 파일을 함께 업로드할 수 있습니다.",
    )

    if not uploaded_login:
        st.warning("은행로그인실적파일(은행)을 업로드하면 청구자료 생성용 실적 데이터를 확인할 수 있습니다.")
        render_billing_source_tables()
        return

    try:
        normalized_files = []
        missing_by_file = []
        for login_file in uploaded_login:
            raw_df = read_billing_login_upload_with_header_scan(login_file)
            file_df, missing = normalize_billing_login_df(raw_df)
            if missing:
                missing_by_file.append(f"{login_file.name}: {', '.join(missing)}")
                continue
            normalized_files.append(file_df)
        if missing_by_file:
            st.error(f"필수 컬럼을 찾을 수 없습니다: {' / '.join(missing_by_file)}")
            st.caption("필요 컬럼: 고객번호, 고객명. 메뉴사용현황 파일은 최근로그인, 로그인 컬럼도 필요합니다.")
            render_billing_source_tables()
            return
        normalized_df = merge_billing_login_dfs(normalized_files)
    except Exception as exc:
        st.error(f"파일을 읽을 수 없습니다: {exc}")
        return

    total_login = int(normalized_df["로그인"].sum()) if not normalized_df.empty else 0
    latest_login = "-"
    parsed_dates = pd.to_datetime(normalized_df["최근로그인"], errors="coerce")
    if parsed_dates.notna().any():
        latest_login = parsed_dates.max().strftime("%Y-%m-%d")

    col_total, col_customer, col_login, col_latest = st.columns(4)
    col_total.metric("업로드 행수", f"{len(normalized_df):,}건")
    col_customer.metric("고객 수", f"{normalized_df['고객번호'].nunique():,}곳")
    col_login.metric("로그인 합계", f"{total_login:,}회")
    col_latest.metric("최신 로그인일", latest_login)

    excel_bytes = dataframe_to_excel_bytes({"은행로그인실적파일": normalized_df})
    st.download_button(
        "정규화된 은행로그인실적파일 다운로드",
        data=excel_bytes,
        file_name="은행로그인실적파일_정규화.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True,
    )
    render_billing_source_tables(login_df=normalized_df)


def show_main():
    apply_global_table_css()
    inject_theme_toggle()
    show_sidebar()
    persist_current_menu()

    menu = st.session_state.current_menu
    allowed_menus = {
        "대시보드",
        "업로드 및 실적 확인",
        "이번달 활동 대상고객 추천",
        OPERATION_TARGET_MENU,
        "주간보고 이력 작성",
        "관리자용 실적 확인",
        "실적 분석/계산",
        "실적 보고서",
        BILLING_MENU,
        "청구자료 작성",
        "주간보고 취합",
        "운영계획",
        "직원 및 권한설정",
        "구글 스트레드시트 연동",
        ACTIVITY_TEMPLATE_CONVERT_MENU,
    }
    if menu not in allowed_menus:
        st.session_state.current_menu = "업로드 및 실적 확인"
        persist_current_menu()
        st.rerun()
    user_menus = {"업로드 및 실적 확인", "이번달 활동 대상고객 추천", OPERATION_TARGET_MENU, "주간보고 이력 작성", ACTIVITY_TEMPLATE_CONVERT_MENU, BILLING_MENU}
    if menu not in user_menus and st.session_state.user_role != "관리자":
        st.session_state.current_menu = "업로드 및 실적 확인"
        persist_current_menu()
        st.rerun()

    render_page_title(menu)

    if menu == "대시보드":
        show_dashboard()
    elif menu == "업로드 및 실적 확인":
        show_user_history()
    elif menu == "이번달 활동 대상고객 추천":
        show_target_customers()
    elif menu == OPERATION_TARGET_MENU:
        show_operation_activity_targets()
    elif menu == "주간보고 이력 작성":
        show_weekly_report_user()
    elif menu == "관리자용 실적 확인":
        show_admin_performance()
    elif menu == "실적 분석/계산":
        show_admin_analysis()
    elif menu == "실적 보고서":
        show_report()
    elif menu == BILLING_MENU:
        show_billing_generation()
    elif menu == "청구자료 작성":
        show_billing_materials()
    elif menu == "주간보고 취합":
        show_weekly_report_admin()
    elif menu == "운영계획":
        show_operation_plan()
    elif menu == "직원 및 권한설정":
        show_staff_admin()
    elif menu == "구글 스트레드시트 연동":
        show_google_sync()
    elif menu == ACTIVITY_TEMPLATE_CONVERT_MENU:
        show_activity_template_converter()
    else:
        st.session_state.current_menu = "업로드 및 실적 확인"
        st.rerun()


show_main()
